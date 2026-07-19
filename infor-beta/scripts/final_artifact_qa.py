"""Fail-closed QA for the exact final PowerPoint and Excel artefacts."""

from __future__ import annotations

from dataclasses import asdict, dataclass
from pathlib import Path
from typing import Any, Callable, Mapping

from openpyxl import load_workbook
from pptx import Presentation


@dataclass(frozen=True)
class FinalArtifactQAResult:
    final_deck_path: str
    final_combined_workbook_path: str
    ready_for_delivery: bool
    blocking_issues: tuple[str, ...]
    warnings: tuple[str, ...]
    remaining_placeholders: tuple[str, ...]
    financial_chart_insertion_status: str
    visual_qa_render_status: str
    visual_qa_rendered_slides: tuple[int, ...]
    aggregation_backend: str
    aggregation_degraded: bool
    aggregation_verified: bool
    recalculation_status: str
    sources_deleted: bool | None
    kept_source_workbooks: tuple[str, ...]
    manual_completion_steps: tuple[str, ...]

    def to_dict(self) -> dict[str, Any]:
        """Return a JSON-serializable stage-output mapping."""
        value = asdict(self)
        for key, item in tuple(value.items()):
            if isinstance(item, tuple):
                value[key] = list(item)
        return value


def _remaining_placeholders(deck: Path) -> tuple[str, ...]:
    findings: list[str] = []
    presentation = Presentation(deck)
    for slide_number, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            texts: list[str] = []
            if getattr(shape, "has_text_frame", False):
                texts.append(shape.text)
            if getattr(shape, "has_table", False):
                texts.extend(cell.text for row in shape.table.rows for cell in row.cells)
            for text in texts:
                normalized = " ".join(text.split())
                if "placeholder" in normalized.casefold():
                    findings.append(f"slide {slide_number}: {normalized}")
    return tuple(findings)


def run_final_artifact_qa(
    *,
    deck_path: Path | str,
    combined_workbook_path: Path | str,
    output_dir: Path | str,
    chart_status: Mapping[str, Any] | None = None,
    aggregation_status: Mapping[str, Any] | None = None,
    visual_qa_evidence: Mapping[str, Any] | None = None,
    render_fn: Callable[..., list[Path]] | None = None,
) -> FinalArtifactQAResult:
    """Inspect the exact final artefact paths and return a fail-closed decision."""
    deck = Path(deck_path).expanduser().resolve()
    workbook = Path(combined_workbook_path).expanduser().resolve()
    issues: list[str] = []
    warnings: list[str] = []

    if not deck.is_file():
        issues.append(f"final deck does not exist: {deck}")
    if not workbook.is_file():
        issues.append(f"final combined workbook does not exist: {workbook}")
    else:
        try:
            inspected_workbook = load_workbook(workbook, read_only=True, data_only=False)
            if not inspected_workbook.sheetnames:
                issues.append("final combined workbook contains no worksheets")
            inspected_workbook.close()
        except Exception as exc:
            issues.append(f"final combined workbook could not be inspected: {exc}")

    chart_status = chart_status or {}
    chart_required = bool(chart_status.get("required", False))
    chart_ok = bool(chart_status.get("financial_summary_inserted")) and bool(
        chart_status.get("ltm_pie_inserted")
    )
    chart_insertion_status = "not_applicable"
    manual_steps: list[str] = []
    if chart_required:
        chart_insertion_status = "passed" if chart_ok else "failed"
        if not chart_ok:
            issues.append("financial chart insertion failed or was incomplete")
            manual_steps.append(
                "Run chart insertion on a supported render backend and rerun final QA."
            )

    aggregation = aggregation_status or {}
    aggregation_backend = str(aggregation.get("backend", "unknown"))
    aggregation_degraded = bool(aggregation.get("degraded", False))
    merge_integrity_verified = aggregation.get("merge_integrity_verified") is True
    relink_ok = aggregation.get("relink_ok") is True
    recalculation_ok = aggregation.get("recalculation_ok") is True
    external_refs = tuple(aggregation.get("external_refs") or ())
    sources_deleted_value = aggregation.get("sources_deleted")
    sources_deleted = (
        sources_deleted_value if isinstance(sources_deleted_value, bool) else None
    )
    kept_sources = tuple(str(path) for path in aggregation.get("kept_sources") or ())
    warnings.extend(str(item) for item in aggregation.get("warnings") or ())
    recalculation_status = "passed" if recalculation_ok else "failed"
    aggregation_verified = (
        bool(aggregation)
        and not aggregation_degraded
        and merge_integrity_verified
        and relink_ok
        and recalculation_ok
        and not external_refs
    )
    if not aggregation:
        issues.append("workbook aggregation verification status is missing")
    else:
        if aggregation_degraded:
            issues.append("workbook aggregation is degraded")
        if not merge_integrity_verified or not relink_ok or external_refs:
            issues.append("workbook aggregation integrity is not verified")
        if not recalculation_ok:
            issues.append("workbook recalculation is not verified")
        if not aggregation_verified:
            manual_steps.append(
                "Retry workbook aggregation with a full-fidelity backend and verified recalculation."
            )

    placeholders: tuple[str, ...] = ()
    if deck.is_file():
        try:
            placeholders = _remaining_placeholders(deck)
        except Exception as exc:
            issues.append(f"final deck could not be inspected: {exc}")
        if placeholders:
            issues.append(
                f"final deck contains {len(placeholders)} prohibited placeholder(s)"
            )
            manual_steps.append("Replace every prohibited placeholder and rerun final QA.")

    visual_status = "unavailable"
    rendered_slide_indices: tuple[int, ...] = ()
    if deck.is_file():
        try:
            slide_count = len(Presentation(deck).slides)
            required_indices = list(range(slide_count))
            if not required_indices:
                raise RuntimeError("the final deck contains no slides")
            if render_fn is None:
                from slide_render import render_deck_to_png

                render_fn = render_deck_to_png
            rendered = render_fn(deck, Path(output_dir), slide_indices=required_indices)
            if len(rendered) != len(required_indices) or not all(Path(path).is_file() for path in rendered):
                raise RuntimeError(
                    f"renderer produced {len(rendered)} of {len(required_indices)} required slide image(s)"
                )
            rendered_slide_indices = tuple(required_indices)
            evidence = visual_qa_evidence or {}
            inspected = tuple(evidence.get("inspected_slide_indices", ()))
            overflow_free = evidence.get("overflow_free") is True
            if inspected == rendered_slide_indices and overflow_free:
                visual_status = "passed"
            else:
                visual_status = "rendered_unverified"
                issues.append(
                    "visual QA evidence is missing or incomplete for the rendered final deck"
                )
                manual_steps.append(
                    "Inspect every rendered slide for overflow/overlap and record complete visual QA evidence."
                )
        except Exception as exc:
            issues.append(f"visual QA could not run; deck remains draft/non-final: {exc}")
            manual_steps.append(
                "Run final QA where PowerPoint COM or LibreOffice slide rendering is available."
            )

    return FinalArtifactQAResult(
        final_deck_path=str(deck),
        final_combined_workbook_path=str(workbook),
        ready_for_delivery=not issues,
        blocking_issues=tuple(issues),
        warnings=tuple(warnings),
        remaining_placeholders=placeholders,
        financial_chart_insertion_status=chart_insertion_status,
        visual_qa_render_status=visual_status,
        visual_qa_rendered_slides=rendered_slide_indices,
        aggregation_backend=aggregation_backend,
        aggregation_degraded=aggregation_degraded,
        aggregation_verified=aggregation_verified,
        recalculation_status=recalculation_status,
        sources_deleted=sources_deleted,
        kept_source_workbooks=kept_sources,
        manual_completion_steps=tuple(manual_steps),
    )
