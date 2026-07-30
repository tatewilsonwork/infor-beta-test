"""Financial Summary charts — build the deck's metric charts and place them.

The `financial-charts` stage. One clustered-column chart per metric row — four for
the default single Financial Summary slide, eight when the deck spec asked for two
FS slides — is built on the deal workbook's `financial-summary` tab, where each
flow metric's ``=INDEX('ltm-metrics'!…)`` LTM link is an ordinary internal
reference to a sibling tab. The charts are then rendered and dropped into the
Financial Summary slides' chart placeholders (four per slide, discovered by
scanning the deck for the Metric #1 placeholder), stretched to each placeholder's
box — the same picture-into-placeholder pattern as the cap-table / ownership
insertions.

Before Phase D this had to be a *post-aggregation* stage: the two tabs lived in
separate standalone files until `workbook-aggregation` merged them, so the LTM
links read ``#N/A`` and could not be charted. The deal owns one workbook from
stage one now, so this stage's only ordering constraint is that it follows
`deck` — the deck it edits.

This stage MUTATES THE ALREADY-ASSEMBLED DECK IN PLACE — it must never re-run the
`deck-assembler` (or any other skill). The orchestrators below only open
``deck_path``, insert pictures into existing placeholders, and save; nothing here
dispatches a sub-agent or rebuilds the deck.

One backend, everywhere: **openpyxl + LibreOffice headless**. openpyxl writes
the native charts on the tab; each PNG is rendered from a single-chart temp
workbook built off LibreOffice-recalculated values.

Phase D deleted the Excel-COM alternative. It only ever ran on a Windows dev box
with Excel installed — production is Cowork/Linux, which has neither — so it was
a second implementation of every chart-formatting decision that no deliverable
depended on, and the platform whose output actually ships was the one exercised
less. Its measured cost was real: a wedged automation instance is unkillable
in-process (Phase I), and the four tests that drove it had to be opt-in behind
`INFOR_EXCEL_COM_TESTS` because the Cap IQ add-in raises a modal dialog that
`DisplayAlerts = False` cannot suppress.

INFOR chart formatting (the only formatting that matters): Palatino Linotype 9 pt
black (data labels + category-axis labels); no chart title; no chart border; the
category (horizontal) axis rendered as a **solid black baseline line of visible
width** (a hairline / width-less line is dropped by the LibreOffice PNG render, so
the width is explicit); no major gridlines; the value (vertical) axis hidden — no
line, no label; gap width 50%; data labels on every bar at Outside End, in the
same ``$`` currency format as the tab's value cells (so a bar reads ``$102.7``
exactly like its cell); all bars filled RGB(70, 86, 110) = hex ``46566E``.

The same module also builds the overview slide's **LTM revenue pie**
(``render_ltm_revenue_pie_into_deck``): a by-segment pie over the deal
workbook's ``ltm-metrics`` tab "LTM Revenue Overview" block. The pie series is the
**"% of Total" column** (the ``=B/Btotal`` fraction), so its data labels show the
segment share (value-only, ``#,##0.0%`` format) rather than the dollar amounts —
and only slices whose share is **above 3%** carry a label (smaller ones overlap
each other in the short overview box). The legend is docked on the RIGHT with the
segment names, and the pie's plot area is pinned to the left of the chart box so
it sits clear of the legend; no title/border, and slice fills from the INFOR theme
accent palette (``pptx_helpers.INFOR_ACCENTS``). It rides the same
post-aggregation stage and is dropped into the overview slide's "Rectangle 4"
placeholder.
"""

from __future__ import annotations

import io
import sys
import tempfile
from pathlib import Path

from pptx import Presentation

from deal_workbook import TAB_FINANCIAL_SUMMARY, TAB_LTM_METRICS
from pptx_helpers import INFOR_ACCENTS
from template_layout import (
    MARKER_BUILT_FINANCIAL_SUMMARY,
    MARKER_BUILT_OVERVIEW,
    find_slide_by_marker,
    find_slides_by_marker,
)

# --- INFOR chart constants ---------------------------------------------------
_BAR_RGB_HEX = "46566E"  # RGB(70, 86, 110) as openpyxl RRGGBB
_FONT_NAME = "Palatino Linotype"
_FONT_SIZE_PT = 9
_FONT_SIZE_HUNDREDTHS = _FONT_SIZE_PT * 100  # openpyxl drawing fonts use 1/100 pt
_GAP_WIDTH = 50
# Data-label number format for the four metric charts — the SAME currency format
# as the financial-summary tab's value cells (financial_summary_workbook
# ._VALUE_FORMAT), so a bar's label reads "$102.7" exactly like its cell.
_VALUE_FORMAT = '$#,##0.0_);($#,##0.0);"--"'

# Category (horizontal) axis baseline: an explicit, visible solid-black line.
# Without an explicit width the openpyxl ``<a:ln>`` defaults to a hairline that the
# LibreOffice PNG render drops entirely, leaving the bars with no baseline (Issue
# 2). 12700 EMU = 1.0 pt.
_AXIS_LINE_WIDTH_EMU = 12700

# --- tab / slide geometry ----------------------------------------------------
# From the constants, not spelled out: this stage reads the DEAL workbook, so its
# tab names are `deal_workbook`'s. Both were literals duplicating those constants
# until v0.5.45 — the same shape as that release's bug, where a tab name recorded
# in one place and assumed in another drifted silently.
_SHEET_DEFAULT = TAB_FINANCIAL_SUMMARY
_HEADER_ROW = 5
_FIRST_DATA_ROW = _HEADER_ROW + 1
# Chart placeholder shape names on EVERY Financial Summary slide (the deck can
# carry one or two — a cloned FS slide keeps the same shape names). In tile
# order: #1 top-left, #2 top-right, #3 / #4 below. Slide k charts the tab's
# metric rows 4k..4k+3 (rows 6-9 on the first slide, 10-13 on the second).
#
# ORDERING: the names are the contract, and this list's ORDER pairs them to the
# tab's metric rows, so a stale mapping would put each chart under the wrong
# label with nothing failing. `test_the_chart_placeholders_are_named_in_reading_order`
# checks the library's geometry against it.
_PLACEHOLDER_NAMES: list[str] = [
    "Rectangle 17",  # Metric #1 (label tile Rectangle 13)
    "Rectangle 7",   # Metric #2 (label tile Rectangle 12)
    "Rectangle 19",  # Metric #3 (label tile Rectangle 15)
    "Rectangle 18",  # Metric #4 (label tile Rectangle 14)
]
_CHARTS_PER_SLIDE = len(_PLACEHOLDER_NAMES)
# The shape whose presence identifies a Financial Summary slide when the deck is
# scanned (the Metric #1 chart placeholder, present until this stage fills it).
_FS_MARKER_PLACEHOLDER = _PLACEHOLDER_NAMES[0]
# Placeholder box: 4.53" x 2.51". Used to size the exported chart so its aspect
# matches the slide box (the picture is stretched to the box either way).
_CHART_W_PT = 4.53 * 72
_CHART_H_PT = 2.51 * 72
_CHART_W_CM = 4.53 * 2.54
_CHART_H_CM = 2.51 * 2.54

# --- LTM revenue pie (overview slide) ----------------------------------------
# The pie is built on the deal workbook's `ltm-metrics` tab — its "LTM Revenue
# Overview" block carries literal segment × LTM-revenue values — and dropped into
# the overview slide's wide/short "[Pie Chart Placeholder]" (Rectangle 4). The
# data live on the same deal workbook the FS charts use, so this rides the
# post-aggregation `financial-charts` stage rather than a parallel path.
_PIE_SHEET_DEFAULT = TAB_LTM_METRICS
_PIE_SECTION_LABEL = "LTM Revenue Overview"
_PIE_TOTAL_LABEL = "Total"
_PIE_SEGMENT_COL = 1  # column A — segment names
_PIE_VALUE_COL = 2    # column B — LTM revenue $ amounts (Python grouping/fractions read this)
# The pie charts at most 5 slices: the 4 largest segments (descending) plus an
# "Other" slice grouping the remainder — more slices overflowed the right-docked
# legend into the pie in the wide/short overview box. Both builders write a
# "Pie Chart Source" block in columns E:G beside the overview block (name /
# `=B{r}` $ amount / `=F/Btotal` fraction — Excel charts literal cells, and the
# in-cell formulas keep the block live when the analyst edits a segment $), and
# the chart series references that block's fraction column.
_PIE_MAX_SLICES = 5
_PIE_OTHER_LABEL = "Other"
_PIE_SRC_TITLE = "Pie Chart Source"
_PIE_SRC_NAME_COL = 5   # column E — slice names (categories / legend)
_PIE_SRC_VALUE_COL = 6  # column F — grouped $ amounts (=B refs; Other = total − top 4)
_PIE_SRC_FRAC_COL = 7   # column G — grouped "% of Total" fractions; the chart series
_PIE_SRC_FRAC_FORMAT = "0.0%"
# Data-label number format for the pie: the series is the source block's fraction
# column, so a value-only "%" format renders e.g. 0.452 as "45.2%".
_PIE_LABEL_FORMAT = '#,##0.0%_);(#,##0.0%);"--"'
# The pie legend runs one point smaller than the 9 pt chart text: five entries
# with long segment names wrap to two lines each, and at 9 pt Excel drops the
# entries that no longer fit the 1.77"-tall overview box (the "Other" entry
# vanished in the live-run render).
_PIE_LEGEND_FONT_SIZE_PT = 8
# Only slices whose share of the total is ABOVE this threshold carry a data
# label — the tiny-slice labels overlap each other in the short overview box.
_PIE_LABEL_MIN_FRACTION = 0.03
# Pie data labels sit INSIDE their slice in white: Excel's Best Fit placed the
# small-slice labels outside the pie (the live-run "8.8% / 3.4% float outside
# their sections" report), and a label pinned inside a dark accent slice needs
# white to read. The clustered-column labels stay black/Outside End.
_PIE_LABEL_COLOR = "FFFFFF"
# Pie plot-area manual layout, as fractions of the chart area: pin the pie to
# the LEFT of the chart box so it sits clear of the right-docked legend (the
# auto-layout centers the pie in this wide/short box, where it can collide with
# the legend).
_PIE_PLOT_X = 0.02
_PIE_PLOT_Y = 0.05
_PIE_PLOT_W = 0.58
_PIE_PLOT_H = 0.90
# Legend manual layout: the full remaining right side of the chart box. Excel's
# auto legend in this wide/short box comes out narrower and shorter than the
# space allows, wraps every entry to two lines and then silently DROPS the
# entries that no longer fit (the "Other" entry vanished in the live-run
# render). Pinning the legend to the full height/width keeps all five entries.
_PIE_LEGEND_X = 0.61
_PIE_LEGEND_Y = 0.02
_PIE_LEGEND_W = 0.38
_PIE_LEGEND_H = 0.96
# The overview slide is FOUND, not indexed. Through v0.5.39 this module and
# `pitch_deck_assembler` shared a `OVERVIEW_SLIDE_INDEX = 6` constant, which was
# only correct because the pitch flow happens to delete a slide after it. The
# slide is now located by the very placeholder this stage is about to fill —
# the same self-discovery the Financial Summary slides already used.
_PIE_PLACEHOLDER = MARKER_BUILT_OVERVIEW.shape_name  # 'Rectangle 4'
# Placeholder box: 4.51" x 1.77" (wide and short) — size the exported pie to the
# same aspect so the stretched picture is not distorted.
_PIE_W_PT = 4.51 * 72
_PIE_H_PT = 1.77 * 72
_PIE_W_CM = 4.51 * 2.54
_PIE_H_CM = 1.77 * 2.54

def period_axis_columns(ws) -> tuple[int, int]:
    """Return the (first, last) 1-based column of the period axis on row 5.

    The period header runs ``B5 .. <col before "Units">``: ``B5:G5`` when the LTM
    column is shown, ``B5:F5`` when it is suppressed. Reading it dynamically lets
    a 5- or 6-column tab chart correctly. ``ws`` is an openpyxl worksheet.
    """
    units_col = None
    for col in range(2, ws.max_column + 1):
        if ws.cell(row=_HEADER_ROW, column=col).value == "Units":
            units_col = col
            break
    if units_col is None or units_col < 3:
        raise ValueError(
            "could not locate the 'Units' header on row 5; the Financial Summary "
            "tab is not in the expected chart-ready layout"
        )
    return 2, units_col - 1


def metric_data_rows(ws) -> list[int]:
    """Return the 1-based metric data rows of the Financial Summary tab.

    One metric per row from row 6 down, contiguous — the block ends at the
    first empty column-A cell. Four rows for the default single-slide deck
    (6-9), eight for the two-slide deck (6-13). ``ws`` is an openpyxl worksheet.
    """
    rows: list[int] = []
    r = _FIRST_DATA_ROW
    while True:
        label = ws.cell(row=r, column=1).value
        if label is None or (isinstance(label, str) and not label.strip()):
            break
        rows.append(r)
        r += 1
    if not rows:
        raise ValueError(
            "no metric rows found below the row-5 header; the Financial Summary "
            "tab is not in the expected chart-ready layout"
        )
    return rows


def _find_label_row_openpyxl(ws, prefixes) -> int | None:
    """Row (1-based) of the first col-A cell whose text starts with any prefix.

    The established way
    to locate a labelled block on these tabs without hardcoding row numbers.
    """
    for row in ws.iter_rows(min_col=1, max_col=1):
        value = row[0].value
        if isinstance(value, str) and any(value.strip().startswith(p) for p in prefixes):
            return row[0].row
    return None


def ltm_revenue_overview_range(ws) -> tuple[int, int] | None:
    """Return the (first, last) 1-based data row of the "LTM Revenue Overview" block.

    The block runs: section title row, header row (section + 1), then one row per
    segment starting at section + 2, ending just above the "Total" row. Segment
    names live in column A and LTM revenue values in column B. The Total row is
    excluded. ``ws`` is an openpyxl worksheet. Returns ``None`` when the block (or
    its Total row) is not found.
    """
    section = _find_label_row_openpyxl(ws, (_PIE_SECTION_LABEL,))
    if section is None:
        return None
    first_data = section + 2  # skip the section title + the header row
    total = _find_label_row_openpyxl_from(ws, (_PIE_TOTAL_LABEL,), start=first_data)
    if total is None or total <= first_data:
        return None
    return first_data, total - 1


def _find_label_row_openpyxl_from(ws, prefixes, *, start: int) -> int | None:
    """Like ``_find_label_row_openpyxl`` but only scanning rows >= ``start``."""
    for row in range(start, ws.max_row + 1):
        value = ws.cell(row=row, column=1).value
        if isinstance(value, str) and any(value.strip().startswith(p) for p in prefixes):
            return row
    return None
def _fractions_from_amounts(amounts: list) -> list:
    """Convert a list of $ amounts to fractions of their total (None-safe).

    Used two ways: the off-Windows PNG render charts literal cells, but the pie's
    "% of Total" column on the real tab is the formula ``=B/Btotal`` that openpyxl
    cannot evaluate, so the throwaway render workbook charts fractions recomputed
    here from the column-B ``$`` literals, and the same recomputed fractions decide
    which slices are above the 3% data-label threshold — deterministic regardless
    of recalc state.
    """
    numeric = [a for a in amounts if isinstance(a, (int, float))]
    total = sum(numeric)
    if not total:
        return [None for _ in amounts]
    return [(a / total if isinstance(a, (int, float)) else None) for a in amounts]


def _suppressed_pie_label_indices(fractions: list) -> list[int]:
    """0-based slice indices whose share is NOT above ``_PIE_LABEL_MIN_FRACTION``.

    Only slices larger than 3% of the total carry a data label — smaller ones
    overlap each other in the short overview box. A non-numeric fraction (``None``
    from an empty/zero total) is suppressed too.
    """
    return [
        i
        for i, f in enumerate(fractions)
        if not (isinstance(f, (int, float)) and f > _PIE_LABEL_MIN_FRACTION)
    ]


def _pie_grouping(amounts: list) -> tuple[list[int], bool]:
    """How the overview pie groups the revenue segments into at most 5 slices.

    Returns ``(kept, has_other)``: ``kept`` is the 0-based offsets (into the
    overview block's segment rows) of the directly-charted segments, largest $
    first; ``has_other`` is True when the remaining segments collapse into a
    trailing "Other" slice (i.e. there are more segments than
    ``_PIE_MAX_SLICES``). Five or fewer segments chart as-is (descending), with
    no "Other". Non-numeric amounts rank last, so they land in "Other".
    """
    order = sorted(
        range(len(amounts)),
        key=lambda i: amounts[i] if isinstance(amounts[i], (int, float)) else float("-inf"),
        reverse=True,
    )
    if len(order) <= _PIE_MAX_SLICES:
        return order, False
    return order[: _PIE_MAX_SLICES - 1], True


def _grouped_pie_labels_amounts(names: list, amounts: list) -> tuple[list, list]:
    """The pie's slice labels and $ amounts after Top-4 + "Other" grouping.

    Mirrors the source block both builders write on the tab — used for the
    off-Windows PNG render (which charts literals) and for the >3% data-label
    threshold on both builders (deterministic regardless of recalc state).
    """
    kept, has_other = _pie_grouping(amounts)
    labels = [names[i] for i in kept]
    grouped = [amounts[i] for i in kept]
    if has_other:
        total = sum(a for a in amounts if isinstance(a, (int, float)))
        kept_total = sum(a for a in grouped if isinstance(a, (int, float)))
        labels.append(_PIE_OTHER_LABEL)
        grouped.append(total - kept_total)
    return labels, grouped


# ---------------------------------------------------------------------------
# Public orchestrator
# ---------------------------------------------------------------------------
def render_financial_summary_charts_into_deck(
    *,
    deck_path: Path | str,
    deal_workbook: Path | str,
    sheet_name: str = _SHEET_DEFAULT,
    slide_index: int | None = None,
    output_path: Path | str | None = None,
) -> Path | None:
    """Build the Financial Summary charts and place them into the deck.

    One chart per metric row on the tab — four for the default single Financial
    Summary slide, eight for the two-slide deck spec. The deck's FS slides are
    discovered by scanning for the Metric #1 chart placeholder (a cloned FS
    slide keeps the same shape names), in deck order: slide *k* receives the
    tab's metric rows ``4k..4k+3``. Pass ``slide_index`` to target a single,
    known slide instead of scanning (legacy call shape).

    Returns the output deck path, or ``None`` when the slide is left with its
    placeholders — either because the deal workbook has no ``financial-summary``
    tab (the financial-summary stage produced nothing) **or** because the native
    charts were persisted to the workbook but their PNGs could not be rendered for
    the deck (LibreOffice unavailable — the graceful-degradation path, Issue 1).
    In both cases the slide keeps its placeholders, like the ownership null path.

    Raises ValueError when the tab's metric-row count does not match the deck's
    FS slide count × 4 — a wireframe/financial-summary stage mismatch.

    Modifies the deck at ``deck_path`` IN PLACE (or writes ``output_path``); it does
    not re-assemble it. Never call this from a flow that also re-runs the
    deck-assembler — see the module docstring.
    """
    deck = Path(deck_path).resolve()
    workbook = Path(deal_workbook).resolve()
    if not deck.exists():
        raise FileNotFoundError(f"deck not found: {deck}")
    if not workbook.exists():
        raise FileNotFoundError(f"deal workbook not found: {workbook}")

    geometry = _resolve_geometry(workbook, sheet_name)
    if geometry is None:
        return None
    first_col, last_col, data_rows = geometry

    pngs = _build_charts_openpyxl_libreoffice(
        workbook, sheet_name, first_col, last_col, data_rows
    )

    if not pngs or any(row not in pngs for row in data_rows):
        # Native charts were persisted to the workbook, but their PNGs could not be
        # rendered (LibreOffice unavailable). Leave the slide placeholders rather
        # than crashing — the durable artefact (workbook charts) is already saved.
        return None

    slide_indexes = (
        [slide_index] if slide_index is not None else _find_financial_summary_slides(deck)
    )
    if len(data_rows) != _CHARTS_PER_SLIDE * len(slide_indexes):
        raise ValueError(
            f"the '{sheet_name}' tab carries {len(data_rows)} metric rows but the deck "
            f"has {len(slide_indexes)} Financial Summary slide(s) — expected "
            f"{_CHARTS_PER_SLIDE * len(slide_indexes)} (four charts per slide)"
        )

    out = Path(output_path).resolve() if output_path is not None else deck
    inserts: dict[int, dict[str, bytes]] = {}
    for k, idx in enumerate(slide_indexes):
        slide_rows = data_rows[_CHARTS_PER_SLIDE * k : _CHARTS_PER_SLIDE * (k + 1)]
        inserts[idx] = {
            name: pngs[row] for name, row in zip(_PLACEHOLDER_NAMES, slide_rows, strict=True)
        }
    return _insert_pngs_into_slides(deck_path=deck, inserts=inserts, output_path=out)


def _resolve_geometry(workbook: Path, sheet_name: str) -> tuple[int, int, list[int]] | None:
    """Read the period-axis columns + metric data rows from the deal
    workbook, or None if the ``financial-summary`` tab is absent."""
    from openpyxl import load_workbook

    wb = load_workbook(workbook, read_only=True)
    try:
        if sheet_name not in wb.sheetnames:
            return None
        ws = wb[sheet_name]
        first_col, last_col = period_axis_columns(ws)
        return first_col, last_col, metric_data_rows(ws)
    finally:
        wb.close()


def _find_financial_summary_slides(deck_path: Path) -> list[int]:
    """Zero-based indices of the deck's Financial Summary slides, in order.

    A Financial Summary slide is identified by its Metric #1 chart placeholder
    shape (``Rectangle 17``), which is present until this stage replaces it with
    the chart picture. Raises when none is found — the deck was already charted
    or is not an assembled pitch deck.

    The scan itself now lives in `template_layout.find_slides_by_marker`, which
    generalised this function in Phase C; the local wrapper keeps the
    deck-specific error message.
    """
    indexes = find_slides_by_marker(Presentation(deck_path), MARKER_BUILT_FINANCIAL_SUMMARY)
    if not indexes:
        raise ValueError(
            f"no Financial Summary chart placeholders ({_FS_MARKER_PLACEHOLDER!r}) found "
            f"in {deck_path.name} — the deck was already charted or is not an assembled "
            f"pitch deck"
        )
    return indexes


# ---------------------------------------------------------------------------
# LTM revenue pie (overview slide) — public orchestrator
# ---------------------------------------------------------------------------
def render_ltm_revenue_pie_into_deck(
    *,
    deck_path: Path | str,
    deal_workbook: Path | str,
    sheet_name: str = _PIE_SHEET_DEFAULT,
    slide_index: int | None = None,
    placeholder_name: str = _PIE_PLACEHOLDER,
    output_path: Path | str | None = None,
) -> Path | None:
    """Build the LTM-revenue-by-segment pie and place it on the overview slide.

    The pie is built on the deal workbook's ``ltm-metrics`` tab (its "LTM
    Revenue Overview" block carries literal segment × LTM-revenue values) and
    dropped into the overview slide's "[Pie Chart Placeholder]" (``Rectangle 4``).

    ``slide_index`` defaults to ``None``, meaning "find the overview slide by
    that placeholder" — the deck tells this stage where its overview slide is,
    rather than the two agreeing on a number.

    Returns the output deck path, or ``None`` when the slide keeps its placeholder
    — either because the deal workbook has no ``ltm-metrics`` tab / no "LTM
    Revenue Overview" block, **or** because the native pie was persisted to the
    workbook but its PNG could not be rendered for the deck (LibreOffice unavailable
    — the graceful-degradation path). Mirrors the FS / ownership null paths.

    Modifies the deck at ``deck_path`` IN PLACE (or writes ``output_path``); it does
    not re-assemble it. Never call this from a flow that also re-runs the
    deck-assembler — see the module docstring.
    """
    deck = Path(deck_path).resolve()
    workbook = Path(deal_workbook).resolve()
    if not deck.exists():
        raise FileNotFoundError(f"deck not found: {deck}")
    if not workbook.exists():
        raise FileNotFoundError(f"deal workbook not found: {workbook}")

    rng = _resolve_pie_range(workbook, sheet_name)
    if rng is None:
        return None
    first_row, last_row = rng

    png = _build_pie_openpyxl_libreoffice(workbook, sheet_name, first_row, last_row)

    if png is None:
        # Native pie persisted to the workbook, but its PNG could not be rendered
        # (LibreOffice unavailable). Leave the overview placeholder rather than
        # crashing — the durable artefact (workbook pie) is already saved.
        return None

    out = Path(output_path).resolve() if output_path is not None else deck
    target = (
        slide_index
        if slide_index is not None
        else find_slide_by_marker(
            Presentation(deck), MARKER_BUILT_OVERVIEW, template=deck.name
        )
    )
    return insert_pngs_into_placeholders(
        deck_path=deck,
        slide_index=target,
        pngs_by_placeholder={placeholder_name: png},
        output_path=out,
    )


def _resolve_pie_range(workbook: Path, sheet_name: str) -> tuple[int, int] | None:
    """Read the "LTM Revenue Overview" data-row span from the deal workbook,
    or None if the ``ltm-metrics`` tab / block is absent."""
    from openpyxl import load_workbook

    wb = load_workbook(workbook, read_only=True)
    try:
        if sheet_name not in wb.sheetnames:
            return None
        return ltm_revenue_overview_range(wb[sheet_name])
    finally:
        wb.close()


# ---------------------------------------------------------------------------
# Placeholder insertion (mirrors excel_to_powerpoint.insert_excel_into_placeholder)
# ---------------------------------------------------------------------------
def _insert_pngs_into_slides(
    *,
    deck_path: Path | str,
    inserts: dict[int, dict[str, bytes]],
    output_path: Path | str | None = None,
) -> Path:
    """Replace named placeholders across slides with pictures, in one open/save.

    ``inserts`` maps zero-based slide index -> {placeholder name -> PNG bytes}.
    Each picture is added at the placeholder's exact left/top/width/height (the
    picture is stretched to the box), and the placeholder shape is removed.
    """
    deck = Path(deck_path).resolve()
    prs = Presentation(deck)
    for slide_index, pngs_by_placeholder in inserts.items():
        slide = prs.slides[slide_index]
        for name, png in pngs_by_placeholder.items():
            placeholder = next((s for s in slide.shapes if s.name == name), None)
            if placeholder is None:
                raise KeyError(f"placeholder {name!r} not found on slide {slide_index + 1}")
            left, top, width, height = (
                placeholder.left,
                placeholder.top,
                placeholder.width,
                placeholder.height,
            )
            placeholder._element.getparent().remove(placeholder._element)
            buf = png if hasattr(png, "read") else io.BytesIO(png)
            slide.shapes.add_picture(buf, left, top, width=width, height=height)
    out = Path(output_path).resolve() if output_path is not None else deck
    prs.save(out)
    return out


def insert_pngs_into_placeholders(
    *,
    deck_path: Path | str,
    slide_index: int,
    pngs_by_placeholder: dict[str, bytes],
    output_path: Path | str | None = None,
) -> Path:
    """Single-slide wrapper over :func:`_insert_pngs_into_slides`."""
    return _insert_pngs_into_slides(
        deck_path=deck_path,
        inserts={slide_index: pngs_by_placeholder},
        output_path=output_path,
    )


def insert_png_into_placeholder(
    *,
    deck_path: Path | str,
    slide_index: int,
    placeholder_name: str,
    png_bytes: bytes,
    output_path: Path | str | None = None,
) -> Path:
    """Single-placeholder convenience wrapper over :func:`insert_pngs_into_placeholders`."""
    return insert_pngs_into_placeholders(
        deck_path=deck_path,
        slide_index=slide_index,
        pngs_by_placeholder={placeholder_name: png_bytes},
        output_path=output_path,
    )
# ---------------------------------------------------------------------------
# openpyxl + LibreOffice backend (off-Windows) — best-effort
# ---------------------------------------------------------------------------
def _fs_chart_anchor(idx: int, last_data_row: int) -> str:
    """On-tab anchor cell for the idx-th persisted FS chart (2-per-row grid).

    The grid starts one blank row below the last metric row (B11/I11 then
    B27/I27 for the default four metrics; B15/I15 down for eight) — 16 rows per
    grid row matches the chart height.
    """
    col = "B" if idx % 2 == 0 else "I"
    row = (last_data_row + 2) + 16 * (idx // 2)
    return f"{col}{row}"


def _persist_native_charts_openpyxl(
    workbook: Path,
    *,
    rebuild: str,
    fs_sheet: str = _SHEET_DEFAULT,
    fs_cols: tuple[int, int] | None = None,
    fs_rows: list[int] | None = None,
    pie_sheet: str = _PIE_SHEET_DEFAULT,
    pie_rows: tuple[int, int] | None = None,
):
    """Load the deal workbook, (re)create native charts, save. Returns the wb.

    The FS-chart and pie steps run as separate orchestrator calls against the
    SAME deal workbook, so each save must neither lose the other step's
    charts nor let a re-run accumulate duplicates:

      - The ``rebuild`` side (``"fs"`` | ``"pie"``) has its sheet's charts
        cleared and re-created unconditionally, so re-running a step replaces
        its charts instead of adding a second set next to the stale one.
        (openpyxl has no public chart-removal API; resetting the private
        per-sheet list is the accepted idiom.) Raises ``KeyError`` before any
        mutation when that side's sheet is missing.
      - The sibling side's charts normally survive untouched: openpyxl 3.x
        reads existing chart parts on load and round-trips them on save. As a
        belt-and-braces for any openpyxl build that drops chart parts on load,
        a sibling tab that is present and chart-ready but carries NO charts
        gets its set re-created (at openpyxl fidelity) rather than silently
        lost on this save.

    Geometry is taken from the explicit ``fs_cols`` / ``pie_rows`` when the
    caller already resolved it, and re-derived from the tab otherwise (the
    sibling side on a re-create).
    """
    from openpyxl import load_workbook

    wb = load_workbook(workbook)
    need_sheet = fs_sheet if rebuild == "fs" else pie_sheet
    if need_sheet not in wb.sheetnames:
        raise KeyError(f"sheet {need_sheet!r} not found in {workbook}")

    if fs_sheet in wb.sheetnames and (rebuild == "fs" or not wb[fs_sheet]._charts):
        ws = wb[fs_sheet]
        cols = fs_cols
        rows = fs_rows
        try:
            if cols is None:
                cols = period_axis_columns(ws)
            if rows is None:
                rows = metric_data_rows(ws)
        except ValueError:
            cols = rows = None  # tab not in the chart-ready layout — skip its charts
        if cols is not None and rows is not None:
            ws._charts = []
            first_col, last_col = cols
            for idx, data_row in enumerate(rows):
                ws.add_chart(
                    _make_openpyxl_chart(ws, data_row, first_col, last_col),
                    _fs_chart_anchor(idx, rows[-1]),
                )

    if pie_sheet in wb.sheetnames and (rebuild == "pie" or not wb[pie_sheet]._charts):
        ws = wb[pie_sheet]
        rows = pie_rows if pie_rows is not None else ltm_revenue_overview_range(ws)
        if rows is not None:
            ws._charts = []
            first_row, last_row = rows
            # _make_openpyxl_pie also (re)writes the Top-4 + "Other" source
            # block the chart references, so the sibling re-create path gets
            # the block too.
            ws.add_chart(
                _make_openpyxl_pie(ws, first_row, last_row),
                ws.cell(row=last_row + 3, column=1).coordinate,
            )

    wb.save(workbook)
    return wb


def _build_charts_openpyxl_libreoffice(
    workbook: Path, sheet_name: str, first_col: int, last_col: int, data_rows: list[int]
) -> dict[int, bytes]:
    """Persist native openpyxl charts on the tab, then render each PNG.

    The deal workbook off-Windows was already built by openpyxl (CapIQ links
    do not survive that path), so loading + re-saving it here costs nothing extra.
    PNGs are rendered from single-chart temp workbooks built off
    LibreOffice-recalculated values so the LTM bar is correct.

    Chart persistence on the tab must **not** depend on LibreOffice being present:
    the native charts are added and the workbook saved **first**, then the PNG
    render is attempted. If ``soffice``/``libreoffice`` (or ``pypdfium2``) is
    missing, the workbook charts have already been saved and an empty ``{}`` is
    returned so the caller degrades gracefully (leaves the deck placeholders)
    rather than aborting the whole stage (Issue 1). Persistence goes through
    :func:`_persist_native_charts_openpyxl` so a re-run replaces the charts
    instead of accumulating a duplicate set, and the pie step's chart on the
    sibling ``ltm-metrics`` tab survives this save.
    """
    _persist_native_charts_openpyxl(
        workbook,
        rebuild="fs",
        fs_sheet=sheet_name,
        fs_cols=(first_col, last_col),
        fs_rows=data_rows,
    )

    try:
        resolved = _libreoffice_recalc_values(workbook, sheet_name, first_col, last_col, data_rows)
        labels = resolved["labels"]
        pngs: dict[int, bytes] = {}
        for data_row in data_rows:
            pngs[data_row] = _render_single_chart_png(labels, resolved[data_row])
        return pngs
    except RuntimeError as exc:
        # LibreOffice / pypdfium2 unavailable for the deck-image step. The native
        # charts are already saved on the tab; degrade gracefully instead of
        # aborting the stage — the caller leaves the slide placeholders.
        print(
            f"[financial-charts] deck-image render skipped ({exc}); the native "
            f"charts are saved on the '{sheet_name}' tab of {workbook.name}.",
            file=sys.stderr,
        )
        return {}


def _make_openpyxl_chart(ws, data_row: int, first_col: int, last_col: int):
    """Build one INFOR-formatted clustered-column BarChart for a metric row."""
    from openpyxl.chart import BarChart, Reference
    from openpyxl.chart.label import DataLabelList

    chart = BarChart()
    chart.type = "col"
    chart.grouping = "clustered"
    chart.gapWidth = _GAP_WIDTH
    chart.title = None
    chart.legend = None
    chart.width = _CHART_W_CM
    chart.height = _CHART_H_CM

    data = Reference(ws, min_col=first_col, max_col=last_col, min_row=data_row, max_row=data_row)
    cats = Reference(ws, min_col=first_col, max_col=last_col, min_row=_HEADER_ROW, max_row=_HEADER_ROW)
    chart.add_data(data, from_rows=True, titles_from_data=False)
    chart.set_categories(cats)

    # All bars filled RGB(70, 86, 110).
    chart.series[0].graphicalProperties.solidFill = _BAR_RGB_HEX

    # Data labels on every bar, Outside End, Palatino 9 black. Show the VALUE
    # only — the other label flags must be turned OFF explicitly, or LibreOffice
    # renders "<category>; <series>; <value>" (e.g. "FY2025; Row 2; 589.808")
    # instead of "589.8" (Issue P1.2).
    labels = DataLabelList()
    labels.showVal = True
    labels.showCatName = False
    labels.showSerName = False
    labels.showLegendKey = False
    labels.showPercent = False
    labels.dLblPos = "outEnd"
    labels.numFmt = _VALUE_FORMAT
    labels.txPr = _palatino_text()
    chart.dataLabels = labels

    # Category axis: Palatino 9 black; value axis + gridlines hidden.
    chart.x_axis.delete = False
    chart.x_axis.txPr = _palatino_text()
    chart.y_axis.delete = True
    chart.y_axis.majorGridlines = None
    _openpyxl_no_border_black_axis(chart)

    return chart


def _openpyxl_no_border_black_axis(chart) -> None:
    """No chart-area border + a visible solid-black category-axis baseline.

    Clears the chart-area outline and sets the category-axis line explicitly
    black **with a visible width**. A width-less ``<a:ln>`` renders as a hairline that LibreOffice drops
    on export, so the bars float with no baseline (Issue 2) — pin the width.
    """
    from openpyxl.chart.shapes import GraphicalProperties
    from openpyxl.drawing.line import LineProperties

    chart.graphical_properties = GraphicalProperties(ln=LineProperties(noFill=True))
    chart.x_axis.spPr = GraphicalProperties(
        ln=LineProperties(w=_AXIS_LINE_WIDTH_EMU, solidFill="000000")
    )


def _palatino_text(size_hundredths: int = _FONT_SIZE_HUNDREDTHS, color: str = "000000"):
    """A RichText carrying Palatino Linotype default run properties (9 pt black
    unless overridden — the pie legend passes 8 pt; the pie data labels pass
    white, since they sit inside the accent-filled slices)."""
    from openpyxl.chart.text import RichText
    from openpyxl.drawing.text import (
        CharacterProperties,
        Font as DrawingFont,
        Paragraph,
        ParagraphProperties,
    )

    cp = CharacterProperties(
        latin=DrawingFont(typeface=_FONT_NAME),
        sz=size_hundredths,
        solidFill=color,
    )
    return RichText(p=[Paragraph(pPr=ParagraphProperties(defRPr=cp), endParaRPr=cp)])


def _libreoffice_recalc_values(
    workbook: Path, sheet_name: str, first_col: int, last_col: int, data_rows: list[int]
) -> dict:
    """Recalc the deal workbook with LibreOffice and read the period labels +
    each metric row's resolved values (so the LTM cell is no longer a formula)."""
    from openpyxl import load_workbook

    # Shared locator + recalc-on-load helper (never a bare shutil.which — see
    # find_soffice: the Windows MSI puts nothing on PATH).
    from excel_to_powerpoint import _soffice_convert, find_soffice

    soffice = find_soffice()
    if soffice is None:
        raise RuntimeError(
            "LibreOffice (soffice/libreoffice) not found on PATH or in the standard "
            "install locations; required for the non-Windows Financial Summary chart "
            "renderer. Install LibreOffice or run the conductor on a Windows machine "
            "with Excel."
        )
    with tempfile.TemporaryDirectory() as tmp_dir:
        _soffice_convert(soffice, workbook, "xlsx:Calc MS Excel 2007 XML", Path(tmp_dir))
        recalced = Path(tmp_dir) / f"{workbook.stem}.xlsx"
        if not recalced.exists():
            raise RuntimeError("LibreOffice produced no recalculated workbook")
        wb = load_workbook(recalced, data_only=True)
        ws = wb[sheet_name]
        out: dict = {
            "labels": [
                ws.cell(row=_HEADER_ROW, column=c).value for c in range(first_col, last_col + 1)
            ]
        }
        for data_row in data_rows:
            out[data_row] = [
                ws.cell(row=data_row, column=c).value for c in range(first_col, last_col + 1)
            ]
        wb.close()
        return out


def _render_single_chart_png(labels: list, values: list) -> bytes:
    """Render one INFOR-formatted chart to PNG from literal values via LibreOffice.

    Builds a throwaway single-chart workbook (data parked above the print area,
    print area set to the chart footprint), converts it to PDF, and trims the
    white margins so only the chart survives.
    """
    from openpyxl import Workbook

    from excel_to_powerpoint import _soffice_convert, _trim_white_margins, find_soffice

    soffice = find_soffice()
    if soffice is None:
        raise RuntimeError("LibreOffice (soffice/libreoffice) not found")
    try:
        import pypdfium2 as pdfium
    except ImportError as exc:
        raise RuntimeError("pypdfium2 is required for the non-Windows chart renderer") from exc

    wb = Workbook()
    ws = wb.active
    n = len(labels)
    for j in range(n):
        ws.cell(row=1, column=2 + j, value=labels[j])
        cell = ws.cell(row=2, column=2 + j, value=values[j])
        # Match the data-label format on the source cells too: a renderer that
        # source-links label formats (numFmt without sourceLinked="0") then still
        # shows "$102.7", not the raw literal.
        cell.number_format = _VALUE_FORMAT
    # Anchor the chart well below the data and print only its footprint, so the
    # data cells in rows 1-2 do not bleed into the rendered image.
    ws.add_chart(_make_single_value_chart(ws, n), "A5")
    ws.print_area = "A5:N28"
    ws.page_setup.fitToWidth = 1
    ws.page_setup.fitToHeight = 1
    ws.sheet_properties.pageSetUpPr.fitToPage = True
    for side in ("left", "right", "top", "bottom"):
        setattr(ws.page_margins, side, 0.1)

    with tempfile.TemporaryDirectory() as tmp_dir:
        src = Path(tmp_dir) / "chart.xlsx"
        wb.save(src)
        _soffice_convert(soffice, src, "pdf", Path(tmp_dir))
        pdf_path = Path(tmp_dir) / "chart.pdf"
        if not pdf_path.exists():
            raise RuntimeError("LibreOffice produced no PDF for the chart")
        pdf = pdfium.PdfDocument(str(pdf_path))
        try:
            pil = pdf[0].render(scale=200 / 72).to_pil().convert("RGB")
            pil = _trim_white_margins(pil)
            buf = io.BytesIO()
            pil.save(buf, format="PNG")
            return buf.getvalue()
        finally:
            pdf.close()


def _make_single_value_chart(ws, n: int):
    """A chart over a literal one-row data block (row 2, header row 1) at B..."""
    from openpyxl.chart import BarChart, Reference
    from openpyxl.chart.label import DataLabelList

    chart = BarChart()
    chart.type = "col"
    chart.grouping = "clustered"
    chart.gapWidth = _GAP_WIDTH
    chart.title = None
    chart.legend = None
    chart.width = _CHART_W_CM
    chart.height = _CHART_H_CM
    data = Reference(ws, min_col=2, max_col=1 + n, min_row=2, max_row=2)
    cats = Reference(ws, min_col=2, max_col=1 + n, min_row=1, max_row=1)
    chart.add_data(data, from_rows=True, titles_from_data=False)
    chart.set_categories(cats)
    chart.series[0].graphicalProperties.solidFill = _BAR_RGB_HEX
    # Value-only data labels (see _make_openpyxl_chart): the other flags must be
    # off or LibreOffice renders "<category>; <series>; <value>" (Issue P1.2).
    labels = DataLabelList()
    labels.showVal = True
    labels.showCatName = False
    labels.showSerName = False
    labels.showLegendKey = False
    labels.showPercent = False
    labels.dLblPos = "outEnd"
    labels.numFmt = _VALUE_FORMAT
    labels.txPr = _palatino_text()
    chart.dataLabels = labels
    chart.x_axis.delete = False
    chart.x_axis.txPr = _palatino_text()
    chart.y_axis.delete = True
    chart.y_axis.majorGridlines = None
    _openpyxl_no_border_black_axis(chart)
    return chart


# ---------------------------------------------------------------------------
# openpyxl + LibreOffice backend for the LTM revenue pie — best-effort
# ---------------------------------------------------------------------------
def _build_pie_openpyxl_libreoffice(
    workbook: Path, sheet_name: str, first_row: int, last_row: int
) -> bytes | None:
    """Persist a native openpyxl pie on the ltm-metrics tab, then render its PNG.

    The native pie's series is the Top-4 + "Other" **source block's fraction
    column** (``=F/Btotal`` formulas, written beside the overview block).
    openpyxl cannot evaluate those formulas, so the throwaway render workbook
    charts the equivalent grouped fractions recomputed in Python from the
    column-B ``$`` literals — the slices and value labels then match the native
    pie without a LibreOffice recalc of the source tab.

    Like the FS chart path, the native pie is added and the workbook saved **first**
    so its persistence does not depend on LibreOffice; the PNG render is attempted
    after. If ``soffice``/``libreoffice`` (or ``pypdfium2``) is missing, the pie has
    already been saved on the tab and ``None`` is returned so the caller leaves the
    overview placeholder instead of aborting the stage (Issue 3 parity with Issue 1).
    Persistence goes through :func:`_persist_native_charts_openpyxl` so a re-run
    replaces the pie instead of accumulating a duplicate, and the FS step's four
    charts on the sibling ``financial-summary`` tab survive this save.
    """
    wb = _persist_native_charts_openpyxl(
        workbook, rebuild="pie", pie_sheet=sheet_name, pie_rows=(first_row, last_row)
    )
    ws = wb[sheet_name]
    names = [ws.cell(row=r, column=_PIE_SEGMENT_COL).value for r in range(first_row, last_row + 1)]
    amounts = [ws.cell(row=r, column=_PIE_VALUE_COL).value for r in range(first_row, last_row + 1)]
    labels, grouped = _grouped_pie_labels_amounts(names, amounts)

    try:
        # The render workbook charts literal cells, but the real tab's source-block
        # fractions are formulas openpyxl can't evaluate — recompute the grouped
        # Top-4 + "Other" fractions from the column-B $ literals so the rendered
        # slices/labels match the native pie's source-block series.
        return _render_single_pie_png(labels, _fractions_from_amounts(grouped))
    except RuntimeError as exc:
        print(
            f"[financial-charts] overview pie render skipped ({exc}); the native "
            f"pie is saved on the '{sheet_name}' tab of {workbook.name}.",
            file=sys.stderr,
        )
        return None


def _style_openpyxl_pie(chart, n_points: int, fractions: list | None = None) -> None:
    """INFOR pie formatting: no title, accent slice fills, right-docked legend with
    the pie's plot area pinned to the left of the chart box (clear of the legend),
    and value-only "%" labels (the series is the source block's "% of Total" fraction, so
    a value label with a "%" number format reads e.g. "45.2%") — white, pinned
    Inside End so no label floats outside its slice — on the slices whose
    share is above the 3% threshold — ``fractions`` decides which; smaller slices
    get a per-point all-show-flags-off override (openpyxl 3.1 does not model
    CT_DLbl's ``<c:delete>``, and unlike that shape this one survives the
    load→save round-trip in :func:`_persist_native_charts_openpyxl`). The label
    config lives on the SERIES-level ``dLbls`` — where Excel itself writes it and
    the only level where per-point overrides are honored."""
    from openpyxl.chart.label import DataLabel, DataLabelList
    from openpyxl.chart.layout import Layout, ManualLayout
    from openpyxl.chart.series import DataPoint
    from openpyxl.chart.shapes import GraphicalProperties
    from openpyxl.drawing.line import LineProperties

    chart.title = None
    chart.width = _PIE_W_CM
    chart.height = _PIE_H_CM
    # No chart-area border — the openpyxl mirror of _com_strip_chart_border (the
    # default chart outline otherwise frames the pie picture on the slide).
    chart.graphical_properties = GraphicalProperties(ln=LineProperties(noFill=True))
    chart.series[0].data_points = [
        DataPoint(idx=i, spPr=GraphicalProperties(solidFill=INFOR_ACCENTS[i % len(INFOR_ACCENTS)]))
        for i in range(n_points)
    ]
    # Legend docked on the RIGHT (segment names), Palatino 8 (one point under
    # the chart text), pinned to the full remaining right side of the chart box
    # (see _PIE_LEGEND_*); the plot area is pinned to the left via a manual
    # layout so the pie cannot overlap the legend.
    chart.legend.position = "r"
    chart.legend.overlay = False
    chart.legend.txPr = _palatino_text(_PIE_LEGEND_FONT_SIZE_PT * 100)
    chart.legend.layout = Layout(
        manualLayout=ManualLayout(
            xMode="edge",
            yMode="edge",
            x=_PIE_LEGEND_X,
            y=_PIE_LEGEND_Y,
            w=_PIE_LEGEND_W,
            h=_PIE_LEGEND_H,
        )
    )
    chart.layout = Layout(
        manualLayout=ManualLayout(
            layoutTarget="inner",
            xMode="edge",
            yMode="edge",
            x=_PIE_PLOT_X,
            y=_PIE_PLOT_Y,
            w=_PIE_PLOT_W,
            h=_PIE_PLOT_H,
        )
    )
    labels = DataLabelList()
    labels.showVal = True
    labels.showPercent = False
    labels.showCatName = False
    labels.showSerName = False
    labels.showLegendKey = False
    labels.numFmt = _PIE_LABEL_FORMAT
    # Inside End in white — Best Fit (the pie default) floats small-slice
    # labels outside the pie, and inside a dark accent slice black is unreadable.
    labels.dLblPos = "inEnd"
    labels.txPr = _palatino_text(color=_PIE_LABEL_COLOR)
    if fractions is not None:
        labels.dLbl = [
            DataLabel(
                idx=i,
                showVal=False,
                showPercent=False,
                showCatName=False,
                showSerName=False,
                showLegendKey=False,
            )
            for i in _suppressed_pie_label_indices(fractions)
        ]
    chart.series[0].dLbls = labels


def _write_pie_source_block_openpyxl(ws, first_row: int, last_row: int) -> tuple[int, int, list]:
    """Write the pie's Top-4 + "Other" source block in columns E:G, in place.

    Sits beside the "LTM Revenue Overview" block, row-aligned with it (title on
    the section-title row, headers on the header row, slices from the first
    segment row). The kept slices reference their segment rows live (``=A{r}`` /
    ``=B{r}``) and "Other" is ``=B{total} - SUM(top slices)``, so the block —
    and the chart on it — follows analyst edits to the segment $ after an Excel
    recalc. The block's footprint is cleared first so a re-run (or a shrunk
    segment list) never leaves stale slice rows behind.

    Returns ``(src_first_row, src_last_row, grouped_amounts)`` — the chart range
    and the Python-side grouped $ amounts (for the >3% label threshold).
    """
    from openpyxl.styles import Font
    from openpyxl.utils import get_column_letter

    names = [ws.cell(row=r, column=_PIE_SEGMENT_COL).value for r in range(first_row, last_row + 1)]
    amounts = [ws.cell(row=r, column=_PIE_VALUE_COL).value for r in range(first_row, last_row + 1)]
    kept, has_other = _pie_grouping(amounts)
    total_row = last_row + 1  # the overview block's "Total" row (excluded from the range)

    body_font = Font(name=_FONT_NAME, size=11)
    title_font = Font(name=_FONT_NAME, size=11, bold=True, color="1F3864")
    header_font = Font(name=_FONT_NAME, size=11, bold=True)
    value_format = ws.cell(row=first_row, column=_PIE_VALUE_COL).number_format

    for r in range(first_row - 2, total_row + 1):
        for c in (_PIE_SRC_NAME_COL, _PIE_SRC_VALUE_COL, _PIE_SRC_FRAC_COL):
            ws.cell(row=r, column=c).value = None

    title = _PIE_SRC_TITLE + (" (Top 4 + Other)" if has_other else "")
    ws.cell(row=first_row - 2, column=_PIE_SRC_NAME_COL, value=title).font = title_font
    for col, text in (
        (_PIE_SRC_NAME_COL, "Segment"),
        (_PIE_SRC_VALUE_COL, ws.cell(row=first_row - 1, column=_PIE_VALUE_COL).value),
        (_PIE_SRC_FRAC_COL, "% of Total"),
    ):
        ws.cell(row=first_row - 1, column=col, value=text).font = header_font

    seg_col = get_column_letter(_PIE_SEGMENT_COL)
    val_col = get_column_letter(_PIE_VALUE_COL)
    src_val_col = get_column_letter(_PIE_SRC_VALUE_COL)
    rows: list[tuple[int, str, str]] = [
        (first_row + j, f"={seg_col}{first_row + i}", f"={val_col}{first_row + i}")
        for j, i in enumerate(kept)
    ]
    if has_other:
        r = first_row + len(kept)
        rows.append(
            (r, _PIE_OTHER_LABEL, f"={val_col}{total_row}-SUM({src_val_col}{first_row}:{src_val_col}{r - 1})")
        )
    for r, name, value in rows:
        ws.cell(row=r, column=_PIE_SRC_NAME_COL, value=name).font = body_font
        v = ws.cell(row=r, column=_PIE_SRC_VALUE_COL, value=value)
        v.font = body_font
        v.number_format = value_format
        p = ws.cell(
            row=r,
            column=_PIE_SRC_FRAC_COL,
            value=f"={src_val_col}{r}/${val_col}${total_row}",
        )
        p.font = body_font
        p.number_format = _PIE_SRC_FRAC_FORMAT

    # Mirror the source columns' widths so the block reads cleanly on the tab.
    for src_col, dst_col in (
        (_PIE_SEGMENT_COL, _PIE_SRC_NAME_COL),
        (_PIE_VALUE_COL, _PIE_SRC_VALUE_COL),
        (3, _PIE_SRC_FRAC_COL),
    ):
        width = ws.column_dimensions[get_column_letter(src_col)].width
        if width:
            ws.column_dimensions[get_column_letter(dst_col)].width = width

    return first_row, rows[-1][0], _grouped_pie_labels_amounts(names, amounts)[1]


def _make_openpyxl_pie(ws, first_row: int, last_row: int):
    """Build an INFOR-formatted PieChart over the Top-4 + "Other" source block.

    Writes/refreshes the source block (columns E:G) beside the "LTM Revenue
    Overview" block, then charts its **fraction column** (``=F/Btotal``) so the
    data labels read the slice share; Excel evaluates the formulas. Categories
    (legend) = the block's name column. openpyxl cannot evaluate the fraction
    formulas, so the 3% data-label threshold is decided from the grouped
    amounts recomputed in Python off the column-B ``$`` literals — the same
    shares Excel computes into the block.
    """
    from openpyxl.chart import PieChart, Reference

    src_first, src_last, grouped = _write_pie_source_block_openpyxl(ws, first_row, last_row)
    chart = PieChart()
    data = Reference(
        ws, min_col=_PIE_SRC_FRAC_COL, max_col=_PIE_SRC_FRAC_COL, min_row=src_first, max_row=src_last
    )
    cats = Reference(
        ws, min_col=_PIE_SRC_NAME_COL, max_col=_PIE_SRC_NAME_COL, min_row=src_first, max_row=src_last
    )
    chart.add_data(data, titles_from_data=False)
    chart.set_categories(cats)
    _style_openpyxl_pie(
        chart, src_last - src_first + 1, fractions=_fractions_from_amounts(grouped)
    )
    return chart


def _make_single_pie_chart(ws, n: int):
    """A pie over a literal block: categories A1:A{n}, values B1:B{n}.

    The literal values ARE the fractions (recomputed by the caller), so they
    decide the 3% data-label threshold directly.
    """
    from openpyxl.chart import PieChart, Reference

    chart = PieChart()
    data = Reference(ws, min_col=2, max_col=2, min_row=1, max_row=n)
    cats = Reference(ws, min_col=1, max_col=1, min_row=1, max_row=n)
    chart.add_data(data, titles_from_data=False)
    chart.set_categories(cats)
    fractions = [ws.cell(row=r, column=2).value for r in range(1, n + 1)]
    _style_openpyxl_pie(chart, n, fractions=fractions)
    return chart


def _render_single_pie_png(labels: list, values: list) -> bytes:
    """Render the INFOR pie to PNG from literal labels/values via LibreOffice.

    Mirrors :func:`_render_single_chart_png`: data parked above the print area, the
    print area set to the chart footprint, converted to PDF, white margins trimmed.
    """
    from openpyxl import Workbook

    from excel_to_powerpoint import _soffice_convert, _trim_white_margins, find_soffice

    soffice = find_soffice()
    if soffice is None:
        raise RuntimeError("LibreOffice (soffice/libreoffice) not found")
    try:
        import pypdfium2 as pdfium
    except ImportError as exc:
        raise RuntimeError("pypdfium2 is required for the non-Windows chart renderer") from exc

    wb = Workbook()
    ws = wb.active
    n = len(labels)
    for i in range(n):
        ws.cell(row=1 + i, column=1, value=labels[i])
        cell = ws.cell(row=1 + i, column=2, value=values[i])
        # Match the "%" label format on the source cells (see the single-chart
        # renderer): a source-linking renderer then still shows "60.0%".
        cell.number_format = _PIE_LABEL_FORMAT
    # Anchor the chart away from the data and print only its footprint.
    ws.add_chart(_make_single_pie_chart(ws, n), "D2")
    ws.print_area = "D2:N28"
    ws.page_setup.fitToWidth = 1
    ws.page_setup.fitToHeight = 1
    ws.sheet_properties.pageSetUpPr.fitToPage = True
    for side in ("left", "right", "top", "bottom"):
        setattr(ws.page_margins, side, 0.1)

    with tempfile.TemporaryDirectory() as tmp_dir:
        src = Path(tmp_dir) / "pie.xlsx"
        wb.save(src)
        _soffice_convert(soffice, src, "pdf", Path(tmp_dir))
        pdf_path = Path(tmp_dir) / "pie.pdf"
        if not pdf_path.exists():
            raise RuntimeError("LibreOffice produced no PDF for the pie chart")
        pdf = pdfium.PdfDocument(str(pdf_path))
        try:
            pil = pdf[0].render(scale=200 / 72).to_pil().convert("RGB")
            pil = _trim_white_margins(pil)
            buf = io.BytesIO()
            pil.save(buf, format="PNG")
            return buf.getvalue()
        finally:
            pdf.close()
