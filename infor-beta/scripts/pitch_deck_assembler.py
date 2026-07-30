"""Assembler for the INFOR slide-library pitch deck.

The blank library is 16 slides (incl. the insider-ownership slide, which
follows the Financial Summary slide, and the precedent-transactions slide,
which follows the comparable-companies slide). The deck's slide mix is
configurable, so slide indices are computed from the SlidePlan rather than
hardcoded:

- the market-entry section grows across multiple slides — two targets per
  slide — by cloning the library's market-entry slide (8 targets → 4
  market-entry slides → 19-slide deck);
- the Financial Summary section grows the same way — four metrics per slide —
  when the SlidePlan carries more than one ``financial-summary`` entry (the
  deck spec's "2 slides / 8 metrics" option);
- the Key Investment Highlights slide is deleted when the SlidePlan omits its
  entry (the deck spec's "omit" option).

Every slide *after* the Financial Summary section shifts with the section's
size; disclaimer/contact remain the last two slides.
"""

from __future__ import annotations

import math
from dataclasses import dataclass
from functools import partial
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu, Inches

from deal_workbook import TAB_CAPTABLE, TAB_OWNERSHIP
from deck_repair import assert_converged, converge_deck
from excel_to_powerpoint import (
    assert_range_pictures_are_distinct,
    insert_excel_into_placeholder,
)
from naming import safe_filename
from pptx_helpers import (
    clone_slide_after,
    delete_shape,
    delete_slide,
    fill_footnote_currency,
    find_shape,
    find_table_shape,
    fit_overview_textbox,
    iter_all_shapes,
    set_cell_text,
    set_table_height,
    set_text,
    shapes_in_reading_order,
    write_bullets_or_plain,
)
from schemas import PitchDeckContent, SlidePlan
from template_layout import (
    CAP_TABLE_OUTPUT_CCY_CELL,
    CAP_TABLE_OUTPUT_CCY_NAMES,
    CAP_TABLE_PICTURE_NAMES,
    CAP_TABLE_PICTURE_RANGE,
    MARKER_COMPS,
    MARKER_CONTACT,
    MARKER_COVER,
    MARKER_EARNINGS_SUMMARY,
    MARKER_FINANCIAL_SUMMARY,
    MARKER_KIH,
    MARKER_MARKET_ENTRY,
    MARKER_OVERVIEW,
    MARKER_OWNERSHIP,
    MARKER_PRECEDENTS,
    MARKER_RISKS,
    NAME_CAP_OUTPUT_CCY,
    NAME_CAP_PICTURE_RANGE,
    NAME_FX_RATE,
    NAME_OWN_INSIDERS_PICTURE,
    NAME_OWN_INSTITUTIONS_PICTURE,
    OWNERSHIP_INSIDERS_PICTURE_NAMES,
    OWNERSHIP_INSIDERS_PICTURE_RANGE,
    OWNERSHIP_INSTITUTIONS_PICTURE_NAMES,
    OWNERSHIP_INSTITUTIONS_PICTURE_RANGE,
    TemplateLayoutError,
    find_optional_slide_by_marker,
    find_slide_by_marker,
    find_slides_by_marker,
    resolve_workbook_range,
    verify_workbook_names,
)

# Financial Summary slide title shape + the four metric-label tiles it carries
# (each cloned FS slide has the same shape names).
#
# ORDERING: the tiles are addressed BY NAME, and this list is in the reading
# order the `financial-summary` stage emits its labels in — Rectangle 13 is the
# top-left tile, 12 the top-right, 15 the bottom-left, 14 the bottom-right. The
# names are the contract (Phase C: never a position), and because they read out
# of order alphabetically, `test_the_financial_summary_tiles_are_named_in_reading_order`
# checks the library's geometry still agrees with this list. That test is the
# whole safeguard: nothing else here would notice the mapping going stale.
_FS_TITLE_SHAPE = "Title 6"
_FS_METRIC_TILES = ["Rectangle 13", "Rectangle 12", "Rectangle 15", "Rectangle 14"]
_FS_TILES_PER_SLIDE = len(_FS_METRIC_TILES)

# The footnote shape that carries each slide's source/note lines. Every one of
# these slides is populated with figures, so every one is labelled with a
# currency — see `_label_slide_currency`. Two figure-bearing slides are
# deliberately absent: the ownership slide states share counts and percentages
# (its "Source: Bloomberg" line is complete as shipped), and the considerations /
# mitigants slide is prose.
_FOOTNOTE_OVERVIEW = "Text Placeholder 1"
_FOOTNOTE_FINANCIAL_SUMMARY = "Text Placeholder 11"
_FOOTNOTE_COMPS = "Text Placeholder 1"
_FOOTNOTE_PRECEDENTS = "Text Placeholder 1"
_FOOTNOTE_KIH = "Text Placeholder 13"
_FOOTNOTE_MARKET_ENTRY = "Text Placeholder 3"

# Contact slide. The library ships a 2x2 grid of banker cards — three rows of one
# column each: "<name>, <title>", phone, email — plus a one-row address table it
# must not be confused with. Cards are identified by that shape, never by name or
# index, and filled in reading order.
_CONTACT_CARD_ROWS = 3
_CONTACT_CARD_COLS = 1
_CONTACT_NAME_SIZE = 12   # the library's own sizes on its filled card
_CONTACT_DETAIL_SIZE = 10

# Every workbook this module reads is the DEAL workbook, so its sheets are
# addressed by `deal_workbook.TAB_*` — never by a `template_layout`
# `*_SOURCE_SHEET`, which is the corresponding sheet name in the *source*
# template (the cap table's is `Cap with Links` there, `captable` here). Reaching
# for the source constant is the v0.5.45 bug; see the sheet-name note in
# `template_layout`.

# Overview-slide cap-table placeholder; the picture covers the capitalization
# summary plus the Financial/Valuation metric rows (same range as the earnings
# overview). The range is resolved from the workbook's
# `infor_cap_picture_range` defined name — the constant is only the fallback
# for a cap table built before the templates were named.
_CAP_TABLE_PLACEHOLDER = "Rectangle 3"

# Insider-ownership slide. The left "Insiders" placeholder ('Rectangle 1') is
# replaced by a picture of the ownership workbook's Select-Insiders block. The
# right "Institutions" placeholder ('Rectangle 3') is replaced by the
# Select-Institutions block (top-12 + subtotal + Other Shareholders + Total)
# when the ownership stage ingested a Bloomberg export (Bloomberg Output C14
# populated); otherwise it stays a Bloomberg placeholder. Both picture ranges
# are resolved by defined name, like the cap table's.
_OWNERSHIP_PLACEHOLDER = "Rectangle 1"
_INSTITUTIONS_PLACEHOLDER = "Rectangle 3"


def _bullet_tuple(bullet) -> tuple[str, int]:
    return (bullet.text, bullet.level)


# The library's own "an analyst fills this in" token. Any of these left on a
# delivered deck is a build failure — see `_unfilled_placeholders`.
_PLACEHOLDER_TOKEN = "[x]"


def _iter_slide_text(slide):
    """Yield ``(shape_name, text)`` for every text-bearing thing on `slide`.

    Recurses into groups and into table cells. Both matter: the Key Investment
    Highlights bodies are nested two groups deep and the Contact cards are table
    cells, so a top-level-only walk — which this was — could not see either of
    the two slides that shipped the worst placeholder text.
    """
    for shape in iter_all_shapes(slide.shapes):
        if getattr(shape, "has_text_frame", False):
            yield shape.name, shape.text
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    yield shape.name, cell.text


def _all_text(prs: Presentation) -> str:
    return "\n".join(text for slide in prs.slides for _, text in _iter_slide_text(slide))


def _unfilled_placeholders(prs: Presentation) -> list[str]:
    """Every library `[x]` the fill left behind, named by slide and shape.

    Swept across EVERY slide rather than the handful known to be at risk. The
    three defects this closes were found by reading a delivered deck, not by a
    check: two market-entry slides carried an unsubstituted `[x]$MM` footnote,
    the ownership slide a bare `[x]` where a takeaway belonged, and the Contact
    slide three cards of `[x]`. None of them is a rendering problem, so
    `deck_contract` had nothing to say about them, and each was a review finding
    on a deck that had already gone out.

    The deck's *deferred* placeholders — `[Pie Chart Placeholder]`,
    `[Placeholder for Comps Chart]`, `[Placeholder for <target> Logo]` — are a
    different thing and deliberately ship: they name work the analyst finishes in
    PowerPoint. They carry no `[x]`, which is what separates the two.
    """
    return [
        f"slide {index}: {name} — {text.strip()!r}"
        for index, slide in enumerate(prs.slides, start=1)
        for name, text in _iter_slide_text(slide)
        if _PLACEHOLDER_TOKEN in text
    ]


def _shape_text(shape) -> str:
    return shape.text if getattr(shape, "has_text_frame", False) else ""


def _replace_first_line(shape, first_line: str, remaining_lines: list[str] | None = None) -> None:
    lines = [first_line]
    if remaining_lines:
        lines.extend(remaining_lines)
    set_text(shape, lines)


def _section_boxes(slide):
    """The divider's four label boxes, TOP TO BOTTOM.

    ORDERING: geometric, and it has to be. All four ship as "Rounded Rectangle
    19", so a name lookup cannot tell them apart, and `slide.shapes` hands them
    back in document order — which on this library is the 1st, 3rd, 2nd and 4th
    box down the slide. Zipping `section_labels` onto that put "Valuation" in the
    second box and "Financial Summary" in the third on every deck INFOR shipped;
    the labels were right and the wireframe was right, and the reader still saw
    the deck's contents in the wrong order.
    """
    return shapes_in_reading_order(
        shape for shape in slide.shapes if shape.name == "Rounded Rectangle 19"
    )


def _write_flexible_bullets(shape, bullets) -> None:
    """Write bullets, falling back to plain paragraphs if a library shape has no bullet glyph template."""
    write_bullets_or_plain(shape, [_bullet_tuple(bullet) for bullet in bullets])


def _fill_investment_highlights(slide, content: PitchDeckContent) -> None:
    """Fill the four numbered highlight quadrants, blanking any the content skips.

    ORDERING: by the number the library PRINTS in each quadrant's oval, which is
    the strongest ordering available — it is what the reader sees, so it cannot
    disagree with the render the way a document-order index can. The four groups
    arrive in document order 1, 4, 3, 2, so the sort is doing real work. The
    library's numbering also runs top-to-bottom, which
    `test_the_highlight_quadrants_are_numbered_in_reading_order` pins, so the sort
    agrees with geometry rather than merely being self-consistent.

    A quadrant the content has no highlight for is BLANKED, not left alone: the
    library ships it as "[x]", and a deck must never deliver that.
    """
    quadrants: list[tuple[int, object, object]] = []
    for group in slide.shapes:
        if group.shape_type != MSO_SHAPE_TYPE.GROUP or not group.name.startswith("Group"):
            continue
        number = header_shape = body_shape = None
        for sub in iter_all_shapes(group.shapes):
            if sub.name.startswith("Oval") and getattr(sub, "has_text_frame", False) and sub.text.strip().isdigit():
                number = int(sub.text.strip())
            elif sub.name.startswith("Arrow: Pentagon"):
                header_shape = sub
            elif sub.name.startswith("Rectangle") and getattr(sub, "has_text_frame", False) and "[x]" in sub.text:
                body_shape = sub
        if number is not None and header_shape is not None and body_shape is not None:
            quadrants.append((number, header_shape, body_shape))
    quadrants.sort(key=lambda q: q[0])
    for idx, (_, header_shape, body_shape) in enumerate(quadrants):
        if idx < len(content.investment_highlights):
            highlight = content.investment_highlights[idx]
            set_text(header_shape, [highlight.header])
            set_text(body_shape, list(highlight.bullets))
        else:
            set_text(header_shape, [""])
            set_text(body_shape, [""])
    # The tagline box ships as "[x]"; with no tagline it is blanked, not skipped.
    for shape in slide.shapes:
        if shape.name == "Text Placeholder 2" and getattr(shape, "has_text_frame", False):
            set_text(shape, [content.investment_highlights_tagline or ""])


# Market-entry cell sizing: the label column is white at `_ME_LABEL_SIZE`, the
# target value columns 9 pt. The value cells carry the long Overview / Strategic
# Rationale copy, and the layout engine grows a table row to fit its text (a stored
# row height is only a MINIMUM). At the library's 10 pt the real per-target copy
# wrapped tall enough to re-expand the whole table to ~6.3" on open — past the
# 5.71" clamp `set_table_height` writes. 9 pt keeps that copy inside the clamped
# rows; pair it with concise Overview / Strategic Rationale cells (see
# pitch-content) for headroom.
_ME_VALUE_SIZE = 9
_ME_LABEL_COLOR = "FFFFFF"  # scheme bg1 (white) in the library

# Target total height for a filled market-entry (acquisition-target) table.
# Real content can grow rows past the library's ~5.7" so the table runs off the
# slide; after the cells are filled we clamp the table back to this height
# (mirrors dragging the table's resize handle in PowerPoint). 5.71" keeps the
# table's bottom above the slide edge (table top is 1.2" on a 7.5" slide).
_ME_TABLE_HEIGHT = Inches(5.71)

# Row-label size in the market-entry table's narrow (1.457" usable) label column.
# A label too wide for the column wraps and re-grows its row, which is one of the
# two mechanisms behind PRL17's 5.91"-rendered table. Until v0.5.39 the assembler
# pre-empted that by measuring each label against a Palatino advance-width table
# and stepping the over-wide ones down individually; the converge loop now measures
# the rendered table and caps the label column when it actually overruns, which
# also keeps the column uniform. Pair with concise labels from pitch-content
# (≤ ~18 chars) so the cap is rarely needed.
_ME_LABEL_SIZE = 11


# Slide 10 Considerations/Mitigants table sizing. The library ships the header
# row at 12 pt and the body cells at 10 pt; the old code hardcoded 9 pt / 8 pt,
# which rendered noticeably smaller than the template.
_RISK_HEADER_SIZE = 12
_RISK_BODY_SIZE = 10


def _fill_risk_table(slide, content: PitchDeckContent) -> None:
    """Fill the Considerations/Mitigants table and clamp it to the library height.

    Writes the cells at the library's own sizes — 12 pt header, 10 pt body — and
    clamps the declared rows back to the height the library ships the table at.

    It does **not** pre-shrink the body font. Until v0.5.39 it estimated every
    row's rendered height from its widest cell and stepped the body down 10 -> 9 ->
    8 pt until the estimates fit, which is how PRL18 still shipped a 5.36" table
    against a 5.1715" library height: the estimate said it fit. The converge loop
    now measures the rendered table and steps the font down only when the render
    actually overruns — and only as far as it has to, so a deck whose copy fits at
    10 pt keeps the template's size instead of being shrunk by a cautious model.
    """
    table_frame = find_table_shape(slide)
    table = table_frame.table
    library_height = table_frame.height  # the library ships this table at 5.18"

    header = ("Considerations", "Mitigants")
    max_rows = min(len(content.risk_mitigants), len(table.rows) - 1)
    body = [
        (row.risk, "\n".join(row.mitigants)) for row in content.risk_mitigants[:max_rows]
    ]
    body += [("", "")] * (len(table.rows) - 1 - len(body))

    set_cell_text(table.cell(0, 0), header[0], size_pt=_RISK_HEADER_SIZE)
    set_cell_text(table.cell(0, 1), header[1], size_pt=_RISK_HEADER_SIZE)
    for idx, (risk, mitigants) in enumerate(body):
        set_cell_text(table.cell(idx + 1, 0), risk, size_pt=_RISK_BODY_SIZE)
        set_cell_text(table.cell(idx + 1, 1), mitigants, size_pt=_RISK_BODY_SIZE)
    set_table_height(table_frame, library_height)


# ─── Currency labelling ──────────────────────────────────────────────────────
#
# A slide's currency label is a property of WHAT POPULATES THE SLIDE, and there
# is more than one answer per deck. Three sources fill this deck:
#
#   - the content bundle's authored copy, in `content.figure_currency` — the
#     target's filing reporting currency, which the `financial-summary` and
#     `ltm-metrics` tabs are locked to as well;
#   - the cap table, in its own OUTPUT currency, pasted onto the overview slide;
#   - deferred chart placeholders, which carry no figures yet.
#
# Until v0.5.51 one letter was read off the cap table and applied to the overview
# slide, the highlights slide and the market-entry slides regardless of what those
# slides said — while the comps, precedents and Financial Summary slides were
# reached by no call at all. On the Project Open Text deck that shipped four
# distinct failures in one file: two slides carrying the literal `[x]$MM`, one
# carrying no note, five labelled C$ over US$ content, and — the expensive one —
# LTM revenue stated as C$7,344.2 on the overview slide and $5,207.9 on the
# Financial Summary slide, the same figure 1.4102x apart, with nothing on the deck
# to reconcile them. Hence: label from the content, and where a slide is populated
# from two currencies, say both and print the rate between them.


@dataclass(frozen=True)
class _SecondCurrency:
    """A second currency on one slide: what it populates, its code, and the rate."""

    what: str                 # 'capitalization table'
    code: str                 # ISO code of that element
    fx_rate: float | None     # multiplier from the slide's base currency, if known


def _currency_letter(code: str) -> str:
    """Map an ISO currency code to the footnote's ``[x]$MM`` token replacement.

    USD / CAD map explicitly to the ``'US'`` / ``'C'`` dollar letters (so the
    token resolves to ``US$MM`` / ``C$MM``); **any other code is returned as-is**
    and `fill_footnote_currency` renders the ISO code without a dollar sign
    (``GBP MM``) — a non-dollar filer is never silently mislabelled as C$.
    """
    code = code.strip().upper()
    if code in {"USD", "US$", "US"}:
        return "US"
    if code in {"CAD", "C$", "C"}:
        return "C"
    return code


def _currency_phrase(code: str) -> str:
    """How a currency reads in a note: ``'US$MM'``, ``'C$MM'``, ``'GBP MM'``.

    The same rendering `fill_footnote_currency` produces from the library's
    token, so a note this module writes from scratch (the Financial Summary
    slide, which the library ships with no note at all) is indistinguishable
    from one the library templated.
    """
    letter = _currency_letter(code)
    return f"{letter}$MM" if letter in ("US", "C") else f"{letter} MM"


def _currency_unit(code: str) -> str:
    """One unit of `code` as it reads in an FX quote: ``'US$1.00'`` / ``'GBP 1.00'``."""
    letter = _currency_letter(code)
    return f"{letter}$1.00" if letter in ("US", "C") else f"{letter} 1.00"


def _currency_amount(code: str, value: float) -> str:
    """`value` of `code` as it reads in an FX quote: ``'C$1.4102'`` / ``'GBP 1.4102'``."""
    letter = _currency_letter(code)
    return f"{letter}${value:,.4f}" if letter in ("US", "C") else f"{letter} {value:,.4f}"


def _read_output_currency(workbook_path) -> str:
    """The deal workbook's cap-table OUTPUT currency, as an ISO code.

    Reads the ``captable`` tab's output-currency cell, located by its
    ``infor_cap_output_ccy`` defined name (``F5`` as shipped). An empty or
    non-text cell falls back to ``CAD`` (the template default), so a footnote
    never ships the literal ``[x]``.

    A **missing tab** is a hard failure, not a fallback. Until v0.5.45 this read
    ``wb[CAP_TABLE_SHEET] if … in wb.sheetnames else wb.active`` — and since
    Phase D renamed the sheet, that condition was always false, so every pitch
    build derived the client-facing footnote currency from whichever tab happened
    to be active. Silently labelling a US filer's figures ``C$MM`` is worse than
    failing the deck stage, so there is no sheet fallback here at all.
    """
    from openpyxl import load_workbook

    from template_layout import defined_name_ref

    wb = load_workbook(workbook_path, data_only=True)
    try:
        if TAB_CAPTABLE not in wb.sheetnames:
            raise TemplateLayoutError(
                f"{Path(workbook_path).name}: cannot read the footnote output "
                f"currency — expected the {TAB_CAPTABLE!r} tab (have {wb.sheetnames}). "
                f"The deal workbook is created at deal-init from "
                f"INFOR Deal Workbook Template.xlsx; note the source template names "
                f"this sheet 'Cap with Links' and the deal workbook renames it."
            )
        ws = wb[TAB_CAPTABLE]
        addr = defined_name_ref(ws, NAME_CAP_OUTPUT_CCY) or CAP_TABLE_OUTPUT_CCY_CELL
        code = str(ws[addr].value or "").strip().upper()
    finally:
        wb.close()
    return code or "CAD"  # template default — the cell is empty/unreadable


def _output_currency_letter(workbook_path) -> str:
    """The cap table's output currency as a footnote letter (``'US'`` / ``'C'`` / ISO)."""
    return _currency_letter(_read_output_currency(workbook_path))


def _read_fx_rate(workbook_path) -> float | None:
    """The cap table's FX rate — the multiplier from its input to its output currency.

    This is the number that reconciles the deck's two currencies: the cap table
    converts its filing-currency inputs by it (``Add: Debt = F122*F7``), so a
    figure stated in both currencies differs by exactly this. It belongs on the
    deck for the same reason the figures do — a reader comparing LTM revenue on
    two slides is otherwise looking at an unexplained 41% gap.

    Returns None when the cell is empty or holds a formula rather than a value
    (CapIQ ships un-refreshed here, as everywhere): the caller then names both
    currencies without a rate rather than printing something it made up.

    The tab and the name are both established before this runs — the pre-flight
    verifies `infor_fx_rate` and `_read_output_currency` raises on a missing
    `captable` — so there is no fallback here either.
    """
    from openpyxl import load_workbook

    from template_layout import resolve_name_cell

    wb = load_workbook(workbook_path, data_only=True)
    try:
        ws = wb[TAB_CAPTABLE]
        value = ws[resolve_name_cell(ws, NAME_FX_RATE, template=Path(workbook_path).name)].value
    finally:
        wb.close()
    return float(value) if isinstance(value, (int, float)) else None


def _label_slide_currency(shape, code: str, *, second: _SecondCurrency | None = None) -> None:
    """State on `shape` which currency the slide's figures are in.

    `code` is the ISO currency of whatever populates the slide. Where the library
    ships the ``[x]$MM`` token this substitutes it, preserving the rest of the
    standardized source/note lines; where the library ships **no note at all** —
    the Financial Summary footnote reads "Sources: Company filings" and nothing
    else — this appends one, because a slide of figures with no currency on it
    is as wrong as one labelled in the wrong currency, and harder to notice.

    `second` is a second currency populating the SAME slide: the overview slide
    carries a C$ cap-table picture above US$ overview bullets. The note then names
    both and the rate between them, instead of picking one and leaving the reader
    to discover the gap against another slide.
    """
    lines = [p.text for p in shape.text_frame.paragraphs]
    if any("[x]" in line for line in lines):
        fill_footnote_currency(shape, _currency_letter(code))
    else:
        lines.append(f"Note: All figures in {_currency_phrase(code)} unless indicated otherwise")
        set_text(shape, lines)

    if second is None or _currency_letter(second.code) == _currency_letter(code):
        return

    clause = f"{second.what} in {_currency_phrase(second.code)}"
    if second.fx_rate:
        clause += (
            f" at {_currency_unit(code)} = "
            f"{_currency_amount(second.code, second.fx_rate)}"
        )
    lines = [p.text for p in shape.text_frame.paragraphs]
    for idx, line in enumerate(lines):
        if "All figures in" in line:
            lines[idx] = f"{line}; {clause}"
            break
    else:  # pragma: no cover — every branch above leaves an 'All figures in' line
        lines.append(f"Note: {clause}")
    set_text(shape, lines)


def _ownership_has_bloomberg(workbook_path) -> bool:
    """True when the ownership workbook carries Bloomberg institutional data.

    The ownership skill only populates the ``Bloomberg Output`` tab (first
    holder in ``C14``) when the analyst attached a Bloomberg export, so an
    empty C14 means the Select-Institutions block is uncomputed and the
    slide's right side must stay a placeholder.
    """
    try:
        from openpyxl import load_workbook

        wb = load_workbook(workbook_path, read_only=True)
        try:
            if "Bloomberg Output" not in wb.sheetnames:
                return False
            value = wb["Bloomberg Output"]["C14"].value
        finally:
            wb.close()
        return value is not None and str(value).strip() != ""
    except Exception:
        return False


def _fill_market_entry_targets(
    slide,
    *,
    row_labels: list[str],
    targets: list,
    market: str | None,
    currency: str,
    slide_number: int,
    total_slides: int,
) -> None:
    """Fill one market-entry slide with up to two targets.

    `targets` holds the 1-2 targets for THIS slide. The table is the fixed
    12-row structure (Overview / HQ / Year Founded → 7 consistent industry
    metrics → Scale KPIs / Strategic Rationale): the label column (col 0) is
    written white at 11 pt (stepping down per-label so no label wraps in the
    column) and the target value columns at `_ME_VALUE_SIZE` (9 pt —
    deliberately below the library's 10 pt so the 5.71" table clamp holds; see
    the constant's comment).

    `currency` is the content bundle's, because the cells on this slide are the
    content bundle's. It used to be the cap table's output currency, which is the
    currency of a different slide's picture.
    """
    title = "Potential " + (f"{market} " if market else "") + "Market Entry Targets"
    if total_slides > 1:
        title += f" ({slide_number} of {total_slides})"
    set_text(find_shape(slide, "Title 1"), [title])
    _label_slide_currency(find_shape(slide, _FOOTNOTE_MARKET_ENTRY), currency)

    # No early return on an empty `targets`: the library ships this table as 24
    # cells of "[x]", and returning here left every one of them on the deck. The
    # loops below blank what they do not fill, so a target-less slide comes out
    # empty rather than covered in placeholders.
    table_frame = find_table_shape(slide)
    table = table_frame.table
    n_cols = len(table.columns)  # label column + target columns (3 in the library)
    # Table row 0 is the blank logo/header row; data labels start at row 1. With
    # the fixed 12-row structure every data row is populated — no blank rows.
    for i, label in enumerate(row_labels):
        row = i + 1
        if row >= len(table.rows):
            break
        set_cell_text(table.cell(row, 0), label, size_pt=_ME_LABEL_SIZE, color_hex=_ME_LABEL_COLOR)
        for col in range(1, n_cols):
            target = targets[col - 1] if (col - 1) < len(targets) else None
            value = target.cells[i] if target is not None else ""
            set_cell_text(table.cell(row, col), value, size_pt=_ME_VALUE_SIZE)
    for row in range(len(row_labels) + 1, len(table.rows)):
        for col in range(n_cols):
            set_cell_text(table.cell(row, col), "", size_pt=_ME_VALUE_SIZE)

    # ORDERING: geometric. The logo boxes must line up left→right with the target
    # columns they sit above, and they are one row, so reading order is `.left`.
    #
    # The library ships each box as '[Placeholder for Logo]'; for a populated
    # column, name the target inside that phrase so the box says WHOSE logo goes
    # there while still reading as a placeholder. It used to become
    # '[Glean Logo]', which is guidance shaped like a caption — the delivered deck
    # printed "[Glean Logo]", "[Pinecone Logo]" across four slides and they read as
    # content rather than as something to clear. Keeping the library's own
    # '[Placeholder for …]' vocabulary (the same wording as the deferred chart
    # boxes the deck deliberately ships) makes it unmistakable, and the whole box
    # is one shape the analyst clears in one action. Blank any box whose column has
    # no target (odd final slide).
    logos = shapes_in_reading_order(
        s for s in slide.shapes
        if getattr(s, "has_text_frame", False) and "[Placeholder for Logo]" in s.text
    )
    for col_idx, logo in enumerate(logos):
        if col_idx < len(targets):
            name = getattr(targets[col_idx], "name", None)
            set_text(logo, [f"[Placeholder for {name} Logo]" if name else "[Placeholder for Logo]"])
        else:
            set_text(logo, [""])

    # Clamp the now-filled table to a fixed height so long content can't run it
    # off the slide. Done last, after every cell is populated. The clamp alone does
    # not stop render-time row growth — a declared row height is only a minimum —
    # and the converge loop measures the render and steps the body font down when
    # it happens. That replaced per-row content-height estimates here.
    set_table_height(table_frame, _ME_TABLE_HEIGHT)


def _contact_cards(slide):
    """The Contact slide's banker cards, in reading order.

    ORDERING: geometric, over a shape-identified set. A card is a 3x1 table, which
    is what distinguishes it from the 1x1 address table beside it; the four cards
    are a 2x2 grid, so document order says nothing about which is top-left.
    """
    return shapes_in_reading_order(
        shape for shape in slide.shapes
        if getattr(shape, "has_table", False)
        and len(shape.table.rows) == _CONTACT_CARD_ROWS
        and len(shape.table.columns) == _CONTACT_CARD_COLS
    )


def _fill_contact_cards(slide, contacts) -> None:
    """Fill the Contact slide's cards, and DELETE every card left unfilled.

    The library ships four cards, one of them already carrying INFOR's own
    details and three reading "[x]" down all three rows. The assembler used to
    treat the whole slide as static, so a delivered deck ended on three empty
    contact cards — which reads as an unfinished deck, and is the one slide a
    client is most likely to act on.

    `contacts` supplied: the first N cards are written in reading order and the
    rest deleted. `contacts` empty — the default: whatever the library ships
    filled is kept and the placeholder cards are deleted. That default is
    deliberately the template's own card rather than a name written down here;
    INFOR is single-tenant, so the deal team is knowable, but it is knowable from
    the artefact the analyst maintains, not from this module.
    """
    cards = _contact_cards(slide)
    for idx, card in enumerate(cards):
        table = card.table
        if idx < len(contacts):
            contact = contacts[idx]
            set_cell_text(
                table.cell(0, 0), f"{contact.name}, {contact.title}", size_pt=_CONTACT_NAME_SIZE
            )
            set_cell_text(table.cell(1, 0), contact.phone, size_pt=_CONTACT_DETAIL_SIZE)
            set_cell_text(table.cell(2, 0), contact.email, size_pt=_CONTACT_DETAIL_SIZE)
        elif contacts or any("[x]" in row.cells[0].text for row in table.rows):
            # Unfilled: either the content overrode the cards and this one is
            # surplus, or it is a library placeholder nothing filled.
            delete_shape(card)


class _PitchLayout:
    """Zero-based deck indices for a configurable pitch slide mix, **discovered**.

    Until v0.5.39 this class did arithmetic — ``self.ownership = 7 + nfs``,
    ``self.market_entry_first = 11 + nfs + (1 if kih else 0)`` — over a fixed
    prefix, so the shared library growing by one slide moved every constant.
    Now each slide is located by its marker shape on the finished slide mix.
    The arithmetic is gone; the class is the discovered result.

    Constructed AFTER the clones and deletes, while the markers are still
    intact (nothing is filled until the layout is known). Every lookup insists
    on exactly one match, so this doubles as the verification pass it replaced:
    a re-ordered library raises `TemplateLayoutError` here rather than
    misplacing content.
    """

    def __init__(self, prs, *, template_name: str, expected_total: int):
        find = partial(find_slide_by_marker, prs, template=template_name)
        self.cover = find(MARKER_COVER)
        self.overview = find(MARKER_OVERVIEW)
        self.financial_summary = find_slides_by_marker(prs, MARKER_FINANCIAL_SUMMARY)
        if not self.financial_summary:
            raise TemplateLayoutError(
                f"{template_name}: the deck needs at least one Financial Summary slide, "
                f"but none carries the {MARKER_FINANCIAL_SUMMARY.shape_name!r} marker"
            )
        self.ownership = find(MARKER_OWNERSHIP)
        self.contact = find(MARKER_CONTACT)
        self.risks = find(MARKER_RISKS)
        self.comps = find(MARKER_COMPS)
        self.precedents = find(MARKER_PRECEDENTS)
        self.investment_highlights = find_optional_slide_by_marker(
            prs, MARKER_KIH, template=template_name
        )
        self.market_entry = find_slides_by_marker(prs, MARKER_MARKET_ENTRY)
        self.n_financial_summary = len(self.financial_summary)
        self.n_market_entry = len(self.market_entry)
        # The expected count still comes from the caller's own clone/delete
        # bookkeeping, derived from the library's live slide count — so the
        # check stays meaningful (it catches a clone or delete that went wrong)
        # without re-hardcoding a deck size.
        self.total = expected_total


def assemble_pitch_deck(
    *,
    slide_plan_path: Path | str,
    content_path: Path | str,
    template_path: Path | str,
    output_dir: Path | str,
    captable_workbook_path: Path | str | None = None,
    ownership_workbook_path: Path | str | None = None,
    financial_metric_labels: list[str] | None = None,
    converge: bool = True,
) -> Path:
    """Fill the INFOR slide-library pitch deck.

    The blank library is 16 slides; the SlidePlan drives the slide mix:

    - the market-entry section expands across multiple slides (two targets per
      slide) based on ``content.market_entry_targets``;
    - the Financial Summary section expands to one slide per four metrics when
      the SlidePlan carries repeated ``financial-summary`` entries (the extra
      slides are cloned from the library's FS slide and retitled ``(k of n)``);
    - the Key Investment Highlights slide is deleted when the SlidePlan omits
      its ``key-investment-highlights`` entry.

    The cap table is pasted into slide 7 when ``captable_workbook_path`` is
    supplied, and the insider-ownership table into the ownership slide when
    ``ownership_workbook_path`` is supplied (both via the ``excel_to_powerpoint``
    insertion helper). When the ownership workbook also carries Bloomberg
    institutional data, the Select-Institutions block is pasted into the slide's
    right placeholder too; other chart/table insertions remain deferred
    placeholders.

    ``financial_metric_labels`` are the Financial Summary tile labels, selected
    by and handed off from the ``financial-summary`` stage (no longer on
    ``PitchDeckContent``). When supplied it must hold exactly four names per
    Financial Summary slide in the plan (4 for the default single slide, 8 for
    two); when None the FS tiles keep their template placeholder text.
    """
    slide_plan = SlidePlan.model_validate_json(Path(slide_plan_path).read_text(encoding="utf-8"))
    content = PitchDeckContent.model_validate_json(Path(content_path).read_text(encoding="utf-8"))
    if slide_plan.deliverable_type != "pitch":
        raise ValueError("pitch deck assembler only supports pitch SlidePlan objects")
    if len(slide_plan.slides) < 15:
        raise ValueError("pitch deck expects at least the base INFOR Slide Library entries (incl. precedent transactions)")

    template = Path(template_path)
    if not template.exists():
        raise FileNotFoundError(f"pitch library template not found: {template}")

    out_dir = Path(output_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    output_path = out_dir / f"Pitch Deck - {safe_filename(content.client_name, default='Client')}.pptx"

    # Layout pre-flight on the companion workbooks whose cells and picture
    # ranges are resolved below (the output currency and the FX rate behind the
    # slide-7 note, then the cap-table / insiders picture ranges) — a workbook
    # that lost those names raises here, naming all of them at once, instead of
    # part-way through the fill.
    if captable_workbook_path is not None:
        verify_workbook_names(
            captable_workbook_path,
            sheet=TAB_CAPTABLE,
            names=(NAME_FX_RATE, *CAP_TABLE_OUTPUT_CCY_NAMES, *CAP_TABLE_PICTURE_NAMES),
        )
    if ownership_workbook_path is not None:
        verify_workbook_names(
            ownership_workbook_path,
            sheet=TAB_OWNERSHIP,
            names=OWNERSHIP_INSIDERS_PICTURE_NAMES,
        )

    # Every slide's currency label comes from what populates that slide. The
    # content bundle states its own (`figure_currency`, the filing reporting
    # currency), and it labels every slide the bundle's copy fills. The cap table
    # states its OUTPUT currency, which may differ — it converts — so it labels
    # only the one slide its picture lands on, alongside the content's.
    content_currency = content.figure_currency
    cap_currency = (
        _read_output_currency(captable_workbook_path)
        if captable_workbook_path is not None
        else None
    )
    # The rate is only read when the two currencies actually differ — there is
    # nothing for it to reconcile otherwise, and the read is a second open of a
    # workbook on the analyst's (cloud-synced) deal directory.
    cap_table_second_currency = None
    if cap_currency is not None and _currency_letter(cap_currency) != _currency_letter(
        content_currency
    ):
        cap_table_second_currency = _SecondCurrency(
            what="capitalization table",
            code=cap_currency,
            fx_rate=_read_fx_rate(captable_workbook_path),
        )

    # The SlidePlan drives the deck's slide mix: the Financial Summary slide
    # count and the Key Investment Highlights toggle come from the wireframe's
    # entries; the market-entry slide count comes from the true content-bundle
    # target count (the wireframe count is only a default).
    plan_entry_ids = [entry.library_entry_id for entry in slide_plan.slides]
    n_financial_summary = max(1, plan_entry_ids.count("financial-summary"))
    include_kih = "key-investment-highlights" in plan_entry_ids
    n_market_entry = (
        max(1, math.ceil(len(content.market_entry_targets) / 2))
        if content.market_entry_targets
        else 1
    )
    if financial_metric_labels and len(financial_metric_labels) != _FS_TILES_PER_SLIDE * n_financial_summary:
        raise ValueError(
            f"financial_metric_labels holds {len(financial_metric_labels)} names but the "
            f"SlidePlan carries {n_financial_summary} Financial Summary slide(s) — expected "
            f"{_FS_TILES_PER_SLIDE * n_financial_summary} (four tiles per slide)"
        )

    prs = Presentation(template)

    # Locate the clone/delete targets by marker. Each lookup insists on exactly
    # one match in the blank library, so a re-ordered or re-saved library raises
    # TemplateLayoutError here instead of cloning or deleting the wrong slide.
    market_entry_at = find_slide_by_marker(prs, MARKER_MARKET_ENTRY, template=template.name)
    financial_summary_at = find_slide_by_marker(
        prs, MARKER_FINANCIAL_SUMMARY, template=template.name
    )
    earnings_at = find_slide_by_marker(prs, MARKER_EARNINGS_SUMMARY, template=template.name)

    # The deck's final slide count, derived from the library's own size rather
    # than a written-down total: every library entry survives except the
    # earnings slide (pitch does not use it) and — when the plan omits it — Key
    # Investment Highlights, plus one extra slide per additional market-entry
    # and Financial Summary section slide.
    expected_total = (
        len(prs.slides)
        - 1                                 # the earnings-summary slide, always dropped
        - (0 if include_kih else 1)
        + (n_market_entry - 1)
        + (n_financial_summary - 1)
    )

    # Grow the market-entry section (two targets per slide) and the Financial
    # Summary section (four metrics per slide) by cloning their library slides.
    # Clone BEFORE dropping the earnings slide so python-pptx allocates fresh,
    # non-colliding slide part names; market-entry first, because the FS clones
    # are inserted earlier in the deck and would shift its index.
    for _ in range(n_market_entry - 1):
        clone_slide_after(prs, market_entry_at)
    for _ in range(n_financial_summary - 1):
        clone_slide_after(prs, financial_summary_at)

    # The shared library carries the earnings-update slide; the pitch deck does
    # not use it. (Its index is unaffected by the clones above, which both land
    # after it.)
    delete_slide(prs, earnings_at)

    # Drop the Key Investment Highlights slide when the plan omits it — located
    # by marker on the post-delete deck, so no index arithmetic is involved.
    if not include_kih:
        delete_slide(prs, find_slide_by_marker(prs, MARKER_KIH, template=template.name))

    # The slide mix is final: discover every index the fill needs. Constructing
    # the layout IS the verification pass — each lookup requires exactly one
    # matching slide.
    layout = _PitchLayout(prs, template_name=template.name, expected_total=expected_total)

    # Cover — client name/date only.
    slide1 = prs.slides[layout.cover]
    title = find_shape(slide1, "Title 1")
    _replace_first_line(title, content.client_name, ["Internal Discussion Materials", ""])
    for shape in slide1.shapes:
        if shape.name == "Subtitle 2" and "[Date]" in _shape_text(shape):
            set_text(shape, [content.presentation_date])

    # Executive summary — flexible bullets. Fixed offsets from the cover: the
    # exec-summary, credentials and divider slides carry no marker of their own
    # (nothing distinguishes them but position in the library's front matter).
    slide2 = prs.slides[layout.cover + 1]
    _write_flexible_bullets(
        find_shape(slide2, "Content Placeholder 7"),
        content.executive_summary_bullets,
    )

    # The three credentials slides after it are static. Do not touch.

    # Section divider labels — the last front-matter slide before the overview.
    slide6 = prs.slides[layout.overview - 1]
    for idx, rect in enumerate(_section_boxes(slide6)):
        if idx < len(content.section_labels):
            set_text(rect, [content.section_labels[idx]])
        else:
            set_text(rect, [""])

    # Public company overview. The cap table is pasted into the 'Rectangle 3'
    # placeholder after save (when a workbook is supplied); the revenue pie
    # stays a deferred placeholder. The bullets box is sized to the band above
    # the 'LTM Revenue Breakdown' header and gets an explicit autofit fontScale
    # when over-long, so the copy cannot render into the pie section
    # (PowerPoint ignores a scale-less autofit on open).
    slide7 = prs.slides[layout.overview]
    set_text(find_shape(slide7, "Title 6"), [f"Introduction to {content.client_name}"])
    overview_shape = find_shape(slide7, "TextBox 9")
    _write_flexible_bullets(overview_shape, content.company_overview_bullets)
    fit_overview_textbox(slide7, overview_shape)
    # THE mixed-currency slide: the bullets are the content bundle's, the pasted
    # picture is the cap table's, and on a cross-border deal those are different
    # currencies. Name both, and the rate between them.
    _label_slide_currency(
        find_shape(slide7, _FOOTNOTE_OVERVIEW),
        content_currency,
        second=cap_table_second_currency,
    )

    # Financial Summary slide(s) — metric labels only (from the financial-summary
    # stage); charts remain placeholders. Left as template placeholders when no
    # labels. With two FS slides, tiles fill four labels per slide in order and
    # each slide is retitled '(k of n)'.
    for k, fs_index in enumerate(layout.financial_summary):
        fs_slide = prs.slides[fs_index]
        # The library ships this footnote as "Sources: Company filings" and
        # nothing else, so the note is APPENDED rather than substituted. The tab
        # behind these tiles is locked to millions in the filing currency, which
        # is the content bundle's — the same figures the overview slide restates
        # in the cap table's currency, which is why the two disagreed by 1.4102x
        # with no note on either to say so.
        _label_slide_currency(
            find_shape(fs_slide, _FOOTNOTE_FINANCIAL_SUMMARY), content_currency
        )
        if n_financial_summary > 1:
            set_text(
                find_shape(fs_slide, _FS_TITLE_SHAPE),
                [f"Financial Summary ({k + 1} of {n_financial_summary})"],
            )
        if financial_metric_labels:
            slide_labels = financial_metric_labels[
                _FS_TILES_PER_SLIDE * k : _FS_TILES_PER_SLIDE * (k + 1)
            ]
            for shape_name, label in zip(_FS_METRIC_TILES, slide_labels, strict=True):
                set_text(find_shape(fs_slide, shape_name), [label])

    # Acquirer considerations/mitigants — concise risks + tagline; the table is
    # clamped back to the library's 5.18" after fill (see _fill_risk_table).
    slide_risks = prs.slides[layout.risks]
    set_text(find_shape(slide_risks, "Text Placeholder 6"), [content.risks_tagline])
    _fill_risk_table(slide_risks, content)

    # Ownership takeaway. The insider/institution tables are pasted after save;
    # the takeaway line is the content bundle's, and until v0.5.51 there was no
    # field for it, so the slide shipped a bare '[x]' below 24 insiders and 118
    # institutions. No currency note: this slide states share counts and
    # percentages, and its "Source: Bloomberg" line is complete as shipped.
    slide_own = prs.slides[layout.ownership]
    set_text(find_shape(slide_own, "Text Placeholder 4"), [content.ownership_takeaway])

    # Comps takeaway; chart placeholder remains unless insertion later replaces it.
    # The footnote's currency is the content bundle's, because the takeaway is the
    # only thing on this slide that carries figures — the comps table is a deferred
    # placeholder. When that table lands, the label must come from the tab that
    # populates it; `infor_comps_output_ccy` is stamped for exactly that and is
    # written by nothing today, so it cannot be the source yet.
    slide_comps = prs.slides[layout.comps]
    set_text(find_shape(slide_comps, "Text Placeholder 5"), [content.comps_takeaway])
    _label_slide_currency(find_shape(slide_comps, _FOOTNOTE_COMPS), content_currency)

    # Precedent-transactions takeaway; chart placeholder remains (no
    # Excel→PowerPoint while Capital IQ can't be refreshed), mirroring comps —
    # including where its footnote currency comes from and why.
    slide_prec = prs.slides[layout.precedents]
    set_text(find_shape(slide_prec, "Text Placeholder 5"), [content.precedents_takeaway])
    _label_slide_currency(find_shape(slide_prec, _FOOTNOTE_PRECEDENTS), content_currency)

    # Key investment highlights — only when the plan carries the slide. Every box
    # on it is the content bundle's, so the footnote is the content's currency;
    # it used to be the cap table's, which labelled US$ quadrants C$MM.
    if layout.investment_highlights is not None:
        slide_kih = prs.slides[layout.investment_highlights]
        _fill_investment_highlights(slide_kih, content)
        _label_slide_currency(find_shape(slide_kih, _FOOTNOTE_KIH), content_currency)

    # Market-entry targets, two per slide. The section was grown above; fill
    # each slide with its pair and title it '(N of M)'.
    for j, me_index in enumerate(layout.market_entry):
        pair = content.market_entry_targets[2 * j : 2 * j + 2]
        _fill_market_entry_targets(
            prs.slides[me_index],
            row_labels=content.market_entry_row_labels,
            targets=pair,
            market=content.market_entry_market,
            currency=content_currency,
            slide_number=j + 1,
            total_slides=n_market_entry,
        )

    # The disclaimer is a static library entry — left untouched. Contact is NOT:
    # it ships three of its four cards as '[x]'.
    _fill_contact_cards(prs.slides[layout.contact], content.contacts)

    prs.save(output_path)

    # Write -> verify -> repair -> re-verify against the deck contract, BEFORE the
    # Excel pictures go in: the repair loop re-saves through python-pptx, and doing
    # that after the COM insertions would round-trip the pasted pictures for no
    # reason. A picture cannot overflow its box — only text and tables re-grow — so
    # nothing the insertions add needs fitting.
    #
    # No `out_dir`: the loop's ~170 renders per deck stage under `tempfile` and are
    # deleted on the way out. This used to pass `out_dir=out_dir / ".qa"`, which put
    # them permanently in the analyst's cloud-synced deal directory and cost the
    # per-file sync overhead that killed eight consecutive live pitch runs. Only a
    # FAILING converge leaves anything here, and `keep_on_failure` is where.
    if converge:
        assert_converged(
            output_path,
            converge_deck(output_path, library=template, keep_on_failure=out_dir / ".qa"),
        )

    # Paste the generated cap table into the overview slide's placeholder
    # (mirrors the earnings overview insertion). Done after save so the picture
    # write re-opens and re-saves the finished deck. The picture range comes
    # from the workbook's own `infor_cap_picture_range` name.
    if captable_workbook_path is not None:
        insert_excel_into_placeholder(
            deck_path=output_path,
            workbook_path=captable_workbook_path,
            output_path=output_path,
            slide_index=layout.overview,
            placeholder_name=_CAP_TABLE_PLACEHOLDER,
            sheet_name=TAB_CAPTABLE,
            source_range=resolve_workbook_range(
                captable_workbook_path,
                sheet=TAB_CAPTABLE,
                name=NAME_CAP_PICTURE_RANGE,
                fallback=CAP_TABLE_PICTURE_RANGE,
            ),
        )

    # Paste the insider-ownership table into the ownership slide's left
    # "Insiders" placeholder (mirrors the cap-table insertion). The ownership
    # slide follows the Financial Summary section at its computed layout index.
    # When the ownership stage also ingested a Bloomberg export (Bloomberg
    # Output C14 populated), the right "Institutions" placeholder gets the
    # Select-Institutions block too; otherwise it stays a Bloomberg placeholder.
    institutions_inserted = False
    if ownership_workbook_path is not None:
        insert_excel_into_placeholder(
            deck_path=output_path,
            workbook_path=ownership_workbook_path,
            output_path=output_path,
            slide_index=layout.ownership,
            placeholder_name=_OWNERSHIP_PLACEHOLDER,
            sheet_name=TAB_OWNERSHIP,
            source_range=resolve_workbook_range(
                ownership_workbook_path,
                sheet=TAB_OWNERSHIP,
                name=NAME_OWN_INSIDERS_PICTURE,
                fallback=OWNERSHIP_INSIDERS_PICTURE_RANGE,
            ),
        )
        if _ownership_has_bloomberg(ownership_workbook_path):
            # Check the Select-Institutions block's name is there before
            # pasting, then take the range from it.
            verify_workbook_names(
                ownership_workbook_path,
                sheet=TAB_OWNERSHIP,
                names=OWNERSHIP_INSTITUTIONS_PICTURE_NAMES,
            )
            insert_excel_into_placeholder(
                deck_path=output_path,
                workbook_path=ownership_workbook_path,
                output_path=output_path,
                slide_index=layout.ownership,
                placeholder_name=_INSTITUTIONS_PLACEHOLDER,
                sheet_name=TAB_OWNERSHIP,
                source_range=resolve_workbook_range(
                    ownership_workbook_path,
                    sheet=TAB_OWNERSHIP,
                    name=NAME_OWN_INSTITUTIONS_PICTURE,
                    fallback=OWNERSHIP_INSTITUTIONS_PICTURE_RANGE,
                ),
            )
            institutions_inserted = True

    _verify_pitch_output(
        output_path,
        cap_table_inserted=captable_workbook_path is not None,
        ownership_inserted=ownership_workbook_path is not None,
        institutions_inserted=institutions_inserted,
        expected_slides=layout.total,
    )
    return output_path


def _verify_pitch_output(
    path: Path,
    *,
    cap_table_inserted: bool = False,
    ownership_inserted: bool = False,
    institutions_inserted: bool = False,
    expected_slides: int | None = None,
) -> None:
    prs = Presentation(path)
    if expected_slides is not None and len(prs.slides) != expected_slides:
        raise ValueError(
            f"assembled pitch deck has {len(prs.slides)} slides; the slide-mix "
            f"layout expected {expected_slides}"
        )
    # The ownership slide takes TWO range pictures (insiders + institutions), so it
    # is where a range renderer selecting the wrong sheet shows up as one picture
    # pasted twice. Checked on every slide, not just that one.
    assert_range_pictures_are_distinct(prs)
    text = _all_text(prs)
    forbidden = ["[CLIENT NAME]", "[Date]"]
    leftovers = [token for token in forbidden if token in text]
    if leftovers:
        raise ValueError(f"assembled pitch deck still contains required-field placeholders: {leftovers}")
    unfilled = _unfilled_placeholders(prs)
    if unfilled:
        raise ValueError(
            "assembled pitch deck still contains unfilled '[x]' placeholders — a "
            "deck that reaches a client with these on it reads as unfinished:\n  "
            + "\n  ".join(unfilled)
        )
    required_placeholders = [
        "[Pie Chart Placeholder]",
        "[Placeholder for Metric #1 Chart]",
        "[Placeholder for Comps Chart]",
        # The precedent-transactions slide stays a chart placeholder, like comps
        # (no Excel→PowerPoint while Capital IQ can't be refreshed here).
        "[Placeholder for Precedents Chart]",
    ]
    missing = [token for token in required_placeholders if token not in text]
    if missing:
        raise ValueError(f"deferred placeholders were unexpectedly removed: {missing}")
    has_cap_placeholder = "[Cap Table Placeholder]" in text
    if cap_table_inserted and has_cap_placeholder:
        raise ValueError("slide 7 cap-table placeholder was not replaced by the Excel insertion stage")
    if not cap_table_inserted and not has_cap_placeholder:
        raise ValueError("slide 7 cap-table placeholder must remain when no workbook is supplied")
    has_insider_placeholder = "[Placeholder for Insider Ownership]" in text
    if ownership_inserted and has_insider_placeholder:
        raise ValueError("ownership slide insider placeholder was not replaced by the Excel insertion stage")
    if not ownership_inserted and not has_insider_placeholder:
        raise ValueError("ownership slide insider placeholder must remain when no workbook is supplied")
    # The ownership slide's institutional side is filled only when the ownership
    # stage ingested a Bloomberg export; otherwise it must stay a placeholder.
    has_institutions_placeholder = "[Placeholder for Institutional Ownership]" in text
    if institutions_inserted and has_institutions_placeholder:
        raise ValueError("ownership slide institutions placeholder was not replaced by the Excel insertion stage")
    if not institutions_inserted and not has_institutions_placeholder:
        raise ValueError("ownership slide institutions placeholder must remain without Bloomberg data")
