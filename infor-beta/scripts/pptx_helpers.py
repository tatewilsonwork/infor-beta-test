"""python-pptx helpers shared across INFOR skills.

These helpers encode the formatting rules that recur across every INFOR deck
work and that have a history of regressing when re-derived inline in a skill:

  - `set_text(shape, lines)` preserves run-level rPr (font/size/bold/italic/color)
    by mutating runs[0].text in place rather than wiping and re-adding runs.
  - `write_bulleted_shape(shape, items)` harvests bullet pPr templates from the
    template's seed paragraphs BEFORE wiping, so new bullets keep the square /
    dash glyphs that python-pptx would otherwise drop.
  - `set_cell_text(cell, text, size_pt, color_hex)` forces Palatino on table
    cells (PowerPoint's default fallback is Calibri, which has been observed
    to slip in when cells are rewritten).
  - `delete_slide(prs, index)` removes a slide and drops its presentation-part
    relationship, so reducing a cloned library deck to the wanted entries
    leaves no orphaned slide parts behind.

Skills can load these via:

    import sys, os
    sys.path.insert(0, os.environ.get("CLAUDE_PLUGIN_ROOT", "./infor-beta") + "/scripts")
    from pptx_helpers import set_text, write_bulleted_shape, set_cell_text, delete_slide, find_shape

Tests live in tests/test_pptx_helpers.py and build fresh in-memory
decks so they don't depend on the INFOR template files.
"""

from copy import deepcopy

from lxml import etree
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt


# ─── Brand constants ─────────────────────────────────────────────────────────

PALATINO = "Palatino Linotype"
COLOR_UP = "00B050"    # green — positive delta / beat
COLOR_DOWN = "C00000"  # red   — negative delta / miss

# INFOR theme accent palette (INFORFG.thmx "INFOR (New)", accent1–6) — categorical
# fills for pie / segment charts, used in theme order and cycled past six. accent2
# (46566E) is also the clustered-column bar colour used by the Financial Summary charts.
INFOR_ACCENTS = ["0E213F", "46566E", "ADB9CA", "A4844B", "767171", "E5E3E3"]


# ─── Shape lookup ────────────────────────────────────────────────────────────

def find_shape(slide, name):
    """Return the first shape on the slide whose .name matches."""
    for s in slide.shapes:
        if s.name == name:
            return s
    raise KeyError(f"Shape {name!r} not found on slide")


def find_shape_in_group(group, name):
    """Return the first child of the group whose .name matches."""
    for s in group.shapes:
        if s.name == name:
            return s
    raise KeyError(f"Shape {name!r} not found in group {group.name!r}")


# ─── XML helpers (private) ───────────────────────────────────────────────────

def _pPr_of(paragraph):
    """Return the paragraph-level <a:pPr> element, or None if absent."""
    for child in paragraph._p:
        if child.tag.endswith("}pPr"):
            return child
    return None


def _first_run_rPr(paragraph):
    """Return a deepcopy of the first run's <a:rPr>, or None if absent."""
    for r in paragraph.runs:
        for child in r._r:
            if child.tag.endswith("}rPr"):
                return deepcopy(child)
        return None
    return None


def _replace_run_rPr(run, template_rPr):
    """Strip the run's existing rPr children and graft a fresh deepcopy in."""
    if template_rPr is None:
        return
    for child in [c for c in run._r if c.tag.endswith("}rPr")]:
        run._r.remove(child)
    run._r.insert(0, deepcopy(template_rPr))


# ─── set_text ────────────────────────────────────────────────────────────────

def set_text(shape, lines, size_pt=None, color_hex=None):
    """Replace shape text while preserving the template's run formatting.

    Strategy:
      - For existing paragraphs (i < len(tf.paragraphs)): mutate runs[0].text
        in place to keep its rPr (font/size/bold/italic/color). Remove later
        runs on the same paragraph.
      - For new paragraphs beyond the existing count: add_paragraph, copy
        pPr from paragraph 0, copy first run's rPr from paragraph 0 so the
        new run inherits the template's formatting.

    `size_pt` and `color_hex` are explicit overrides applied AFTER the
    template formatting is restored. Pass them only when you intentionally
    want to override — e.g., the delta boxes on the earnings deck need
    forced 10 pt + green/red. Title bars / quote boxes should NOT receive
    overrides (the template's bold/italic/color would be wiped).
    """
    tf = shape.text_frame
    template_pPr = _pPr_of(tf.paragraphs[0])
    template_rPr = _first_run_rPr(tf.paragraphs[0])

    for i, line in enumerate(lines):
        if i < len(tf.paragraphs):
            p = tf.paragraphs[i]
            for r in list(p.runs[1:]):
                r._r.getparent().remove(r._r)
            if p.runs:
                # In-place mutation preserves the run's rPr
                p.runs[0].text = line
                run = p.runs[0]
            else:
                run = p.add_run()
                run.text = line
                _replace_run_rPr(run, template_rPr)
        else:
            p = tf.add_paragraph()
            if template_pPr is not None:
                for child in list(p._p):
                    if child.tag.endswith("}pPr"):
                        p._p.remove(child)
                p._p.insert(0, deepcopy(template_pPr))
            run = p.add_run()
            run.text = line
            _replace_run_rPr(run, template_rPr)

        if size_pt is not None:
            run.font.name = PALATINO
            run.font.size = Pt(size_pt)
        if color_hex is not None:
            run.font.color.rgb = RGBColor.from_string(color_hex)

    while len(tf.paragraphs) > len(lines):
        p = tf.paragraphs[-1]
        p._p.getparent().remove(p._p)


# ─── set_cell_text ───────────────────────────────────────────────────────────

def _declare_paragraph_size(paragraph, size_pt):
    """Put an explicit size on a run-less paragraph, so it reserves only that line.

    A cell blanked with a zero-length run (`<a:t/>`) and no declared size takes its
    line height from the table style's default — which on the INFOR library is far
    larger than the body copy. Every row then renders taller than the assembler
    declared, and no amount of font stepping helps because the empty cell's size is
    not what was stepped.

    Measured on the market-entry table's unused third column (odd target count):
    12 rows grew a little each and the table rendered **0.587" past the slide
    edge**; declaring the size here takes it to 0.007", i.e. clean. The size itself
    does not matter (1 pt, 4 pt and 9 pt all measure identically) — what matters is
    that the paragraph declares one at all.
    """
    pPr = paragraph._p.find(qn("a:pPr"))
    if pPr is None:
        pPr = etree.Element(qn("a:pPr"))
        paragraph._p.insert(0, pPr)  # pPr must be the first child of <a:p>
    defRPr = pPr.find(qn("a:defRPr"))
    if defRPr is None:
        defRPr = etree.SubElement(pPr, qn("a:defRPr"))
    endParaRPr = paragraph._p.find(qn("a:endParaRPr"))
    if endParaRPr is None:
        endParaRPr = etree.SubElement(paragraph._p, qn("a:endParaRPr"))  # must be last
    for node in (defRPr, endParaRPr):
        node.set("sz", str(int(round(size_pt * 100))))


def set_cell_text(cell, text, size_pt=9, color_hex=None):
    """Overwrite a table cell as a single Palatino run at size_pt.

    Unlike `set_text`, this DOES set font.name + font.size explicitly because
    PowerPoint's table-cell default fallback is Calibri — inheriting from the
    template has been observed to slip back to Calibri across rewrites.

    Blanking a cell (``text=""``) leaves a genuinely empty paragraph that declares
    `size_pt`, rather than a zero-length run — see `_declare_paragraph_size` for
    the row-growth defect that fixes.
    """
    tf = cell.text_frame
    while len(tf.paragraphs) > 1:
        last = tf.paragraphs[-1]
        last._p.getparent().remove(last._p)
    p = tf.paragraphs[0]
    for r in list(p.runs):
        r._r.getparent().remove(r._r)
    if not text:
        _declare_paragraph_size(p, size_pt)
        return
    run = p.add_run()
    run.text = text
    run.font.name = PALATINO
    run.font.size = Pt(size_pt)
    if color_hex is not None:
        run.font.color.rgb = RGBColor.from_string(color_hex)


# ─── write_bulleted_shape ────────────────────────────────────────────────────

def _has_bullet_glyph(pPr):
    """True if the paragraph properties carry a real bullet glyph (not buNone)."""
    return any(
        child.tag.endswith("}buChar") or child.tag.endswith("}buAutoNum")
        for child in pPr
    )


def _harvest_bullet_templates(shape):
    """Capture pPr + rPr templates from the shape's bulleted seed paragraphs.

    Returns {level_index: (pPr_copy, rPr_copy)} keyed 0, 1, 2... by ascending
    marL (smallest indent = main bullet = level 0).

    Only paragraphs that carry a real bullet glyph (`<a:buChar>` / `<a:buAutoNum>`)
    define a level. Empty spacer paragraphs and explicit `<a:buNone>` paragraphs
    are skipped, and multiple seed paragraphs sharing an indent collapse to one
    level — keyed by distinct marL, not one-entry-per-paragraph. (The INFOR
    library's exec-summary placeholder ships ~14 paragraphs, most empty or
    buNone; the old per-paragraph enumeration mapped level 0 to the marL=0
    buNone spacer, so main bullets lost their glyph and fell back to the
    placeholder's navy list colour.) When several seed paragraphs share a marL,
    the first wins, preferring one whose run carries an rPr so the level keeps
    the template's colour / size.
    """
    tf = shape.text_frame
    by_marL: dict[int, tuple] = {}
    for para in tf.paragraphs:
        pPr = _pPr_of(para)
        if pPr is None or not _has_bullet_glyph(pPr):
            continue
        marL = int(pPr.get("marL") or "0")
        rPr = _first_run_rPr(para)
        if marL not in by_marL:
            by_marL[marL] = (deepcopy(pPr), deepcopy(rPr) if rPr is not None else None)
        elif by_marL[marL][1] is None and rPr is not None:
            by_marL[marL] = (deepcopy(pPr), deepcopy(rPr))
    ordered = sorted(by_marL.items(), key=lambda kv: kv[0])
    return {i: tpl for i, (_marL, tpl) in enumerate(ordered)}


def _emit_bullet_run(paragraph, text, tmpl_rPr, size, *, bold):
    """Append a run, grafting the harvested template rPr so it keeps the
    template's colour (and italic / etc.), then re-assert Palatino/size/bold.

    Grafting the harvested rPr is what keeps body copy the template's intended
    colour: without it the run carries no `<a:solidFill>` and inherits the
    placeholder list-style `defRPr`, which on the INFOR library is navy
    (`1B2759`) and renders as blue body text. name/size/bold are re-applied on
    top because PowerPoint's list/table fallback typeface is Calibri."""
    run = paragraph.add_run()
    run.text = text
    if tmpl_rPr is not None:
        _replace_run_rPr(run, tmpl_rPr)
    run.font.name = PALATINO
    run.font.size = size
    run.font.bold = bold
    return run


def write_bulleted_shape(shape, items):
    """Wipe the shape and rewrite bullets with pPr + rPr correctly preserved.

    `items` is a list of tuples:
        (text, level)                       — single-run bullet
        (prefix_bold, rest_regular, level)  — two-run bullet (bold prefix +
                                              regular tail, used for segment
                                              names like 'easyfinancial: ...')

    `level` 0 = main (square glyph, larger font); 1 = sub (dash, smaller).

    Harvests the seed pPr templates BEFORE wiping so the bullet characters
    survive. Sets font.name = Palatino Linotype and font.size explicitly on
    every run. After writing, asserts every paragraph has a buChar element —
    raises RuntimeError if any bullet is missing its glyph, so a broken deck
    fails at write-time instead of shipping silently.
    """
    tf = shape.text_frame
    templates = _harvest_bullet_templates(shape)  # must happen BEFORE we wipe
    if not templates:
        raise RuntimeError(
            f"Shape {shape.name!r} has no bullet templates to harvest; "
            f"write_bulleted_shape requires the template to ship at least one "
            f"seed paragraph with a pPr (square or dash bullet)."
        )

    def _size_for(level):
        _, rPr = templates.get(level, (None, None))
        if rPr is not None:
            sz = rPr.get("sz")
            if sz is not None:
                return Pt(int(sz) / 100)
        return Pt(10.5 if level == 0 else 10.0)

    # Wipe: leave one paragraph behind, clear its runs and pPr
    while len(tf.paragraphs) > 1:
        last = tf.paragraphs[-1]
        last._p.getparent().remove(last._p)
    first = tf.paragraphs[0]
    for r in list(first.runs):
        r._r.getparent().remove(r._r)
    for child in list(first._p):
        if child.tag.endswith("}pPr"):
            first._p.remove(child)

    for i, item in enumerate(items):
        if len(item) == 2:
            prefix, rest, level = "", item[0], item[1]
        elif len(item) == 3:
            prefix, rest, level = item
        else:
            raise ValueError(
                "items must be (text, level) or (prefix_bold, rest_regular, level)"
            )

        p = first if i == 0 else tf.add_paragraph()
        if i != 0:
            for child in list(p._p):
                if child.tag.endswith("}pPr"):
                    p._p.remove(child)

        tmpl_pPr, tmpl_rPr = templates.get(level, templates[0])
        p._p.insert(0, deepcopy(tmpl_pPr))

        size = _size_for(level)
        if prefix:
            _emit_bullet_run(p, prefix, tmpl_rPr, size, bold=True)
            _emit_bullet_run(p, rest, tmpl_rPr, size, bold=False)
        else:
            _emit_bullet_run(p, rest, tmpl_rPr, size, bold=False)

    # Post-write: every paragraph must have a bullet character
    for i, para in enumerate(tf.paragraphs):
        has_bu = False
        for elem in para._p.iter():
            if elem.tag.endswith("}buChar") or elem.tag.endswith("}buAutoNum"):
                has_bu = True
                break
        if not has_bu:
            raise RuntimeError(
                f"Shape {shape.name!r} paragraph {i} has no bullet character — "
                f"pPr template was not propagated. Refusing to ship a broken deck."
            )


# ─── Autofit (shrink text on overflow) ───────────────────────────────────────

def enable_normal_autofit(shape, font_scale=None, line_space_reduction=None):
    """Set the text frame's body autofit to 'Shrink text on overflow'.

    PowerPoint's `<a:normAutofit/>` autofit makes the renderer scale the font
    down at display time when the text would overflow the frame. Use this on
    any shape whose copy length varies from deck to deck — bulleted overview
    blocks, business-updates lists, etc. — so an over-budget LLM run doesn't
    silently render text past the divider line.

    Pass `font_scale` (0-100, percent of original) and `line_space_reduction`
    (0-100, percent reduction in line spacing) to pre-apply a specific
    scaling instead of letting PowerPoint compute one on first render. The
    OOXML stores both as integer thousandths-of-a-percent.
    """
    tf = shape.text_frame
    txBody = tf._txBody
    bodyPr = txBody.find(qn("a:bodyPr"))
    if bodyPr is None:
        bodyPr = etree.SubElement(txBody, qn("a:bodyPr"))
        txBody.insert(0, bodyPr)
    for child in list(bodyPr):
        if child.tag in (qn("a:noAutofit"), qn("a:normAutofit"), qn("a:spAutoFit")):
            bodyPr.remove(child)
    norm = etree.SubElement(bodyPr, qn("a:normAutofit"))
    if font_scale is not None:
        norm.set("fontScale", str(int(round(font_scale * 1000))))
    if line_space_reduction is not None:
        norm.set("lnSpcReduction", str(int(round(line_space_reduction * 1000))))


def normal_autofit_scale(shape):
    """Stored `normAutofit` fontScale in percent, or None when there is no normAutofit.

    A `<a:normAutofit/>` with no `fontScale` attribute returns 100.0: PowerPoint
    ignores a scale-less autofit on open and renders at full size (the v0.5.23
    finding), so full size is what it means in practice.
    """
    bodyPr = shape.text_frame._txBody.find(qn("a:bodyPr"))
    if bodyPr is None:
        return None
    norm = bodyPr.find(qn("a:normAutofit"))
    if norm is None:
        return None
    raw = norm.get("fontScale")
    return 100.0 if raw is None else int(raw) / 1000.0


def strip_autofit(shape):
    """Replace whatever autofit the shape carries with an explicit `<a:noAutofit/>`.

    Used to build a render probe that shows the shape at its *stored* text size.
    LibreOffice treats ANY `<a:normAutofit>` — explicit `fontScale` included — as
    "shrink to fit" and recomputes its own scale, so a rendered measurement of an
    autofit shape can never show overflow. Stripping the autofit (after baking the
    stored scale into the run sizes with `apply_text_scale`) is what makes the
    render measure what PowerPoint will actually draw.
    """
    txBody = shape.text_frame._txBody
    bodyPr = txBody.find(qn("a:bodyPr"))
    if bodyPr is None:
        bodyPr = etree.SubElement(txBody, qn("a:bodyPr"))
        txBody.insert(0, bodyPr)
    for child in list(bodyPr):
        if child.tag in (qn("a:noAutofit"), qn("a:normAutofit"), qn("a:spAutoFit")):
            bodyPr.remove(child)
    etree.SubElement(bodyPr, qn("a:noAutofit"))


# Text-size attributes carried by a run, an empty-paragraph mark, and a
# paragraph-level default. All three drive rendered line height.
_SIZED_TAGS = ("a:rPr", "a:endParaRPr", "a:defRPr")


def apply_text_scale(shape, scale_pct):
    """Multiply every EXPLICIT text size in the shape by `scale_pct` percent.

    The counterpart to `strip_autofit`: together they turn "this shape carries a
    75% autofit" into "this shape's runs are 75% of their nominal size and nothing
    will shrink them further", which is renderable and therefore measurable.

    Sizes that are *inherited* (a run with no `sz`, taking the placeholder list
    style's) are left alone, because the effective size is not in this shape's
    XML. That errs toward rendering the text LARGER than it will ship, which is
    the safe direction for an overflow check. Every INFOR assembler writes
    explicit sizes (`set_cell_text`, `write_bulleted_shape`), so in practice
    nothing is missed.
    """
    factor = scale_pct / 100.0
    for tag in _SIZED_TAGS:
        for node in shape.text_frame._txBody.iter(qn(tag)):
            raw = node.get("sz")
            if raw is not None:
                node.set("sz", str(max(100, int(round(int(raw) * factor)))))


# ─── Variable-length text-block fit (both deck assemblers) ───────────────────
# Gap left between a text block and the section header that closes its band, so
# the two do not touch. Measured off the library, not derived.
_FIT_BAND_GAP_IN = 0.12
_OVERVIEW_BAND_PREFIX = "LTM Revenue"


def fit_overview_textbox(slide, shape, *, band_prefix=_OVERVIEW_BAND_PREFIX):
    """Size a variable-length text block to the free band above its next header.

    The shared library ships these blocks **one line tall** and relies on autofit,
    but PowerPoint ignores a scale-less ``<a:normAutofit/>`` on open, so an
    over-long run renders full size and spills into whatever is below — the
    recurring "overview text overlaps the LTM revenue breakdown" analyst report,
    and (found by the Phase B contract) the same thing on the earnings-update
    Business Updates block. Sizing the box to the band is the fix that belongs
    here: it is pure geometry, read off the neighbouring shape's position.

    **Deciding how far to shrink the text is deliberately NOT done here.** Until
    v0.5.39 this function solved for a ``fontScale`` using hand-calibrated
    Palatino em constants — an average character width, a line height, a paragraph
    spacing allowance — i.e. a reimplementation of the layout engine in Python.
    Those constants were recalibrated ~15% in v0.5.23 after a "fitted" block spilled
    anyway, which is the argument against having them at all. `deck_repair` now
    measures the real render and steps the scale down until the ink fits.

    So this writes a plain shrink-on-overflow autofit and leaves the scale alone.
    Returns the available band height in inches.
    """
    top_in = Emu(shape.top).inches
    band_bottom = None
    for other in slide.shapes:
        if getattr(other, "has_text_frame", False) and band_prefix in other.text_frame.text:
            band_bottom = Emu(other.top).inches
            break
    if band_bottom is not None and band_bottom - top_in > 0.5:
        avail_in = band_bottom - _FIT_BAND_GAP_IN - top_in
        shape.height = Inches(avail_in)
    else:
        avail_in = Emu(shape.height).inches
    enable_normal_autofit(shape)
    return avail_in


# ─── delete_slide ────────────────────────────────────────────────────────────

def delete_slide(prs, index):
    """Remove a slide from a presentation by zero-based index.

    Drops both the `<p:sldId>` entry AND the presentation-part relationship that
    points at the slide. Removing only the sldId leaves the slide part orphaned
    in the package, which python-pptx still serialises — producing duplicate
    part-name warnings and stray slides on re-open. Dropping the relationship
    makes the part unreachable so it is not written out.

    Used to reduce a cloned library deck to the slides you want: open the full
    `INFOR Slide Library.pptx`, then delete the entries you don't need.
    """
    xml_slides = prs.slides._sldIdLst
    sldId = list(xml_slides)[index]
    rId = sldId.get(qn("r:id"))
    xml_slides.remove(sldId)
    if rId is not None:
        prs.part.drop_rel(rId)


# ─── clone_slide_after ────────────────────────────────────────────────────────

def clone_slide_after(prs, index):
    """Duplicate the slide at `index`, inserting the copy immediately after it.

    Returns the new slide. Creates a fresh slide from the source's layout, drops
    the placeholders `add_slide` seeds, deep-copies every shape element from the
    source `spTree`, then moves the new `sldId` from the tail to position
    `index + 1`.

    Used to grow a fixed-size library section to a variable count — e.g. the
    pitch market-entry slide, which the deck repeats (two targets per slide) when
    a deal has more than two targets.

    Caveat: this copies shape XML only. Slide-level relationships (embedded
    pictures, charts, hyperlinks) are NOT re-created, so the source slide must
    not depend on them. The INFOR market-entry slide carries only a table, text,
    and shape-based logo placeholders, so an spTree copy is sufficient.
    """
    source = prs.slides[index]
    new_slide = prs.slides.add_slide(source.slide_layout)
    for shape in list(new_slide.shapes):
        shape._element.getparent().remove(shape._element)
    for shape in source.shapes:
        new_slide.shapes._spTree.append(deepcopy(shape._element))
    sldIdLst = prs.slides._sldIdLst
    new_sldId = list(sldIdLst)[-1]
    sldIdLst.remove(new_sldId)
    sldIdLst.insert(index + 1, new_sldId)
    return new_slide


# ─── Shape traversal / lookup shared by the deck assemblers ──────────────────

def iter_all_shapes(shapes):
    """Yield every shape in `shapes`, recursing into groups depth-first."""
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_all_shapes(shape.shapes)


def find_table_shape(slide):
    """Return the first shape on `slide` that holds a table; raise KeyError if none."""
    for shape in slide.shapes:
        if getattr(shape, "has_table", False):
            return shape
    raise KeyError("no table shape found on slide")


def set_table_height(table_frame, total_height):
    """Resize a table to an exact total height by scaling its row heights.

    Mirrors setting a table's height in PowerPoint (drag the bottom handle / set
    Height in the Size pane): the graphic-frame extent and every row height are
    scaled to `total_height` proportionally, so the rows still sum to exactly the
    target (rounding remainder folded into the last row). Call AFTER the cells are
    filled — row heights are only meaningful once content is in place.

    A declared row height is only a render-time MINIMUM, so this clamp does not by
    itself stop the layout engine re-growing a row whose text needs more.

    Until v0.5.39 it tried to: callers passed per-row `min_heights` estimated from
    a Palatino character-width table and a 1.2-em line height, and a waterfall
    pinned any row below its estimate. That is a layout engine reimplemented in
    Python to guess an answer the renderer already knows, and it guessed wrong often
    enough to need recalibrating. What actually stops row re-growth is smaller text,
    and what decides how much smaller is the measured render — see `deck_repair`,
    which steps the body font down until the rendered table fits and re-clamps
    through this function.
    """
    table = table_frame.table
    rows = list(table.rows)
    base = sum(r.height for r in rows)
    if not rows or base <= 0:
        return
    target = int(total_height)
    heights = [max(0, int(round(r.height * target / base))) for r in rows]
    heights[-1] = max(0, heights[-1] + target - sum(heights))  # rows sum exactly
    for row, h in zip(rows, heights):
        row.height = h
    table_frame.height = sum(heights)


# ─── Footnote currency token ─────────────────────────────────────────────────

def fill_footnote_token(shape, letter, *, token="[x]"):
    """Substitute the currency-letter token in a library footnote, in place.

    The shared INFOR library footnotes ship a 'All figures in [x]$MM' line where
    `[x]` is the currency letter ('US' / 'C'). Replace the token across every
    paragraph and rewrite via `set_text`, so the rest of the standardized
    source/note lines (and their formatting) are preserved rather than
    re-hardcoded.
    """
    lines = [p.text.replace(token, letter) for p in shape.text_frame.paragraphs]
    set_text(shape, lines)


def fill_footnote_currency(shape, letter):
    """Fill the footnote's '[x]$MM' token for a dollar letter OR a bare ISO code.

    `letter` is what the assemblers' currency helpers return: ``"US"`` / ``"C"``
    for the two dollar currencies (substituted into the token as before, giving
    ``US$MM`` / ``C$MM``), or a bare ISO code (``"GBP"``, ``"EUR"``, …) for any
    other currency — the code then replaces the letter AND the dollar sign
    (``[x]$MM`` -> ``GBP MM``), so a non-dollar filer's footnote renders its
    ISO code instead of mislabelling the figures with a dollar sign.
    """
    if letter in ("US", "C"):
        fill_footnote_token(shape, letter)
    elif any("[x]$" in p.text for p in shape.text_frame.paragraphs):
        fill_footnote_token(shape, f"{letter} ", token="[x]$")
    else:
        fill_footnote_token(shape, letter)


# ─── Bullet writer with plain-text fallback ──────────────────────────────────

def write_bullets_or_plain(shape, items, *, autofit=False):
    """Write bullets, falling back to plain paragraphs when the shape has no
    bullet-glyph template to harvest.

    `items` is the tuple list `write_bulleted_shape` expects — `(text, level)` or
    `(prefix_bold, rest, level)`. If the template ships no seed bullet (so
    `write_bulleted_shape` raises RuntimeError), degrade to a flat `set_text` of
    each item's text. Pass `autofit=True` to also enable shrink-on-overflow
    autofit, for variable-length blocks (overview bullets, business updates).
    """
    try:
        write_bulleted_shape(shape, items)
    except RuntimeError:
        set_text(shape, [item[0] for item in items])
    if autofit:
        enable_normal_autofit(shape)
