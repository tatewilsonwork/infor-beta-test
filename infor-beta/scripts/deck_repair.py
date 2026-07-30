"""Write → verify → repair → re-verify for a built INFOR deck (Phase B step 3).

`deck_contract.verify_deck` can see a deck. This module acts on what it sees, and
it is what lets the estimation code go away.

For thirteen releases the assemblers **predicted** the renderer: a per-character
Palatino advance-width table, em constants for line height and paragraph spacing,
per-row content-height floors, and a font-stepping ladder driven by all three.
Every one of those was a hand-calibrated model of a layout engine, recalibrated
whenever an analyst opened a file and found it wrong (the v0.5.23 note "the old
constants under-estimated rendered height ~15%" is the shape of the problem).

Here the renderer answers for itself. To choose a font size, this module builds a
probe slide per candidate size, renders them all in one LibreOffice pass, and
reads back how far the ink actually falls past the box. Nothing is predicted.

Two properties make that affordable and sound:

  - **One conversion per ladder, not per step.** `deck_contract.build_probe_deck`
    packs every candidate of every shape into a single deck, so measuring a
    four-step ladder on six shapes costs one conversion.
  - **The measurement is conservative.** LibreOffice wraps Palatino about one line
    taller than PowerPoint, so a size that fits under LibreOffice fits in
    PowerPoint. This deliberately does not correct toward PowerPoint's metrics.

Autofit is the subtle part
--------------------------
LibreOffice treats any `<a:normAutofit>` as "shrink to fit" and recomputes its own
scale, so an autofit shape rendered as-is **always** measures clean no matter how
much text it holds. PowerPoint does not: it applies the stored `fontScale` and
nothing more. Measuring the shipped XML would therefore certify every overview
block as fine — which is exactly the v0.5.23 defect, undetected.

So every text probe bakes the candidate scale into the run sizes and replaces the
autofit with `<a:noAutofit/>` (`pptx_helpers.apply_text_scale` / `strip_autofit`).
`lnSpcReduction` is deliberately not modelled, which renders slightly taller than
ships — again the safe direction.

What gets repaired
------------------
Only the three geometric kinds, and only by making content smaller — never by
moving or growing a box, because the box is the brand design:

  - `rendered-overflow` / `masked-overflow` on a **text** shape → step the
    autofit `fontScale` down the measured ladder.
  - the same on a **table**, or `table-taller-than-library` → step the body cells'
    point size down and re-clamp the declared rows through
    `pptx_helpers.set_table_height`.

Header rows are left alone. The library ships the Considerations/Mitigants header
at 12 pt against 10 pt body copy, and an earlier revision that shrank everything
uniformly "rendered noticeably smaller than the template" (v0.5.14). Row 0 of
every INFOR table is either a header or, on market-entry, an empty logo row.

Failure is loud
---------------
The loop is bounded (default 3 iterations) and **fails the stage** when it cannot
converge — `assert_converged` raises `DeckNotConvergedError`. A deck that still
overflows is never saved quietly, because "quietly" is how thirteen of these
shipped. Vision findings never enter the loop: they are advisory, non-deterministic
and belong at the `deck` checkpoint where an agent or the analyst reads them.

QA scratch is ephemeral unless the stage fails
---------------------------------------------
The loop renders a lot: per measured pass, 19 slide PNGs, 63 attribution PNGs and
a 1.7 MB probe deck for the verify, plus 8 PNGs and another probe deck for the
repair — ~170 files and ~10 MB for one deck. **None of it is a deliverable.** So
`out_dir` defaults to a `tempfile` root this function deletes on the way out, and
the happy path leaves nothing behind. The assemblers used to point it at the
deal's artefacts directory, which is cloud-synced: same code and the same warm
render cache measured 11.7 s for the first repair on local disk and 26.9 s on the
mount, and eight consecutive live attempts were killed at ~43 s without ever
converging. Bulk I/O to the mount is *not* slow (60 × 256 KB in 36 ms) — it is
per-file metadata and sync overhead, which is exactly what ~170 small files buys.

Failure is the one case where the renders are worth keeping, because they are the
evidence for *why* it would not fit. Pass `keep_on_failure=<artefacts>/.qa` and the
last pass's renders are copied there before the scratch root goes away, with the
path named in `DeckNotConvergedError`.

The font the ladder is calibrated against
-----------------------------------------
Every size this module chooses is measured in whatever face the render host
resolves `pptx_helpers.PALATINO` to, and the ladder was calibrated against real
Palatino Linotype. `font_probe.probe_font_resolution` runs once per converge and
the resolution goes into the stage log either way, so a silent substitution — the
worst case, and the one that leaves no trace on disk — becomes a warning instead
of a mystery.
"""

from __future__ import annotations

import shutil
import sys
import tempfile
from dataclasses import dataclass, field
from functools import partial
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Inches, Pt

from deck_contract import (
    Finding,
    LibraryBaseline,
    ProbeRequest,
    default_library_path,
    library_baseline,
    match_library_slide,
    measure_probe_overflow,
    verify_deck,
)
from font_probe import FontResolution, probe_font_resolution
from pptx_helpers import (
    apply_text_scale,
    enable_normal_autofit,
    normal_autofit_scale,
    set_table_height,
    strip_autofit,
)

# ─── What the loop may act on ────────────────────────────────────────────────

OVERFLOW_KINDS = ("rendered-overflow", "masked-overflow")
TABLE_HEIGHT_KIND = "table-taller-than-library"
REPAIRABLE_KINDS = (*OVERFLOW_KINDS, TABLE_HEIGHT_KIND)

# Autofit fontScale ladder, in percent. The 70% floor is inherited from the
# retired `fit_text_scale`: below it Palatino body copy stops being legible on a
# projected slide, so a deck needing more must lose words, not points — which is
# why exhausting the ladder fails the stage instead of shrinking further.
TEXT_SCALES = (100.0, 95.0, 90.0, 85.0, 80.0, 75.0, 70.0)

# Line-space reduction paired with a reduced scale, matching what the retired
# estimator wrote. Not modelled in the probe, so it only ever buys headroom.
_LINE_SPACE_REDUCTION = 8.0

# Points taken off the LARGEST body size in a table. Each step caps every body run
# at `max_size - k`, so the shrink lands where there is most to give and a column
# that is already small is left alone until the larger ones have caught up.
#
# That matters on the market-entry table, whose label column is 11 pt against 9 pt
# values: step 1 caps at 10 pt, taking the labels down and leaving the values at
# the 9 pt they were deliberately set to. Subtracting a point from every run
# instead would drop the values to 8 pt for a defect the labels caused — which is
# what the retired `_me_label_size_pt` avoided by measuring label widths against a
# Palatino advance-width table. Capping gets the same targeted outcome with no
# character metrics, and as a bonus keeps the column uniform rather than mixing
# 11 pt and 10 pt labels in one column.
#
# On a uniform body (the Considerations table's 10 pt, the broker table's 9 pt) a
# cap is identical to a subtraction.
TABLE_SIZE_DROPS = (0, 1, 2, 3)
_MIN_TABLE_PT = 6.0

# Ink this far below a box is render noise, not overflow: one pixel row at the
# contract's 150 dpi is 0.007", and antialiasing on a table's bottom border
# routinely lights one. Matches `deck_contract._OVERFLOW_TOL_IN`.
_ACCEPT_TOL_IN = 0.05

DEFAULT_MAX_ITERATIONS = 3

#: Prefix of the ephemeral scratch root, so a killed run's leftovers are
#: recognisable as ours in the system temp directory.
SCRATCH_PREFIX = "deck-repair-"


class DeckNotConvergedError(RuntimeError):
    """The deck still violates the contract after the repair budget is spent."""


@dataclass
class ConvergeResult:
    """Outcome of the loop, for the stage log and the `deck` checkpoint."""

    converged: bool
    iterations: int
    findings: list[Finding] = field(default_factory=list)
    actions: list[str] = field(default_factory=list)
    #: What the render host resolved the deck's primary latin typeface to. `None`
    #: only on a hand-built result; the loop always probes.
    font: FontResolution | None = None
    #: Where the failing pass's renders were kept, when the caller asked for that
    #: and the deck did not converge. `None` on every converged run — the happy
    #: path leaves no QA scratch anywhere.
    kept_dir: Path | None = None

    @property
    def blocking(self) -> list[Finding]:
        return [f for f in self.findings if f.blocking]

    @property
    def unrepaired(self) -> list[Finding]:
        """Blocking geometric findings that survived the loop."""
        return [f for f in self.blocking if f.kind in REPAIRABLE_KINDS]

    @property
    def advisory(self) -> list[Finding]:
        return [f for f in self.findings if not f.blocking]

    def summary(self) -> str:
        state = "converged" if self.converged else "DID NOT CONVERGE"
        return (
            f"{state} after {self.iterations} repair iteration(s); "
            f"{len(self.blocking)} blocking / {len(self.advisory)} advisory finding(s)"
        )


# ─── Targets ─────────────────────────────────────────────────────────────────


@dataclass
class _Target:
    """One shape to shrink, and the constraint it has to satisfy."""

    slide: int
    shape: str
    is_table: bool
    accept_in: float
    clamp_height: int | None = None  # EMU; set by a table-height finding

    @property
    def key(self) -> tuple[int, str]:
        return (self.slide, self.shape)


def _targets_from(findings, prs, baseline) -> list[_Target]:
    """Collapse the repairable findings to one target per shape.

    A shape can attract two findings at once — PRL18's table is both taller than
    the library and (once clamped) at risk of rendering past it — so the strictest
    acceptance threshold and any clamp height are merged into a single target.

    **The acceptance threshold comes from the library's ATTRIBUTED baseline, not
    from the finding's `limit_in`.** The two are different quantities and mixing
    them silently stalls the loop. A probe measures a shape's own ink with nothing
    masking it; `rendered-overflow` measures only ink no other shape has claimed,
    which stops at the first neighbour below. On the market-entry table those read
    0.59" and 0.18" for the same slide — so accepting the finding's 0.14" limit
    against a probe's 0.59" is unsatisfiable, and the loop spends its whole budget
    re-applying the deepest step. Both sides must be measured the same way.
    """
    targets: dict[tuple[int, str], _Target] = {}
    for finding in findings:
        if finding.kind not in REPAIRABLE_KINDS or not finding.shape:
            continue
        if finding.slide >= len(prs.slides):
            continue
        slide = prs.slides[finding.slide]
        shape = next((s for s in slide.shapes if s.name == finding.shape), None)
        if shape is None:
            continue
        is_table = bool(getattr(shape, "has_table", False))
        if not is_table and not getattr(shape, "has_text_frame", False):
            continue  # a picture cannot be made to fit by shrinking text

        matched = match_library_slide(slide, baseline.signatures)
        allowed = (
            baseline.attributed.get(matched, {}).get(finding.shape, 0.0)
            if matched is not None
            else 0.0
        )
        accept = allowed + _ACCEPT_TOL_IN
        clamp = (
            Inches(finding.limit_in)
            if finding.kind == TABLE_HEIGHT_KIND and finding.limit_in is not None
            else None
        )

        existing = targets.get((finding.slide, finding.shape))
        if existing is None:
            targets[(finding.slide, finding.shape)] = _Target(
                finding.slide, finding.shape, is_table, accept, clamp
            )
        else:
            existing.accept_in = min(existing.accept_in, accept)
            if clamp is not None:
                existing.clamp_height = (
                    clamp if existing.clamp_height is None else min(existing.clamp_height, clamp)
                )
    return sorted(targets.values(), key=lambda t: t.key)


# ─── Probe mutations ─────────────────────────────────────────────────────────


def _probe_text_at_scale(shape, scale: float) -> None:
    """Show a text shape as PowerPoint will draw it at `scale` percent."""
    apply_text_scale(shape, scale)
    strip_autofit(shape)


def _body_runs(table):
    """Every run in every cell below the header row."""
    for row in list(table.rows)[1:]:
        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    yield run


def _shrink_table_body(shape, drop_pt: float, clamp_height: int | None) -> None:
    """Cap the body cells at (largest size - `drop_pt`) and re-clamp the declared rows.

    A cap rather than a subtraction, so a mixed-size body shrinks its largest
    column first — see `TABLE_SIZE_DROPS`. Runs with no explicit size are left
    alone: nothing in this shape's XML says what they render at. Every INFOR table
    is written through `set_cell_text`, which always sets one.
    """
    table = shape.table
    sizes = [run.font.size.pt for run in _body_runs(table) if run.font.size is not None]
    if drop_pt and sizes:
        cap = max(_MIN_TABLE_PT, max(sizes) - drop_pt)
        for run in _body_runs(table):
            if run.font.size is not None and run.font.size.pt > cap:
                run.font.size = Pt(cap)
    target = clamp_height if clamp_height is not None else shape.height
    set_table_height(shape, target)


def _table_probe_bottom(shape, clamp_height: int | None) -> float:
    """The y a clamped table is aiming for, which is what its probe is measured against."""
    height = clamp_height if clamp_height is not None else shape.height
    return Emu(shape.top).inches + Emu(height).inches


# ─── One measured repair pass ────────────────────────────────────────────────


def _ladder(target: _Target, shape) -> list:
    """Candidate settings for a target, least-shrink first."""
    if target.is_table:
        return list(TABLE_SIZE_DROPS)
    current = normal_autofit_scale(shape)
    ceiling = 100.0 if current is None else current
    return [scale for scale in TEXT_SCALES if scale <= ceiling]


def _repair_pass(deck: Path, findings, work: Path, baseline) -> list[str]:
    """Measure every target's ladder in one render, apply the best step, save.

    Returns a human-readable action per shape actually changed; an empty list means
    nothing could be improved, which the caller reads as "not converging".
    """
    prs = Presentation(deck)
    targets = _targets_from(findings, prs, baseline)
    if not targets:
        return []

    requests: list[ProbeRequest] = []
    ladders: dict[tuple[int, str], list] = {}
    for target in targets:
        shape = next(s for s in prs.slides[target.slide].shapes if s.name == target.shape)
        steps = _ladder(target, shape)
        ladders[target.key] = steps
        for step in steps:
            if target.is_table:
                mutate = partial(
                    _shrink_table_body, drop_pt=step, clamp_height=target.clamp_height
                )
                bottom = _table_probe_bottom(shape, target.clamp_height)
            else:
                mutate = partial(_probe_text_at_scale, scale=step)
                bottom = None
            requests.append(
                ProbeRequest(
                    target.slide,
                    target.shape,
                    key=(target.slide, target.shape, step),
                    mutate=mutate,
                    bottom=bottom,
                )
            )

    depths = measure_probe_overflow(deck, requests, work)

    actions: list[str] = []
    for target in targets:
        steps = ladders[target.key]
        if not steps:
            continue
        fitting = [
            step
            for step in steps
            if depths.get((target.slide, target.shape, step), 0.0) <= target.accept_in
        ]
        # Least shrink that fits; if nothing fits, the deepest step, so the loop
        # has done everything it can before reporting failure.
        chosen = fitting[0] if fitting else steps[-1]
        measured = depths.get((target.slide, target.shape, chosen), 0.0)
        shape = next(s for s in prs.slides[target.slide].shapes if s.name == target.shape)

        if target.is_table:
            if chosen == 0 and target.clamp_height is None:
                continue  # nothing to change
            _shrink_table_body(shape, chosen, target.clamp_height)
            clamped = (
                f", clamped to {Emu(target.clamp_height).inches:.3f}\""
                if target.clamp_height is not None
                else ""
            )
            actions.append(
                f"slide {target.slide + 1} {target.shape!r}: body cells "
                f"-{chosen:g} pt{clamped} (measured {measured:.3f}\" overflow, "
                f"accepting <= {target.accept_in:.3f}\")"
            )
        else:
            current = normal_autofit_scale(shape)
            if current is not None and chosen >= current:
                continue  # already at or below the best the ladder offers
            enable_normal_autofit(
                shape,
                font_scale=chosen,
                line_space_reduction=0.0 if chosen >= 100.0 else _LINE_SPACE_REDUCTION,
            )
            actions.append(
                f"slide {target.slide + 1} {target.shape!r}: autofit fontScale "
                f"{'unset' if current is None else f'{current:g}%'} -> {chosen:g}% "
                f"(measured {measured:.3f}\" overflow, accepting <= {target.accept_in:.3f}\")"
            )

    if actions:
        prs.save(deck)
    return actions


# ─── Ephemeral scratch ───────────────────────────────────────────────────────


def _latest_pass(root: Path, prefix: str) -> Path | None:
    """The highest-numbered `<prefix>-N` directory under `root`, if any.

    Read off disk rather than computed from the iteration count, because the loop
    breaks *before* re-verifying when a pass could apply nothing — so the newest
    verify directory is not always `verify-<iterations>`.
    """
    numbered: list[tuple[int, Path]] = []
    for path in root.glob(f"{prefix}-*"):
        suffix = path.name.rsplit("-", 1)[-1]
        if suffix.isdigit() and path.is_dir():
            numbered.append((int(suffix), path))
    if not numbered:
        return None
    return max(numbered, key=lambda pair: pair[0])[1]


def _keep_failing_artefacts(root: Path, dest: Path | str) -> Path | None:
    """Copy the LAST verify/repair pass out of the scratch root, for the analyst.

    Only the last pass. Its renders are the ones showing the overflow that
    survived, and copying every pass is how ~170 files and ~10 MB per deck ended
    up in the deal directory permanently in the first place.
    """
    dest = Path(dest)
    kept = False
    for prefix in ("verify", "repair"):
        src = _latest_pass(root, prefix)
        if src is None:
            continue
        try:
            shutil.copytree(src, dest / src.name, dirs_exist_ok=True)
            kept = True
        except OSError:
            pass  # a kept render is a diagnostic, never the reason a stage fails
    return dest if kept else None


# ─── The loop ────────────────────────────────────────────────────────────────


def converge_deck(
    deck: Path | str,
    *,
    library: Path | str | None = None,
    max_iterations: int = DEFAULT_MAX_ITERATIONS,
    out_dir: Path | str | None = None,
    keep_on_failure: Path | str | None = None,
    vision: bool = False,
    log=None,
) -> ConvergeResult:
    """Verify the deck, repair what can be repaired, and re-verify.

    Repairs the deck **in place**. `vision=False` keeps the advisory tier out of
    the loop's renders; the `deck` checkpoint runs it once on the converged file.

    Returns the result rather than raising, so a caller can log it and decide.
    `assert_converged` is the raising form the assemblers use.

    `out_dir` is the QA scratch root. **Leave it unset in production**: the
    default is a `tempfile` directory that this function deletes before
    returning, which is what keeps ~10 MB of renders per deck out of the
    analyst's cloud-synced deal directory (see the module docstring). Pass
    `keep_on_failure` — the assemblers pass `<artefacts>/.qa` — and the last
    pass's renders survive into it *only* when the deck does not converge, with
    the path carried on the result and named in `DeckNotConvergedError`.

    Pass an explicit `out_dir` when you want the whole tree to persist (tests, a
    hands-on debugging session) or when `vision=True`: the advisory findings name
    render paths, and an ephemeral root deletes what they point at.
    """
    deck = Path(deck)
    emit = log if log is not None else partial(print, file=sys.stderr)
    ephemeral = out_dir is None
    root = Path(out_dir) if out_dir else Path(tempfile.mkdtemp(prefix=SCRATCH_PREFIX))
    try:
        result = _converge_into(
            deck,
            root,
            library=library,
            max_iterations=max_iterations,
            vision=vision,
            emit=emit,
        )
        if not result.converged and keep_on_failure is not None:
            result.kept_dir = _keep_failing_artefacts(root, keep_on_failure)
            if result.kept_dir is not None:
                emit(
                    f"deck_repair: kept the failing pass's renders in "
                    f"{result.kept_dir}"
                )
        return result
    finally:
        if ephemeral:
            shutil.rmtree(root, ignore_errors=True)


def _converge_into(
    deck: Path,
    root: Path,
    *,
    library: Path | str | None,
    max_iterations: int,
    vision: bool,
    emit,
) -> ConvergeResult:
    """The loop proper, writing every render under `root`."""
    font = probe_font_resolution()
    emit(f"deck_repair: {font.log_line()}")

    def _verify(tag: str) -> list[Finding]:
        return verify_deck(deck, library=library, vision=vision, out_dir=root / tag)

    findings = _verify("verify-0")
    # Same cached baseline `verify_deck` just used; the repair pass needs its
    # attributed map to set acceptance thresholds the probes can actually meet.
    library_path = Path(library) if library else default_library_path()
    baseline = (
        library_baseline(library_path) if library_path and library_path.is_file()
        else LibraryBaseline()
    )
    actions: list[str] = []
    iterations = 0
    while iterations < max_iterations:
        repairable = [f for f in findings if f.blocking and f.kind in REPAIRABLE_KINDS]
        if not repairable:
            break
        iterations += 1
        applied = _repair_pass(deck, repairable, root / f"repair-{iterations}", baseline)
        for action in applied:
            emit(f"deck_repair: {action}")
        actions.extend(applied)
        if not applied:
            emit(
                f"deck_repair: {len(repairable)} blocking geometric finding(s) remain and "
                f"no further shrink is available — the content is over budget, not the "
                f"layout"
            )
            break
        findings = _verify(f"verify-{iterations}")

    result = ConvergeResult(
        converged=not [f for f in findings if f.blocking and f.kind in REPAIRABLE_KINDS],
        iterations=iterations,
        findings=findings,
        actions=actions,
        font=font,
    )
    emit(f"deck_repair: {result.summary()}")
    # And WHAT they were. The summary's counts are the headline, but "18 blocking
    # finding(s)" on its own is unactionable: the stage that reads this log — since
    # v0.5.51 the driver captures it to `stages/<id>/log.txt` — has no other route to
    # the list, because a converged deck deletes its own QA scratch on the way out.
    # One line each, `Finding.__str__`'s wording, so there is no second rendering.
    for finding in result.findings:
        emit(f"deck_repair:   {finding}")
    return result


def assert_converged(deck: Path | str, result: ConvergeResult) -> None:
    """Raise `DeckNotConvergedError` when the deck still breaks the contract.

    The stage fails here rather than returning a deck the analyst has to catch.
    The message carries every surviving finding and every repair that was tried,
    because the remedy is usually editorial — shorter mitigants, fewer overview
    bullets — and the author needs to know which shape and by how much.

    It also names the kept renders and, when the render host substituted the
    typeface, says so: a deck that will not fit measured in the wrong metrics is
    not necessarily over budget in the right ones, so that has to be ruled out
    before anyone rewrites copy.
    """
    if result.converged:
        return
    lines = [
        f"{Path(deck).name} did not converge against the deck contract after "
        f"{result.iterations} repair iteration(s).",
        "",
        "Surviving findings:",
        *(f"  {finding}" for finding in result.unrepaired),
    ]
    if result.actions:
        lines += ["", "Repairs applied:", *(f"  {action}" for action in result.actions)]
    if result.kept_dir is not None:
        lines += [
            "",
            f"The failing pass's renders were kept here: {result.kept_dir}",
            "(the rest of the QA scratch was ephemeral and has been removed)",
        ]
    if result.font is not None and not result.font.ok:
        lines += ["", result.font.log_line()]
    lines += [
        "",
        "Every automatic remedy shrinks text, and the ladder is spent. The content "
        "is over budget for the box: shorten the copy at the content stage.",
    ]
    raise DeckNotConvergedError("\n".join(lines))
