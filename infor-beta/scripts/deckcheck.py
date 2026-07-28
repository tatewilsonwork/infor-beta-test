"""Falsification pass over a built deck — Phase G's `deckcheck` stage.

`deck_contract` asks whether the deck *looks* right. This asks whether its
**numbers are true**, and it asks by trying to disprove them: for every figure on
a slide, find the provenance record that claims to support it, open the filing that
record names, and look. A figure that survives that is confirmed; one that does
not is a finding.

The split of labour is the same one Phase B settled on, for the same reason.

**Mechanical, here.** Pull every figure out of the deck's text shapes and table
cells, normalise it, and join it to the run's provenance ledger
(`provenance.read_run_provenance`). That produces an *agenda*: figures that trace
to a record (verify the record), figures that trace to nothing (find a source or
report there isn't one), and rasterised pictures no string scan can read at all.

**Judgement, in the SKILL.md.** Whether the cited page actually says 4,520.3,
whether the period label matches the statement it came from, whether a bridge's
components are the right three quarters. A model has to read the filing for that.

Advisory, always
----------------
Every finding here is `SEVERITY_ADVISORY` and :class:`CheckFinding` refuses to be
constructed otherwise. This is a review pass surfaced at a checkpoint, not a gate:
the plans' one `required` gate is on `deck`, and it stays there. A falsification
pass that could halt a run would be a pass that has to be right about a target's
financial statements, and nothing here is that confident.

What is NOT a finding
---------------------
Error values in CapIQ-dependent cells are the **normal** state of a shipped
artefact in this environment — CapIQ cannot be refreshed here, so the forward
estimates, the comps and precedents array formulas, and the pre-resolution LTM
link all ship un-evaluated for the analyst to refresh in Excel.
:data:`EXPECTED_ERROR_CONTEXTS` is that list, and it is rendered into the agenda
and the report rather than only stated in the skill's prose, so it sits in front of
whoever is reading at the moment they would otherwise report one. Re-flagging them
is how a review gets ignored.

The library is the baseline for "whose figure is this?"
------------------------------------------------------
A pitch deck carries three static credential slides whose tombstone values are
INFOR's own and have nothing to do with this target — plus a disclaimer and a
contact page. They have no provenance in this run and never will, so reporting them
would bury the real findings under boilerplate.

They are excluded the same way `deck_contract` excludes the library's own
overhanging footnotes: **a figure that already appears in the blank slide library,
on a shape of the same name, is not this run's figure.** Static slides are copied
from the library verbatim, so they drop out automatically; a filled slide's figures
cannot match, because what the library holds there is a `[x]` token. No slide list
to migrate when the library gains an entry — the same self-maintaining property
Phase C bought everywhere else.
"""

from __future__ import annotations

import re
from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from deck_contract import default_library_path, write_picture_crops
from provenance import FigureProvenance, ProvenanceLedger

SEVERITY_ADVISORY = "advisory"

#: The verdicts a reviewer may return for one figure.
#:
#: `contradicted` is the one that matters — the source was read and says something
#: else. `unsupported` means no source could be found at all; `unverifiable` means
#: the source was named but not available to read (an un-attached filing, a CapIQ
#: value that cannot be refreshed here), which is a different and much weaker
#: statement than "wrong".
VERDICTS = ("confirmed", "contradicted", "unsupported", "unverifiable")

#: Error values that are correct behaviour in a shipped artefact, never defects.
#: Rendered into the agenda and the report verbatim (:func:`expected_error_note`)
#: so the rule travels with the evidence instead of living only in SKILL.md prose.
EXPECTED_ERROR_CONTEXTS = (
    "The cap table's forward-estimate columns — the CapIQ `SP_REV_EST` / "
    "`SP_EBITDA_EST` calls and the EV/metric rows that divide by them. CapIQ cannot "
    "be refreshed in this environment, so `#VALUE!` / `n/a` there is the normal "
    "state of a shipped cap table, and it is visible in the rasterised cap-table "
    "picture on the overview / earnings-summary slide.",
    "The comps tab's CapIQ array formulas, which ship un-evaluated for the analyst "
    "to refresh in Excel — which is also why the deck's comps slide is a placeholder.",
    "The precedents tab's column-C FX array formula and every converted TEV, "
    "$-metric ratio and statistic row that depends on it.",
    "The financial-summary tab's LTM link before the `ltm-metrics` tab is written.",
    "A `[Placeholder for …]` / `[Pie Chart Placeholder]` region that the plan "
    "deliberately defers — `deck_contract` owns unsubstituted tokens, and it has "
    "already run at the `deck` stage.",
)


def expected_error_note() -> str:
    """The do-not-report list as a markdown block."""
    lines = ["**Not defects — do not report these:**", ""]
    lines += [f"- {item}" for item in EXPECTED_ERROR_CONTEXTS]
    return "\n".join(lines)


# ─── Figure extraction ───────────────────────────────────────────────────────
#
# A figure is worth auditing when it is written like a figure: it carries a
# currency symbol, a magnitude suffix, a percent or a multiple, or it is a decimal
# / thousands-separated number. Bare integers are excluded, and that exclusion is
# doing real work — a deck is full of them (slide numbers, "Q3 2026", "2024 Annual
# Report", years founded, "10-K") and not one is a claim about the target's
# financials. The cost is a bare share count in a tile whose scale lives in the
# label; the benefit is an agenda a human will actually read.

_SCALE_TO_MILLIONS = {
    "MM": 1.0,
    "M": 1.0,  # "$500M senior notes" — this desk means millions
    "B": 1000.0,
    "bn": 1000.0,
    "K": 0.001,
    "k": 0.001,
}

_KIND_CURRENCY = "currency"
_KIND_PERCENT = "percent"
_KIND_MULTIPLE = "multiple"
_KIND_PLAIN = "plain"

_FIGURE_RE = re.compile(
    r"""
    (?<![A-Za-z0-9])                                  # not mid-identifier
    (?P<open>\()?
    (?P<sign>-|−)?                               # hyphen-minus or U+2212
    (?P<currency>US\$|C\$|CAD\$|A\$|\$|€|£|¥)?
    \s?
    (?P<number>\d{1,3}(?:,\d{3})+(?:\.\d+)?|\d+(?:\.\d+)?)
    \s?
    (?:(?P<scale>MM|bn|B|M|K|k|%|x|X)(?![A-Za-z]))?
    (?P<close>\))?
    """,
    re.VERBOSE,
)


@dataclass(frozen=True)
class DeckFigure:
    """One number as the deck writes it, plus where it sits and what it means.

    ``slide`` is zero-based, matching `deck_contract.Finding`; add 1 for the number
    PowerPoint shows. ``millions`` is the currency value normalised to millions —
    the scale the deal workbook is locked to — so a deck rendering ``$1.2B`` and a
    record holding ``1200.0`` compare directly. None for a percent or a multiple.
    """

    slide: int
    shape: str
    raw: str
    value: float
    kind: str
    decimals: int
    context: str
    scale: str | None = None
    currency: str | None = None
    millions: float | None = None


def _iter_shapes(shapes):
    """Every shape, flattening groups (mirrors `deck_contract._iter_shapes`)."""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes)
        else:
            yield shape


def _shape_texts(slide):
    """(shape_name, text) for every text-bearing shape and table cell on a slide.

    Same scope as `deck_contract`'s error scan, and deliberately so: a rasterised
    range picture is not reached, because an error value inside one is usually
    correct (see :data:`EXPECTED_ERROR_CONTEXTS`) and its figures need a reader,
    not a regex. Those become :attr:`DeckAudit.pictures` entries instead.
    """
    for shape in _iter_shapes(slide.shapes):
        if shape.has_text_frame and shape.text_frame.text.strip():
            yield shape.name, shape.text_frame.text
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        yield shape.name, cell.text


def _decimals(number_text: str) -> int:
    return len(number_text.split(".")[1]) if "." in number_text else 0


def _classify(currency: str | None, scale: str | None, number_text: str) -> str | None:
    """The figure's kind, or None when it is not worth auditing."""
    if scale == "%":
        return _KIND_PERCENT
    if scale in ("x", "X"):
        return _KIND_MULTIPLE
    if currency or scale in _SCALE_TO_MILLIONS:
        return _KIND_CURRENCY
    if "." in number_text or "," in number_text:
        return _KIND_PLAIN
    return None


#: Characters of surrounding text kept either side of a figure. A window, not the
#: whole shape: a bullet block holds a dozen figures, and repeating its full text
#: once per figure made the agenda table unreadable — which is the same failure as
#: not reporting at all.
_CONTEXT_RADIUS = 70


def _context(text: str, start: int, end: int) -> str:
    """The sentence around a figure, collapsed to one line and window-trimmed."""
    left = max(0, start - _CONTEXT_RADIUS)
    right = min(len(text), end + _CONTEXT_RADIUS)
    window = " ".join(text[left:right].split())
    return ("…" if left > 0 else "") + window + ("…" if right < len(text) else "")


def _figures_in(text: str, *, slide: int, shape: str):
    """Yield every auditable :class:`DeckFigure` in one run of text."""
    for match in _FIGURE_RE.finditer(text):
        number_text = match.group("number")
        currency = match.group("currency")
        scale = match.group("scale")
        kind = _classify(currency, scale, number_text)
        if kind is None:
            continue
        magnitude = float(number_text.replace(",", ""))
        negative = bool(match.group("sign")) or bool(match.group("open") and match.group("close"))
        value = -magnitude if negative else magnitude
        factor = _SCALE_TO_MILLIONS.get(scale or "", 1.0) if kind == _KIND_CURRENCY else None
        yield DeckFigure(
            slide=slide,
            shape=shape,
            raw=match.group(0).strip(),
            value=value,
            kind=kind,
            decimals=_decimals(number_text),
            context=_context(text, match.start(), match.end()),
            scale=scale,
            currency=currency,
            millions=None if factor is None else value * factor,
        )


def _library_keys(library: Path | str | None) -> set[tuple[str, str]]:
    """`(shape name, figure text)` pairs that exist in the blank slide library.

    Anything matching one of these came with the library — a credential
    tombstone, a disclaimer figure — not from this run's data. Returns an empty
    set when the library cannot be resolved, which means boilerplate is *reported*
    rather than silently dropped: over-reporting is recoverable, and a silent
    filter that failed open would be a review that quietly stopped covering the
    filled slides too.
    """
    if library is None:
        return set()
    path = Path(library)
    if not path.is_file():
        return set()
    keys: set[tuple[str, str]] = set()
    for slide in Presentation(path).slides:
        for shape_name, text in _shape_texts(slide):
            for figure in _figures_in(text, slide=0, shape=shape_name):
                keys.add((shape_name, figure.raw))
    return keys


def extract_deck_figures(
    deck: Path | str,
    *,
    library: Path | str | None = None,
) -> list[DeckFigure]:
    """Every auditable figure in the deck's text shapes and table cells.

    Figures the blank slide library also carries on a shape of the same name are
    excluded as boilerplate (see the module docstring). ``library`` defaults to the
    shipped library via `deck_contract.default_library_path`.
    """
    path = Path(deck)
    if not path.is_file():
        raise FileNotFoundError(f"deck not found: {path}")
    reference = library if library is not None else default_library_path()
    boilerplate = _library_keys(reference)

    figures: list[DeckFigure] = []
    for index, slide in enumerate(Presentation(path).slides):
        for shape_name, text in _shape_texts(slide):
            for figure in _figures_in(text, slide=index, shape=shape_name):
                if (shape_name, figure.raw) in boilerplate:
                    continue
                figures.append(figure)
    return figures


def _picture_names(deck: Path | str) -> set[str]:
    """Shape names of every picture in a deck — the library baseline for pictures."""
    names: set[str] = set()
    for slide in Presentation(Path(deck)).slides:
        for shape in _iter_shapes(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                names.add(shape.name)
    return names


def deck_pictures(
    deck: Path | str,
    *,
    library: Path | str | None = None,
) -> list[tuple[int, str]]:
    """`(zero-based slide, shape name)` for every picture THIS RUN put on the deck.

    A pasted range or a rendered chart — the figures no string scan reaches, so
    the reviewer reads the crop (:func:`write_evidence`) instead.

    The library baseline applies here too, and it earns its keep: the pitch deck
    carries 49 pictures, of which 44 are the library's own logos and decorative
    graphics. A run's pastes are the ones whose shape name is not a library picture
    — the cap-table range, the Financial Summary charts, the ownership blocks —
    which is 5, and a list of 5 gets read where a list of 49 does not.
    """
    reference = library if library is not None else default_library_path()
    boilerplate = (
        _picture_names(reference) if reference is not None and Path(reference).is_file() else set()
    )
    out: list[tuple[int, str]] = []
    for index, slide in enumerate(Presentation(Path(deck)).slides):
        for shape in _iter_shapes(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and shape.name not in boilerplate:
                out.append((index, shape.name))
    return out


# ─── Joining the deck to the ledger ──────────────────────────────────────────

#: Floor on the match tolerance, in the figure's own units. A deck figure is a
#: rounded rendering of a workbook value, so the tolerance is normally half the
#: last shown decimal place; this floor covers a whole-number rendering, and the
#: relative term below covers rounding inside a derived record's own arithmetic.
_TOLERANCE_FLOOR = 0.05
_TOLERANCE_RELATIVE = 0.001


def _tolerance(figure: DeckFigure) -> float:
    half_place = 0.5 * (10.0**-figure.decimals)
    if figure.kind == _KIND_CURRENCY and figure.scale in _SCALE_TO_MILLIONS:
        half_place *= _SCALE_TO_MILLIONS[figure.scale]
    reference = abs(figure.millions if figure.millions is not None else figure.value)
    return max(half_place, _TOLERANCE_FLOOR, reference * _TOLERANCE_RELATIVE)


def _candidate_values(figure: DeckFigure, record: FigureProvenance) -> list[float]:
    """The record values worth comparing this figure against.

    A record holding an Excel formula (a combined figure written as
    ``"=9000+800"``) has no number to compare — openpyxl never evaluates, so the
    cached value does not exist here either. Such a record can still be *cited* by
    a reviewer; it just cannot be auto-matched.
    """
    value = record.value
    if not isinstance(value, (int, float)) or isinstance(value, bool):
        return []
    units = (record.units or "").lower()
    if figure.kind == _KIND_PERCENT:
        # A percent record may be stored as 15.3 or as 0.153.
        return [float(value), float(value) * 100.0] if "%" in units else []
    if figure.kind == _KIND_MULTIPLE:
        return [float(value)] if "x" in units else []
    return [float(value)]


def _agrees(figure: DeckFigure, record: FigureProvenance) -> bool:
    target = figure.millions if figure.millions is not None else figure.value
    tolerance = _tolerance(figure)
    return any(abs(target - candidate) <= tolerance for candidate in _candidate_values(figure, record))


@dataclass(frozen=True)
class FigureMatch:
    """One deck figure and the provenance record that supports it, if any."""

    figure: DeckFigure
    record: FigureProvenance | None = None

    @property
    def traced(self) -> bool:
        return self.record is not None


@dataclass
class DeckAudit:
    """The machine half of a deck check: the join, and what could not be joined."""

    deck: Path
    matches: list[FigureMatch] = field(default_factory=list)
    pictures: list[tuple[int, str]] = field(default_factory=list)
    ledger: ProvenanceLedger = field(default_factory=ProvenanceLedger)

    @property
    def traced(self) -> list[FigureMatch]:
        return [m for m in self.matches if m.traced]

    @property
    def untraced(self) -> list[FigureMatch]:
        return [m for m in self.matches if not m.traced]

    @property
    def slides(self) -> tuple[int, ...]:
        return tuple(sorted({m.figure.slide for m in self.matches}))


def audit_deck(
    deck: Path | str,
    ledger: ProvenanceLedger,
    *,
    library: Path | str | None = None,
) -> DeckAudit:
    """Join every figure on the deck to the run's provenance ledger.

    Purely deterministic and rendering-free — no LibreOffice, no model. A figure
    is `traced` when some record's value agrees with it within the tolerance its
    own rendering implies; the reviewer still has to read the record's filing,
    because agreeing with a record only means the deck copied the workbook
    faithfully, not that the workbook is right.
    """
    figures = extract_deck_figures(deck, library=library)
    matches = [
        FigureMatch(figure, next((r for r in ledger.figures if _agrees(figure, r)), None))
        for figure in figures
    ]
    return DeckAudit(
        deck=Path(deck),
        matches=matches,
        pictures=deck_pictures(deck, library=library),
        ledger=ledger,
    )


# ─── Evidence ────────────────────────────────────────────────────────────────


def write_evidence(deck: Path | str, out_dir: Path | str) -> tuple[dict[int, Path], list[tuple[int, str, Path]]]:
    """Render every slide to PNG and extract every embedded picture.

    Returns `({slide index: png}, [(slide index, shape name, crop)])`. The renders
    are what the reviewer reads figures off; the crops carry a pasted range at its
    native resolution, which 150 dpi of slide does not (`write_picture_crops`).

    Raises `RuntimeError` when LibreOffice is absent — the caller must say the
    review could not run rather than report a clean deck. The picture crops need no
    renderer, so they are written either way.
    """
    from slide_render import render_deck_to_png  # local: keeps import cost off the audit path

    deck = Path(deck)
    root = Path(out_dir)
    root.mkdir(parents=True, exist_ok=True)
    crops = write_picture_crops(Presentation(deck), root / "pictures")
    pngs = render_deck_to_png(deck, root / "slides")
    renders = {int(p.stem.split("_")[-1]) - 1: p for p in pngs}
    return renders, crops


# ─── Findings and the report ─────────────────────────────────────────────────


@dataclass(frozen=True)
class CheckFinding:
    """One reviewer verdict on one figure. Advisory by construction.

    ``severity`` exists only to be pinned: `deckcheck` surfaces a review at a
    checkpoint and never gates a run, so a finding that claimed to be blocking
    would be a category error. Constructing one raises.
    """

    slide: int
    figure: str
    verdict: str
    detail: str
    source: str | None = None
    severity: str = SEVERITY_ADVISORY

    def __post_init__(self) -> None:
        if self.verdict not in VERDICTS:
            raise ValueError(
                f"unknown verdict {self.verdict!r}; expected one of {', '.join(VERDICTS)}"
            )
        if self.severity != SEVERITY_ADVISORY:
            raise ValueError(
                f"a deckcheck finding is always {SEVERITY_ADVISORY!r} — this stage is a "
                f"review surfaced at a checkpoint, not a gate. The run's only required "
                f"gate is on `deck`."
            )

    @property
    def slide_number(self) -> int:
        """The 1-based number the analyst sees in PowerPoint."""
        return self.slide + 1


def _figure_row(match: FigureMatch) -> str:
    figure = match.figure
    record = match.record
    where = f"{figure.slide + 1}"
    if record is None:
        return f"| {where} | `{figure.raw}` | {figure.shape} | — | {figure.context} |"
    citation = "; ".join(record.citation_lines) or (record.derivation or "—")
    return (
        f"| {where} | `{figure.raw}` | {figure.shape} | {record.figure}"
        f"{f' ({record.location})' if record.location else ''} | {citation} |"
    )


def render_agenda(audit: DeckAudit) -> str:
    """The machine half of the report: what to check, and what not to report.

    Written for a reader (the checkpoint agent, or the analyst) rather than for a
    parser. The untraced figures come first because they are where the work is.
    """
    lines = [
        "## Agenda",
        "",
        f"- Deck: `{audit.deck}`",
        f"- Figures found: {len(audit.matches)} across {len(audit.slides)} slide(s)",
        f"- Provenance records: {len(audit.ledger)}"
        + (f" from stage(s) {', '.join(audit.ledger.stages)}" if audit.ledger.stages else ""),
        f"- Traced to a record: {len(audit.traced)}; untraced: {len(audit.untraced)}",
        f"- Rasterised pictures (no string scan reaches these): {len(audit.pictures)}",
        "",
    ]

    lines += [
        f"### Untraced figures — {len(audit.untraced)}",
        "",
        "No provenance record supports these. Find the source and confirm the value,",
        "or record that there is none. Analyst-supplied figures (a valuation range, a",
        "market size from the CIM) legitimately land here — say so and name where they",
        "came from.",
        "",
    ]
    if audit.untraced:
        lines += ["| Slide | Figure | Shape | Record | Source |", "|---|---|---|---|---|"]
        lines += [_figure_row(m) for m in audit.untraced]
    else:
        lines.append("(none)")
    lines.append("")

    lines += [
        f"### Traced figures — {len(audit.traced)}",
        "",
        "A record's value agrees with what the slide shows. That proves the deck copied",
        "the workbook faithfully — NOT that the workbook is right. Open each record's",
        "filing at the statement and page it names, and try to disprove the figure.",
        "",
    ]
    if audit.traced:
        lines += ["| Slide | Figure | Shape | Record | Source |", "|---|---|---|---|---|"]
        lines += [_figure_row(m) for m in audit.traced]
    else:
        lines.append("(none)")
    lines.append("")

    lines += [
        f"### Rasterised pictures — {len(audit.pictures)}",
        "",
        "Read the native-resolution crop, not the slide render — a pasted range carries",
        "far more pixels than 150 dpi of slide gives it.",
        "",
    ]
    lines += (
        [f"- slide {index + 1}: `{name}`" for index, name in audit.pictures]
        if audit.pictures
        else ["(none)"]
    )
    lines += ["", expected_error_note(), ""]
    return "\n".join(lines)


def render_report(
    audit: DeckAudit,
    findings: "list[CheckFinding] | tuple[CheckFinding, ...]" = (),
    *,
    company: str,
    provenance_path: Path | str | None = None,
    notes: "list[str] | tuple[str, ...]" = (),
) -> str:
    """The analyst-facing deck check: verdicts first, then the agenda behind them."""
    by_verdict = {v: [f for f in findings if f.verdict == v] for v in VERDICTS}
    lines = [
        f"# Deck check — {company}",
        "",
        "**Advisory review, not a gate.** Every item below is a figure someone should",
        "look at before this deck leaves the building; nothing here halts a run, and the",
        "plan's one required approval is on the `deck` stage, which has already passed.",
        "",
        f"- Deck: `{audit.deck}`",
        f"- Provenance record: `{provenance_path}`" if provenance_path else "- Provenance record: (not written)",
        f"- Figures checked: {len(audit.matches)} ({len(audit.traced)} traced, "
        f"{len(audit.untraced)} untraced)",
        "- Verdicts: "
        + ", ".join(f"{len(by_verdict[v])} {v}" for v in VERDICTS),
        "",
        "## Findings",
        "",
    ]
    if findings:
        lines += ["| Slide | Figure | Verdict | Detail | Source |", "|---|---|---|---|---|"]
        # Contradictions first — they are the only verdict that says a number is wrong.
        for verdict in VERDICTS:
            for f in sorted(by_verdict[verdict], key=lambda f: (f.slide, f.figure)):
                lines.append(
                    f"| {f.slide_number} | {f.figure} | **{f.verdict}** | {f.detail} "
                    f"| {f.source or '—'} |"
                )
    else:
        lines.append("(no figure could be faulted)")
    lines.append("")
    if notes:
        lines += ["## Notes", ""] + [f"- {n}" for n in notes] + [""]
    lines.append(render_agenda(audit))
    return "\n".join(lines)


def write_report(
    path: Path | str,
    audit: DeckAudit,
    findings: "list[CheckFinding] | tuple[CheckFinding, ...]" = (),
    *,
    company: str,
    provenance_path: Path | str | None = None,
    notes: "list[str] | tuple[str, ...]" = (),
) -> Path:
    """Render the report and write it. Returns the path."""
    target = Path(path)
    target.parent.mkdir(parents=True, exist_ok=True)
    target.write_text(
        render_report(
            audit,
            findings,
            company=company,
            provenance_path=provenance_path,
            notes=notes,
        ),
        encoding="utf-8",
    )
    return target
