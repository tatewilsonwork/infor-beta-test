"""Earnings-update deck assembler built on the shared INFOR slide library.

The earnings update no longer ships its own template. It clones the relevant
entries out of `INFOR Slide Library.pptx` — cover (1), public-company overview
(7), earnings summary (8), disclaimer (15), contact (16) — into a fresh
five-slide deck and fills them from a typed `EarningsUpdateContent` bundle.
"""

from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

from deal_workbook import TAB_CAPTABLE
from deck_repair import assert_converged, converge_deck
from excel_to_powerpoint import (
    assert_range_pictures_are_distinct,
    insert_excel_into_placeholder,
)
from naming import safe_filename
from pptx_helpers import (
    COLOR_DOWN,
    COLOR_UP,
    delete_slide,
    fill_footnote_currency,
    find_shape,
    find_shape_in_group,
    find_table_shape,
    fit_overview_textbox,
    iter_all_shapes,
    set_cell_text,
    set_text,
    write_bullets_or_plain,
)
from schemas import EarningsUpdateContent, SlidePlan
from template_layout import (
    CAP_TABLE_PICTURE_NAMES,
    CAP_TABLE_PICTURE_RANGE,
    MARKER_CONTACT,
    MARKER_COVER,
    MARKER_DISCLAIMER,
    MARKER_EARNINGS_SUMMARY,
    MARKER_OVERVIEW,
    NAME_CAP_PICTURE_RANGE,
    find_slide_by_marker,
    resolve_workbook_range,
    verify_workbook_names,
)

# The five library entries the earnings deck keeps, in final deck order. Their
# INDICES are discovered per-run with `find_slide_by_marker`, not written down:
# through v0.5.39 this was the literal tuple `(0, 6, 7, 15, 16)`, and it had
# needed three manual migrations — (0,6,7,13,14) before v0.5.8's ownership
# slide, (0,6,7,14,15) before v0.5.14's precedents slide — because every
# insertion ahead of the closers silently shifted them. Locating each slide by
# its marker shape ends that: an inserted library slide moves the indices and
# nothing here changes.
_KEEP_MARKERS = (
    MARKER_COVER,
    MARKER_OVERVIEW,
    MARKER_EARNINGS_SUMMARY,
    MARKER_DISCLAIMER,
    MARKER_CONTACT,
)

# Earnings-summary slide cap-table placeholder. The picture range is resolved
# from the workbook's `infor_cap_picture_range` defined name (shared with the
# pitch assembler — same picture); the constant is only the shipped address, and
# the pre-flight above requires the name, so it is not reachable in practice.
_CAP_TABLE_PLACEHOLDER = "Rectangle 3"


def cap_table_picture_range(workbook_path) -> str:
    """The cap table's slide-picture range, by name.

    ``workbook_path`` is the **deal workbook**, so the sheet is ``TAB_CAPTABLE``
    and not the source template's ``Cap with Links`` — see the sheet-name note in
    `template_layout`.
    """
    return resolve_workbook_range(
        workbook_path,
        sheet=TAB_CAPTABLE,
        name=NAME_CAP_PICTURE_RANGE,
        fallback=CAP_TABLE_PICTURE_RANGE,
    )

# Section header that closes the Business Updates band on the earnings-summary
# slide — the bullets above it are sized to stop short of it.
_BROKER_TABLE_BAND = "Broker Estimates"


def _bullet_tuple(bullet) -> tuple[str, int]:
    text = f"{bullet.bold_prefix or ''}{bullet.text}"
    return (text, bullet.level)


def _strip_currency_unit(label: str) -> str:
    """Strip "(C$MM)" / "(US$MM)" / "(MM)" suffixes from labels.

    Leaves non-MM markers such as "(US$)" and "(C$)" intact — per-share
    metrics need them since the table header announces the MM scope only.
    """
    return re.sub(r"\s*\([A-Z]{0,3}\$?MM\)\s*$", "", label).strip()


def _money(value: str) -> str:
    """Prefix a broker-table figure with `$`, respecting sign/paren convention.

    Leaves percent values and already-prefixed values untouched.
    """
    s = value.strip()
    if not s or s[0] == "$" or "%" in s:
        return s
    if s.startswith("(") and s.endswith(")"):
        return f"(${s[1:-1].lstrip('-').strip()})"
    if s[0] in "+-":
        return f"{s[0]}${s[1:].strip()}"
    return f"${s}"


def _fmt_mm(value: str) -> str:
    """Format a plain figure (in MM) for a financial-highlights tile.

    Convention: whole millions as ``$XMM`` (no decimals); ``>= 1,000`` MM as
    ``$X.XB`` (one decimal). Percent values are returned untouched so a non-
    dollar KPI tile is not mangled. Unparseable input is returned as-is.
    """
    s = value.strip()
    if not s or "%" in s:
        return s
    neg = s.startswith("(") and s.endswith(")")
    raw = s[1:-1] if neg else s
    raw = raw.lstrip("+-").replace("$", "").replace(",", "").strip()
    try:
        num = abs(float(raw))
    except ValueError:
        return value
    sign = "-" if (neg or s.lstrip().startswith("-")) else ""
    body = f"${num / 1000:.1f}B" if num >= 1000 else f"${num:,.0f}MM"
    return f"{sign}{body}"


def _currency_letter(currency: str) -> str:
    """Map a footnote currency scope to its token replacement.

    ``'US$MM'`` / ``'USD'`` -> ``'US'`` and ``'C$MM'`` / ``'CAD'`` -> ``'C'``
    (explicit map — the only two dollar renderings the footnote supports). Any
    other scope returns its bare code with the ``'MM'`` suffix / ``'$'``
    stripped (``'GBPMM'`` -> ``'GBP'``), which `fill_footnote_currency` renders
    as the ISO code without a dollar sign — never silently mislabelled as C$.
    """
    code = currency.strip()
    if code.upper().endswith("MM"):
        code = code[:-2].rstrip()
    code = code.rstrip("$").strip()
    upper = code.upper()
    if upper in {"US", "USD"}:
        return "US"
    if upper in {"C", "CAD"}:
        return "C"
    return code or currency.strip()


def _fill_footnote(shape, currency: str) -> None:
    """Swap the '[x]$MM' currency token in the library footnote.

    Derives the currency letter (or bare ISO code) from the deck's currency
    scope, then defers the in-place token swap to the shared
    `fill_footnote_currency` helper so the rest of the standardized source/note
    lines are preserved verbatim.
    """
    fill_footnote_currency(shape, _currency_letter(currency))


def _write_flexible_bullets(shape, bullets, items=None) -> None:
    """Write bullets (with shrink-on-overflow autofit), falling back to plain
    paragraphs if the shape lacks a glyph template."""
    tuples = items if items is not None else [_bullet_tuple(b) for b in bullets]
    write_bullets_or_plain(shape, tuples, autofit=True)


# The overview-bullets fit (shared with the pitch assembler — the two decks'
# overview slides are the same library entry) lives in
# `pptx_helpers.fit_overview_textbox`: it sizes TextBox 9 to the band above the
# 'LTM Revenue Breakdown' header and writes an explicit normAutofit fontScale
# when the copy would overflow, since PowerPoint ignores a scale-less
# `<a:normAutofit/>` on open.


def _metric_name_pt(name: str) -> float | None:
    """Point size for an over-long KPI tile name, or None to keep the template 9 pt.

    The tile is ~1.44" wide; a name past ~24 chars wraps to a second line and,
    together with the value line, overflows the fixed tile height. Shrink only
    the name (not the value) so the values stay uniform across tiles.
    """
    n = len(name)
    if n <= 24:
        return None
    if n <= 32:
        return 8.0
    return 7.0


def _set_metric_box(box, value_str: str, name: str) -> None:
    """Write the value + name into a tile, shrinking only an over-long name."""
    set_text(box, [value_str, name])
    size = _metric_name_pt(name)
    if size is not None:
        for run in box.text_frame.paragraphs[1].runs:
            run.font.size = Pt(size)


def _set_cover(slide, content: EarningsUpdateContent) -> None:
    title = find_shape(slide, "Title 1")
    set_text(title, [content.company_name, f"{content.reporting_quarter} Earnings Update"])
    for shape in slide.shapes:
        if (
            shape.name.startswith("Subtitle")
            and getattr(shape, "has_text_frame", False)
            and "[Date]" in shape.text_frame.text
        ):
            set_text(shape, [content.cover_date])


def _set_overview(slide, content: EarningsUpdateContent) -> None:
    set_text(find_shape(slide, "Title 6"), [f"Introduction to {content.company_name}"])
    overview = find_shape(slide, "TextBox 9")
    _write_flexible_bullets(overview, content.company_overview_bullets)
    fit_overview_textbox(slide, overview)
    _fill_footnote(find_shape(slide, "Text Placeholder 1"), content.currency)


# (group_name, prior_box, current_box, variance_box) for the four metric rows.
_METRIC_GROUPS = (
    ("Group 12", "Rectangle 1032", "Rectangle 1034", "Rectangle 1041"),
    ("Group 9", "Rectangle 1043", "Rectangle 1037", "Rectangle 1042"),
    ("Group 8", "Rectangle 1035", "Rectangle 1036", "Rectangle 1061"),
    ("Group 2", "Rectangle 1057", "Rectangle 1058", "Rectangle 1064"),
)


def _set_earnings_summary(slide, content: EarningsUpdateContent) -> None:
    set_text(
        find_shape(slide, "Title 1"),
        [f"{content.company_name} {content.reporting_quarter} Earnings Summary"],
    )
    _fill_footnote(find_shape(slide, "Text Placeholder 1"), content.currency)

    # Period header bar (mid-blue) below the Financial Highlights title bar:
    # comparison quarter | Variance | reporting quarter.
    set_text(find_shape(slide, "Rectangle 16"), [content.comparison_quarter])
    set_text(find_shape(slide, "Rectangle 21"), [content.reporting_quarter])

    # Business updates (left column) and performance summary. The library ships
    # TextBox 6 one line tall (0.28") with 2.4" of empty band beneath it and
    # relies on autofit — which PowerPoint ignores when it carries no explicit
    # scale, so the bullets rendered full-size straight through the Broker
    # Estimates header. Same defect `fit_overview_textbox` was written for in
    # v0.5.23, on a shape nobody had rendered; found by the Phase B contract.
    updates = find_shape(slide, "TextBox 6")
    _write_flexible_bullets(updates, None, items=[(b, 0) for b in content.business_updates])
    fit_overview_textbox(slide, updates, band_prefix=_BROKER_TABLE_BAND)
    set_text(find_shape(slide, "Rectangle 1111"), [content.performance_summary])

    # Four metric rows — metric name + value only; the period lives in the bar.
    for (group_name, prior_box, current_box, var_box), kpi in zip(
        _METRIC_GROUPS, content.kpi_rows, strict=True
    ):
        group = find_shape(slide, group_name)
        name = _strip_currency_unit(kpi.name)
        _set_metric_box(find_shape_in_group(group, prior_box), _fmt_mm(kpi.prior_value), name)
        _set_metric_box(find_shape_in_group(group, current_box), _fmt_mm(kpi.current_value), name)
        color = COLOR_UP if kpi.delta_sign > 0 else (COLOR_DOWN if kpi.delta_sign < 0 else None)
        set_text(find_shape_in_group(group, var_box), [kpi.delta_str], size_pt=10, color_hex=color)

    # Broker estimates vs actuals — $ on reported / estimate / variance.
    tbl = find_table_shape(slide).table
    set_cell_text(tbl.cell(0, 0), f"Figures in {content.currency_short}", size_pt=9)
    set_cell_text(tbl.cell(0, 1), "Reported", size_pt=9)
    set_cell_text(tbl.cell(0, 2), "Bloomberg Estimate", size_pt=9)
    set_cell_text(tbl.cell(0, 3), "Variance", size_pt=9)
    for i, row in enumerate(content.broker_rows, start=1):
        set_cell_text(tbl.cell(i, 0), _strip_currency_unit(row.label), size_pt=9)
        set_cell_text(tbl.cell(i, 1), _money(row.reported), size_pt=9)
        set_cell_text(tbl.cell(i, 2), _money(row.estimate), size_pt=9)
        color = COLOR_UP if row.variance_sign > 0 else (COLOR_DOWN if row.variance_sign < 0 else None)
        set_cell_text(tbl.cell(i, 3), _money(row.variance), size_pt=9, color_hex=color)

    # Management quotes.
    q1, q2 = content.management_quotes
    g1 = find_shape(slide, "Group 1070")
    set_text(find_shape_in_group(g1, "TextBox 1072"), [f"“{q1.quote}”"])
    set_text(find_shape_in_group(g1, "TextBox 1073"), [f"{q1.speaker} – {q1.role}"])
    g2 = find_shape(slide, "Group 1086")
    set_text(find_shape_in_group(g2, "TextBox 1088"), [f"“{q2.quote}”"])
    set_text(find_shape_in_group(g2, "TextBox 1089"), [f"{q2.speaker} – {q2.role}"])


def assemble_earnings_update_deck(
    *,
    slide_plan_path: Path | str,
    content_path: Path | str,
    template_path: Path | str,
    output_dir: Path | str,
    captable_workbook_path: Path | str | None = None,
    converge: bool = True,
) -> Path:
    """Build the earnings-update deck by cloning library entries and filling them.

    `template_path` must point at the shared `INFOR Slide Library.pptx`.
    """
    slide_plan = SlidePlan.model_validate_json(Path(slide_plan_path).read_text(encoding="utf-8"))
    content = EarningsUpdateContent.model_validate_json(Path(content_path).read_text(encoding="utf-8"))
    if slide_plan.deliverable_type != "earnings-update":
        raise ValueError("earnings assembler only supports earnings-update SlidePlan objects")

    template = Path(template_path)
    if not template.exists():
        raise FileNotFoundError(f"slide-library template not found: {template}")

    out_dir = Path(output_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    output_path = out_dir / f"Earnings Update - {safe_filename(content.company_name, default='Company')}.pptx"

    prs = Presentation(template)

    # Locate each kept entry by its marker shape. `find_slide_by_marker` insists
    # on exactly one match, so a library that lost, duplicated or renamed one of
    # these raises TemplateLayoutError here — before any clone or delete — rather
    # than assembling the wrong slides.
    keep_indices = [find_slide_by_marker(prs, m, template=template.name) for m in _KEEP_MARKERS]

    # Reduce the library to the five earnings entries (delete from the tail so
    # earlier indices stay valid).
    keep = set(keep_indices)
    for idx in range(len(prs.slides) - 1, -1, -1):
        if idx not in keep:
            delete_slide(prs, idx)

    # Re-find the three slides that get filled, now that the deck is five slides
    # long. Re-finding rather than assuming the survivors kept their listed order
    # costs one cheap scan and makes the fill independent of the library's slide
    # order entirely — the markers are still intact at this point, since nothing
    # has been filled yet.
    cover_at = find_slide_by_marker(prs, MARKER_COVER, template=template.name)
    overview_at = find_slide_by_marker(prs, MARKER_OVERVIEW, template=template.name)
    summary_at = find_slide_by_marker(prs, MARKER_EARNINGS_SUMMARY, template=template.name)

    _set_cover(prs.slides[cover_at], content)
    _set_overview(prs.slides[overview_at], content)
    _set_earnings_summary(prs.slides[summary_at], content)
    # Disclaimer + contact are static library entries — untouched.

    prs.save(output_path)

    # Write -> verify -> repair -> re-verify against the deck contract, BEFORE the
    # cap-table picture goes in: the repair loop re-saves the deck through
    # python-pptx, and doing that after the COM insertion would round-trip the
    # pasted picture for no reason. A picture cannot overflow its box anyway —
    # only text and tables re-grow — so nothing the insertion adds needs fitting.
    #
    # No `out_dir`: the loop's renders are scratch and stage under `tempfile`, so a
    # successful converge leaves nothing in the deal directory (see deck_repair's
    # module docstring). `keep_on_failure` preserves the failing pass's renders —
    # and only those — for the analyst.
    if converge:
        assert_converged(
            output_path,
            converge_deck(output_path, library=template, keep_on_failure=out_dir / ".qa"),
        )

    if captable_workbook_path is not None:
        # The picture range is read through its defined name — a cap table that
        # lost the name raises instead of pasting whatever the shipped address
        # happens to hold now.
        verify_workbook_names(
            captable_workbook_path, sheet=TAB_CAPTABLE, names=CAP_TABLE_PICTURE_NAMES
        )
        insert_excel_into_placeholder(
            deck_path=output_path,
            workbook_path=captable_workbook_path,
            output_path=output_path,
            # The OVERVIEW slide, not the earnings summary. `Rectangle 3` is the
            # overview's "Capitalization Summary" placeholder (library slide 6 ->
            # assembled index 1); the summary slide has no such shape, so
            # targeting it raises KeyError and fails the whole deck stage.
            #
            # v0.5.40 (Phase C) replaced the hardcoded `slide_index=1` with
            # `summary_at` while converting indices to marker lookups — the right
            # migration, the wrong marker. It shipped because the only two tests
            # covering this insertion were both invisible on a dev box: the Excel
            # one behind the opt-in `excel_com` gate, the LibreOffice one behind
            # `skipif win32`. Deleting COM removed the reason for both guards, and
            # un-gating them surfaced this immediately.
            slide_index=overview_at,
            placeholder_name=_CAP_TABLE_PLACEHOLDER,
            sheet_name=TAB_CAPTABLE,
            source_range=cap_table_picture_range(captable_workbook_path),
        )
    _verify_output(
        output_path,
        overview_index=overview_at,
        summary_index=summary_at,
        cap_table_inserted=captable_workbook_path is not None,
    )
    return output_path


def _slide_text(slide) -> str:
    parts: list[str] = []
    for shape in iter_all_shapes(slide.shapes):
        if getattr(shape, "has_text_frame", False):
            parts.append(shape.text)
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    parts.append(cell.text)
    return "\n".join(parts)


def _verify_output(
    path: Path,
    *,
    overview_index: int,
    summary_index: int,
    cap_table_inserted: bool = False,
) -> None:
    # The slide indices are passed in rather than re-discovered: this runs on
    # the FILLED deck, where the library markers have been overwritten with the
    # company's own copy. The caller resolved them by marker before the fill.
    prs = Presentation(path)
    if len(prs.slides) != len(_KEEP_MARKERS):
        raise ValueError(
            f"earnings deck must have {len(_KEEP_MARKERS)} slides, got {len(prs.slides)}"
        )
    # This deck takes one range picture today, so the check cannot bite yet — but
    # it is the same guard the pitch assembler runs, and a second range paste here
    # (an LTM tab, say) must not be able to silently duplicate the first.
    assert_range_pictures_are_distinct(prs)
    overview_text = _slide_text(prs.slides[overview_index])
    summary_text = _slide_text(prs.slides[summary_index])

    forbidden = ["[x]", "[Client Name]", "[Company]", "[Quarter]", "[Name]", "[Role]", "[Date]"]
    leftovers = sorted({t for t in forbidden if t in (overview_text + "\n" + summary_text)})
    if leftovers:
        raise ValueError(f"assembled earnings deck still contains placeholders: {leftovers}")

    # The LTM revenue pie remains a deferred placeholder.
    if "[Pie Chart Placeholder]" not in overview_text:
        raise ValueError("LTM revenue pie placeholder was unexpectedly removed from the overview slide")

    has_cap_placeholder = "[Cap Table Placeholder]" in overview_text
    if cap_table_inserted and has_cap_placeholder:
        raise ValueError("cap-table placeholder was not replaced by the Excel insertion stage")
    if not cap_table_inserted and not has_cap_placeholder:
        raise ValueError("cap-table placeholder was modified; it must remain when no workbook is supplied")
