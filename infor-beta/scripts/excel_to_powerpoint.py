"""Reusable Excel-to-PowerPoint insertion helpers.

The cap-table insertion path renders the Excel source range as a picture and
pastes it into the deck placeholder at the placeholder's exact width and height.
**One render backend, on every platform: LibreOffice headless** — set the print
area to the source range via openpyxl, convert the workbook to PDF with
`soffice --headless --convert-to pdf`, render PDF page 1 to PIL via `pypdfium2`,
and feed the bytes to python-pptx.

Phase D deleted the Excel-COM backend this preferred on Windows. It was the more
faithful renderer in principle — a real Excel recalc, then
`Range.CopyPicture(xlScreen, xlPicture)` into a temporary `ChartObject` exported
via `Chart.Export` — but production is Cowork/Linux with no Excel, so it never
rendered anything that shipped, while being the branch a Windows dev box
exercised instead of the one that does. That is the inverse of the dev/prod parity
Phase A set out to create, and Phase A had already flipped the slide renderer for
the same reason.

Deleted with it: the `excel_com_app` instance owner, the visible-but-parked
off-screen instance (an invisible one renders a blank picture after a recalc), the
`Chart.Paste`-pasted-nothing clipboard race with its retry loop and its distinct
`_ClipboardPasteError` — which, once the retries were spent, had to degrade to
LibreOffice anyway.

Tune the workbook's column widths and row heights so the natural aspect ratio of
the source range matches the target placeholder; the picture is stretched to fit
either way.

Two standing rules govern the temp copy this renders from:

- **A renderer must render a private copy.** Both engines hold the caller's file
  open (LibreOffice drops a `.~lock`), and `zipfile.is_zipfile` swallows the
  resulting `OSError` into a bogus `PackageNotFoundError`.
- **The private copy must print ONE sheet, and the PDF must have exactly one
  page.** Hiding a sheet does not stop LibreOffice exporting it — neither its
  print range, nor its whole used range when it is the workbook's active tab. A
  deal workbook trips both: it inherits print areas from its source templates
  (`captable` -> `$A$1:$T$187`, `comps` -> `$A$1:$BI$45`) and it opens on
  `precedents`. So the export was a multi-page PDF whose page 1 was the entire
  captable sheet, and rendering `pdf[0]` silently returned the cap table for every
  caller — the wrong-sheet bug that shipped a pitch deck with two cap-table
  pictures on the ownership slide, byte-identical and therefore deduped by
  python-pptx into one shared `r:embed`. `_prepare_print_copy` handles both
  mechanisms; the page-count assertion below is the backstop that makes a third
  one loud instead of silent.
"""

from __future__ import annotations

import io
import shutil
import subprocess
import tempfile
from collections import Counter
from pathlib import Path

from pptx import Presentation

# The name given to a picture that replaced a placeholder, so an Excel-range
# picture stays identifiable in the finished deck (`Excel Range: Rectangle 1`)
# rather than an anonymous `Picture 14`. `assert_range_pictures_are_distinct`
# scopes itself by this prefix.
RANGE_PICTURE_PREFIX = "Excel Range: "


def range_picture_name(placeholder_name: str) -> str:
    """The shape name for the picture that replaces ``placeholder_name``."""
    return f"{RANGE_PICTURE_PREFIX}{placeholder_name}"


def _range_picture_embeds(slide) -> "list[tuple[str, str]]":
    """``(shape name, r:embed id)`` for each Excel-range picture on ``slide``.

    Flat scan: `add_picture` appends to the slide's own shape tree, so a range
    picture is never inside a group.
    """
    return [
        (shape.name, shape._element.blip_rId)
        for shape in slide.shapes
        if shape.name.startswith(RANGE_PICTURE_PREFIX)
        and getattr(shape._element, "blip_rId", None)
    ]


def assert_range_pictures_are_distinct(prs: Presentation) -> None:
    """Fail when two Excel-range pictures on one slide share an embedded image.

    The symptom-level guard on the wrong-sheet bug class in the module docstring.
    python-pptx dedupes image parts by SHA1, so two ranges that rendered
    byte-identical PNGs ship as *one* picture referenced twice — which is what
    the ownership slide did when both of its ranges rendered the cap table. The
    slide is valid XML and python-pptx reports two pictures, so nothing short of
    comparing the embeds notices.

    Scoped to range pictures on purpose: the credentials slides legitimately show
    one client logo twice, so an unscoped assertion would fail every pitch deck.
    """
    for index, slide in enumerate(prs.slides, start=1):
        embeds = _range_picture_embeds(slide)
        counts = Counter(rid for _, rid in embeds)
        shared = {
            rid: sorted(name for name, r in embeds if r == rid)
            for rid, n in counts.items()
            if n > 1
        }
        if shared:
            raise ValueError(
                f"slide {index} carries Excel-range pictures sharing one embedded "
                f"image: {shared}. python-pptx dedupes identical image bytes to a "
                f"single r:embed, so two ranges rendered the same picture — check "
                f"that the range renderer is not selecting the wrong sheet's page "
                f"(see excel_to_powerpoint's one-sheet-one-page rule)."
            )


# LibreOffice recalculates OOXML formulas on load only when told to. openpyxl
# writes the cap table with no cached values and the template is manual-calc, so
# without this the headless export prints every formula cell (market cap, net
# debt, Enterprise Value, the whole Financial/Valuation block) as 0/blank. The
# profile below sets Calc's "recalc OOXML/ODF on load" mode to 0 = Always, so the
# conversion recomputes them — and 0 = Always (not 2 = Prompt) is essential in
# headless mode, where a prompt would silently skip the recalc.
def insert_excel_into_placeholder(
    *,
    deck_path: Path | str,
    workbook_path: Path | str,
    placeholder_name: str,
    source_range: str,
    sheet_name: str,
    output_path: Path | str | None = None,
    slide_index: int = 0,
) -> Path:
    """Replace a deck placeholder with a picture of an Excel range.

    Generic over (slide, placeholder, sheet, range): the cap table pastes
    ``Cap with Links!B15:F40`` into a slide-7 ``Rectangle 3``; the ownership
    slide pastes ``Ownership!B4:G17`` into its ``Rectangle 1`` insider
    placeholder. The picture is stretched to the placeholder's exact box.
    """
    deck = Path(deck_path).resolve()
    workbook = Path(workbook_path).resolve()
    if not deck.exists():
        raise FileNotFoundError(f"deck not found: {deck}")
    if not workbook.exists():
        raise FileNotFoundError(f"workbook not found: {workbook}")

    prs = Presentation(deck)
    slide = prs.slides[slide_index]
    placeholder = next((shape for shape in slide.shapes if shape.name == placeholder_name), None)
    if placeholder is None:
        raise KeyError(f"placeholder {placeholder_name!r} not found on slide {slide_index + 1}")
    left, top, width, height = placeholder.left, placeholder.top, placeholder.width, placeholder.height

    png_buffer = _render_range_to_png(workbook, sheet_name, source_range)

    placeholder._element.getparent().remove(placeholder._element)
    picture = slide.shapes.add_picture(png_buffer, left, top, width=width, height=height)
    # Name it after the placeholder it replaced. python-pptx would call it
    # `Picture <next shape id>`, which says nothing about which range it holds and
    # gives `assert_range_pictures_are_distinct` nothing to scope itself by.
    picture.name = range_picture_name(placeholder_name)

    out = Path(output_path) if output_path is not None else deck
    prs.save(out)
    return out


def _render_range_to_png(workbook: Path, sheet_name: str, source_range: str) -> io.BytesIO:
    """Render an Excel range as PNG with headless LibreOffice.

    One backend, on every platform. Phase D deleted the Excel-COM alternative
    this used to prefer on Windows: production is Cowork/Linux with no Excel, so
    the COM path rendered nothing that ships while being the branch a dev box
    exercised — the inverse of the dev/prod parity Phase A set out to create. It
    also carried the whole clipboard-race apparatus (`Chart.Paste` silently
    pasting nothing, a retry loop, and a distinct `_ClipboardPasteError` that had
    to degrade to LibreOffice anyway once the retries were spent).

    Fidelity is close but not pixel perfect; tune the workbook to render cleanly
    under LibreOffice, which is what the analyst receives.
    """
    return _libreoffice_range_to_png(workbook, sheet_name, source_range)


_LO_RECALC_XCU = """\
<?xml version="1.0" encoding="UTF-8"?>
<oor:items xmlns:oor="http://openoffice.org/2001/registry" xmlns:xs="http://www.w3.org/2001/XMLSchema" xmlns:xsi="http://www.w3.org/2001/XMLSchema-instance">
 <item oor:path="/org.openoffice.Office.Calc/Formula/Load">
  <prop oor:name="OOXMLRecalcMode" oor:op="fuse"><value>0</value></prop>
 </item>
 <item oor:path="/org.openoffice.Office.Calc/Formula/Load">
  <prop oor:name="ODFRecalcMode" oor:op="fuse"><value>0</value></prop>
 </item>
</oor:items>
"""


def _write_lo_recalc_profile(base_dir: Path) -> str:
    """Create a throwaway LibreOffice user profile that forces a formula recalc
    on load, and return its ``file://`` URI for ``-env:UserInstallation``.

    Self-contained under ``base_dir`` (LibreOffice reads
    ``<UserInstallation>/user/registrymodifications.xcu``), so the user's global
    LibreOffice profile is never touched.
    """
    profile = base_dir / "lo_profile"
    (profile / "user").mkdir(parents=True, exist_ok=True)
    (profile / "user" / "registrymodifications.xcu").write_text(
        _LO_RECALC_XCU, encoding="utf-8"
    )
    return profile.as_uri()


# Standard install locations for platforms whose LibreOffice installer does not
# put `soffice` on PATH. The Windows MSI never does — so a Windows dev box with
# LibreOffice correctly installed still fails `shutil.which("soffice")`, which
# would leave the LibreOffice-by-default slide renderer unusable on the very
# machine it exists to keep honest.
_SOFFICE_FALLBACK_PATHS = (
    r"C:\Program Files\LibreOffice\program\soffice.exe",
    r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    "/Applications/LibreOffice.app/Contents/MacOS/soffice",
)


def find_soffice() -> str | None:
    """Return a usable ``soffice`` command, or None when LibreOffice is absent.

    Prefers PATH (how Cowork / Linux prod resolves it), then falls back to the
    standard per-platform install locations.

    **The single LibreOffice locator — never call ``shutil.which("soffice")``
    directly.** A bare PATH lookup is what shipped in v0.5.35: the renderer was
    flipped to LibreOffice-by-default on every platform while five other call
    sites (this module's range renderer, three in ``financial_charts``, one in
    the aggregator, since deleted) still resolved through PATH only, so on a Windows
    dev box (MSI install, no PATH entry) they failed or silently degraded —
    inverting the dev/prod parity the flip existed to create. The drift lock in
    ``test_excel_to_powerpoint.py`` fails if a bare lookup reappears.
    """
    found = shutil.which("soffice") or shutil.which("libreoffice")
    if found:
        return found
    return next((p for p in _SOFFICE_FALLBACK_PATHS if Path(p).is_file()), None)


def _soffice_convert(soffice: str, src: Path, out_fmt: str, out_dir: Path) -> None:
    """Convert ``src`` to ``out_fmt`` with headless LibreOffice, recalculating
    in-workbook formulas on load via a throwaway recalc profile.

    ``out_fmt`` is a ``--convert-to`` target (e.g. ``"pdf"`` or
    ``"xlsx:Calc MS Excel 2007 XML"``). Raises RuntimeError on failure. CapIQ
    ``_xll.*`` cells are unknown to LibreOffice and resolve to ``#NAME?`` (the
    template's IFERROR wrappers degrade those to ``n/a``); the recalc does not
    crash on them and every in-workbook arithmetic cell still computes.
    """
    with tempfile.TemporaryDirectory() as prof_base:
        profile_uri = _write_lo_recalc_profile(Path(prof_base))
        try:
            subprocess.run(
                [
                    soffice,
                    f"-env:UserInstallation={profile_uri}",
                    "--headless",
                    "--nologo",
                    "--nodefault",
                    "--nofirststartwizard",
                    "--convert-to",
                    out_fmt,
                    "--outdir",
                    str(out_dir),
                    str(src),
                ],
                check=True,
                capture_output=True,
                timeout=180,
            )
        except subprocess.CalledProcessError as exc:
            raise RuntimeError(
                f"LibreOffice {out_fmt!r} conversion failed: "
                f"{exc.stderr.decode(errors='replace')}"
            ) from exc
        except subprocess.TimeoutExpired as exc:
            # A wedged soffice must degrade like a missing one: every caller's
            # graceful-degradation net catches RuntimeError only, so a raw
            # TimeoutExpired would abort the whole stage even though the durable
            # artefacts (e.g. the native workbook charts) are already saved.
            raise RuntimeError(
                f"LibreOffice {out_fmt!r} conversion timed out after "
                f"{exc.timeout:.0f}s"
            ) from exc


# Stem of the private copy and therefore of LibreOffice's PDF output (it names the
# output after the input). One constant so the two paths cannot drift apart.
_PRINT_COPY_STEM = "range_print"


def _prepare_print_copy(
    workbook: Path, sheet_name: str, source_range: str, dest: Path
) -> None:
    """Save a private copy of ``workbook`` whose only print area is ``source_range``.

    Makes ``sheet_name`` the only visible sheet **and the selected/active one**,
    **clears every other sheet's print area**, points the target sheet's at
    ``source_range`` and forces fit-to-1-page with tight margins.

    Two independent reasons a non-target sheet gets exported anyway, both of which
    put a foreign page ahead of the requested range — and hiding the sheet stops
    neither:

    1. **Its print area.** LibreOffice exports a hidden sheet's print range. A deal
       workbook inherits print areas from its source templates (`captable` ->
       `'captable'!$A$1:$T$187`, `comps` -> `'comps'!$A$1:$BI$45`), so those two
       tabs print in full unless cleared.
    2. **The stored active tab.** A workbook whose active tab is some other sheet
       exports that sheet's used range in full, hidden or not. The shipped deal
       workbook template opens on `precedents`, which is how the cap table stayed
       page 1 even after every stray print area was cleared.

    Together they are why the caller's `pdf[0]` was the cap table for every
    requested range. See the one-sheet-one-page rule in the module docstring; the
    page-count assertion in `_libreoffice_range_to_png` is the backstop for a
    third mechanism nobody has found yet.
    """
    from openpyxl import load_workbook

    wb = load_workbook(workbook)
    if sheet_name not in wb.sheetnames:
        raise KeyError(
            f"sheet {sheet_name!r} not found in workbook (available: {wb.sheetnames})"
        )
    for name in wb.sheetnames:
        sheet = wb[name]
        sheet.sheet_state = "visible" if name == sheet_name else "hidden"
        if name != sheet_name:
            sheet.print_area = None
        for view in sheet.views.sheetView:
            view.tabSelected = name == sheet_name
    wb.active = wb.sheetnames.index(sheet_name)
    ws = wb[sheet_name]
    ws.print_area = source_range
    ws.page_setup.fitToWidth = 1
    ws.page_setup.fitToHeight = 1
    ws.sheet_properties.pageSetUpPr.fitToPage = True
    ws.page_margins.left = 0.25
    ws.page_margins.right = 0.25
    ws.page_margins.top = 0.25
    ws.page_margins.bottom = 0.25
    wb.save(dest)


def _libreoffice_range_to_png(workbook: Path, sheet_name: str, source_range: str) -> io.BytesIO:
    """Render an Excel range as PNG via LibreOffice headless PDF export.

    Strategy: save a temp copy in which the target sheet is the only visible and
    active one, no other sheet has a print area, and the target's print area is
    the source range, fit to one page (`_prepare_print_copy`). Convert that copy to PDF with
    headless LibreOffice — **forcing a recalculation on load** so the
    openpyxl-authored (cache-less), manual-calc cap-table formulas actually
    compute (see `_soffice_convert` / `_write_lo_recalc_profile`) — then render
    PDF page 1 via pypdfium2 at 200 DPI and return as PNG bytes.

    CapIQ `_xll.*` cells (forward consensus estimates) are unknown to LibreOffice
    and resolve to `#NAME?`, which the template's IFERROR wrappers degrade to
    `n/a`; the in-workbook arithmetic still computes, and so does the LTM column,
    whose cells are `=INDEX('ltm-metrics'!…)` links. That is why the non-target
    sheets are **hidden, not removed**: a cross-tab link needs its sibling tab in
    the copy, and recalc-on-load resolves it there.

    The export must be exactly one page. A second page means another sheet was
    exported too, in which case page 1 is not necessarily the requested range — so
    this raises rather than rendering whatever page 1 happens to be.

    Requires LibreOffice (resolved by `find_soffice`) and the `pypdfium2`
    package. Raises RuntimeError with a clear message if either is
    missing — the conductor surfaces this to the analyst.
    """
    soffice = find_soffice()
    if soffice is None:
        raise RuntimeError(
            "LibreOffice (soffice/libreoffice) not found on PATH or in the "
            "standard install locations; required for the non-Windows "
            "cap-table renderer. Install LibreOffice or run the conductor on "
            "a Windows machine with Excel."
        )
    try:
        import pypdfium2 as pdfium
    except ImportError as exc:
        raise RuntimeError(
            "pypdfium2 is required for the non-Windows cap-table renderer; "
            "run `pip install pypdfium2`."
        ) from exc

    with tempfile.TemporaryDirectory() as tmp_dir:
        tmp_xlsx = Path(tmp_dir) / f"{_PRINT_COPY_STEM}.xlsx"
        _prepare_print_copy(workbook, sheet_name, source_range, tmp_xlsx)

        # Convert to PDF with recalc-on-load forced (the saved xlsx has no cached
        # formula values), so the EV cascade and metric blocks print computed.
        _soffice_convert(soffice, tmp_xlsx, "pdf", Path(tmp_dir))

        pdf_path = Path(tmp_dir) / f"{_PRINT_COPY_STEM}.pdf"
        if not pdf_path.exists():
            raise RuntimeError(
                f"LibreOffice produced no PDF output for {workbook.name} "
                f"({sheet_name}!{source_range})"
            )

        pdf = pdfium.PdfDocument(str(pdf_path))
        try:
            # One sheet in, one page out. More than one page means another sheet
            # was exported too, and page 1 is then whichever sheet comes first in
            # the workbook rather than the requested range — the wrong-sheet bug
            # class, which is silent unless asserted here.
            if len(pdf) != 1:
                raise RuntimeError(
                    f"LibreOffice exported {len(pdf)} PDF pages for "
                    f"{sheet_name}!{source_range} in {workbook.name}; the range "
                    f"renderer requires exactly 1 because it renders page 1, and "
                    f"an extra page means page 1 may be a different sheet. Look for "
                    f"a sheet LibreOffice exports despite being hidden: a stray "
                    f"print area, or the stored active tab (_prepare_print_copy "
                    f"clears both)."
                )
            page = pdf[0]
            pil_image = page.render(scale=200 / 72).to_pil().convert("RGB")
            pil_image = _trim_white_margins(pil_image)
            buf = io.BytesIO()
            pil_image.save(buf, format="PNG")
            buf.seek(0)
            return buf
        finally:
            pdf.close()


def _trim_white_margins(image, threshold: int = 250):
    """Crop solid-white borders left by LibreOffice's print-area PDF export."""
    from PIL import ImageChops, Image

    bg = Image.new("RGB", image.size, (255, 255, 255))
    diff = ImageChops.difference(image, bg)
    # Treat near-white as background by quantising
    diff = diff.point(lambda p: 0 if p < (255 - threshold) else 255)
    bbox = diff.getbbox()
    return image.crop(bbox) if bbox else image
