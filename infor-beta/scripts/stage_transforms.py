"""In-process stage transforms — the deterministic half of a conductor plan.

Phase F sorts every plan stage into one of two kinds.

**Judgment** — `captable`, `comps`, `precedents`, `ownership`,
`financial-summary`, `ltm-metrics`, `pitch-content`, `earningsupdate-content`,
`deckread`, `deckcheck`. These read filings, search the web, choose peer sets, draft
copy, look at rendered slides, and argue about whether a figure on a slide is true.
They stay sub-agents with a real tool allow-list, dispatched by the model.

**Transform** — `pitch-wireframe`, `earningsupdate-wireframe`, `deck-assembler`,
`financial-charts`. Each is one call to a function that already exists in
`scripts/`, on inputs the conductor already resolved. Before Phase F each was a
sub-agent whose SKILL.md ended in a "Reference command" the model was asked to
reproduce: a Task round trip, a fresh context, and a language model retyping a
call whose arguments were sitting in `inputs.json`. Nothing in that loop could go
*right* that the direct call does not, and three things could go wrong — a
paraphrased argument, a skipped step, or a sub-agent that reports success without
writing `outputs.json`.

So the driver calls them. The functions here are ports of those reference
commands, and they keep the same handoff object the dispatched form had —
:class:`stage_io.StageIO`, built from the same three paths and reading the same
`inputs.json` off disk — so a transform's inputs, its artefact locations and its
`outputs.json` are identical to what the sub-agent produced.

**A transform is not a shortcut past the checkpoint.** Every stage still writes
`inputs.json` and `outputs.json`, still carries its `$stages` dependency edges
into `plan_schedule.compute_waves`, and still reports through
`conductor.complete_wave` — so a checkpoint (or a failure) on a transform behaves
exactly as it would on a sub-agent. What changed is who issues the call, not what
the run looks like.

**What the reclassification could have quietly dropped, and where it went.** Two
of these skills ended in a *mandatory* analyst-facing QA section — the
deck-assembler's "read the slides" pass and `financial-charts`' chart check — and
those are judgment, not transform. Deleting the SKILL.md would have deleted them
with nothing failing.

`financial-charts` came out whole: it renders the overview + Financial Summary
slides and hands back `charts_inserted` / `pie_inserted`, which makes a skipped chart
visible at the wave boundary where it used to depend on a sub-agent remembering to
mention it.

The deck's did not, and it took four releases to notice. `deck-assembler` writes a
vision review (the agenda `deck_contract` already builds, plus the renders and
picture crops) and hands back `vision_review_path` — but an agenda is only half of a
QA pass, and the half that was deleted with the SKILL.md was the *reader*. Nothing
consumed `vision_review_path`, so 19 KB of questions went out with every run and not
one of them was ever answered; a run's own picture defect sat on a slide that file
listed by name and the run reported clean. Reading slides is judgment and a transform
cannot do it, so the reader came back as a stage: `deckread`, dispatched, final wave,
consuming this checklist plus the finished deck and returning findings. What this
module writes is the input to that stage — which is what it always was. The finished
artefact goes to `deckread` and `deckcheck` together, both judgment, both dispatched,
both advisory.
"""

from __future__ import annotations

from collections.abc import Callable
from pathlib import Path

from stage_io import StageIO

#: Filename of the vision-review checklist the `deck` transform writes into its
#: stage directory. Reported through the `vision_review_path` output, which the
#: wave-boundary surface names automatically (every declared `Path` output is) and
#: which the `deckread` stage consumes — it is that stage's agenda, not a review.
VISION_REVIEW_NAME = "vision_review.md"


class StageTransformError(RuntimeError):
    """A transform was asked for a skill that has none, or for a shape it cannot build."""


# ---------------------------------------------------------------------------
# Wireframes
# ---------------------------------------------------------------------------
def pitch_wireframe(io: StageIO) -> dict:
    """`pitch-wireframe` — the typed pitch `SlidePlan`, driven by the slide mix."""
    from pitch_deck_wireframe import build_pitch_deck_slide_plan
    from schemas import Company
    from wireframe_common import write_slide_plan

    inputs = io.inputs
    slide_plan = build_pitch_deck_slide_plan(
        company=Company.model_validate(inputs["company"]),
        section_labels=inputs.get("section_labels"),
        current_section=inputs.get("current_section"),
        market_entry_target_count=inputs.get("market_entry_target_count"),
        financial_metric_count=inputs.get("financial_metric_count"),
        include_investment_highlights=inputs.get("include_investment_highlights"),
    )
    path = write_slide_plan(slide_plan, io.stage_dir / "slide_plan.json")
    return {"slide_plan_path": str(path)}


def earnings_update_wireframe(io: StageIO) -> dict:
    """`earningsupdate-wireframe` — the fixed five-slide earnings-update `SlidePlan`."""
    from earnings_update_wireframe import build_earnings_update_slide_plan
    from schemas import Company
    from wireframe_common import write_slide_plan

    inputs = io.inputs
    slide_plan = build_earnings_update_slide_plan(
        company=Company.model_validate(inputs["company"]),
        reporting_quarter=inputs["reporting_quarter"],
        comparison_quarter=inputs["comparison_quarter"],
    )
    path = write_slide_plan(slide_plan, io.stage_dir / "slide_plan.json")
    return {"slide_plan_path": str(path)}


# ---------------------------------------------------------------------------
# Deck assembly
# ---------------------------------------------------------------------------
def _vision_checklist() -> str:
    """`deck_contract`'s checklist, not a second copy of it — one wording, one place."""
    from deck_contract import VISION_CHECKLIST

    return VISION_CHECKLIST


def render_vision_review(deck: Path | str, vision) -> str:
    """Render the `deck` stage's vision review as markdown.

    The agenda is `deck_contract.vision_pass`'s, unchanged — one entry per thing
    worth a close look, each naming the slide, the question, the slide render and
    (for a rasterised range or chart) the picture blob at native resolution.

    It is the **agenda** for the reading half of deck QA, and only that: every entry
    is a question, and this function answers none of them. Geometry was already
    measured and repaired inside the assembler, so what is left is whether the slides
    read correctly — which no measurement answers, and which the `deckread` stage
    answers by looking. `deckread` re-renders the same agenda against the *finished*
    artefact (`deckread.render_worklist` delegates here), because in the pitch plan
    the charts land after assembly and these crops describe a file that no longer
    exists on disk. One wording, one place.

    It opens by naming the typeface the geometry was measured in, because that is
    the caveat on the sentence before it: `deck_repair` chose every font size from a
    render, and a substituted face means those sizes were measured against advance
    widths the analyst will not see. `converge_deck` logs the same line, but a log
    is not on disk — this is, at `vision_review_path`.
    """
    from font_probe import probe_font_resolution

    deck = Path(deck)
    lines = [
        f"# Deck vision review — {deck.name}",
        "",
        (
            "Geometry is already measured and repaired: `deck_repair.converge_deck` ran "
            "inside the assembler, and the stage would have failed if the deck still "
            "broke the contract. What is left is the part a measurement cannot do — "
            "**reading the slides**. Open each render below and check for "
            f"{_vision_checklist()}."
        ),
        "",
        probe_font_resolution().log_line(),
        "",
    ]

    by_slide: dict[int, list] = {}
    for target in vision.targets:
        by_slide.setdefault(target.slide, []).append(target)

    if not by_slide:
        lines += ["Nothing was flagged for a close look.", ""]
    for index in sorted(by_slide):
        lines.append(f"## Slide {index + 1}")
        lines.append("")
        for target in by_slide[index]:
            shape = f"`{target.shape}` — " if target.shape else ""
            lines.append(f"- {shape}{target.question}")
            if target.crop:
                lines.append(f"  - picture (native resolution): {target.crop}")
        render = next((t.render for t in by_slide[index] if t.render), None)
        if render:
            lines.append(f"- render: {render}")
        lines.append("")

    advisory = [f for f in vision.findings if f.kind != "vision-review"]
    if advisory:
        lines += ["## Also reported", ""]
        lines += [f"- slide {f.slide + 1}: {f.kind} — {f.message}" for f in advisory]
        lines.append("")

    renders = vision.review_images
    lines += [
        "## Every slide's render",
        "",
        *(f"- slide {i + 1}: {renders[i]}" for i in sorted(renders)),
        "" if renders else "- (the deck could not be rendered — QA did not run)",
        "",
    ]
    return "\n".join(lines)


def _write_vision_review(io: StageIO, deck: Path) -> str | None:
    """Build the vision review for `deck` and return the path to it.

    A failure here must not fail the stage — the deck exists and the run carries on
    — but it must not vanish either, so the reason is written into the file
    `vision_review_path` points at, where the analyst will look for the review.
    """
    path = io.stage_dir / VISION_REVIEW_NAME
    try:
        from deck_contract import vision_pass

        vision = vision_pass(deck, out_dir=io.stage_dir / "vision")
        body = render_vision_review(deck, vision)
    except Exception as exc:  # noqa: BLE001 — reported, not swallowed
        body = (
            f"# Deck vision review — {deck.name}\n\n"
            f"**The review could not be built: {type(exc).__name__}: {exc}**\n\n"
            f"The deck itself assembled and converged. Read every slide of "
            f"`{deck}` by hand before it goes out.\n"
        )
    path.write_text(body, encoding="utf-8")
    return str(path)


def deck_assembler(io: StageIO) -> dict:
    """`deck-assembler` — clone the shared library into the deliverable's deck.

    Dispatches on the `SlidePlan`'s own `deliverable_type` rather than on the
    plan, so the two deliverables cannot be crossed. Both assemblers verify the
    library and workbook layouts by marker shape / defined name before touching
    anything, and both run the converge loop internally — a deck that will not
    converge raises `DeckNotConvergedError`, which the driver records as a stage
    failure with the shape and depth named in the message.
    """
    from earnings_update_assembler import assemble_earnings_update_deck
    from pitch_deck_assembler import assemble_pitch_deck
    from schemas import SlidePlan

    inputs = io.inputs
    slide_plan_path = Path(inputs["slide_plan_path"])
    slide_plan = SlidePlan.model_validate_json(slide_plan_path.read_text(encoding="utf-8"))
    template_path = io.plugin_root / "templates" / inputs["template_name"]
    output_dir = io.artefacts_dir

    if slide_plan.deliverable_type == "earnings-update":
        deck_path = assemble_earnings_update_deck(
            slide_plan_path=inputs["slide_plan_path"],
            content_path=inputs["content_bundle_path"],
            template_path=template_path,
            output_dir=output_dir,
            captable_workbook_path=inputs.get("captable_workbook_path"),
        )
    elif slide_plan.deliverable_type == "pitch":
        deck_path = assemble_pitch_deck(
            slide_plan_path=inputs["slide_plan_path"],
            content_path=inputs["content_bundle_path"],
            template_path=template_path,
            output_dir=output_dir,
            captable_workbook_path=inputs.get("captable_workbook_path"),
            ownership_workbook_path=inputs.get("ownership_workbook_path"),
            financial_metric_labels=inputs.get("financial_metric_labels"),
        )
    else:
        raise StageTransformError(
            f"unsupported deliverable_type {slide_plan.deliverable_type!r} in {slide_plan_path}"
        )

    return {
        "deck_path": str(deck_path),
        "vision_review_path": _write_vision_review(io, Path(deck_path)),
    }


# ---------------------------------------------------------------------------
# Financial Summary charts
# ---------------------------------------------------------------------------
#: The overview slide and the first Financial Summary slide, zero-based. The QA
#: render adds each extra FS slide the deck carries.
_CHART_QA_BASE_SLIDES = (6, 7)


def financial_charts(io: StageIO) -> dict:
    """`financial-charts` — the FS metric charts and the overview LTM revenue pie.

    Runs after `deck`, on the deck `deck` produced, and edits it in place. It never
    re-assembles: a re-assembly would write a clean deck over these charts, which
    is why this was the one skill whose allow-list deliberately excluded `Task`.
    As a transform it cannot dispatch anything at all.

    A `None` return from either helper means the workbook had nothing to chart
    (no `financial-summary` tab, no "LTM Revenue Overview" block) or LibreOffice
    could not render the charts it did build. The placeholders stay, the native
    charts are still on the workbook, and the two booleans say so at the wave
    boundary — where the dispatched form relied on the sub-agent mentioning it.
    """
    from financial_charts import (
        render_financial_summary_charts_into_deck,
        render_ltm_revenue_pie_into_deck,
    )

    inputs = io.inputs
    deck_path = str(inputs["deck_path"])
    deal_workbook = inputs["deal_workbook"]

    fs_deck = render_financial_summary_charts_into_deck(
        deck_path=deck_path,
        deal_workbook=deal_workbook,
    )
    if fs_deck is not None:
        deck_path = str(fs_deck)

    pie_deck = render_ltm_revenue_pie_into_deck(
        deck_path=deck_path,
        deal_workbook=deal_workbook,
    )
    if pie_deck is not None:
        deck_path = str(pie_deck)

    return {
        "deck_path": deck_path,
        "charts_inserted": fs_deck is not None,
        "pie_inserted": pie_deck is not None,
        "chart_qa_dir": _render_chart_qa(io, Path(deck_path)),
    }


def _render_chart_qa(io: StageIO, deck: Path) -> str | None:
    """Render the charted slides so the finished artefact can be looked at.

    The overview slide (the pie) plus every Financial Summary slide the deck
    carries — counted from the deck rather than assumed, because the deck spec's
    two-slide option shifts everything after it. `deckcheck` reads the final
    artefact next; these are the renders of the slides this stage changed.
    """
    try:
        from pptx import Presentation

        from slide_render import render_deck_to_png

        slide_count = len(Presentation(deck).slides)
        indices = [i for i in _CHART_QA_BASE_SLIDES if i < slide_count]
        extra = max(_CHART_QA_BASE_SLIDES) + 1
        if extra < slide_count:  # a second Financial Summary slide shifts the rest
            indices.append(extra)
        out_dir = io.stage_dir / "chart-qa"
        render_deck_to_png(deck, out_dir, slide_indices=indices)
        return str(out_dir)
    except Exception:  # noqa: BLE001 — no renderer is not a stage failure
        return None


# ---------------------------------------------------------------------------
# The registry
# ---------------------------------------------------------------------------
#: skill name -> the in-process call. A skill in here is a **transform**: the
#: conductor runs it itself and never dispatches a sub-agent for it. A skill NOT
#: in here is **judgment** and is dispatched. There is no third state, and the
#: plans need no annotation — the classification lives in exactly one place.
TRANSFORMS: dict[str, Callable[[StageIO], dict]] = {
    "pitch-wireframe": pitch_wireframe,
    "earningsupdate-wireframe": earnings_update_wireframe,
    "deck-assembler": deck_assembler,
    "financial-charts": financial_charts,
}


def is_transform(skill: str) -> bool:
    """True when `skill` is executed in-process rather than dispatched."""
    return skill in TRANSFORMS


def run_transform(skill: str, io: StageIO) -> dict:
    """Execute one transform and return its outputs dict (the caller persists it)."""
    try:
        fn = TRANSFORMS[skill]
    except KeyError:
        raise StageTransformError(
            f"{skill!r} is not a transform — it is a judgment stage and must be "
            f"dispatched. Known transforms: {', '.join(sorted(TRANSFORMS))}"
        ) from None
    outputs = fn(io)
    if not isinstance(outputs, dict):
        raise StageTransformError(
            f"transform {skill!r} returned {type(outputs).__name__}, not a dict of outputs"
        )
    return outputs
