"""Tests for pptx_helpers.

Builds tiny in-memory decks with python-pptx — no INFOR template fixtures
required. Run with `python -m unittest infor-beta/scripts/tests/test_pptx_helpers.py`
or `python infor-beta/scripts/tests/test_pptx_helpers.py`.
"""

import os
import sys
import unittest
from copy import deepcopy

# Make pptx_helpers importable when run directly: the helper modules live in
# scripts/, one level up from this tests/ directory.
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from pptx_helpers import (
    COLOR_DOWN,
    COLOR_UP,
    PALATINO,
    delete_shape,
    delete_slide,
    fill_footnote_currency,
    find_shape,
    find_shape_in_group,
    set_cell_text,
    set_text,
    shapes_in_reading_order,
    write_bulleted_shape,
)


# ─── Fixture builders ────────────────────────────────────────────────────────

def _make_shape_with_bold_italic_run(slide, name, text):
    """Add a textbox whose paragraph 0 has a bold-italic Palatino run — mimics
    the template's quote / title shapes that set_text must preserve."""
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    tb.name = name
    p = tb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = PALATINO
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.italic = True
    return tb


def _make_bulleted_shape(slide, name, seed_paragraphs):
    """Add a textbox with seed paragraphs that include explicit pPr + buChar,
    mimicking the template's bulleted shapes.

    seed_paragraphs: list of (text, marL_emu, font_size_pt, bullet_char)
    e.g. [("seed main", 180975, 10.5, "■"), ("seed sub", 360000, 10.0, "-")]
    """
    from lxml import etree

    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(4))
    tb.name = name
    tf = tb.text_frame

    # First paragraph already exists; reuse it for the first seed
    for i, (text, marL, size_pt, bu) in enumerate(seed_paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

        # Wipe runs + pPr the python-pptx default left
        for r in list(p.runs):
            r._r.getparent().remove(r._r)
        for child in list(p._p):
            if child.tag.endswith("}pPr"):
                p._p.remove(child)

        # Build a <a:pPr marL="..." indent="..."><a:buChar char="..."/></a:pPr>
        a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        pPr = etree.SubElement(p._p, f"{{{a_ns}}}pPr")
        pPr.set("marL", str(marL))
        pPr.set("indent", f"-{marL}")
        buChar = etree.SubElement(pPr, f"{{{a_ns}}}buChar")
        buChar.set("char", bu)
        # pPr must come BEFORE any run children — move it to the front
        p._p.remove(pPr)
        p._p.insert(0, pPr)

        run = p.add_run()
        run.text = text
        run.font.name = PALATINO
        run.font.size = Pt(size_pt)
    return tb


def _make_exec_summary_shape(slide, name, body_color="203864"):
    """Mimic the INFOR exec-summary placeholder: a leading marL=0 buNone spacer
    and a trailing glyph-less spacer bracketing a coloured square-main + dash-sub
    seed. The old harvest keyed templates per-paragraph after a marL sort, so the
    marL=0 buNone spacer became 'level 0' and main bullets lost their glyph and
    inherited the placeholder list colour. This guards the level-deduped harvest
    + run-colour preservation."""
    from lxml import etree

    a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(4))
    tb.name = name
    tf = tb.text_frame

    def _seed(p, marL, *, glyph, size_pt=None, no_bullet=False):
        for r in list(p.runs):
            r._r.getparent().remove(r._r)
        for child in list(p._p):
            if child.tag.endswith("}pPr"):
                p._p.remove(child)
        pPr = etree.SubElement(p._p, f"{{{a_ns}}}pPr")
        pPr.set("marL", str(marL))
        pPr.set("indent", f"-{marL}")
        if no_bullet:
            etree.SubElement(pPr, f"{{{a_ns}}}buNone")
        elif glyph is not None:
            etree.SubElement(pPr, f"{{{a_ns}}}buChar").set("char", glyph)
        p._p.remove(pPr)
        p._p.insert(0, pPr)
        if size_pt is not None:
            run = p.add_run()
            run.text = "[x]"
            run.font.name = PALATINO
            run.font.size = Pt(size_pt)
            run.font.color.rgb = RGBColor.from_string(body_color)

    _seed(tf.paragraphs[0], 0, glyph=None, no_bullet=True)          # buNone spacer
    _seed(tf.add_paragraph(), 180000, glyph="§", size_pt=12)        # main (square)
    _seed(tf.add_paragraph(), 360000, glyph="-", size_pt=11)        # sub (dash)
    _seed(tf.add_paragraph(), 180000, glyph=None)                   # glyph-less spacer
    return tb


# ─── fill_footnote_currency tests ────────────────────────────────────────────

class TestFillFootnoteCurrency(unittest.TestCase):
    """The library footnote ships an 'All figures in [x]$MM' token line. US / C
    substitute the dollar letter; any other (ISO) code must also swallow the
    dollar sign so a GBP filer's footnote never reads like a dollar scope."""

    def setUp(self):
        self.prs = Presentation()
        self.slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])

    def _footnote(self, text="All figures in [x]$MM unless otherwise noted"):
        return _make_shape_with_bold_italic_run(self.slide, "Footnote", text)

    def test_us_letter_keeps_dollar_sign(self):
        shape = self._footnote()
        fill_footnote_currency(shape, "US")
        self.assertEqual(shape.text_frame.text, "All figures in US$MM unless otherwise noted")

    def test_c_letter_keeps_dollar_sign(self):
        shape = self._footnote()
        fill_footnote_currency(shape, "C")
        self.assertEqual(shape.text_frame.text, "All figures in C$MM unless otherwise noted")

    def test_iso_code_replaces_dollar_sign_too(self):
        shape = self._footnote()
        fill_footnote_currency(shape, "GBP")
        self.assertEqual(shape.text_frame.text, "All figures in GBP MM unless otherwise noted")

    def test_iso_code_without_dollar_token_falls_back_to_plain_swap(self):
        shape = self._footnote("Figures in [x]MM")
        fill_footnote_currency(shape, "GBP")
        self.assertEqual(shape.text_frame.text, "Figures in GBPMM")


# ─── set_text tests ──────────────────────────────────────────────────────────

class TestSetText(unittest.TestCase):
    def setUp(self):
        self.prs = Presentation()
        self.slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])

    def test_preserves_bold_italic_when_no_overrides(self):
        """The whole point of set_text — mutating runs[0].text must keep the rPr.
        This is what would have shipped Calibri-default text if we'd wiped runs."""
        shape = _make_shape_with_bold_italic_run(self.slide, "Quote", "[original]")
        set_text(shape, ["the new quote"])
        run = shape.text_frame.paragraphs[0].runs[0]
        self.assertEqual(run.text, "the new quote")
        self.assertEqual(run.font.name, PALATINO)
        self.assertEqual(run.font.size, Pt(11))
        self.assertTrue(run.font.bold)
        self.assertTrue(run.font.italic)

    def test_size_override_applies(self):
        shape = _make_shape_with_bold_italic_run(self.slide, "Delta", "[x]")
        set_text(shape, ["+$4.2MM"], size_pt=10)
        run = shape.text_frame.paragraphs[0].runs[0]
        self.assertEqual(run.text, "+$4.2MM")
        self.assertEqual(run.font.size, Pt(10))

    def test_color_override_applies(self):
        shape = _make_shape_with_bold_italic_run(self.slide, "Delta", "[x]")
        set_text(shape, ["+$4.2MM"], size_pt=10, color_hex=COLOR_UP)
        run = shape.text_frame.paragraphs[0].runs[0]
        self.assertEqual(str(run.font.color.rgb), COLOR_UP)

    def test_multi_line_creates_new_paragraphs(self):
        shape = _make_shape_with_bold_italic_run(self.slide, "Footnote", "[seed]")
        set_text(shape, ["line 1", "line 2", "line 3"])
        tf = shape.text_frame
        self.assertEqual(len(tf.paragraphs), 3)
        self.assertEqual(tf.paragraphs[0].runs[0].text, "line 1")
        self.assertEqual(tf.paragraphs[1].runs[0].text, "line 2")
        self.assertEqual(tf.paragraphs[2].runs[0].text, "line 3")
        # New paragraph must inherit template formatting (italic + bold)
        new_run = tf.paragraphs[1].runs[0]
        self.assertEqual(new_run.font.name, PALATINO)
        self.assertTrue(new_run.font.italic)

    def test_shrinks_when_fewer_lines_than_paragraphs(self):
        shape = _make_shape_with_bold_italic_run(self.slide, "X", "[seed]")
        set_text(shape, ["a", "b", "c"])
        self.assertEqual(len(shape.text_frame.paragraphs), 3)
        set_text(shape, ["only one"])
        self.assertEqual(len(shape.text_frame.paragraphs), 1)
        self.assertEqual(shape.text_frame.paragraphs[0].runs[0].text, "only one")


# ─── set_cell_text tests ─────────────────────────────────────────────────────

class TestSetCellText(unittest.TestCase):
    def setUp(self):
        self.prs = Presentation()
        self.slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        self.table = self.slide.shapes.add_table(
            rows=2, cols=2, left=Inches(1), top=Inches(1), width=Inches(4), height=Inches(1)
        ).table

    def test_forces_palatino_at_size(self):
        set_cell_text(self.table.cell(0, 0), "Revenue", size_pt=9)
        run = self.table.cell(0, 0).text_frame.paragraphs[0].runs[0]
        self.assertEqual(run.text, "Revenue")
        self.assertEqual(run.font.name, PALATINO)
        self.assertEqual(run.font.size, Pt(9))

    def test_color_hex_applies_to_variance_cell(self):
        set_cell_text(self.table.cell(1, 1), "+$1.2", size_pt=9, color_hex=COLOR_UP)
        run = self.table.cell(1, 1).text_frame.paragraphs[0].runs[0]
        self.assertEqual(str(run.font.color.rgb), COLOR_UP)

    def test_overwrite_replaces_previous_content(self):
        set_cell_text(self.table.cell(0, 0), "Old")
        set_cell_text(self.table.cell(0, 0), "New")
        self.assertEqual(self.table.cell(0, 0).text_frame.text, "New")


# ─── write_bulleted_shape tests ──────────────────────────────────────────────

class TestWriteBulletedShape(unittest.TestCase):
    def setUp(self):
        self.prs = Presentation()
        self.slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])

    def test_main_level_bullets_get_square_glyph(self):
        shape = _make_bulleted_shape(
            self.slide,
            "TextBox 16",
            [("seed main", 180975, 10.5, "■"), ("seed sub", 360000, 10.0, "-")],
        )
        write_bulleted_shape(shape, [("bullet 1", 0), ("bullet 2", 0)])

        paras = shape.text_frame.paragraphs
        self.assertEqual(len(paras), 2)
        for p in paras:
            buChars = [e for e in p._p.iter() if e.tag.endswith("}buChar")]
            self.assertEqual(len(buChars), 1, "every paragraph must have a buChar")
            self.assertEqual(buChars[0].get("char"), "■", "level-0 square glyph")

    def test_sub_level_bullets_get_dash_glyph(self):
        shape = _make_bulleted_shape(
            self.slide,
            "TextBox 16",
            [("seed main", 180975, 10.5, "■"), ("seed sub", 360000, 10.0, "-")],
        )
        write_bulleted_shape(shape, [("main", 0), ("sub one", 1), ("sub two", 1)])

        paras = shape.text_frame.paragraphs
        glyphs = []
        for p in paras:
            buChar = next(e for e in p._p.iter() if e.tag.endswith("}buChar"))
            glyphs.append(buChar.get("char"))
        self.assertEqual(glyphs, ["■", "-", "-"])

    def test_bold_prefix_two_run_paragraph(self):
        shape = _make_bulleted_shape(
            self.slide,
            "TextBox 16",
            [("seed main", 180975, 10.5, "■"), ("seed sub", 360000, 10.0, "-")],
        )
        write_bulleted_shape(
            shape,
            [("easyfinancial:", " direct-to-consumer unsecured loans", 1)],
        )
        p = shape.text_frame.paragraphs[0]
        self.assertEqual(len(p.runs), 2)
        self.assertEqual(p.runs[0].text, "easyfinancial:")
        self.assertTrue(p.runs[0].font.bold)
        self.assertEqual(p.runs[1].text, " direct-to-consumer unsecured loans")
        self.assertFalse(p.runs[1].font.bold)

    def test_explicit_palatino_on_every_run(self):
        shape = _make_bulleted_shape(
            self.slide,
            "TextBox 16",
            [("seed main", 180975, 10.5, "■"), ("seed sub", 360000, 10.0, "-")],
        )
        write_bulleted_shape(shape, [("a", 0), ("b", 1), ("c:", " tail", 0)])
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                self.assertEqual(
                    r.font.name, PALATINO,
                    f"every run must have Palatino set explicitly (got {r.font.name!r})"
                )

    def test_raises_when_template_has_no_pPr(self):
        """Defensive: a shape with no seed bullets cannot be safely rewritten."""
        tb = self.slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        tb.name = "Empty"
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "no bullet here"

        with self.assertRaises(RuntimeError):
            write_bulleted_shape(tb, [("x", 0)])

    def test_bad_item_shape_raises(self):
        shape = _make_bulleted_shape(
            self.slide, "TextBox 16",
            [("seed", 180975, 10.5, "■")],
        )
        with self.assertRaises(ValueError):
            write_bulleted_shape(shape, [("a", "b", "c", "d")])  # wrong arity

    def test_exec_summary_keeps_template_colour_and_buchar(self):
        """Slide-2-style placeholder: main/sub bullets must keep the template
        run colour (not fall back to the navy list default) and every paragraph
        must carry a bullet glyph — the regression behind 'blue body text'."""
        shape = _make_exec_summary_shape(self.slide, "Content Placeholder 7")
        write_bulleted_shape(
            shape,
            [("Main one", 0), ("Sub under it", 1), ("Second main", 0)],
        )
        paras = shape.text_frame.paragraphs
        self.assertEqual(len(paras), 3)
        glyphs = []
        for p in paras:
            buChars = [e for e in p._p.iter() if e.tag.endswith("}buChar")]
            self.assertEqual(len(buChars), 1, "every paragraph must keep a bullet glyph")
            glyphs.append(buChars[0].get("char"))
            for run in p.runs:
                self.assertEqual(
                    str(run.font.color.rgb), "203864",
                    "run must keep the harvested template colour, not the list default",
                )
        # buNone spacer must NOT hijack level 0: main->square, sub->dash, main->square.
        self.assertEqual(glyphs, ["§", "-", "§"])

    def test_harvest_dedupes_levels_and_skips_bunone(self):
        """Harvest must collapse the many polluted paragraphs to two levels."""
        from pptx_helpers import _harvest_bullet_templates

        shape = _make_exec_summary_shape(self.slide, "Content Placeholder 7")
        templates = _harvest_bullet_templates(shape)
        self.assertEqual(sorted(templates), [0, 1], "exactly two bullet levels")
        # level 0 = smallest indent = the square main bullet
        lvl0_pPr = templates[0][0]
        self.assertEqual(lvl0_pPr.get("marL"), "180000")
        self.assertTrue(
            any(c.tag.endswith("}buChar") and c.get("char") == "§" for c in lvl0_pPr.iter())
        )


# ─── find_shape / find_shape_in_group tests ──────────────────────────────────

class TestFindShape(unittest.TestCase):
    def test_find_shape_by_name(self):
        prs = Presentation()
        # Layout 6 is "Blank" — no auto-added placeholders to collide with our textbox name
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        tb.name = "Rectangle 1032"
        found = find_shape(slide, "Rectangle 1032")
        # python-pptx returns a fresh proxy each iteration — equality is the right check
        self.assertEqual(found.name, "Rectangle 1032")
        self.assertEqual(found.shape_id, tb.shape_id)

    def test_find_shape_missing_raises(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        with self.assertRaises(KeyError):
            find_shape(slide, "Nonexistent")


# ─── delete_slide tests ──────────────────────────────────────────────────────

class TestDeleteSlide(unittest.TestCase):
    def test_delete_removes_slide(self):
        prs = Presentation()
        s0 = prs.slides.add_slide(prs.slide_layouts[6])
        s0.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1)).name = "keep-me"
        s1 = prs.slides.add_slide(prs.slide_layouts[6])
        s1.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1)).name = "delete-me"
        self.assertEqual(len(prs.slides), 2)
        delete_slide(prs, 1)
        self.assertEqual(len(prs.slides), 1)
        remaining = [s.name for s in prs.slides[0].shapes]
        self.assertIn("keep-me", remaining)


# ─── shapes_in_reading_order / delete_shape ──────────────────────────────────

def _box(slide, name, top_in, left_in):
    box = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(1), Inches(0.4))
    box.name = name
    return box


class TestReadingOrder(unittest.TestCase):
    """`slide.shapes` is document order; these pin the sort that fixes it.

    The pitch deck's section divider ships four identically-named boxes in
    document order 1st, 3rd, 2nd, 4th down the slide, and pairing labels against
    that order put two of them in the wrong box on every deck INFOR shipped.
    """

    def _slide(self):
        prs = Presentation()
        return prs.slides.add_slide(prs.slide_layouts[6])

    def test_single_column_sorts_top_to_bottom(self):
        slide = self._slide()
        # Added out of visual order, exactly as the library authors the divider.
        added = [_box(slide, "a", 1.0, 1.0), _box(slide, "c", 3.0, 1.0),
                 _box(slide, "b", 2.0, 1.0), _box(slide, "d", 4.0, 1.0)]
        self.assertEqual([s.name for s in added], ["a", "c", "b", "d"])
        self.assertEqual(
            [s.name for s in shapes_in_reading_order(added)], ["a", "b", "c", "d"]
        )

    def test_grid_sorts_rows_then_left_to_right(self):
        slide = self._slide()
        added = [_box(slide, "br", 3.0, 5.0), _box(slide, "tl", 1.0, 1.0),
                 _box(slide, "bl", 3.0, 1.0), _box(slide, "tr", 1.0, 5.0)]
        self.assertEqual(
            [s.name for s in shapes_in_reading_order(added)], ["tl", "tr", "bl", "br"]
        )

    def test_a_row_survives_a_small_vertical_nudge(self):
        """Row members are banded, not compared exactly.

        A template re-save that shifts one tile of a row by a hair must not
        reshuffle the row — that would be the same silent defect in a new place.
        """
        slide = self._slide()
        added = [_box(slide, "right", 1.05, 5.0), _box(slide, "left", 1.0, 1.0)]
        self.assertEqual(
            [s.name for s in shapes_in_reading_order(added)], ["left", "right"]
        )

    def test_a_real_row_gap_is_not_banded_together(self):
        slide = self._slide()
        added = [_box(slide, "below", 1.5, 1.0), _box(slide, "above", 1.0, 5.0)]
        self.assertEqual(
            [s.name for s in shapes_in_reading_order(added)], ["above", "below"]
        )

    def test_delete_shape_removes_it_from_the_slide(self):
        slide = self._slide()
        keep, drop = _box(slide, "keep", 1.0, 1.0), _box(slide, "drop", 2.0, 1.0)
        delete_shape(drop)
        self.assertEqual([s.name for s in slide.shapes], ["keep"])
        self.assertEqual(keep.name, "keep")


# ─── Constants ───────────────────────────────────────────────────────────────

class TestConstants(unittest.TestCase):
    def test_palatino_full_name(self):
        self.assertEqual(PALATINO, "Palatino Linotype")

    def test_brand_color_hex_uppercase_no_hash(self):
        self.assertEqual(COLOR_UP, "00B050")
        self.assertEqual(COLOR_DOWN, "C00000")


if __name__ == "__main__":
    unittest.main(verbosity=2)


# ─── fit_overview_textbox ────────────────────────────────────────────────────

from pptx.oxml.ns import qn as _qn
from pptx.util import Emu as _Emu

from pptx_helpers import fit_overview_textbox, normal_autofit_scale


def _make_overview_slide():
    """A synthetic overview slide: TextBox 9 up top, the 'LTM Revenue Breakdown'
    header at 4.747" — the geometry of the shared library's overview slide."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.375), Inches(1.471), Inches(4.507), Inches(0.575))
    box.name = "TextBox 9"
    box.text_frame.margin_left = 0
    box.text_frame.margin_right = 0
    header = slide.shapes.add_textbox(Inches(0.365), Inches(4.747), Inches(4.528), Inches(0.341))
    header.text_frame.paragraphs[0].add_run().text = "LTM Revenue Breakdown"
    return prs, slide, box


def _font_scale_of(shape):
    bodyPr = shape.text_frame._txBody.find(_qn("a:bodyPr"))
    autofit = bodyPr.find(_qn("a:normAutofit"))
    if autofit is None:
        return None
    scale = autofit.get("fontScale")
    return None if scale is None else int(scale) / 1000.0


class FitOverviewTextboxTests(unittest.TestCase):
    """What is left after v0.5.39: geometry only.

    Choosing a font scale used to live here, solved from hand-calibrated Palatino
    em constants. That is now measured off a real render in `deck_repair`, and the
    tests for it are in `test_deck_repair.py` — there is deliberately nothing here
    that predicts how tall text will be.
    """

    def test_fit_sizes_box_to_the_band_and_enables_autofit(self):
        prs, slide, box = _make_overview_slide()
        for _ in range(8):
            p = box.text_frame.add_paragraph()
            p.add_run().text = "Propel is a Toronto-based fintech that facilitates access to credit for underbanked consumers across three operating brands and multiple funding programs"
        avail = fit_overview_textbox(slide, box)

        # Box sized to the band above the LTM header (4.747 - 0.12 - 1.471).
        expected = 4.747 - 0.12 - 1.471
        self.assertAlmostEqual(avail, expected, places=2)
        self.assertAlmostEqual(_Emu(box.height).inches, expected, places=2)
        # Shrink-on-overflow is enabled, at full size: the measured repair loop is
        # what lowers the scale, and it must not be pre-empted by a guess here.
        self.assertEqual(normal_autofit_scale(box), 100.0)
        self.assertIsNone(_font_scale_of(box))

    def test_fit_keeps_short_copy_at_full_size(self):
        prs, slide, box = _make_overview_slide()
        box.text_frame.paragraphs[0].add_run().text = "One short line"
        fit_overview_textbox(slide, box)
        # autofit enabled, but no explicit downscale
        self.assertIsNone(_font_scale_of(box))
        self.assertEqual(normal_autofit_scale(box), 100.0)

    def test_fit_without_band_header_keeps_existing_height(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(0.4), Inches(1.5), Inches(4.5), Inches(2.0))
        box.text_frame.paragraphs[0].add_run().text = "No band on this slide"
        fit_overview_textbox(slide, box)
        self.assertAlmostEqual(_Emu(box.height).inches, 2.0, places=3)
