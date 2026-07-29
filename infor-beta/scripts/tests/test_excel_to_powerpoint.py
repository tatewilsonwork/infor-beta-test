"""Tests for the cap-table Excel→PowerPoint renderer's LibreOffice recalc path,
plus the shared ``find_soffice`` locator every LibreOffice caller resolves through.

The Windows COM path forces ``excel.CalculateFull()``; the non-Windows
(Cowork/Linux) path must force LibreOffice to recalculate the openpyxl-authored,
manual-calc workbook on load, or every formula cell prints as 0/blank.

The one-sheet-one-page section is the wrong-sheet bug class: the renderer takes
`pdf[0]`, so anything that puts a foreign page first silently renders the wrong
sheet for every caller. Those tests are deliberately renderer-free (the print-area
preparation is pure openpyxl; the page-count assertion is exercised with a stubbed
conversion), so they run everywhere rather than skipping where LibreOffice is
absent — the bug shipped precisely because the tests that could have caught it
were invisible on the dev box.
"""

import io
import re
from pathlib import Path

import pytest
from openpyxl import Workbook, load_workbook
from pptx import Presentation
from pptx.util import Inches

import excel_to_powerpoint as e2p
from excel_to_powerpoint import (
    _PRINT_COPY_STEM,
    _prepare_print_copy,
    _soffice_convert,
    _write_lo_recalc_profile,
    assert_range_pictures_are_distinct,
    find_soffice,
    insert_excel_into_placeholder,
    range_picture_name,
)

_SCRIPTS_DIR = Path(__file__).resolve().parents[1]
_TEMPLATES_DIR = _SCRIPTS_DIR.parent / "templates"
_DEAL_WORKBOOK_TEMPLATE = _TEMPLATES_DIR / "INFOR Deal Workbook Template.xlsx"
_REPO_ROOT = _SCRIPTS_DIR.parents[1]


def test_soffice_convert_timeout_raises_runtime_error(tmp_path: Path, monkeypatch):
    """v0.5.21: a wedged LibreOffice (TimeoutExpired) must surface as RuntimeError
    like every other soffice failure, so callers' graceful-degradation nets
    (``except RuntimeError``) engage instead of the stage aborting raw."""
    import subprocess

    import excel_to_powerpoint as e2p

    def _hang(*_args, **_kwargs):
        raise subprocess.TimeoutExpired(cmd="soffice", timeout=180)

    monkeypatch.setattr(e2p.subprocess, "run", _hang)
    with pytest.raises(RuntimeError, match="timed out"):
        _soffice_convert("soffice", tmp_path / "x.xlsx", "pdf", tmp_path)


# ─── The shared LibreOffice locator (v0.5.35 flip / v0.5.36 consolidation) ────


def test_find_soffice_returns_path_hit_before_install_locations(tmp_path: Path, monkeypatch):
    """PATH first — that is how Cowork / Linux prod resolves the binary."""
    monkeypatch.setattr("shutil.which", lambda name: "/usr/bin/soffice" if name == "soffice" else None)
    monkeypatch.setattr(e2p, "_SOFFICE_FALLBACK_PATHS", (str(tmp_path),))

    assert find_soffice() == "/usr/bin/soffice"


def test_find_soffice_falls_back_to_standard_install_location(tmp_path: Path, monkeypatch):
    """The Windows MSI puts nothing on PATH, so PATH-miss must not mean absent."""
    monkeypatch.setattr("shutil.which", lambda *_a, **_k: None)
    installed = tmp_path / "soffice.exe"
    installed.write_bytes(b"")
    monkeypatch.setattr(
        e2p, "_SOFFICE_FALLBACK_PATHS", (str(tmp_path / "not-here.exe"), str(installed))
    )

    assert find_soffice() == str(installed)


def test_find_soffice_is_none_when_libreoffice_absent(monkeypatch):
    monkeypatch.setattr("shutil.which", lambda *_a, **_k: None)
    monkeypatch.setattr(e2p, "_SOFFICE_FALLBACK_PATHS", ())

    assert find_soffice() is None


def test_no_bare_libreoffice_path_lookups_outside_the_locator():
    """Every LibreOffice caller resolves through ``find_soffice``.

    v0.5.35 made LibreOffice the default renderer on every platform but wired the
    locator into ``slide_render`` only; five other call sites still probed PATH
    directly, which the Windows MSI never satisfies — so they failed or silently
    degraded on the very dev box the flip existed to keep honest. This locks the
    consolidation: the PATH probe lives in exactly one module.
    """
    probe = re.compile(r"""which\(\s*["'](soffice|libreoffice)""")
    offenders = [
        f"{py.relative_to(_SCRIPTS_DIR).as_posix()}:{n}"
        for py in sorted(_SCRIPTS_DIR.rglob("*.py"))
        if py.name != "excel_to_powerpoint.py"  # the locator's own module
        for n, line in enumerate(py.read_text(encoding="utf-8").splitlines(), 1)
        if probe.search(line)
    ]

    assert not offenders, (
        "resolve LibreOffice through excel_to_powerpoint.find_soffice() rather than a "
        f"bare PATH lookup (it also probes the standard install locations): {offenders}"
    )


def test_write_lo_recalc_profile_forces_always_recalc(tmp_path: Path):
    """The throwaway profile must set OOXML recalc-on-load to 0 (Always)."""
    uri = _write_lo_recalc_profile(tmp_path)

    xcu = tmp_path / "lo_profile" / "user" / "registrymodifications.xcu"
    assert xcu.exists(), "registrymodifications.xcu must live under <profile>/user/"
    text = xcu.read_text(encoding="utf-8")
    assert "OOXMLRecalcMode" in text
    # 0 = Always; 1 = Never and 2 = Prompt would both skip the recalc headless.
    assert "<value>0</value>" in text
    assert uri.startswith("file:") and uri.endswith("lo_profile")


def _build_cacheless_cap_table(path: Path) -> None:
    """A manual-calc cap table with formula cells and NO cached values, plus a
    CapIQ ``_xll.*`` cell LibreOffice cannot resolve.

    The tab is named for the deal workbook's, though nothing here depends on it:
    `insert_excel_into_placeholder` is passed `sheet_name` explicitly, so this
    fixture tests recalc and rendering, not tab resolution. It read
    ``Cap with Links`` until v0.5.45 — harmless, but that is the SOURCE template's
    sheet name and reading it here as "what the pipeline produces" is the
    misreading that release was about.
    """
    from deal_workbook import TAB_CAPTABLE

    wb = Workbook()
    ws = wb.active
    ws.title = TAB_CAPTABLE
    ws["F16"] = 30.0  # share price (hardcoded)
    ws["F17"] = 100.0  # basic shares (hardcoded)
    ws["F18"] = "=F16*F17"  # basic market cap -> 3000
    ws["F28"] = 500.0  # net debt (hardcoded)
    ws["F31"] = "=F18+F28"  # Enterprise Value -> 3500
    ws["D47"] = 1200.0  # LTM revenue, hardcoded by ltm-metrics
    ws["D48"] = 300.0  # LTM Adj. EBITDA, hardcoded
    ws["D34"] = '=IFERROR(D47,"n/a ")'  # -> 1200
    ws["D35"] = '=IFERROR(D48,"n/a ")'  # -> 300
    # A CapIQ add-in cell: unknown to LibreOffice -> #NAME?, must degrade not crash.
    ws["E47"] = '=_xll.SNL.Clients.Office.Excel.Functions.SPG("X","IQ_REV")'
    ws["E34"] = '=IFERROR(E47,"n/a ")'  # -> "n/a "
    wb.calculation.calcMode = "manual"
    wb.calculation.fullCalcOnLoad = True
    wb.save(path)


# Not skipped on Windows any more: the guard said "Windows uses the Excel COM
# path", which stopped being true when Phase D deleted it. Skipping here meant the
# dev platform never exercised the renderer production uses.
def test_libreoffice_recalc_populates_inworkbook_formulas(tmp_path: Path):
    soffice = find_soffice()
    if soffice is None:
        pytest.skip("LibreOffice (soffice/libreoffice) not installed")

    from deal_workbook import TAB_CAPTABLE

    src = tmp_path / "cap_cacheless.xlsx"
    _build_cacheless_cap_table(src)
    # Precondition: openpyxl wrote no cached value, so without recalc this is blank.
    assert load_workbook(src, data_only=True)[TAB_CAPTABLE]["F31"].value is None

    out_dir = tmp_path / "out"
    out_dir.mkdir()
    _soffice_convert(soffice, src, "xlsx:Calc MS Excel 2007 XML", out_dir)

    ws = load_workbook(out_dir / "cap_cacheless.xlsx", data_only=True)[TAB_CAPTABLE]
    # In-workbook arithmetic now computes through Enterprise Value.
    assert ws["F18"].value == pytest.approx(3000.0)
    assert ws["F31"].value == pytest.approx(3500.0)
    # The hardcoded LTM column (fed by ltm-metrics) resolves through IFERROR.
    assert ws["D34"].value == pytest.approx(1200.0)
    assert ws["D35"].value == pytest.approx(300.0)
    # CapIQ cell could not resolve -> degraded to n/a; the recalc did not crash.
    assert "n/a" in str(ws["E34"].value)


# ─── One sheet in, one page out (the wrong-sheet bug class) ──────────────────


def test_preparation_leaves_only_the_target_sheet_printable():
    """Driven against the SHIPPED deal workbook template, which really carries the
    two stray print areas the bug rode in on.

    `scripts/tests/fixtures/pitch-workbook.xlsx` is deliberately NOT used here: it
    has no print area on any sheet, so it renders correctly either way — which is
    exactly why nothing caught this. (If a future deal workbook ever gains print
    areas the fixture would start covering it; do not rely on that, and do not
    "simplify" this test onto the fixture.)

    Both halves of the preparation are asserted, because each alone leaves a
    foreign page ahead of the target: the stray print areas, and the stored active
    tab (this template opens on `precedents`, and LibreOffice exports the active
    sheet's whole used range even when it is hidden).
    """
    from openpyxl import load_workbook as _load

    source = _load(_DEAL_WORKBOOK_TEMPLATE)
    stray = {n: source[n].print_area for n in source.sheetnames if source[n].print_area}
    assert stray == {
        "captable": "'captable'!$A$1:$T$187",
        "comps": "'comps'!$A$1:$BI$45",
    }, (
        f"the shipped deal workbook template no longer carries the stray print areas "
        f"this test exists to reproduce (found {stray}). If a template rebuild removed "
        f"them, the renderer's guard is still required — a re-saved source template can "
        f"bring one back — so find another workbook that has one rather than deleting "
        f"this test."
    )
    assert source.views[0].activeTab != source.sheetnames.index("Ownership")


def test_prepared_copy_prints_only_the_target_range(tmp_path: Path):
    _prepare_print_copy(_DEAL_WORKBOOK_TEMPLATE, "Ownership", "B4:G17", tmp_path / "p.xlsx")

    wb = load_workbook(tmp_path / "p.xlsx")
    printable = {name: wb[name].print_area for name in wb.sheetnames if wb[name].print_area}
    assert printable == {"Ownership": "'Ownership'!$B$4:$G$17"}
    # The target is the only visible sheet AND the active/selected one.
    assert [n for n in wb.sheetnames if wb[n].sheet_state == "visible"] == ["Ownership"]
    assert wb.views[0].activeTab == wb.sheetnames.index("Ownership")
    assert [
        n for n in wb.sheetnames for v in wb[n].views.sheetView if v.tabSelected
    ] == ["Ownership"]


def test_prepare_print_copy_rejects_an_unknown_sheet(tmp_path: Path):
    with pytest.raises(KeyError, match="not found in workbook"):
        _prepare_print_copy(_DEAL_WORKBOOK_TEMPLATE, "nope", "A1:B2", tmp_path / "p.xlsx")


def _stub_multi_page_conversion(monkeypatch, pages: int):
    """Stand in for LibreOffice, emitting a PDF with ``pages`` blank pages."""
    import pypdfium2 as pdfium

    monkeypatch.setattr(e2p, "find_soffice", lambda: "soffice")

    def _convert(_soffice, _src, _fmt, out_dir):
        pdf = pdfium.PdfDocument.new()
        for _ in range(pages):
            pdf.new_page(300, 200)
        pdf.save(str(Path(out_dir) / f"{_PRINT_COPY_STEM}.pdf"))

    monkeypatch.setattr(e2p, "_soffice_convert", _convert)


def _two_sheet_workbook(path: Path) -> Path:
    wb = Workbook()
    wb.active.title = "captable"
    wb.active["A1"] = "cap"
    own = wb.create_sheet("Ownership")
    own["B4"] = "Select Insiders"
    wb.save(path)
    return path


def test_a_multi_page_export_raises_naming_the_sheet_range_and_page_count(
    tmp_path: Path, monkeypatch
):
    """Silent wrong-page selection is the whole bug class, so an extra page halts.

    Renderer-free on purpose: the conversion is stubbed, so this asserts the guard
    rather than LibreOffice's pagination.
    """
    _stub_multi_page_conversion(monkeypatch, pages=2)
    workbook = _two_sheet_workbook(tmp_path / "deal.xlsx")

    with pytest.raises(RuntimeError) as exc:
        e2p._libreoffice_range_to_png(workbook, "Ownership", "B4:G17")

    message = str(exc.value)
    assert "2 PDF pages" in message
    assert "Ownership!B4:G17" in message
    assert "deal.xlsx" in message
    assert "_prepare_print_copy" in message


def test_a_single_page_export_is_rendered(tmp_path: Path, monkeypatch):
    """The other side of the assertion: one page still renders, so the guard is a
    guard and not a blanket refusal."""
    _stub_multi_page_conversion(monkeypatch, pages=1)
    workbook = _two_sheet_workbook(tmp_path / "deal.xlsx")

    buf = e2p._libreoffice_range_to_png(workbook, "Ownership", "B4:G17")

    assert buf.getvalue().startswith(b"\x89PNG")


def test_the_one_sheet_one_page_rule_is_documented_next_to_the_private_copy_rule():
    """The drift lock on the rule itself.

    The clearing and the active-tab reset both look like tidying, so a future edit
    can drop either without a test noticing *why* they were there — which is what
    the docstring is for. Pinned in both places a reader looks: the module that
    implements it, and `CLAUDE.md`'s rendering rules, where the private-copy rule
    it sits beside already lives.
    """
    module_doc = e2p.__doc__
    assert "private copy" in module_doc
    for token in ("print area", "active tab", "exactly one\n  page"):
        assert token in module_doc, f"the module docstring no longer states {token!r}"

    brief = (_REPO_ROOT / "CLAUDE.md").read_text(encoding="utf-8")
    rendering = brief.split("### Rendering and geometry")[1].split("\n### ")[0]
    assert "private copy" in rendering
    assert "print area" in rendering and "one page" in rendering, (
        "CLAUDE.md's Rendering and geometry section no longer carries the "
        "one-sheet-one-page rule; it is the only place a fresh session reads it."
    )


# ─── Range pictures are identifiable, and never silently the same picture ─────


def _png(colour: tuple[int, int, int]) -> io.BytesIO:
    from PIL import Image

    buf = io.BytesIO()
    Image.new("RGB", (40, 20), colour).save(buf, format="PNG")
    buf.seek(0)
    return buf


def _deck_with_placeholders(path: Path, *names: str) -> Path:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for offset, name in enumerate(names):
        box = slide.shapes.add_textbox(
            Inches(1), Inches(1 + offset), Inches(3), Inches(0.8)
        )
        box.name = name
    prs.save(path)
    return path


def test_an_inserted_range_picture_is_named_after_its_placeholder(tmp_path: Path, monkeypatch):
    """`Picture 4` says nothing about which range it holds, and gives the
    duplicate-embed guard nothing to scope itself by."""
    monkeypatch.setattr(e2p, "_render_range_to_png", lambda *_a: _png((10, 20, 30)))
    deck = _deck_with_placeholders(tmp_path / "deck.pptx", "Rectangle 1")

    insert_excel_into_placeholder(
        deck_path=deck,
        workbook_path=_two_sheet_workbook(tmp_path / "deal.xlsx"),
        placeholder_name="Rectangle 1",
        source_range="B4:G17",
        sheet_name="Ownership",
    )

    shapes = {s.name for s in Presentation(deck).slides[0].shapes}
    assert shapes == {range_picture_name("Rectangle 1")}


def test_two_range_pictures_sharing_one_embed_are_rejected(tmp_path: Path, monkeypatch):
    """The symptom-level guard: identical bytes are deduped to one r:embed, so the
    ownership slide shipped one cap-table picture referenced twice."""
    monkeypatch.setattr(e2p, "_render_range_to_png", lambda *_a: _png((10, 20, 30)))
    deck = _deck_with_placeholders(tmp_path / "deck.pptx", "Rectangle 1", "Rectangle 3")
    workbook = _two_sheet_workbook(tmp_path / "deal.xlsx")
    for placeholder in ("Rectangle 1", "Rectangle 3"):
        insert_excel_into_placeholder(
            deck_path=deck,
            workbook_path=workbook,
            placeholder_name=placeholder,
            source_range="B4:G17",
            sheet_name="Ownership",
        )

    with pytest.raises(ValueError, match="sharing one embedded image"):
        assert_range_pictures_are_distinct(Presentation(deck))


def test_distinct_range_pictures_pass(tmp_path: Path, monkeypatch):
    colours = iter([(10, 20, 30), (200, 100, 50)])
    monkeypatch.setattr(e2p, "_render_range_to_png", lambda *_a: _png(next(colours)))
    deck = _deck_with_placeholders(tmp_path / "deck.pptx", "Rectangle 1", "Rectangle 3")
    workbook = _two_sheet_workbook(tmp_path / "deal.xlsx")
    for placeholder in ("Rectangle 1", "Rectangle 3"):
        insert_excel_into_placeholder(
            deck_path=deck,
            workbook_path=workbook,
            placeholder_name=placeholder,
            source_range="B4:G17",
            sheet_name="Ownership",
        )

    assert_range_pictures_are_distinct(Presentation(deck))  # does not raise


def test_the_guard_ignores_repeated_library_logos(tmp_path: Path):
    """Why the guard is scoped to range pictures rather than to every picture.

    The credentials slides in `INFOR Slide Library.pptx` show one client logo more
    than once, so an unscoped assertion would fail every pitch deck. Pinned here so
    a later "tighten this up" reads the reason first.
    """
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    logo = _png((1, 2, 3)).getvalue()
    for offset in range(2):
        slide.shapes.add_picture(
            io.BytesIO(logo), Inches(1), Inches(1 + offset), Inches(1), Inches(1)
        )
    embeds = {s._element.blip_rId for s in slide.shapes}
    assert len(embeds) == 1, "python-pptx no longer dedupes identical image bytes"

    assert_range_pictures_are_distinct(prs)  # does not raise


def test_libreoffice_captable_render_is_covered():
    """The LibreOffice sibling is a DIFFERENT surface — it is what should keep
    running by default, so it must NOT pick up the Excel-COM gate."""
    lines = (
        (Path(__file__).resolve().parent / "test_earnings_update_assembler.py")
        .read_text(encoding="utf-8")
        .splitlines()
    )
    idx = next(
        i
        for i, line in enumerate(lines)
        if line.startswith("def test_assemble_earnings_update_deck_inserts_cap_table_via_libreoffice(")
    )
    assert "@pytest.mark.excel_com" not in lines[idx - 1]
