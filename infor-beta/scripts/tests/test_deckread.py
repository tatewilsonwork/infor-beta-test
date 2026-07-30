"""Unit tests for the read-the-slides pass (scripts/deckread.py) and its stage.

What is testable here is the mechanical half plus the contract around the
judgement. The judgement itself — *is* that label pileup, *is* that table legible
at 4.5" — is a model looking at a PNG, and no assertion stands in for it.

So: the work list singles out a slide with a real defect on it, a finding cannot be
constructed with a gating severity, the report says what was read before it says
what was found, and a pass that rendered nothing says so instead of reading as a
clean deck. That last one is the failure this stage exists to end — an agenda that
nobody answered reported clean for four releases.
"""

from __future__ import annotations

import re
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from deck_contract import SEVERITY_ADVISORY, SEVERITY_BLOCKING, VISION_CHECKLIST
from deckread import (
    NOT_YOURS,
    READ_KINDS,
    DeckReadError,
    SlideFinding,
    existing_renders,
    read_deck,
    render_report,
    render_worklist,
    write_report,
    write_worklist,
)

_SKILLS_DIR = Path(__file__).resolve().parents[2] / "skills"
_SKILL_DOC = _SKILLS_DIR / "deckread" / "SKILL.md"


def _textbox(slide, *, left, top, width, height, text):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    frame = box.text_frame
    frame.text = text
    frame.paragraphs[0].runs[0].font.size = Pt(18)
    return box


@pytest.fixture
def deck_with_a_defect(tmp_path: Path) -> Path:
    """Slide 1 clean, slide 2 carrying two text boxes drawn on top of each other.

    Overlapping declared boxes are exactly what `deck_contract`'s vision tier flags
    and refuses to adjudicate: overlap is routine and legitimate in PowerPoint (a
    label over a filled band), so a reader has to look. That is the defect shape this
    stage exists to put in front of one.
    """
    prs = Presentation()
    blank = prs.slide_layouts[6]

    clean = prs.slides.add_slide(blank)
    _textbox(clean, left=1, top=1, width=4, height=1, text="Nothing wrong here")

    defective = prs.slides.add_slide(blank)
    _textbox(defective, left=1, top=1, width=4, height=1.2, text="Revenue by segment")
    _textbox(defective, left=1.4, top=1.1, width=4, height=1.2, text="US$999.9MM total")

    path = tmp_path / "deck-with-a-defect.pptx"
    prs.save(path)
    return path


# ─── The mechanical half ─────────────────────────────────────────────────────


def test_the_work_list_singles_out_the_slide_with_the_defect(deck_with_a_defect, tmp_path):
    # `renders={}` skips the LibreOffice conversion: the overlap is read off the
    # declared boxes, so the agenda does not depend on a renderer.
    evidence = read_deck(deck_with_a_defect, out_dir=tmp_path / "evidence", renders={})

    assert evidence.slide_count == 2
    assert evidence.flagged_slides == (1,), (
        f"expected only the defective slide on the work list, got "
        f"{[i + 1 for i in evidence.flagged_slides]}"
    )
    questions = " ".join(t.question for t in evidence.vision.targets)
    assert "overlap" in questions

    worklist = render_worklist(evidence)
    assert "## Slide 2" in worklist
    assert "## Slide 1" not in worklist


def test_the_work_list_reuses_the_deck_stages_checklist_and_names_it(
    deck_with_a_defect, tmp_path
):
    """One wording for the checklist, and the pre-chart review named, not re-derived.

    `render_worklist` delegates its body to `stage_transforms.render_vision_review` —
    the same renderer the `deck` transform writes `vision_review_path` with. A second
    copy of the checklist is the drift H1 collapsed on the intake side, and this pass
    exists because that agenda had no reader, not because it was worded wrong.
    """
    evidence = read_deck(deck_with_a_defect, out_dir=tmp_path / "evidence", renders={})
    checklist = tmp_path / "runs" / "stages" / "deck" / "vision_review.md"

    worklist = render_worklist(evidence, checklist_path=checklist)

    assert VISION_CHECKLIST in worklist
    assert str(checklist) in worklist
    # The boundary travels with the work list, so it is in front of the reader at the
    # moment they would otherwise report a deferred placeholder or a CapIQ error.
    for what, _ in NOT_YOURS:
        assert what in worklist


def test_write_worklist_persists_it(deck_with_a_defect, tmp_path):
    evidence = read_deck(deck_with_a_defect, out_dir=tmp_path / "evidence", renders={})
    path = write_worklist(tmp_path / "stage" / "worklist.md", evidence)
    assert path.is_file()
    assert "Deck read — work list" in path.read_text(encoding="utf-8")


def test_read_deck_refuses_a_deck_that_is_not_there(tmp_path):
    with pytest.raises(DeckReadError, match="no deck to read"):
        read_deck(tmp_path / "gone.pptx", out_dir=tmp_path / "evidence", renders={})


def test_read_deck_renders_every_slide_and_reports_that_it_did(
    deck_with_a_defect, tmp_path
):
    """The renderer path, for real — `rendered` is the claim the report leans on.

    Not skippable: LibreOffice is the only render backend on every platform and the
    suite renders decks throughout, so "no renderer" here means the environment is
    broken, which is a failure and not a skip.
    """
    evidence = read_deck(deck_with_a_defect, out_dir=tmp_path / "evidence")

    assert evidence.rendered
    assert sorted(evidence.renders) == [0, 1]
    assert all(p.is_file() and p.stat().st_size > 0 for p in evidence.renders.values())


def test_the_second_snippet_reuses_the_renders_instead_of_converting_again(
    deck_with_a_defect, tmp_path, monkeypatch
):
    """The stage runs as two `python` invocations, i.e. two cold render caches.

    `existing_renders` is how the report-writing snippet rebuilds the same evidence
    for free. On a 19-slide deck over a cloud-synced deal directory, a second
    conversion is not a rounding error — and a re-render would also mean the report
    describes different files from the ones the reader just read.
    """
    out = tmp_path / "evidence"
    first = read_deck(deck_with_a_defect, out_dir=out)

    def _no_second_render(*args, **kwargs):  # pragma: no cover — the point is it isn't called
        raise AssertionError("the deck was converted a second time")

    monkeypatch.setattr("deck_contract._render_slides", _no_second_render)
    second = read_deck(deck_with_a_defect, out_dir=out, renders=existing_renders(out))

    assert second.rendered
    assert second.renders == first.renders


def test_existing_renders_is_empty_when_nothing_was_rendered(tmp_path):
    # Which `read_deck` treats exactly as a host with no LibreOffice does.
    assert existing_renders(tmp_path / "never-written") == {}


# ─── Findings ────────────────────────────────────────────────────────────────


def _finding(**kwargs) -> SlideFinding:
    base = {
        "slide": 1,
        "kind": "text-over-text",
        "issue": "the segment heading is drawn over the total",
        "detail": "'Revenue by segment' and 'US$999.9MM total' occupy the same band",
        "shape": "TextBox 2",
        "evidence": "/deals/runs/stages/deckread/evidence/slides/slide_02.png",
    }
    return SlideFinding(**{**base, **kwargs})


def test_a_finding_defaults_to_advisory():
    assert _finding().severity == SEVERITY_ADVISORY
    assert _finding().slide_number == 2


def test_a_finding_cannot_be_constructed_with_a_gating_severity():
    """The same construction-time refusal `deckcheck.CheckFinding` carries.

    No shipped stage gates (v0.5.49), and this one is the least entitled to: the deck
    it is reviewing already converged against measured renders, so a blocking finding
    here would be an opinion overruling a measurement.
    """
    with pytest.raises(ValueError, match="always 'advisory'"):
        _finding(severity=SEVERITY_BLOCKING)


def test_a_finding_refuses_an_unrecognised_kind():
    # A free-text kind is a field nothing can group, count, or act on.
    with pytest.raises(ValueError, match="unknown finding kind"):
        _finding(kind="looks-a-bit-off")
    assert "label-pileup" in READ_KINDS


# ─── The report ──────────────────────────────────────────────────────────────


def test_the_report_carries_the_findings_and_the_evidence(deck_with_a_defect, tmp_path):
    evidence = read_deck(deck_with_a_defect, out_dir=tmp_path / "evidence")
    findings = [_finding()]

    report = render_report(evidence, findings, company="SampleCo Ltd.")

    assert "# Deck read — SampleCo Ltd." in report
    assert "not a gate" in report
    assert "the segment heading is drawn over the total" in report
    assert "`text-over-text`" in report
    assert "slide_02.png" in report          # the finding's own evidence
    assert "Slides read: 2 of 2" in report
    assert "Slides on the work list: 1 (2)" in report


def test_a_clean_read_says_every_slide_was_read(deck_with_a_defect, tmp_path):
    evidence = read_deck(deck_with_a_defect, out_dir=tmp_path / "evidence")
    report = render_report(evidence, [], company="SampleCo Ltd.")
    assert "Every slide was read" in report
    assert VISION_CHECKLIST in report
    assert "did not run" not in report


def test_a_report_with_no_renders_refuses_to_read_as_a_clean_deck(
    deck_with_a_defect, tmp_path
):
    """The failure this stage exists to end, asserted from the other direction.

    "no findings" and "no renderer" are the same file otherwise — which is how a run
    with an unread 19 KB agenda reported clean. So a pass that rendered nothing says
    so above any finding, says it again where a clean verdict would go, and carries
    `slides_read = 0` for the conductor to surface.
    """
    evidence = read_deck(deck_with_a_defect, out_dir=tmp_path / "evidence", renders={})
    assert not evidence.rendered

    report = render_report(evidence, [], company="SampleCo Ltd.")

    assert "The visual half did not run" in report
    assert "no slide of this deck was read" in report
    assert "Every slide was read" not in report
    assert "Slides read: 0 of 2" in report


def test_write_report_persists_it(deck_with_a_defect, tmp_path):
    evidence = read_deck(deck_with_a_defect, out_dir=tmp_path / "evidence", renders={})
    path = write_report(
        tmp_path / "artefacts" / "deckread-Project Sample.md",
        evidence,
        [_finding()],
        company="SampleCo Ltd.",
        notes=["Slide 1 is a deferred placeholder."],
    )
    text = path.read_text(encoding="utf-8")
    assert path.is_file()
    assert "Slide 1 is a deferred placeholder." in text


# ─── The stage's own contract ────────────────────────────────────────────────


def _frontmatter(doc: str) -> str:
    return doc.split("---", 2)[1]


def test_the_stage_doc_exists_and_names_itself_after_its_directory():
    # A dispatched skill's `name:` must equal its directory, or the plan's `skill:`
    # string resolves to a doc that does not describe it.
    assert _SKILL_DOC.is_file()
    assert re.search(r"^name: deckread$", _frontmatter(_SKILL_DOC.read_text(encoding="utf-8")), re.M)


def test_the_stage_can_neither_dispatch_nor_edit_the_deck():
    """Its allow-list is the enforcement, the way `financial-charts`' used to be.

    No `Task`: a review that could dispatch could re-run the assembly it is
    reviewing. No `Edit`: a repair here would silently undo the converged state the
    assembler measured, and the deck is the deliverable.
    """
    tools = re.search(
        r"^allowed-tools: \[(.+)\]$",
        _frontmatter(_SKILL_DOC.read_text(encoding="utf-8")),
        re.M,
    )
    assert tools, "deckread's SKILL.md no longer declares an allowed-tools list"
    granted = {t.strip() for t in tools.group(1).split(",")}
    assert "Read" in granted, "a reading pass has to be able to read a PNG"
    assert not granted & {"Task", "Edit", "NotebookEdit"}, (
        f"deckread must not be able to dispatch or edit: {sorted(granted)}"
    )


def test_the_stage_doc_states_the_advisory_rule_and_the_unread_rule():
    doc = _SKILL_DOC.read_text(encoding="utf-8")
    # Wrapping is the doc's business, so match on the reflowed text.
    flowed = " ".join(doc.split())
    assert "advisory" in doc.lower()
    assert "Do not report a clean deck when nothing looked at it." in flowed
    # And it must not claim to own what it cannot: the boundary in prose, matching
    # `NOT_YOURS` in code.
    assert "`deckread.NOT_YOURS`" in doc
