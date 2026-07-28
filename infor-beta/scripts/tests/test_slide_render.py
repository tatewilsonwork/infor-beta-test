"""Unit tests for the slide-to-PNG overflow QA renderer.

Backend selection is tested on every platform (it is pure dispatch); the actual
render tests skip when the backend binary is absent.
"""

from pathlib import Path

import pytest

import slide_render
from excel_to_powerpoint import find_soffice
from slide_render import BACKEND_ENV_VAR, BACKEND_LIBREOFFICE, render_deck_to_png

_LIBRARY = Path("infor-beta/templates/INFOR Slide Library.pptx")


def _libreoffice_available() -> bool:
    return find_soffice() is not None


@pytest.fixture
def stub_backends(monkeypatch):
    """Replace the renderer with a recorder so dispatch can be asserted.

    One backend since Phase D deleted the PowerPoint-COM path, so there is only
    one recorder — the fixture name is kept because the assertions read the same.
    """
    calls: list[str] = []
    monkeypatch.setattr(
        slide_render, "_libreoffice_render", lambda *a, **k: calls.append(BACKEND_LIBREOFFICE) or []
    )
    monkeypatch.delenv(BACKEND_ENV_VAR, raising=False)
    return calls


def _dummy_deck(tmp_path: Path) -> Path:
    deck = tmp_path / "deck.pptx"
    deck.write_bytes(b"not a real pptx - backend is stubbed")
    return deck


# ─── Backend selection (v0.5.35 render parity) ───────────────────────────────


def test_default_backend_is_libreoffice_even_on_windows(tmp_path: Path, stub_backends):
    """The v0.5.35 flip: Windows dev must render the way Cowork/Linux prod does.

    Before the flip this dispatched to PowerPoint COM on win32, which is why a
    production rendering bug could not be reproduced locally.
    """
    render_deck_to_png(_dummy_deck(tmp_path), tmp_path / "out")

    assert stub_backends == [BACKEND_LIBREOFFICE]
def test_env_var_naming_libreoffice_still_works(tmp_path: Path, stub_backends, monkeypatch):
    """A leftover `=libreoffice` in an environment keeps working, case-tolerantly."""
    monkeypatch.setenv(BACKEND_ENV_VAR, "LibreOffice")

    render_deck_to_png(_dummy_deck(tmp_path), tmp_path / "out")

    assert stub_backends == [BACKEND_LIBREOFFICE]


def test_the_removed_powerpoint_backend_is_rejected_not_ignored(tmp_path: Path, stub_backends):
    """Asking for the deleted COM backend must RAISE, not quietly use LibreOffice.

    A stale caller (or a stale `INFOR_SLIDE_RENDER_BACKEND=powerpoint` in someone's
    shell) that silently got a different renderer is exactly the dev/prod
    ambiguity Phase A set out to remove.
    """
    with pytest.raises(ValueError, match="Phase D deleted the PowerPoint-COM"):
        render_deck_to_png(_dummy_deck(tmp_path), tmp_path / "out", backend="powerpoint")
    assert stub_backends == []


def test_libreoffice_failure_raises_loudly(tmp_path: Path, monkeypatch):
    """A missing LibreOffice must fail, not degrade — there is nothing to degrade to."""
    def boom(*a, **k):
        raise RuntimeError("LibreOffice not found on PATH")

    monkeypatch.setattr(slide_render, "_libreoffice_render", boom)
    monkeypatch.delenv(BACKEND_ENV_VAR, raising=False)

    with pytest.raises(RuntimeError, match="LibreOffice"):
        render_deck_to_png(_dummy_deck(tmp_path), tmp_path / "out")


def test_unknown_backend_raises(tmp_path: Path, stub_backends):
    with pytest.raises(ValueError, match="unknown slide-render backend"):
        render_deck_to_png(_dummy_deck(tmp_path), tmp_path / "out", backend="ghostscript")


def test_render_missing_deck_raises(tmp_path: Path):
    with pytest.raises(FileNotFoundError):
        render_deck_to_png(tmp_path / "nope.pptx", tmp_path / "out")


# ─── Real renders ────────────────────────────────────────────────────────────


@pytest.mark.skipif(not _libreoffice_available(), reason="LibreOffice not installed")
def test_render_selected_slides_returns_png_paths(tmp_path: Path):
    if not _LIBRARY.exists():
        pytest.skip("slide library template not present")
    pytest.importorskip("pypdfium2", reason="LibreOffice backend needs pypdfium2")

    out_dir = tmp_path / "png"
    paths = render_deck_to_png(_LIBRARY, out_dir, slide_indices=[0, 6])

    assert len(paths) == 2
    for p in paths:
        assert p.exists()
        assert p.suffix == ".png"
        assert p.stat().st_size > 0


# ─── Converted-PDF cache (Phase C) ───────────────────────────────────────────
# Correctness first: the cache must key on CONTENT, must not serve a stale PDF
# for changed content, and must be switchable off.


@pytest.fixture
def private_cache(tmp_path: Path, monkeypatch):
    """Give this test its own render cache directory.

    The suite points every process at ONE shared cache (see `conftest.py`), so a
    test that cleared it would be deleting PDFs the other five workers are
    reading — wasteful at best, and a genuine race at worst. Isolating is both
    safer and a better test: the cache starts empty either way.
    """
    cache = tmp_path / "render-cache"
    monkeypatch.setenv(slide_render.CACHE_DIR_ENV_VAR, str(cache))
    monkeypatch.setattr(slide_render, "_CACHE_DIR", None)
    monkeypatch.setattr(slide_render, "_PDF_CACHE", {})
    return cache


def _one_slide_deck(path: Path, text: str) -> Path:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = text
    prs.save(str(path))
    return path


def test_digest_ignores_zip_timestamps_but_not_content(tmp_path: Path):
    """The key that makes the cache useful.

    python-pptx stamps each zip member with the save time, so two decks built
    from identical content differ byte-for-byte. Hashing the members' names and
    payloads sees them as the same deck — which is what lets a regenerated probe
    deck hit the cache. Different content must still miss.
    """
    import time as _time

    a = _one_slide_deck(tmp_path / "a.pptx", "same words")
    _time.sleep(2)  # push the zip timestamp into a different second
    b = _one_slide_deck(tmp_path / "b.pptx", "same words")
    c = _one_slide_deck(tmp_path / "c.pptx", "different words")

    assert a.read_bytes() != b.read_bytes(), "expected differing zip timestamps"
    assert slide_render._deck_digest(a) == slide_render._deck_digest(b)
    assert slide_render._deck_digest(a) != slide_render._deck_digest(c)
    assert slide_render._deck_digest(tmp_path / "not-a-zip.pptx") is None


@pytest.mark.skipif(not _libreoffice_available(), reason="LibreOffice not installed")
def test_identical_content_is_converted_once(tmp_path: Path, private_cache):
    pytest.importorskip("pypdfium2", reason="LibreOffice backend needs pypdfium2")
    conversions: list[Path] = []
    real = slide_render._convert_to_pdf

    def counting(soffice, deck, tmp_dir):
        conversions.append(deck)
        return real(soffice, deck, tmp_dir)

    original, slide_render._convert_to_pdf = slide_render._convert_to_pdf, counting
    try:
        a = _one_slide_deck(tmp_path / "a.pptx", "cache me")
        b = _one_slide_deck(tmp_path / "b.pptx", "cache me")  # same content, new file
        c = _one_slide_deck(tmp_path / "c.pptx", "do not cache me")
        first = render_deck_to_png(a, tmp_path / "o1")
        second = render_deck_to_png(b, tmp_path / "o2")
        third = render_deck_to_png(c, tmp_path / "o3")
    finally:
        slide_render._convert_to_pdf = original

    assert len(conversions) == 2, "identical content must convert once; new content must convert"
    # Served from cache, but still a real render of the right page.
    for paths in (first, second, third):
        assert len(paths) == 1 and paths[0].stat().st_size > 0
    assert first[0].read_bytes() == second[0].read_bytes()


@pytest.mark.skipif(not _libreoffice_available(), reason="LibreOffice not installed")
def test_cache_can_be_switched_off(tmp_path: Path, monkeypatch, private_cache):
    pytest.importorskip("pypdfium2", reason="LibreOffice backend needs pypdfium2")
    monkeypatch.setenv(slide_render.CACHE_ENV_VAR, "0")
    conversions: list[Path] = []
    real = slide_render._convert_to_pdf

    def counting(soffice, deck, tmp_dir):
        conversions.append(deck)
        return real(soffice, deck, tmp_dir)

    original, slide_render._convert_to_pdf = slide_render._convert_to_pdf, counting
    try:
        deck = _one_slide_deck(tmp_path / "a.pptx", "no cache")
        render_deck_to_png(deck, tmp_path / "o1")
        render_deck_to_png(deck, tmp_path / "o2")
    finally:
        slide_render._convert_to_pdf = original

    assert len(conversions) == 2


def test_cache_is_shared_between_processes(tmp_path: Path, monkeypatch):
    """A PDF one process published is served to the next.

    This is what makes the distributed suite worthwhile: the six workers would
    otherwise each convert the blank library and its attribution probe deck.
    Simulated here by publishing into the shared directory, then clearing only
    the in-process index — which is exactly what a second process sees.
    """
    shared = tmp_path / "shared-cache"
    monkeypatch.setenv(slide_render.CACHE_DIR_ENV_VAR, str(shared))
    monkeypatch.setattr(slide_render, "_CACHE_DIR", None)
    monkeypatch.setattr(slide_render, "_PDF_CACHE", {})

    pdf = tmp_path / "converted.pdf"
    pdf.write_bytes(b"%PDF-1.4 stub")
    slide_render._cache_store("deadbeef", pdf)
    assert (shared / "deadbeef.pdf").is_file()
    assert not list(shared.glob("*.part")), "the staging file must not be left behind"

    # A fresh process: same directory, empty in-process index.
    monkeypatch.setattr(slide_render, "_PDF_CACHE", {})
    conversions = []

    def must_not_convert(soffice, deck, tmp_dir):
        conversions.append(deck)
        raise AssertionError("should have been served from the shared cache")

    monkeypatch.setattr(slide_render, "_convert_to_pdf", must_not_convert)
    monkeypatch.setattr(slide_render, "_deck_digest", lambda deck: "deadbeef")
    assert slide_render._pdf_for("soffice", tmp_path / "any.pptx", tmp_path) == shared / "deadbeef.pdf"
    assert conversions == []


def test_a_shared_cache_dir_is_not_deleted_by_a_process_that_joined_it(tmp_path, monkeypatch):
    # Only the creator cleans up; a worker must not remove the directory its
    # siblings are still reading.
    shared = tmp_path / "shared-cache"
    monkeypatch.setenv(slide_render.CACHE_DIR_ENV_VAR, str(shared))
    monkeypatch.setattr(slide_render, "_CACHE_DIR", None)
    monkeypatch.setattr(slide_render, "_CACHE_DIR_IS_OURS", False)
    assert slide_render._cache_dir() == shared
    assert slide_render._CACHE_DIR_IS_OURS is False

    monkeypatch.setattr(slide_render, "_CACHE_DIR", None)
    monkeypatch.delenv(slide_render.CACHE_DIR_ENV_VAR)
    private = slide_render._cache_dir()
    assert private != shared and slide_render._CACHE_DIR_IS_OURS is True


def test_cache_eviction_is_bounded(tmp_path: Path, monkeypatch, private_cache):
    monkeypatch.setattr(slide_render, "_CACHE_MAX_ENTRIES", 3)
    for i in range(6):
        pdf = tmp_path / f"{i}.pdf"
        pdf.write_bytes(b"%PDF-1.4 stub")
        slide_render._cache_store(f"digest{i}", pdf)
    assert len(slide_render._PDF_CACHE) <= 3
    for path in slide_render._PDF_CACHE.values():
        assert path.is_file(), "a surviving entry must still be on disk"
