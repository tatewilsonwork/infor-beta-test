"""Unit tests for the POC earnings-update deck assembler (slide-library based)."""

import shutil
from pathlib import Path

import pytest
from openpyxl import Workbook
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from earnings_update_assembler import _currency_letter, assemble_earnings_update_deck
from earnings_update_wireframe import build_earnings_update_slide_plan, write_slide_plan
from pptx_helpers import find_shape, find_shape_in_group
from schemas import Company, EarningsUpdateContent

_TEMPLATE = Path("infor-beta/templates/INFOR Slide Library.pptx")


def _sample_content() -> EarningsUpdateContent:
    return EarningsUpdateContent(
        company_name="SampleCo",
        ticker="TSX:SMPL",
        reporting_quarter="Q4 2025",
        comparison_quarter="Q4 2024",
        currency="C$MM",
        currency_short="C$MM",
        cover_date="May 2026",
        company_overview_bullets=[
            {"text": "Leading provider of mission-critical compliance and workflow software serving blue-chip enterprises across Canada", "level": 0},
            {"text": "Recurring revenue model supported by multi-year contracts and a high-retention installed base in regulated markets", "level": 0},
            {"text": "Diversified product suite addressing daily operational needs for finance, legal and compliance teams", "level": 0},
            {"text": "Experienced management team with a record of disciplined execution and successful tuck-in acquisitions", "level": 0},
            {"text": "Strong balance sheet supporting organic growth, selective acquisitions and continued platform investment", "level": 0},
            {"text": "Established go-to-market platform with direct sales coverage and a repeatable land-and-expand motion", "level": 0},
            {"text": "Well-positioned to benefit from the ongoing digitization of enterprise compliance workflows", "level": 0},
        ],
        business_updates=[
            "Revenue growth reflected continued enterprise demand and disciplined customer expansion",
            "Management highlighted improving sales productivity and constructive renewal activity",
            "The Company continued to invest in product development while maintaining margin discipline",
            "Near-term priorities remain focused on enterprise execution and cash conversion",
        ],
        kpi_rows=[
            {"name": "Revenue", "prior_value": "100", "current_value": "125", "delta_str": "+25", "delta_sign": 1},
            {"name": "Adjusted EBITDA", "prior_value": "20", "current_value": "22", "delta_str": "+2", "delta_sign": 1},
            {"name": "Net Income", "prior_value": "8", "current_value": "7", "delta_str": "-1", "delta_sign": -1},
            {"name": "Gross Margin", "prior_value": "60.0%", "current_value": "62.0%", "delta_str": "+2.0%", "delta_sign": 1},
        ],
        broker_rows=[
            {"label": "Revenue", "reported": "125", "estimate": "120", "variance": "+5", "variance_sign": 1},
            {"label": "Adjusted EBITDA", "reported": "22", "estimate": "21", "variance": "+1", "variance_sign": 1},
            {"label": "EPS (C$)", "reported": "0.10", "estimate": "0.12", "variance": "(0.02)", "variance_sign": -1},
            {"label": "Operating income", "reported": "18", "estimate": "17", "variance": "+1", "variance_sign": 1},
            {"label": "Gross Margin", "reported": "62.0%", "estimate": "61.0%", "variance": "+1.0%", "variance_sign": 1},
        ],
        management_quotes=[
            {"quote": "We delivered a strong quarter driven by enterprise execution and continued customer expansion", "speaker": "Jane Doe", "role": "CEO"},
            {"quote": "Our focus remains on profitable growth, cash conversion and disciplined capital allocation", "speaker": "John Smith", "role": "CFO"},
        ],
        performance_summary="SampleCo beat consensus revenue while maintaining disciplined operating execution",
        sources=[{"section": "All", "citation": "Company filings, S&P CapIQ, equity research"}],
        manual_steps=["Refresh CapIQ in the companion cap table before pasting into slide 2"],
    )


def _write_sample_cap_table(path: Path) -> Path:
    wb = Workbook()
    ws = wb.active
    ws.title = "Cap with Links"
    ws["B15"] = "SampleCo Cap Table"
    rows = {
        15: ("Company Ticker:", "TSX:SMPL"),
        16: ("Share Price (21-Apr-26)", "C$12.34"),
        17: ("Basic Shares Outstanding", "100.0"),
        18: ("Basic Market Cap", "C$1,234.0"),
        21: ("Fully-Diluted Shares Outstanding", "104.0"),
        22: ("Fully-Diluted Market Cap", "C$1,283.4"),
        28: ("Net Debt", "C$200.0"),
        31: ("Enterprise Value", "C$1,483.4"),
        # B40 pins the bottom of the B15:F40 picture range — the assembler
        # verifies both sentinel anchors (template_layout) before pasting.
        40: ("EV / Adj. EBITDA", "10.0x"),
    }
    for row, (label, value) in rows.items():
        ws.cell(row=row, column=2).value = label
        ws.cell(row=row, column=6).value = value
    wb.save(path)
    return path


def _assemble_sample_deck(tmp_path: Path, content: EarningsUpdateContent | None = None, **kwargs) -> Path:
    slide_plan = build_earnings_update_slide_plan(
        company=Company(legal_name="SampleCo", ticker="TSX:SMPL"),
        reporting_quarter="Q4 2025",
        comparison_quarter="Q4 2024",
    )
    slide_plan_path = write_slide_plan(slide_plan, tmp_path / "slide_plan.json")
    content = content or _sample_content()
    content_path = tmp_path / "content.json"
    content_path.write_text(content.model_dump_json(indent=2), encoding="utf-8")

    return assemble_earnings_update_deck(
        slide_plan_path=slide_plan_path,
        content_path=content_path,
        template_path=_TEMPLATE,
        output_dir=tmp_path,
        **kwargs,
    )


def _slide_text(slide) -> str:
    parts = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            parts.append(shape.text)
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub in shape.shapes:
                if getattr(sub, "has_text_frame", False):
                    parts.append(sub.text)
    return "\n".join(parts)


def test_assemble_earnings_update_deck_clones_five_library_slides(tmp_path: Path):
    deck_path = _assemble_sample_deck(tmp_path)

    assert deck_path.exists()
    prs = Presentation(deck_path)
    assert len(prs.slides) == 5

    overview_text = _slide_text(prs.slides[1])
    summary_text = _slide_text(prs.slides[2])

    assert "Introduction to SampleCo" in overview_text
    assert "SampleCo Q4 2025 Earnings Summary" in summary_text
    # Pie + cap-table placeholders remain when no workbook is supplied.
    assert "[Pie Chart Placeholder]" in overview_text
    assert "[Cap Table Placeholder]" in overview_text
    for token in ("[Company]", "[Quarter]", "[Date]", "[Name]", "[Role]"):
        assert token not in overview_text + summary_text

    # Regression guard: the v0.5.8 ownership-slide insertion shifted the library's
    # disclaimer/contact closers, and the pre-0.5.9 keep-indices (0,6,7,13,14)
    # dropped the Contact slide and kept a stray Ownership slide. The earnings deck
    # must end cover → overview → earnings summary → disclaimer → contact.
    def _deep_text(slide) -> str:
        parts: list[str] = []

        def walk(shapes):
            for sh in shapes:
                if getattr(sh, "has_text_frame", False):
                    parts.append(sh.text)
                if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
                    walk(sh.shapes)

        walk(slide.shapes)
        return "\n".join(parts)

    all_text = "\n".join(_deep_text(s) for s in prs.slides)
    assert "Contact" in _deep_text(prs.slides[4]), "Contact must be the last slide (dropped pre-0.5.9)"
    assert "[Placeholder for Insider Ownership]" not in all_text, "no stray Ownership slide in the earnings deck"


def test_period_label_lives_in_bar_not_metric_boxes(tmp_path: Path):
    deck_path = _assemble_sample_deck(tmp_path)
    summary = Presentation(deck_path).slides[2]

    # Mid-blue period bar carries the comparison / reporting quarters.
    assert find_shape(summary, "Rectangle 16").text_frame.text == "Q4 2024"
    assert find_shape(summary, "Rectangle 21").text_frame.text == "Q4 2025"

    # Metric boxes carry value + name only — no quarter label.
    group = find_shape(summary, "Group 12")
    current_box = find_shape_in_group(group, "Rectangle 1034")
    box_text = current_box.text_frame.text
    assert "125" in box_text
    assert "Revenue" in box_text
    assert "Q4 2025" not in box_text
    assert "Q4 2024" not in box_text


def test_metric_tiles_format_dollars_mm_and_billions(tmp_path: Path):
    content = _sample_content()
    content.kpi_rows[0].prior_value = "1254"
    content.kpi_rows[0].current_value = "1283"  # >= 1000 MM -> billions, 1 decimal
    content.kpi_rows[1].prior_value = "395"
    content.kpi_rows[1].current_value = "493"  # < 1000 MM -> whole millions

    deck_path = _assemble_sample_deck(tmp_path, content)
    summary = Presentation(deck_path).slides[2]

    g0 = find_shape(summary, "Group 12")
    assert "$1.3B" in find_shape_in_group(g0, "Rectangle 1034").text_frame.text
    assert "$1.3B" in find_shape_in_group(g0, "Rectangle 1032").text_frame.text

    g1 = find_shape(summary, "Group 9")
    assert "$493MM" in find_shape_in_group(g1, "Rectangle 1037").text_frame.text

    # A percent KPI tile is left untouched (no $ / MM forced onto it).
    g3 = find_shape(summary, "Group 2")
    assert "62.0%" in find_shape_in_group(g3, "Rectangle 1058").text_frame.text


def test_footnote_substitutes_currency_letter_and_keeps_library_source(tmp_path: Path):
    deck_path = _assemble_sample_deck(tmp_path)  # sample content is C$MM
    prs = Presentation(deck_path)

    overview_note = find_shape(prs.slides[1], "Text Placeholder 1").text_frame.text
    summary_note = find_shape(prs.slides[2], "Text Placeholder 1").text_frame.text

    # '[x]$MM' currency-letter token is replaced; no placeholder leftover.
    assert "C$MM" in overview_note and "[x]" not in overview_note
    assert "C$MM" in summary_note and "[x]" not in summary_note
    # The standardized library source line is preserved, not re-hardcoded.
    assert "S&P Capital IQ" in summary_note


def test_metric_box_name_strips_currency_unit(tmp_path: Path):
    content = _sample_content()
    content.kpi_rows[0].name = "Cloud Revenue (C$MM)"

    deck_path = _assemble_sample_deck(tmp_path, content)
    summary = Presentation(deck_path).slides[2]
    group = find_shape(summary, "Group 12")
    box_text = find_shape_in_group(group, "Rectangle 1034").text_frame.text
    assert "Cloud Revenue" in box_text
    assert "(C$MM)" not in box_text


def test_broker_values_get_dollar_prefix(tmp_path: Path):
    deck_path = _assemble_sample_deck(tmp_path)
    summary = Presentation(deck_path).slides[2]
    table = next(s for s in summary.shapes if getattr(s, "has_table", False)).table

    # Revenue row: plain "125" -> "$125"; variance "+5" -> "+$5".
    assert table.cell(1, 1).text_frame.text == "$125"
    assert table.cell(1, 2).text_frame.text == "$120"
    assert table.cell(1, 3).text_frame.text == "+$5"

    # Parenthesised negative variance keeps the paren convention.
    assert table.cell(3, 3).text_frame.text == "($0.02)"

    # Percent rows are left untouched (no $).
    assert table.cell(5, 1).text_frame.text == "62.0%"


def test_broker_labels_strip_mm_units(tmp_path: Path):
    content = _sample_content()
    content.broker_rows[0].label = "Revenue (US$MM)"
    content.broker_rows[1].label = "Adj. EBITDA (US$MM)"
    content.broker_rows[2].label = "EPS (US$)"  # per-share — not stripped
    content.broker_rows[3].label = "Operating income (C$MM)"
    content.broker_rows[4].label = "Free cashflow (MM)"

    deck_path = _assemble_sample_deck(tmp_path, content)
    summary = Presentation(deck_path).slides[2]
    table = next(s for s in summary.shapes if getattr(s, "has_table", False)).table
    labels = [table.cell(i, 0).text_frame.text for i in range(1, 6)]
    assert labels == [
        "Revenue",
        "Adj. EBITDA",
        "EPS (US$)",
        "Operating income",
        "Free cashflow",
    ]


def test_assemble_earnings_update_deck_enables_autofit_on_overflow_shapes(tmp_path: Path):
    """Overview bullets and business updates both need shrink-on-overflow."""
    deck_path = _assemble_sample_deck(tmp_path)

    prs = Presentation(deck_path)
    overview = find_shape(prs.slides[1], "TextBox 9")
    business = find_shape(prs.slides[2], "TextBox 6")

    for shape in (overview, business):
        bodyPr = shape.text_frame._txBody.find(qn("a:bodyPr"))
        assert bodyPr is not None, f"{shape.name} must have a bodyPr"
        autofit = bodyPr.find(qn("a:normAutofit"))
        assert autofit is not None, (
            f"{shape.name} must have <a:normAutofit/> to shrink on overflow"
        )


def test_assemble_earnings_update_deck_does_not_bold_overview_headers(tmp_path: Path):
    content = _sample_content()
    content.company_overview_bullets[0].bold_prefix = "Header:"
    content.company_overview_bullets[0].text = " Enterprise software provider with recurring revenue visibility across regulated markets"

    deck_path = _assemble_sample_deck(tmp_path, content)

    overview_shape = find_shape(Presentation(deck_path).slides[1], "TextBox 9")
    header_runs = [
        run
        for para in overview_shape.text_frame.paragraphs
        for run in para.runs
        if "Header:" in run.text
    ]
    assert header_runs, "expected the overview bullet prefix text to be present"
    assert all(run.font.bold is not True for run in header_runs)


def test_assemble_earnings_update_deck_inserts_cap_table_from_workbook(tmp_path: Path):
    pytest.importorskip("win32com.client", reason="picture-based insertion requires pywin32 + Excel")

    workbook_path = _write_sample_cap_table(tmp_path / "cap-table.xlsx")
    deck_path = _assemble_sample_deck(tmp_path, captable_workbook_path=workbook_path)

    prs = Presentation(deck_path)
    slide2 = prs.slides[1]
    assert next((s for s in slide2.shapes if s.name == "Rectangle 3"), None) is None, (
        "overview cap-table placeholder should be removed after picture insertion"
    )
    pictures = [s for s in slide2.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert pictures, "expected a picture shape on slide 2 after cap-table insertion"


def test_assemble_earnings_update_deck_inserts_cap_table_via_libreoffice(tmp_path: Path):
    """LibreOffice fallback path for Cowork / Linux runs."""
    import sys

    if sys.platform == "win32":
        pytest.skip("Windows uses the Excel COM path; LibreOffice fallback exercised on other OSes")
    if shutil.which("soffice") is None and shutil.which("libreoffice") is None:
        pytest.skip("LibreOffice not installed; cannot exercise the fallback renderer")
    pytest.importorskip("pypdfium2", reason="pypdfium2 required for the LibreOffice fallback")

    workbook_path = _write_sample_cap_table(tmp_path / "cap-table.xlsx")
    deck_path = _assemble_sample_deck(tmp_path, captable_workbook_path=workbook_path)

    prs = Presentation(deck_path)
    slide2 = prs.slides[1]
    assert next((s for s in slide2.shapes if s.name == "Rectangle 3"), None) is None
    pictures = [s for s in slide2.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert pictures, "expected a picture shape on slide 2 after LibreOffice insertion"

# ─── Footnote currency mapping (v0.5.34) ─────────────────────────────────────


def test_currency_letter_maps_dollar_scopes_explicitly():
    assert _currency_letter("US$MM") == "US"
    assert _currency_letter("C$MM") == "C"
    assert _currency_letter("USD") == "US"
    assert _currency_letter("CAD") == "C"


def test_currency_letter_returns_iso_code_for_non_dollar_currency():
    # A GBP filer's footnote must render the ISO code, never a wrong dollar
    # sign and never a silent 'C' default.
    assert _currency_letter("GBPMM") == "GBP"
    assert _currency_letter("GBP") == "GBP"
    assert _currency_letter("EUR MM") == "EUR"
    assert _currency_letter("CHF") == "CHF"  # would have read as 'C' pre-v0.5.34
