"""The falsification pass — figure extraction, the ledger join, and the report.

Asserted here is the **machine** half only, and deliberately: what a reviewer
concluded after reading a filing is judgement, exactly as the vision tier's
verdicts are (`test_deck_contract`'s "the vision tier is deliberately not asserted
for verdicts"). What can be pinned is that the agenda points at the right figures
and that nothing in it can ever block a run.

Both frozen fixture decks anchor the extraction, because the two failure modes are
opposite: the pitch deck's figures live in bullet text and market-entry tables, the
earnings-update deck's in a broker table, and the pitch deck carries three static
credential slides whose tombstone values are INFOR's own and must not be reported.
"""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from deckcheck import (
    EXPECTED_ERROR_CONTEXTS,
    MATCH_IDENTITY,
    MATCH_VALUE,
    SEVERITY_ADVISORY,
    VERDICTS,
    CheckFinding,
    audit_deck,
    deck_pictures,
    extract_deck_figures,
    render_agenda,
    render_report,
    write_report,
)
from provenance import DeckPlacement, FigureRef, FigureSource, ProvenanceLedger

_FIXTURES = Path(__file__).parent / "fixtures"
_PITCH = _FIXTURES / "pitch-deck.pptx"
_EARNINGS = _FIXTURES / "earnings-update-deck.pptx"

_SRC = FigureSource(filing="FY2025 10-K", statement="Consolidated Statements of Operations")


def _deck_with(tmp_path: Path, text: str, *, name: str = "check.pptx") -> Path:
    """A one-slide deck whose single textbox holds `text`."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
    box.name = "Figures"
    box.text_frame.text = text
    path = tmp_path / name
    prs.save(path)
    return path


# ─── Extraction ──────────────────────────────────────────────────────────────


def test_currency_figures_are_normalised_to_millions(tmp_path: Path):
    deck = _deck_with(tmp_path, "Revenue of US$589.8MM, EBITDA of $1.2B, cash of $450K")
    by_raw = {f.raw: f for f in extract_deck_figures(deck, library=None)}
    assert by_raw["US$589.8MM"].millions == pytest.approx(589.8)
    assert by_raw["$1.2B"].millions == pytest.approx(1200.0)
    assert by_raw["$450K"].millions == pytest.approx(0.45)


def test_percent_and_multiple_figures_keep_their_own_scale(tmp_path: Path):
    deck = _deck_with(tmp_path, "Margins of 28.2% at 11.5x EBITDA")
    kinds = {f.raw: (f.kind, f.value, f.millions) for f in extract_deck_figures(deck, library=None)}
    assert kinds["28.2%"] == ("percent", pytest.approx(28.2), None)
    assert kinds["11.5x"] == ("multiple", pytest.approx(11.5), None)


def test_an_accounting_negative_is_read_as_negative(tmp_path: Path):
    deck = _deck_with(tmp_path, "Net loss of ($12.3MM) for the period")
    figure = next(f for f in extract_deck_figures(deck, library=None) if "12.3" in f.raw)
    assert figure.value == pytest.approx(-12.3)
    assert figure.millions == pytest.approx(-12.3)


@pytest.mark.parametrize(
    "text",
    [
        "Q3 2026 results per the 2024 Annual Report",  # quarters, years, doc names
        "Founded 1987, page 3 of 16",                  # bare counts
        "See the 10-K and the 10-Q",                   # filing names
        "Presented April 2026",                        # a date
    ],
)
def test_bare_integers_are_not_figures(tmp_path: Path, text: str):
    # This exclusion is what makes the agenda readable — a deck is full of bare
    # integers and not one of them is a claim about the target's financials.
    assert extract_deck_figures(_deck_with(tmp_path, text), library=None) == []


def test_a_figure_inside_a_table_cell_is_found():
    # The earnings-update deck's broker table is the case: every figure it carries
    # is in a table, which a text-frame-only scan would miss entirely.
    figures = extract_deck_figures(_EARNINGS)
    tabled = [f for f in figures if f.shape.startswith("Table")]
    assert tabled, "no table-cell figures found in the earnings-update fixture"
    assert all(f.slide == 2 for f in tabled), "the broker table is on slide 3 (zero-based 2)"


def test_the_blank_library_excludes_the_static_credential_slides():
    # The pitch deck's slides 3-5 are copied from the library verbatim: their
    # tombstone values are INFOR's own and have no provenance in this run. They
    # drop out because the same shape in the blank library carries the same figure
    # — no slide list to migrate when the library gains an entry.
    with_baseline = extract_deck_figures(_PITCH)
    without = extract_deck_figures(_PITCH, library=_FIXTURES / "no-such-library.pptx")

    dropped = {(f.slide, f.shape, f.raw) for f in without} - {
        (f.slide, f.shape, f.raw) for f in with_baseline
    }
    assert dropped, "the library baseline excluded nothing"
    assert {slide for slide, _, _ in dropped} == {2, 3, 4}, "only the credential slides"
    assert {f.slide for f in with_baseline}.isdisjoint({2, 3, 4})


def test_filled_slides_survive_the_library_baseline():
    # The library holds `[x]` tokens where the fill puts figures, so a filled
    # figure can never match the baseline.
    figures = extract_deck_figures(_PITCH)
    assert any(f.slide == 1 for f in figures), "the executive summary lost its figures"
    assert len(figures) > 20


def test_pictures_are_listed_rather_than_scanned():
    # A rasterised range is out of scope for a string scan on purpose (an error
    # value inside one is usually correct), so it becomes a reader's job instead.
    pictures = deck_pictures(_EARNINGS)
    assert pictures
    assert all(isinstance(index, int) and name for index, name in pictures)


def test_the_library_baseline_also_drops_decorative_pictures():
    # 49 pictures on the pitch deck, 44 of them the library's own logos and
    # graphics. What is left is what this run pasted: the cap-table range on the
    # overview slide, the Financial Summary charts, the ownership blocks. A list of
    # 5 gets read; a list of 49 does not.
    kept = deck_pictures(_PITCH)
    everything = deck_pictures(_PITCH, library=_FIXTURES / "no-such-library.pptx")
    assert len(everything) > 40
    assert 0 < len(kept) <= 8, kept
    # Zero-based: the overview (6), the Financial Summary slide (7), ownership (8).
    assert {index for index, _ in kept} == {6, 7, 8}


def test_a_figures_context_is_a_window_not_the_whole_shape(tmp_path: Path):
    # A bullet block holds a dozen figures; repeating its full text once per figure
    # made the agenda table unreadable, which is the same failure as not reporting.
    long_text = "Lead-in. " * 40 + "Revenue of US$589.8MM. " + "Trailing. " * 40
    figure = next(f for f in extract_deck_figures(_deck_with(tmp_path, long_text), library=None)
                  if "589.8" in f.raw)
    assert "US$589.8MM" in figure.context
    assert len(figure.context) < len(long_text) / 3
    assert figure.context.startswith("…") and figure.context.endswith("…")


# ─── The ledger join ─────────────────────────────────────────────────────────
#
# Two tiers, and the distinction is the point. A record that says WHERE it lands
# (a `DeckPlacement`) and whose value agrees is `traced`. A record that only shares
# the number is a `value_matched` lead — reported as one, never as provenance. The
# arithmetic below (tolerance, scale, percent storage) decides whether the values
# agree at all, so it is asserted on the value tier; `_placed` adds the placement
# that promotes an agreement to a trace.


def _ledger(*figures) -> ProvenanceLedger:
    """Records with no placement — they can only ever be value-matched."""
    ledger = ProvenanceLedger(stage="financial-summary")
    for name, value, units in figures:
        ledger.record(name, sources=_SRC, value=value, units=units, location="financial-summary!B6")
    return ledger


def _placed(*figures, slide: int = 1, shape: str | None = None) -> ProvenanceLedger:
    """The same records, each claiming a slide (1-based) — the identity join."""
    ledger = ProvenanceLedger(stage="content")
    for name, value, units in figures:
        ledger.record(
            name,
            sources=_SRC,
            value=value,
            units=units,
            placement=DeckPlacement(slide=slide, field="executive_summary_bullets[0]", shape=shape),
        )
    return ledger


def test_a_record_that_claims_the_figure_is_traced(tmp_path: Path):
    deck = _deck_with(tmp_path, "Revenue of US$589.8MM")
    audit = audit_deck(deck, _placed(("Revenue FY2025", 589.8, "US$MM")), library=None)
    assert len(audit.traced) == 1
    assert audit.traced[0].record.figure == "Revenue FY2025"
    assert audit.traced[0].kind == MATCH_IDENTITY
    assert audit.untraced == [] and audit.value_matched == []


def test_a_placement_on_another_slide_is_only_a_value_match(tmp_path: Path):
    # The record says it lands on slide 8; this figure is on slide 1. Same number,
    # different figure — which is exactly the collision the old join printed as
    # traced.
    deck = _deck_with(tmp_path, "Revenue of US$589.8MM")
    audit = audit_deck(deck, _placed(("Revenue FY2025", 589.8, "US$MM"), slide=8), library=None)
    assert audit.traced == []
    assert len(audit.value_matched) == 1


def test_a_shape_name_identifies_a_figure_too(tmp_path: Path):
    deck = _deck_with(tmp_path, "Revenue of US$589.8MM")
    ledger = ProvenanceLedger(stage="deck")
    ledger.record("Revenue FY2025", sources=_SRC, value=589.8, units="US$MM",
                  placement=DeckPlacement(shape="Figures"))
    assert audit_deck(deck, ledger, library=None).traced


def test_a_rounded_rendering_still_agrees_with_its_record(tmp_path: Path):
    # The deck rounds; the workbook does not. Tolerance is half the last decimal
    # place the deck actually shows.
    deck = _deck_with(tmp_path, "Revenue of US$589.8MM")
    assert audit_deck(deck, _placed(("Revenue", 589.83, "US$MM")), library=None).traced
    assert not audit_deck(deck, _placed(("Revenue", 591.2, "US$MM")), library=None).matches[0].matched


def test_a_scale_change_between_workbook_and_deck_still_agrees(tmp_path: Path):
    # The workbook is locked to millions; a deck may render the same figure in
    # billions.
    deck = _deck_with(tmp_path, "Revenue of $1.2B")
    assert audit_deck(deck, _placed(("Revenue", 1200.0, "US$MM")), library=None).traced


def test_a_figure_with_no_record_is_untraced(tmp_path: Path):
    deck = _deck_with(tmp_path, "Total addressable market of $4.2B")
    audit = audit_deck(deck, _ledger(("Revenue FY2025", 589.8, "US$MM")), library=None)
    assert len(audit.untraced) == 1
    assert audit.untraced[0].record is None
    assert audit.untraced[0].kind is None


def test_a_percent_only_agrees_with_a_percent_record(tmp_path: Path):
    deck = _deck_with(tmp_path, "Margin of 28.2%")
    assert audit_deck(deck, _placed(("EBITDA margin", 28.2, "%")), library=None).traced
    # Same number, but the record is a dollar figure — not the same claim, and the
    # placement does not rescue it.
    assert not audit_deck(deck, _placed(("Revenue", 28.2, "US$MM")), library=None).traced


def test_a_percent_record_stored_as_a_fraction_still_agrees(tmp_path: Path):
    deck = _deck_with(tmp_path, "Margin of 28.2%")
    assert audit_deck(deck, _placed(("EBITDA margin", 0.282, "%")), library=None).traced


def test_a_record_holding_an_excel_formula_cannot_be_auto_matched(tmp_path: Path):
    # A combined figure is written as a formula so the arithmetic stays in the
    # cell; openpyxl never evaluates it, so there is no number to compare. The
    # figure lands in the untraced list for a reviewer, not silently in traced.
    deck = _deck_with(tmp_path, "Combined balances of US$9,800.0MM")
    audit = audit_deck(deck, _placed(("Combined balances", "=9000+800", "US$MM")), library=None)
    assert audit.untraced and not audit.traced


def test_audit_needs_no_renderer():
    # Deliberate: the join is deterministic and must run on a box with no
    # LibreOffice, so a missing renderer degrades the evidence, not the audit.
    audit = audit_deck(_PITCH, ProvenanceLedger())
    assert audit.matches and audit.untraced == audit.matches


# ─── A value match is never reported as traced ───────────────────────────────
#
# The observed case, kept as the fixture: on a real pitch run the executive
# summary's ARR of US$4,190.5MM matched the `financial-summary` record for FY2024
# gross profit, 4,191.0 — 0.5 apart, inside the tolerance a derived figure's own
# rounding needs (0.1% of 4,190.5 is 4.19). Four of that run's twelve "traced"
# joins were coincidences like it, and the report printed them as provenance.

_ARR_ON_THE_SLIDE = "ARR of US$4,190.5MM, 81% of revenue"
_GROSS_PROFIT_RECORD = ("Gross Profit FY2024", 4191.0, "US$MM")


def test_the_arr_gross_profit_collision_is_a_value_match_not_a_trace(tmp_path: Path):
    deck = _deck_with(tmp_path, _ARR_ON_THE_SLIDE)
    audit = audit_deck(deck, _ledger(_GROSS_PROFIT_RECORD), library=None)

    arr = next(m for m in audit.matches if "4,190.5" in m.figure.raw)
    assert arr.matched, "the values do agree — that is what makes the coincidence possible"
    assert arr.kind == MATCH_VALUE
    assert not arr.traced
    assert audit.traced == [], "nothing on this deck is traced"
    assert audit.value_matched == [arr]


def test_the_agenda_labels_a_value_match_and_keeps_it_out_of_traced(tmp_path: Path):
    deck = _deck_with(tmp_path, _ARR_ON_THE_SLIDE)
    agenda = render_agenda(audit_deck(deck, _ledger(_GROSS_PROFIT_RECORD), library=None))

    assert "### Value matches — 1 (NOT traced)" in agenda
    assert "### Traced figures — 0" in agenda
    # The row itself has to say so: a reader skimming rows never reaches the
    # heading's caveat.
    row = next(line for line in agenda.splitlines() if "`US$4,190.5MM`" in line)
    assert "**value match only**" in row
    assert "no placement recorded" in row


def test_a_value_match_reports_how_many_records_share_the_number(tmp_path: Path):
    # Two records with the same value make the value join arbitrary; printing the
    # first one silently is how a coincidence becomes a citation.
    deck = _deck_with(tmp_path, _ARR_ON_THE_SLIDE)
    ledger = _ledger(_GROSS_PROFIT_RECORD, ("Deferred revenue FY2024", 4190.7, "US$MM"))
    agenda = render_agenda(audit_deck(deck, ledger, library=None))
    assert "2 records share this value" in agenda


def test_the_same_figure_traces_once_the_content_stage_names_where_it_lands(tmp_path: Path):
    # The fix on the recording side: pitch-content names the slide and the typed
    # PitchDeckContent field, and the ARR figure joins by identity — while the
    # gross-profit record in the same ledger stays a value match candidate.
    deck = _deck_with(tmp_path, _ARR_ON_THE_SLIDE)
    ledger = _ledger(_GROSS_PROFIT_RECORD)
    ledger.record(
        "ARR",
        sources=FigureSource(filing="Q1 2026 10-Q", statement="MD&A — key metrics", page=7),
        value=4190.5,
        units="US$MM",
        location=None,
        placement=DeckPlacement(slide=1, field="executive_summary_bullets[0]"),
    )
    audit = audit_deck(deck, ledger, library=None)

    arr = next(m for m in audit.matches if "4,190.5" in m.figure.raw)
    assert arr.kind == MATCH_IDENTITY
    assert arr.record.figure == "ARR"
    assert "Q1 2026 10-Q" in render_agenda(audit)


# ─── Findings are advisory, by construction ──────────────────────────────────


def test_a_finding_cannot_claim_to_be_blocking():
    # `deckcheck` is a review surfaced at a checkpoint, and since v0.5.49 no shipped
    # stage gates at all — so this would be the only one. A falsification pass that
    # could halt a run would have to be right about a target's financial statements.
    with pytest.raises(ValueError, match="review surfaced at a checkpoint"):
        CheckFinding(slide=1, figure="Revenue", verdict="contradicted", detail="x",
                     severity="blocking")


def test_a_finding_defaults_to_advisory():
    assert CheckFinding(slide=1, figure="Revenue", verdict="confirmed",
                        detail="x").severity == SEVERITY_ADVISORY


def test_an_unknown_verdict_is_rejected():
    with pytest.raises(ValueError, match="unknown verdict"):
        CheckFinding(slide=1, figure="Revenue", verdict="probably fine", detail="x")


def test_slide_number_is_one_based_for_the_analyst():
    assert CheckFinding(slide=6, figure="Revenue", verdict="confirmed", detail="x").slide_number == 7


# ─── Agenda and report ───────────────────────────────────────────────────────


def test_the_agenda_carries_the_do_not_report_list(tmp_path: Path):
    # The constraint that keeps this review from being ignored: expected CapIQ
    # error values must be in front of the reader at the moment they would
    # otherwise report one, not only in the skill's prose.
    deck = _deck_with(tmp_path, "Revenue of US$589.8MM")
    agenda = render_agenda(audit_deck(deck, _ledger(("Revenue", 589.8, "US$MM")), library=None))
    assert "Not defects" in agenda
    for context in EXPECTED_ERROR_CONTEXTS:
        assert context in agenda


def test_the_agenda_names_untraced_figures_and_their_slides(tmp_path: Path):
    deck = _deck_with(tmp_path, "Market size of $4.2B")
    agenda = render_agenda(audit_deck(deck, ProvenanceLedger(), library=None))
    assert "### Untraced figures — 1" in agenda
    assert "$4.2B" in agenda


def test_the_agenda_shows_a_traced_figure_with_its_citation(tmp_path: Path):
    deck = _deck_with(tmp_path, "Revenue of US$589.8MM")
    ledger = ProvenanceLedger(stage="financial-summary")
    ledger.record("Revenue FY2025", sources=_SRC, value=589.8, units="US$MM",
                  location="financial-summary!B6", placement=DeckPlacement(slide=1))
    agenda = render_agenda(audit_deck(deck, ledger, library=None))
    assert "### Traced figures — 1" in agenda
    assert "FY2025 10-K, Consolidated Statements of Operations" in agenda
    assert "financial-summary!B6" in agenda


def test_the_agenda_follows_a_derived_figure_to_the_filings_underneath(tmp_path: Path):
    # The LTM tile on a deck is a bridge total: its own record has no source, and
    # before the refs existed the agenda printed the derivation sentence, which a
    # reviewer could not tell apart from "unsourced".
    deck = _deck_with(tmp_path, "LTM Adj. EBITDA of US$1,840.4MM")
    ledger = ProvenanceLedger(stage="ltm-metrics")
    ledger.record("LTM Adj. EBITDA — FY2025 Adj. EBITDA", sources=_SRC, value=1700.0,
                  units="US$MM", location="ltm-metrics!B30")
    ledger.record(
        "LTM Adj. EBITDA",
        value=1840.4,
        units="US$MM",
        location="ltm-metrics!B33",
        derivation="FY2025 + Q1 2026 YTD − Q1 2025 YTD",
        derived_from=[FigureRef(location="ltm-metrics!B30")],
        placement=DeckPlacement(slide=1),
    )
    agenda = render_agenda(audit_deck(deck, ledger, library=None))

    assert "LTM Adj. EBITDA — FY2025 Adj. EBITDA" in agenda, "the component is named"
    assert "FY2025 10-K, Consolidated Statements of Operations" in agenda, "and its filing"


def test_the_agenda_flags_a_derivation_whose_component_has_no_record(tmp_path: Path):
    deck = _deck_with(tmp_path, "LTM Adj. EBITDA of US$1,840.4MM")
    ledger = ProvenanceLedger(stage="ltm-metrics")
    ledger.record(
        "LTM Adj. EBITDA",
        value=1840.4,
        units="US$MM",
        derivation="FY2025 + Q1 2026 YTD − Q1 2025 YTD",
        derived_from=[FigureRef(location="ltm-metrics!B30", figure="FY2025 Adj. EBITDA")],
        placement=DeckPlacement(slide=1),
    )
    agenda = render_agenda(audit_deck(deck, ledger, library=None))
    assert "UNRESOLVABLE: ltm-metrics!B30" in agenda


def test_the_agenda_lists_records_no_text_figure_reached(tmp_path: Path):
    # The cap table is a rasterised picture, so Enterprise Value can only reach a
    # reviewer through this list — it is the whole path from the biggest number on
    # the overview slide to a filing page.
    deck = _deck_with(tmp_path, "Revenue of US$589.8MM")
    ledger = ProvenanceLedger(stage="captable")
    ledger.record("Total Debt", sources=_SRC, value=4200.0, units="C$MM",
                  location="captable!F122")
    ledger.record(
        "Enterprise Value",
        value="=F22+F28+F29+F30",
        location="captable!F31",
        derivation="cap-table formula =F22+F28+F29+F30",
        derived_from=[FigureRef(location="captable!F122", figure="Total Debt")],
    )
    agenda = render_agenda(audit_deck(deck, ledger, library=None))

    assert "#### Records no text figure joined to — 2" in agenda
    row = next(line for line in agenda.splitlines() if line.startswith("| Enterprise Value "))
    assert "captable!F31" in row
    assert "Total Debt" in row and "FY2025 10-K" in row


def test_the_report_leads_with_the_verdicts_then_the_agenda(tmp_path: Path):
    deck = _deck_with(tmp_path, "Revenue of US$589.8MM and a market size of $4.2B")
    audit = audit_deck(deck, _ledger(("Revenue FY2025", 589.8, "US$MM")), library=None)
    findings = [
        CheckFinding(slide=0, figure="Revenue FY2025", verdict="confirmed",
                     detail="agrees with the cited statement", source="FY2025 10-K p. 61"),
        CheckFinding(slide=0, figure="market size $4.2B", verdict="unsupported",
                     detail="no record and not in the attached CIM"),
    ]
    report = render_report(audit, findings, company="Example Target Inc.",
                          provenance_path=tmp_path / "provenance.json")

    assert report.index("## Findings") < report.index("## Agenda")
    assert "Advisory review, not a gate" in report
    assert "**confirmed**" in report and "**unsupported**" in report
    assert all(v in report for v in VERDICTS), "the verdict tally names every verdict"


def test_the_report_says_so_when_nothing_could_be_faulted(tmp_path: Path):
    deck = _deck_with(tmp_path, "Revenue of US$589.8MM")
    audit = audit_deck(deck, _ledger(("Revenue", 589.8, "US$MM")), library=None)
    assert "no figure could be faulted" in render_report(audit, company="Example Target Inc.")


def test_write_report_creates_its_parent_directory(tmp_path: Path):
    deck = _deck_with(tmp_path, "Revenue of US$589.8MM")
    audit = audit_deck(deck, ProvenanceLedger(), library=None)
    path = write_report(tmp_path / "artefacts" / "deckcheck-Project Test.md", audit,
                        company="Example Target Inc.")
    assert path.is_file()
    assert path.read_text(encoding="utf-8").startswith("# Deck check — Example Target Inc.")
