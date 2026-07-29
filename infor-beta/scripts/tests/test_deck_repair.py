"""Phase B step 3: the converge loop, wired against the same regression artefacts.

`test_deck_contract.py` asserts the oracle SEES the historical defects. This file
asserts the loop ACTS on them — and, just as importantly, that it fails loudly
when it cannot, since shipping a shrunk-but-still-overflowing deck is what the
retired estimator did for thirteen releases.

Every repair here is measured off a real LibreOffice render, so these tests need
LibreOffice. There is deliberately no estimate-based fallback to test against.
"""

from __future__ import annotations

import shutil
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Emu, Pt

import deck_repair
from deck_contract import default_library_path, verify_deck
from deck_repair import (
    REPAIRABLE_KINDS,
    TEXT_SCALES,
    ConvergeResult,
    DeckNotConvergedError,
    assert_converged,
    converge_deck,
)
from excel_to_powerpoint import find_soffice
from font_probe import FontResolution
from pptx_helpers import PALATINO, normal_autofit_scale

_FIXTURES = Path(__file__).parent / "fixtures"
_REGRESSIONS = _FIXTURES / "regressions"

needs_render = pytest.mark.skipif(
    find_soffice() is None, reason="the repair loop measures on a LibreOffice render"
)


def _copy(src: Path, tmp_path: Path) -> Path:
    """Work on a copy — `converge_deck` repairs in place."""
    dst = tmp_path / src.name
    shutil.copy(src, dst)
    return dst


def _geometric(findings):
    return [f for f in findings if f.blocking and f.kind in REPAIRABLE_KINDS]


# ─── The loop repairs the historical defects ─────────────────────────────────


@needs_render
def test_repairs_prl17_market_entry_table(tmp_path):
    """PRL17's 5.91"-rendered table, fixed by measurement rather than estimation.

    The v0.5.23 fix for this was a per-character Palatino width table plus per-row
    content-height floors. Here the loop renders the table at each candidate body
    size and picks the largest that actually fits — and the table's declared height
    is untouched, because the declaration was never the problem.
    """
    deck = _copy(_REGRESSIONS / "prl17-market-entry-table.pptx", tmp_path)
    before = Emu(
        next(s for s in Presentation(deck).slides[0].shapes if s.name == "Table 1215").height
    ).inches

    result = converge_deck(deck, out_dir=tmp_path / "qa")

    assert result.converged, result.summary() + "\n" + "\n".join(
        str(f) for f in result.unrepaired
    )
    assert result.iterations == 1, f"expected one measured pass, took {result.iterations}"
    assert _geometric(verify_deck(deck, vision=False, out_dir=tmp_path / "after")) == []

    after = Emu(
        next(s for s in Presentation(deck).slides[0].shapes if s.name == "Table 1215").height
    ).inches
    assert after == pytest.approx(before, abs=0.01), (
        "the table's declared height must not move — the repair is smaller text, not "
        "a smaller declaration"
    )


@needs_render
def test_repairs_prl18_risk_table_by_clamping_to_the_library_height(tmp_path):
    """PRL18 declared 5.360" against the library's 5.1715"; the clamp is the fix.

    Distinct from PRL17: nothing renders past anything here, the declaration itself
    is wrong. So the repair clamps and only steps the font if the clamped rows then
    render past their new bottom.
    """
    deck = _copy(_REGRESSIONS / "prl18-risk-table.pptx", tmp_path)

    result = converge_deck(deck, out_dir=tmp_path / "qa")

    assert result.converged, result.summary()
    table = next(s for s in Presentation(deck).slides[0].shapes if s.name == "Table 4")
    assert Emu(table.height).inches == pytest.approx(5.1715, abs=0.01)
    assert _geometric(verify_deck(deck, vision=False, out_dir=tmp_path / "after")) == []


# ─── The loop fails loudly rather than shipping a bad deck ───────────────────


@needs_render
def test_does_not_converge_when_the_content_is_over_budget(tmp_path):
    """PRL14's overview block cannot fit its box at any legible size.

    1,126 characters in a box declaring 0.58" of height, with the box never sized to
    the band (the pre-v0.5.23 defect). The loop spends its ladder, reports the
    remaining overflow honestly, and refuses to call it converged. The retired
    estimator floored at 70% and shipped.
    """
    deck = _copy(_REGRESSIONS / "prl14-overview-bullets.pptx", tmp_path)

    result = converge_deck(deck, out_dir=tmp_path / "qa")

    assert not result.converged
    assert result.iterations <= 3, "the loop must stay bounded"
    surviving = [f for f in result.unrepaired if f.shape == "TextBox 9"]
    assert surviving, f"expected the overview block to still be reported: {result.summary()}"

    # It did shrink as far as it could before giving up.
    box = next(s for s in Presentation(deck).slides[0].shapes if s.name == "TextBox 9")
    assert normal_autofit_scale(box) == pytest.approx(min(TEXT_SCALES))

    with pytest.raises(DeckNotConvergedError) as excinfo:
        assert_converged(deck, result)
    message = str(excinfo.value)
    assert "TextBox 9" in message
    assert "shorten the copy" in message, (
        f"the remedy for over-budget content is editorial, and the error should say so: "
        f"{message}"
    )


def test_assert_converged_is_silent_on_a_converged_result(tmp_path):
    assert_converged(tmp_path / "any.pptx", ConvergeResult(converged=True, iterations=0)) is None


# ─── QA scratch is ephemeral unless the stage fails ──────────────────────────
# The loop writes ~170 files and ~10 MB per deck. The assemblers used to point
# that at the deal's artefacts directory, which is cloud-synced: the same code
# and the same warm render cache took 11.7 s for the first repair on local disk
# and 26.9 s on the mount, and eight consecutive live pitch attempts were killed
# at ~43 s having never converged.


@pytest.fixture
def scratch_roots(monkeypatch):
    """Record every ephemeral scratch root `converge_deck` creates.

    Records rather than redirects, so nothing about where LibreOffice puts its
    profile or its PDF cache changes underneath the test.
    """
    made: list[Path] = []
    real_mkdtemp = deck_repair.tempfile.mkdtemp

    def spy(*args, **kwargs):
        path = real_mkdtemp(*args, **kwargs)
        if str(kwargs.get("prefix", "")).startswith(deck_repair.SCRATCH_PREFIX):
            made.append(Path(path))
        return path

    monkeypatch.setattr(deck_repair.tempfile, "mkdtemp", spy)
    return made


@needs_render
def test_a_converged_run_leaves_no_scratch_anywhere(tmp_path, scratch_roots):
    """The happy path writes nothing durable — not beside the deck, not in temp."""
    deck = _copy(_REGRESSIONS / "prl18-risk-table.pptx", tmp_path)
    beside_the_deck = set(tmp_path.iterdir())

    result = converge_deck(deck, keep_on_failure=tmp_path / ".qa")

    assert result.converged, result.summary()
    assert result.kept_dir is None
    assert not (tmp_path / ".qa").exists(), "a converged run kept QA artefacts"
    assert set(tmp_path.iterdir()) == beside_the_deck, (
        "the loop wrote scratch into the deck's own directory"
    )
    assert scratch_roots, "the loop did not stage its renders under tempfile"
    assert [root for root in scratch_roots if root.exists()] == [], (
        "the ephemeral scratch root was left behind"
    )


@needs_render
def test_a_failing_run_keeps_the_last_pass_and_names_it(tmp_path, scratch_roots):
    """Failure is the one case worth keeping: the renders ARE the diagnosis.

    Only the last pass, and only what the analyst would open — the earlier passes
    are what made this ~10 MB per deck.
    """
    deck = _copy(_REGRESSIONS / "prl14-overview-bullets.pptx", tmp_path)
    qa = tmp_path / "artefacts" / ".qa"

    result = converge_deck(deck, keep_on_failure=qa)

    assert not result.converged
    assert result.kept_dir == qa and qa.is_dir()
    kept = sorted(path.name for path in qa.iterdir())
    assert any(name.startswith("verify-") for name in kept), kept
    assert list(qa.rglob("*.png")), "the kept pass carries no renders to look at"
    assert [root for root in scratch_roots if root.exists()] == [], (
        "keeping the failing pass must not keep the whole scratch tree"
    )

    with pytest.raises(DeckNotConvergedError) as excinfo:
        assert_converged(deck, result)
    assert str(qa) in str(excinfo.value), (
        "the error must name where the failing renders were kept"
    )


@needs_render
def test_an_explicit_out_dir_still_persists(tmp_path):
    """A caller that asks for the whole tree gets it — tests, hands-on debugging."""
    deck = _copy(_REGRESSIONS / "prl18-risk-table.pptx", tmp_path)

    converge_deck(deck, out_dir=tmp_path / "qa")

    assert (tmp_path / "qa" / "verify-0").is_dir()


# ─── The font the ladder is calibrated against ───────────────────────────────


@needs_render
def test_the_converge_records_which_font_the_oracle_resolved(tmp_path):
    """Every size in the deck is measured in this face; the run must say which.

    A silent substitution is the worst case — the deck converges and every point
    size is calibrated against the wrong advance widths — and it left no trace
    anywhere until this line existed.
    """
    deck = _copy(default_library_path(), tmp_path)
    logged: list[str] = []

    result = converge_deck(deck, out_dir=tmp_path / "qa", log=logged.append)

    assert result.font is not None
    assert result.font.requested == PALATINO
    assert any("font" in line.lower() for line in logged), logged
    assert any(result.font.log_line() in line for line in logged), (
        "the resolution itself must reach the stage log, pass or fail"
    )


def test_a_substituted_font_is_named_in_the_failure(tmp_path):
    """A deck that will not fit in the WRONG metrics may fit in the right ones.

    So the substitution has to be in front of whoever is about to rewrite copy,
    not only in the stage log they may not still have.
    """
    substituted = FontResolution(
        requested=PALATINO,
        family="DejaVu Serif",
        file="/usr/share/fonts/truetype/dejavu/DejaVuSerif.ttf",
        method="fontconfig",
        metric_compatible=False,
    )
    result = ConvergeResult(converged=False, iterations=3, font=substituted)

    with pytest.raises(DeckNotConvergedError) as excinfo:
        assert_converged(tmp_path / "any.pptx", result)

    assert "DejaVu Serif" in str(excinfo.value)


# ─── Scope: what the loop must NOT touch ─────────────────────────────────────


@needs_render
def test_string_findings_are_not_treated_as_repairable(tmp_path):
    """The pitch fixture has 12 blocking findings and needs zero repairs.

    All twelve are `[x]` tokens and unsubstituted currency letters — real defects,
    but not ones a font size fixes. The loop must recognise that its remedies do
    not apply and converge immediately rather than shrinking innocent shapes.
    """
    deck = _copy(_FIXTURES / "pitch-deck.pptx", tmp_path)
    before = Presentation(deck).slides[6]
    before_scale = normal_autofit_scale(
        next(s for s in before.shapes if s.name == "TextBox 9")
    )

    result = converge_deck(deck, out_dir=tmp_path / "qa")

    assert result.converged, result.summary()
    assert result.iterations == 0
    assert result.actions == []
    assert len(result.blocking) == 12, [str(f) for f in result.blocking]
    assert all(f.kind not in REPAIRABLE_KINDS for f in result.blocking)
    after_scale = normal_autofit_scale(
        next(s for s in Presentation(deck).slides[6].shapes if s.name == "TextBox 9")
    )
    assert after_scale == before_scale, "a converged shape must not be re-scaled"


@needs_render
def test_vision_findings_never_enter_the_loop(tmp_path):
    """Advisory findings are non-deterministic and belong at the `deck` checkpoint.

    Run with the vision tier ON: it must contribute findings to the result (so the
    checkpoint can show them) while contributing nothing to the repair decision.
    """
    deck = _copy(_FIXTURES / "pitch-deck.pptx", tmp_path)

    result = converge_deck(deck, out_dir=tmp_path / "qa", vision=True)

    assert result.converged
    assert result.actions == []
    assert any(f.kind == "vision-review" for f in result.advisory), (
        "the vision agenda must still reach the caller"
    )
    assert all(not f.blocking for f in result.advisory)


@needs_render
def test_converging_an_already_converged_deck_changes_nothing(tmp_path):
    """Idempotence: the stage can be re-run without cumulative shrinking.

    A loop that shaved a point off every pass would degrade a deck each time the
    stage was retried.
    """
    deck = _copy(_REGRESSIONS / "prl18-risk-table.pptx", tmp_path)
    assert converge_deck(deck, out_dir=tmp_path / "qa1").converged
    first = deck.read_bytes()

    second = converge_deck(deck, out_dir=tmp_path / "qa2")
    assert second.converged
    assert second.iterations == 0
    assert second.actions == []
    assert deck.read_bytes() == first, "a second pass must not modify a converged deck"


@needs_render
def test_table_repair_leaves_the_header_row_alone(tmp_path):
    """The library's 12 pt Considerations header must survive a body-font step.

    Shrinking every cell uniformly is what made an earlier revision "render
    noticeably smaller than the template" (v0.5.14).
    """
    deck = _copy(_REGRESSIONS / "prl18-risk-table.pptx", tmp_path)
    table = next(s for s in Presentation(deck).slides[0].shapes if s.name == "Table 4").table
    header_before = [
        run.font.size
        for cell in table.rows[0].cells
        for para in cell.text_frame.paragraphs
        for run in para.runs
    ]

    converge_deck(deck, out_dir=tmp_path / "qa")

    table = next(s for s in Presentation(deck).slides[0].shapes if s.name == "Table 4").table
    header_after = [
        run.font.size
        for cell in table.rows[0].cells
        for para in cell.text_frame.paragraphs
        for run in para.runs
    ]
    assert header_after == header_before
    assert all(size is None or size >= Pt(12) for size in header_after)


# ─── The library is the reference, and satisfies its own contract ────────────


@needs_render
def test_blank_library_needs_no_repair(tmp_path):
    """The geometric baseline cannot need repairing against itself."""
    deck = _copy(default_library_path(), tmp_path)

    result = converge_deck(deck, out_dir=tmp_path / "qa")

    assert result.converged
    assert result.actions == []
    assert deck.read_bytes() == default_library_path().read_bytes()
