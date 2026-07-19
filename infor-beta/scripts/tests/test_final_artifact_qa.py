from pathlib import Path

from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches


def _minimal_artifacts(tmp_path: Path) -> tuple[Path, Path]:
    deck = tmp_path / "final.pptx"
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.save(deck)
    workbook = tmp_path / "final.xlsx"
    Workbook().save(workbook)
    return deck, workbook


def _render_all(deck_path, output_dir, *, slide_indices):
    paths = []
    for index in slide_indices:
        path = Path(output_dir) / f"slide_{index + 1}.png"
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_bytes(b"png")
        paths.append(path)
    return paths


def test_missing_deck_or_workbook_blocks_delivery(tmp_path: Path):
    from final_artifact_qa import run_final_artifact_qa

    result = run_final_artifact_qa(
        deck_path=tmp_path / "missing.pptx",
        combined_workbook_path=tmp_path / "missing.xlsx",
        output_dir=tmp_path / "qa",
    )

    assert result.ready_for_delivery is False
    assert any("deck does not exist" in issue for issue in result.blocking_issues)
    assert any("combined workbook does not exist" in issue for issue in result.blocking_issues)
    assert result.final_deck_path == str((tmp_path / "missing.pptx").resolve())
    assert result.final_combined_workbook_path == str((tmp_path / "missing.xlsx").resolve())


def test_forced_chart_render_failure_blocks_delivery(tmp_path: Path):
    from final_artifact_qa import run_final_artifact_qa

    deck, workbook = _minimal_artifacts(tmp_path)
    result = run_final_artifact_qa(
        deck_path=deck,
        combined_workbook_path=workbook,
        output_dir=tmp_path / "qa",
        chart_status={
            "required": True,
            "financial_summary_inserted": False,
            "ltm_pie_inserted": True,
        },
    )

    assert result.ready_for_delivery is False
    assert result.financial_chart_insertion_status == "failed"
    assert any("financial chart insertion failed" in issue for issue in result.blocking_issues)


def test_remaining_prohibited_placeholder_blocks_delivery(tmp_path: Path):
    from final_artifact_qa import run_final_artifact_qa

    deck, workbook = _minimal_artifacts(tmp_path)
    prs = Presentation(deck)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text = (
        "[Placeholder for Metric #1 Chart]"
    )
    prs.save(deck)

    result = run_final_artifact_qa(
        deck_path=deck,
        combined_workbook_path=workbook,
        output_dir=tmp_path / "qa",
    )

    assert result.ready_for_delivery is False
    assert result.remaining_placeholders == (
        "slide 2: [Placeholder for Metric #1 Chart]",
    )
    assert any("prohibited placeholder" in issue for issue in result.blocking_issues)


def test_missing_visual_qa_evidence_blocks_delivery_and_renders_every_slide(tmp_path: Path):
    from final_artifact_qa import run_final_artifact_qa

    deck, workbook = _minimal_artifacts(tmp_path)
    prs = Presentation(deck)
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.save(deck)
    calls = []

    def render(deck_path, output_dir, *, slide_indices):
        calls.append(slide_indices)
        paths = []
        for index in slide_indices:
            path = Path(output_dir) / f"slide_{index + 1}.png"
            path.parent.mkdir(parents=True, exist_ok=True)
            path.write_bytes(b"png")
            paths.append(path)
        return paths

    result = run_final_artifact_qa(
        deck_path=deck,
        combined_workbook_path=workbook,
        output_dir=tmp_path / "qa",
        render_fn=render,
    )

    assert calls == [[0, 1]]
    assert result.ready_for_delivery is False
    assert result.visual_qa_render_status == "rendered_unverified"
    assert result.visual_qa_rendered_slides == (0, 1)
    assert any("visual QA evidence" in issue for issue in result.blocking_issues)


def test_degraded_or_unverified_aggregation_blocks_delivery(tmp_path: Path):
    from final_artifact_qa import run_final_artifact_qa

    deck, workbook = _minimal_artifacts(tmp_path)
    result = run_final_artifact_qa(
        deck_path=deck,
        combined_workbook_path=workbook,
        output_dir=tmp_path / "qa",
        aggregation_status={
            "backend": "openpyxl",
            "degraded": True,
            "merge_integrity_verified": True,
            "relink_ok": True,
            "recalculation_ok": False,
            "sources_deleted": False,
            "kept_sources": [str(tmp_path / "source.xlsx")],
        },
        render_fn=_render_all,
        visual_qa_evidence={"inspected_slide_indices": [0], "overflow_free": True},
    )

    assert result.ready_for_delivery is False
    assert result.aggregation_backend == "openpyxl"
    assert result.aggregation_degraded is True
    assert result.aggregation_verified is False
    assert result.recalculation_status == "failed"
    assert result.sources_deleted is False
    assert result.kept_source_workbooks == (str(tmp_path / "source.xlsx"),)
    assert any("aggregation is degraded" in issue for issue in result.blocking_issues)
    assert any("recalculation is not verified" in issue for issue in result.blocking_issues)
