"""TDD tests for the 16-slide INFOR slide-library POC."""

from pathlib import Path

import yaml
from openpyxl import Workbook
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu, Inches, Pt
from pydantic import ValidationError
import pytest

from schemas import Company, PitchDeckContent, Plan
from deck_contract import verify_deck
from excel_to_powerpoint import find_soffice
from deck_repair import DeckNotConvergedError
from pptx_helpers import find_shape, shapes_in_reading_order
from pitch_deck_wireframe import build_pitch_deck_slide_plan, write_slide_plan
from pitch_deck_assembler import (
    _iter_slide_text,
    _output_currency_letter,
    _verify_pitch_output,
    assemble_pitch_deck,
)
from slide_library_registry import load_slide_library_registry


PLUGIN_ROOT = Path(__file__).resolve().parents[2]
TEMPLATE = PLUGIN_ROOT / "templates" / "INFOR Slide Library.pptx"


def _sample_content() -> PitchDeckContent:
    return PitchDeckContent(
        client_name="SampleCo Ltd.",
        presentation_date="April 2026",
        figure_currency="USD",
        executive_summary_bullets=[
            {"text": "SampleCo is a compelling public-company advisory candidate supported by durable market positions and recurring revenue visibility", "level": 0},
            {"text": "The Company benefits from diversified end-market exposure and a management team focused on profitable growth", "level": 0},
            {"text": "Potential acquirors are expected to focus on scale, margin durability and cash conversion", "level": 0},
            {"text": "INFOR is well positioned to help evaluate strategic alternatives and prepare the Company for a targeted process", "level": 0},
        ],
        section_labels=["Overview", "Financial Summary", "Valuation", "Process"],
        current_section="Overview",
        company_overview_bullets=[
            {"text": "Provides mission-critical software and services to enterprise customers across North America", "level": 0},
            {"text": "Recurring revenue model supported by multi-year contracts and high customer retention", "level": 0},
            {"text": "Publicly listed issuer with established reporting history and access to capital markets", "level": 0},
        ],
        risk_mitigants=[
            {
                "risk": "Acquirors may question the durability of organic growth as market conditions normalize",
                "mitigants": [
                    "Highlight recurring revenue base",
                    "Show multi-year retention trends",
                    "Frame growth by segment",
                ],
            },
            {
                "risk": "Scale relative to larger public peers may affect valuation expectations",
                "mitigants": [
                    "Benchmark to focused peers",
                    "Emphasize margin profile",
                    "Position scarcity value",
                ],
            },
            {
                "risk": "Buyers may diligence customer concentration and renewal quality",
                "mitigants": [
                    "Prepare cohort support",
                    "Show renewal history",
                    "Document top-customer exposure",
                ],
            },
        ],
        risks_tagline="INFOR will help proactively frame key diligence topics to support a constructive acquiror dialogue.",
        comps_takeaway="SampleCo trades at a discount to higher-growth peers, creating a clear valuation framing opportunity.",
        precedents_takeaway="Recent precedent transactions in the sector cleared at premium multiples, supporting an attractive valuation case.",
        ownership_takeaway="Institutions hold a clear majority of the register, with insider ownership concentrated in long-tenured management.",
        investment_highlights=[
            {"header": "Durable Recurring Revenue", "bullets": ["High retention across multi-year contracts", "Predictable cash conversion through cycles"]},
            {"header": "Fragmented, Underpenetrated Market", "bullets": ["Few scaled competitors", "Clear whitespace for consolidation"]},
            {"header": "Strategic Buyer Interest", "bullets": ["Active strategic and sponsor consolidation", "Scarcity value for scaled assets"]},
            {"header": "Multiple Expansion Paths", "bullets": ["Organic, greenfield and M&A optionality", "Growth not reliant on any single lever"]},
        ],
        investment_highlights_tagline="SampleCo offers a scarce, scaled platform in a structurally attractive market with multiple value-creation paths.",
        market_entry_market="Canada",
        market_entry_row_labels=[
            "Overview",
            "Headquarters",
            "Year Founded",
            "Lending Product",
            "Customer Segment",
            "Funding Model",
            "Credit Approach",
            "Geographic Footprint",
            "Regulatory Status",
            "Technology Platform",
            "Scale KPIs",
            "Strategic Rationale",
        ],
        market_entry_targets=[
            {"cells": ["Digital lender", "Toronto, ON", "2014", "Unsecured instalment", "Prime / near-prime", "Equity and debt facilities", "Proprietary scoring", "Canada-wide", "Provincially licensed", "Cloud-native", ">1MM customers", "Scaled platform with brand recognition"]},
            {"cells": ["Mobile-first lender", "Vancouver, BC", "2017", "Line of credit", "Thin-file", "Venture capital and debt warehouse", "Alternative data", "Western Canada", "Provincially licensed", "Mobile-first", "~0.5MM customers", "Actionable platform subject to diligence"]},
        ],
        sources=[{"section": "Company Overview", "citation": "Company filings, company website and analyst notes"}],
        manual_steps=["LTM revenue breakdown and financial summary charts remain placeholders in this POC."],
    )


def test_pitch_deck_content_schema_validates_concise_risk_mitigants():
    content = _sample_content()

    assert content.presentation_date == "April 2026"
    assert content.risk_mitigants[0].mitigants == [
        "Highlight recurring revenue base",
        "Show multi-year retention trends",
        "Frame growth by segment",
    ]

    bad = content.model_dump()
    bad["risk_mitigants"][0]["mitigants"] = ["one", "two"]
    with pytest.raises(ValidationError):
        PitchDeckContent.model_validate(bad)


def test_risk_mitigant_length_allows_short_sentence_rejects_paragraph():
    # A one-sentence mitigant (> the old 90-char fragment cap, <= 160) is accepted.
    one_sentence = (
        "Frame organic growth by segment so acquirors see the durable recurring "
        "base rather than one-off contributions"
    )
    assert 90 < len(one_sentence) <= 160
    ok = _sample_content().model_dump()
    ok["risk_mitigants"][0]["mitigants"][0] = one_sentence
    assert PitchDeckContent.model_validate(ok).risk_mitigants[0].mitigants[0] == one_sentence

    # A paragraph-length mitigant (> 160 chars) is rejected.
    bad = _sample_content().model_dump()
    bad["risk_mitigants"][0]["mitigants"][0] = "x" * 161
    with pytest.raises(ValidationError):
        PitchDeckContent.model_validate(bad)


def test_registry_loads_16_blank_slide_library_entries():
    entries = load_slide_library_registry()

    assert len(entries) == 16
    assert entries[0].library_entry_id == "pitch-cover"
    assert entries[6].library_entry_id == "public-company-overview"
    assert entries[7].library_entry_id == "financial-summary"
    # Insider-ownership slide follows Financial Summary (before Considerations).
    assert entries[8].library_entry_id == "insider-ownership"
    assert entries[8].static is False
    assert entries[9].library_entry_id == "acquirer-considerations-mitigants"
    assert entries[10].library_entry_id == "comparable-companies"
    # Precedent-transactions slide follows comparable-companies.
    assert entries[11].library_entry_id == "precedent-transactions"
    assert entries[11].static is False
    assert entries[12].library_entry_id == "key-investment-highlights"
    assert entries[13].library_entry_id == "market-entry-targets"
    assert entries[12].static is False
    assert entries[13].static is False
    assert entries[14].library_entry_id == "disclaimer"
    assert entries[14].static is True
    assert entries[15].library_entry_id == "contact"
    assert entries[15].static is True


def test_pitch_deck_wireframe_uses_blank_library_order():
    # Pin to one market-entry slide (2 targets) so this stays a focused check of
    # the 16-entry blank-library order; the default-8 behaviour is covered by
    # test_pitch_wireframe_expands_market_entry_slides.
    plan = build_pitch_deck_slide_plan(
        company=Company(legal_name="SampleCo Ltd.", ticker="TSX:SMP"),
        section_labels=["Overview", "Financial Summary", "Valuation", "Process"],
        market_entry_target_count=2,
    )

    assert plan.deliverable_type == "pitch"
    assert len(plan.slides) == 16
    # Insider-ownership now follows Financial Summary (index 8).
    assert plan.slides[8].library_entry_id == "insider-ownership"
    assert plan.slides[8].content_block["requires"] == ["ownership_table"]
    # Precedent-transactions follows comparable-companies (index 11).
    assert [slide.library_entry_id for slide in plan.slides[10:14]] == [
        "comparable-companies",
        "precedent-transactions",
        "key-investment-highlights",
        "market-entry-targets",
    ]
    assert plan.slides[11].content_block["requires"] == ["precedents_takeaway"]
    assert plan.slides[12].content_block["requires"] == ["investment_highlights"]
    assert plan.slides[13].content_block["requires"] == ["market_entry_row_labels", "market_entry_targets"]
    assert [slide.library_entry_id for slide in plan.slides[:8]] == [
        "pitch-cover",
        "executive-summary",
        "infor-overview",
        "infor-ma-advisor",
        "infor-key-highlights",
        "section-divider",
        "public-company-overview",
        "financial-summary",
    ]
    assert plan.slides[2].content_block["template_static"] is True
    assert plan.slides[6].content_block["requires"] == ["company_overview_bullets", "cap_table"]


def test_assemble_pitch_deck_preserves_static_slides_and_fills_allowed_fields(tmp_path: Path):
    plan = build_pitch_deck_slide_plan(
        company=Company(legal_name="SampleCo Ltd.", ticker="TSX:SMP"),
        section_labels=["Overview", "Financial Summary", "Valuation", "Process"],
    )
    plan_path = write_slide_plan(plan, tmp_path / "slide_plan.json")
    content = _sample_content()
    content_path = tmp_path / "content.json"
    content_path.write_text(content.model_dump_json(indent=2), encoding="utf-8")

    deck_path = assemble_pitch_deck(
        slide_plan_path=plan_path,
        content_path=content_path,
        template_path=TEMPLATE,
        output_dir=tmp_path,
        financial_metric_labels=["Revenue", "Gross Profit", "Adjusted EBITDA", "Net Income"],
        # This test is about fill logic, so the converge loop is off — the module
        # convention `_assemble` applies, which this call predates and bypassed by
        # going straight to the assembler. Leaving it on cost ~45s of renders for
        # assertions that are entirely about text. Geometry is covered by the
        # `converge=True` tests below and by test_deck_repair.py.
        converge=False,
    )

    assert deck_path.exists()
    prs = Presentation(deck_path)
    assert len(prs.slides) == 16
    # Slide-8 tiles are filled from the financial-summary stage labels.
    assert "Adjusted EBITDA" in _all_slides_text(prs)
    text_parts = []

    def _collect(shapes):
        for shape in shapes:
            if getattr(shape, "has_text_frame", False):
                text_parts.append(shape.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        text_parts.append(cell.text)
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                _collect(shape.shapes)

    for slide in prs.slides:
        _collect(slide.shapes)
    all_text = "\n".join(text_parts)
    assert "[CLIENT NAME]" not in all_text
    assert "[Date]" not in all_text
    assert "SampleCo Ltd." in all_text
    assert "April 2026" in all_text
    assert "[Cap Table Placeholder]" in all_text  # insertion skill fills this later
    assert "[Pie Chart Placeholder]" in all_text  # intentionally deferred
    assert "[Placeholder for Metric #1 Chart]" in all_text  # intentionally deferred
    assert "[Placeholder for Comps Chart]" in all_text  # comps slide stays a placeholder
    # Precedent-transactions slide: takeaway filled, chart stays a placeholder.
    assert "Recent precedent transactions in the sector cleared at premium multiples" in all_text
    assert "[Placeholder for Precedents Chart]" in all_text
    assert "Neil Selfe, Managing Principal" in all_text
    # Ownership slide present; both sides are placeholders when no ownership workbook is supplied.
    assert "[Placeholder for Insider Ownership]" in all_text
    assert "[Placeholder for Institutional Ownership]" in all_text

    # New content slides are filled and their placeholders replaced.
    assert "Durable Recurring Revenue" in all_text
    assert "SampleCo offers a scarce, scaled platform" in all_text
    assert "Potential Canada Market Entry Targets" in all_text
    assert "[Investment Highlight #1]" not in all_text
    assert "Potential [x] Market Entry Targets" not in all_text
    # Market-entry logo boxes name the company INSIDE the library's own
    # '[Placeholder for …]' wording, so the box says whose logo goes there while
    # still reading as something to clear. With no target name supplied it stays
    # the generic '[Placeholder for Logo]'.
    assert "[Placeholder for Logo]" in all_text

    # Static credential slide text remains present.
    assert "INFOR is Canada’s leading provider of innovative, independent, forward thinking financial & strategic advice" in all_text
    assert "These materials are confidential and proprietary" in all_text


def test_earnings_update_plan_runs_captable_before_deck_for_insertion():
    plan_path = PLUGIN_ROOT / "plans" / "earnings-update.yaml"
    plan = Plan.model_validate(yaml.safe_load(plan_path.read_text(encoding="utf-8")))

    assert [stage.id for stage in plan.stages] == [
        "wireframe",
        "content",
        "ltm-metrics",
        "captable",
        "deck",
        "deckcheck",
    ]
    deck_stage = next(s for s in plan.stages if s.id == "deck")
    assert deck_stage.inputs["captable_workbook_path"] == "$stages.captable.workbook_path"
    assert deck_stage.inputs["template_name"] == "INFOR Slide Library.pptx"
    # Since Phase D `deck` produces the final ARTEFACT: there is no aggregation
    # stage, because every producer wrote its own tab of the deal's single
    # workbook. Phase G's `deckcheck` follows it but writes no deliverable — it
    # audits the figures on the deck `deck` produced.
    assert plan.stages[-1].id == "deckcheck"
    assert plan.stages[-1].inputs["deck_path"] == "$stages.deck.deck_path"
    # No stage gates (v0.5.49): `deck` reports its outputs at the wave boundary and
    # the run continues to `deckcheck` without an approval pause.
    assert deck_stage.checkpoint == "informational"
    assert [s.id for s in plan.stages if s.checkpoint != "informational"] == []


def test_pitch_library_poc_plan_stage_order():
    plan_path = PLUGIN_ROOT / "plans" / "pitch.yaml"
    plan = Plan.model_validate(yaml.safe_load(plan_path.read_text(encoding="utf-8")))

    assert [stage.id for stage in plan.stages] == [
        "wireframe",
        "content",
        "financial-summary",
        "ltm-metrics",
        "captable",
        "ownership",
        "comps",
        "precedents",
        "deck",
        "financial-charts",
        "deckcheck",
    ]
    assert plan.stages[0].skill == "pitch-wireframe"
    assert plan.stages[1].skill == "pitch-content"
    assert plan.stages[2].skill == "financial-summary"
    assert plan.stages[3].skill == "ltm-metrics"
    assert plan.stages[4].skill == "captable"
    assert plan.stages[5].skill == "ownership"
    assert plan.stages[6].skill == "comps"
    assert plan.stages[7].skill == "precedents"
    assert plan.stages[8].skill == "deck-assembler"
    assert plan.stages[9].skill == "financial-charts"
    assert plan.stages[10].skill == "deckcheck"
    # financial-charts produces the final artefact: it charts the deal workbook and
    # edits the deck. Since Phase D its only ordering constraint is `deck` — the
    # aggregation stage it used to wait for is gone, along with the combined
    # workbook it produced.
    fc_stage = next(s for s in plan.stages if s.id == "financial-charts")
    assert fc_stage.inputs["deal_workbook"] == "$deal.deal_workbook"
    assert fc_stage.inputs["deck_path"] == "$stages.deck.deck_path"
    # Phase G's `deckcheck` audits that FINAL artefact, so it reads
    # financial-charts' deck_path, not deck's — a review of the pre-chart deck
    # would miss every figure the charts carry, and reading a file
    # financial-charts is editing in place would race it.
    dc_stage = next(s for s in plan.stages if s.id == "deckcheck")
    assert dc_stage.inputs["deck_path"] == "$stages.financial-charts.deck_path"
    # No stage gates (v0.5.49): `deck` no longer holds the charts and `deckcheck`
    # behind an analyst approval — the run carries on to the final artefact.
    assert next(s for s in plan.stages if s.id == "deck").checkpoint == "informational"
    assert [s.id for s in plan.stages if s.checkpoint != "informational"] == []
    deck_stage = next(s for s in plan.stages if s.id == "deck")
    assert deck_stage.inputs["slide_plan_path"] == "$stages.wireframe.slide_plan_path"
    assert deck_stage.inputs["content_bundle_path"] == "$stages.content.content_bundle_path"
    assert deck_stage.inputs["captable_workbook_path"] == "$stages.captable.workbook_path"
    assert deck_stage.inputs["ownership_workbook_path"] == "$stages.ownership.workbook_path"
    # The deck reads the four Financial Summary tile labels from the financial-summary stage.
    assert deck_stage.inputs["financial_metric_labels"] == "$stages.financial-summary.financial_metric_labels"
    # financial-summary is the single source of truth for the deck's four metrics and
    # drives ltm-metrics' extra bridges; it runs before ltm-metrics.
    fs_stage = next(s for s in plan.stages if s.id == "financial-summary")
    assert fs_stage.inputs["company"] == "$deal.subject_company"
    assert {o.name for o in fs_stage.outputs} == {"workbook_path", "financial_metric_labels", "ltm_bridge_specs"}
    ltm_stage = next(s for s in plan.stages if s.id == "ltm-metrics")
    assert ltm_stage.inputs["ltm_bridge_specs"] == "$stages.financial-summary.ltm_bridge_specs"
    # Ownership runs after captable so F35 can be sourced from the cap table's basic shares.
    ownership_stage = next(s for s in plan.stages if s.id == "ownership")
    assert ownership_stage.inputs["captable_workbook_path"] == "$stages.captable.workbook_path"
    # Comps only needs the target's facts; it generates its own companion workbook.
    comps_stage = next(s for s in plan.stages if s.id == "comps")
    assert comps_stage.inputs["company"] == "$deal.subject_company"
    assert [o.name for o in comps_stage.outputs] == ["workbook_path"]
    # Precedents, like comps, only needs the target's facts; own companion workbook.
    precedents_stage = next(s for s in plan.stages if s.id == "precedents")
    assert precedents_stage.inputs["company"] == "$deal.subject_company"
    assert [o.name for o in precedents_stage.outputs] == ["workbook_path"]
    # LTM metrics feed the cap table's LTM valuation column (mirrors earnings update).
    captable_stage = next(s for s in plan.stages if s.id == "captable")
    assert captable_stage.inputs["ltm_revenue"] == "$stages.ltm-metrics.ltm_revenue"
    assert captable_stage.inputs["ltm_adj_ebitda"] == "$stages.ltm-metrics.ltm_adj_ebitda"
    # Every workbook-producing stage writes a tab of the deal's ONE workbook, so
    # each is handed the same `$deal.deal_workbook` instead of emitting a
    # standalone file for a later merge to consolidate.
    for stage_id in ("ltm-metrics", "financial-summary", "ownership", "comps",
                     "precedents", "captable"):
        stage = next(s for s in plan.stages if s.id == stage_id)
        assert stage.inputs["deal_workbook"] == "$deal.deal_workbook", stage_id


# ─── Helpers for the post-review fixes ───────────────────────────────────────

def _assemble(
    tmp_path: Path,
    content: PitchDeckContent,
    *,
    market_entry_target_count=None,
    financial_metric_count=None,
    include_investment_highlights=None,
    **kwargs,
) -> Path:
    # The deck-contract converge loop renders the deck (seconds per call), so it is
    # OFF for the tests that are about fill logic and ON for the geometry tests that
    # exist to exercise it — see test_converges_* below and test_deck_contract.py.
    kwargs.setdefault("converge", False)
    plan = build_pitch_deck_slide_plan(
        company=Company(legal_name="SampleCo Ltd.", ticker="TSX:SMP"),
        section_labels=["Overview", "Financial Summary", "Valuation", "Process"],
        market_entry_target_count=market_entry_target_count,
        financial_metric_count=financial_metric_count,
        include_investment_highlights=include_investment_highlights,
    )
    plan_path = write_slide_plan(plan, tmp_path / "slide_plan.json")
    content_path = tmp_path / "content.json"
    content_path.write_text(content.model_dump_json(indent=2), encoding="utf-8")
    kwargs.setdefault(
        "financial_metric_labels",
        ["Revenue", "Gross Profit", "Adjusted EBITDA", "Net Income"],
    )
    return assemble_pitch_deck(
        slide_plan_path=plan_path,
        content_path=content_path,
        template_path=TEMPLATE,
        output_dir=tmp_path,
        **kwargs,
    )


def _content_with_targets(n: int) -> PitchDeckContent:
    """Clone the sample content but with `n` market-entry targets (12 cells each)."""
    base = _sample_content().model_dump()
    template_cells = base["market_entry_targets"][0]["cells"]
    base["market_entry_targets"] = [
        {"cells": [f"Target {i + 1}"] + template_cells[1:]} for i in range(n)
    ]
    return PitchDeckContent.model_validate(base)


def _write_sample_cap_table(deal_dir: Path, currency: str = "CAD") -> Path:
    """The cap table the assembler is really handed: the deal workbook.

    `pitch.yaml` passes `$stages.captable.workbook_path`, which since Phase D is
    the deal workbook — one file, `captable` tab. So this builds it with
    `init_deal_workbook`, the conductor's own deal-init call, and writes only the
    output currency, resolved through `infor_cap_output_ccy` like the captable
    skill does.

    It used to be a synthetic `Workbook()` titled `Cap with Links` with both
    names stamped on by hand — a copy of the SOURCE template's shape.
    `build_deal_workbook_template.py` renames that sheet to `captable`, so the
    fixture and the artefact disagreed and this test stayed green through the
    whole v0.5.45 outage; see `test_assembler_deal_workbook_inputs`.
    """
    from deal_workbook import TAB_CAPTABLE, TabSpec, init_deal_workbook, write_tab
    from template_layout import NAME_CAP_OUTPUT_CCY, resolve_name_cell

    path = init_deal_workbook(
        deal_dir=deal_dir, deliverable_type="pitch", deal_name="Project Test"
    )

    # The shipped `captable` tab already carries real content across
    # `infor_cap_picture_range` (B15:F40), so only the output currency — which
    # drives the `[x]$MM` footnote letter — has to be set.
    def _fill(_wb, ws):
        ws[resolve_name_cell(ws, NAME_CAP_OUTPUT_CCY)] = currency

    write_tab(
        path, TAB_CAPTABLE, TabSpec(write=_fill, verify_names=(NAME_CAP_OUTPUT_CCY,))
    )
    return path


def _set_fx_rate(deal_workbook: Path, rate: float) -> Path:
    """Write the cap table's FX rate, resolved through `infor_fx_rate` like the skill."""
    from deal_workbook import TAB_CAPTABLE, TabSpec, write_tab
    from template_layout import NAME_FX_RATE, resolve_name_cell

    def _fill(_wb, ws):
        ws[resolve_name_cell(ws, NAME_FX_RATE)] = rate

    write_tab(deal_workbook, TAB_CAPTABLE, TabSpec(write=_fill, verify_names=(NAME_FX_RATE,)))
    return deal_workbook


def _all_slides_text(prs: Presentation) -> str:
    parts: list[str] = []

    def _collect(shapes):
        for shape in shapes:
            if getattr(shape, "has_text_frame", False):
                parts.append(shape.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        parts.append(cell.text)
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                _collect(shape.shapes)

    for slide in prs.slides:
        _collect(slide.shapes)
    return "\n".join(parts)


# The two market-entry logo boxes in the library, ordered left→right by .left
# (matches the assembler's own sort): Rectangle 7 sits above the first target
# column, Rectangle 5 above the second.
_LOGO_SHAPE_NAMES = {"Rectangle 5", "Rectangle 7"}


def _logo_boxes(slide):
    return sorted(
        (s for s in slide.shapes if s.name in _LOGO_SHAPE_NAMES),
        key=lambda s: s.left,
    )


# ─── Fix 1: slide 2 exec summary keeps template colour + bullets ─────────────

def test_slide2_exec_summary_keeps_template_colour_and_bullets(tmp_path: Path):
    deck_path = _assemble(tmp_path, _sample_content())
    shape = find_shape(Presentation(deck_path).slides[1], "Content Placeholder 7")
    paras = [p for p in shape.text_frame.paragraphs if p.text.strip()]
    assert len(paras) == 4  # the four exec-summary bullets
    for p in paras:
        buchars = [e for e in p._p.iter() if e.tag.endswith("}buChar")]
        assert buchars, "every exec-summary bullet must keep its template glyph"
        for run in p.runs:
            xml = run._r.xml
            assert "<a:solidFill>" in xml, "run must carry an explicit template colour"
            assert "1B2759" not in xml, "body text must not fall back to the navy list colour"


# ─── Fix 4a: fixed 12-row market-entry structure ─────────────────────────────

def test_market_entry_row_labels_enforce_fixed_12_row_structure():
    # Wrong count (8 rows) rejected even when targets align 1:1.
    bad = _sample_content().model_dump()
    bad["market_entry_row_labels"] = [
        "Overview", "Headquarters", "Year Founded", "A", "B", "C", "Scale KPIs", "Strategic Rationale",
    ]
    bad["market_entry_targets"] = [{"cells": ["x"] * 8}, {"cells": ["y"] * 8}]
    with pytest.raises(ValidationError):
        PitchDeckContent.model_validate(bad)

    # Wrong fixed top label rejected.
    bad2 = _sample_content().model_dump()
    labels = list(bad2["market_entry_row_labels"])
    labels[0] = "Summary"  # must be 'Overview'
    bad2["market_entry_row_labels"] = labels
    with pytest.raises(ValidationError):
        PitchDeckContent.model_validate(bad2)

    # Wrong fixed bottom label rejected.
    bad3 = _sample_content().model_dump()
    labels = list(bad3["market_entry_row_labels"])
    labels[-1] = "Why Acquire"  # must be 'Strategic Rationale'
    bad3["market_entry_row_labels"] = labels
    with pytest.raises(ValidationError):
        PitchDeckContent.model_validate(bad3)


# ─── Fix 4b / 4c: N targets across N/2 slides, formatting ────────────────────

def test_market_entry_expands_two_targets_per_slide(tmp_path: Path):
    deck_path = _assemble(tmp_path, _content_with_targets(8))
    prs = Presentation(deck_path)
    # 16 base - 1 market-entry + 4 market-entry = 19 slides.
    assert len(prs.slides) == 19
    titles = [find_shape(prs.slides[13 + j], "Title 1").text for j in range(4)]
    assert titles == [
        "Potential Canada Market Entry Targets (1 of 4)",
        "Potential Canada Market Entry Targets (2 of 4)",
        "Potential Canada Market Entry Targets (3 of 4)",
        "Potential Canada Market Entry Targets (4 of 4)",
    ]
    # Static disclaimer + contact preserved at the tail.
    tail = _all_slides_text(prs)
    assert "These materials are confidential and proprietary" in tail
    assert "Neil Selfe, Managing Principal" in tail


def test_market_entry_tables_clamped_to_fixed_height(tmp_path: Path):
    # After the cells are filled, every market-entry (acquisition-target) table is
    # clamped to 5.71" so long content can't run it off the slide — the graphic
    # frame and the row heights both land on the target.
    deck_path = _assemble(tmp_path, _content_with_targets(8))  # 4 market-entry slides
    prs = Presentation(deck_path)
    target = Inches(5.71)
    for j in range(4):
        table_shape = next(
            s for s in prs.slides[13 + j].shapes if getattr(s, "has_table", False)
        )
        assert table_shape.height == target, f"slide {13 + j} frame height not clamped"
        row_total = sum(r.height for r in table_shape.table.rows)
        assert row_total == target, f"slide {13 + j} row heights must sum to the target"


def test_market_entry_table_formatting_no_blank_rows(tmp_path: Path):
    deck_path = _assemble(tmp_path, _sample_content())  # 2 targets -> 1 slide
    table = next(
        s for s in Presentation(deck_path).slides[13].shapes if getattr(s, "has_table", False)
    ).table
    # 12 data rows (rows 1-12); every one populated, label white@11 (a long
    # label may step down one size so it can't wrap the row taller), values @9.
    for row in range(1, 13):
        label_run = table.cell(row, 0).text_frame.paragraphs[0].runs[0]
        assert label_run.font.size in (Pt(11), Pt(10), Pt(9)), "label column must be 9-11 pt"
        assert str(label_run.font.color.rgb) == "FFFFFF", "label column must be white"
        for col in (1, 2):
            cell = table.cell(row, col)
            assert cell.text.strip(), f"data cell ({row},{col}) must be populated"
            assert cell.text_frame.paragraphs[0].runs[0].font.size == Pt(9), "values must be 9 pt"
    # Short labels keep the template 11 pt; only an over-wide label steps down.
    assert table.cell(2, 0).text_frame.paragraphs[0].runs[0].font.size == Pt(11)  # 'Headquarters'


def test_market_entry_labels_are_written_at_one_uniform_size(tmp_path: Path):
    """The fill writes every row label at 11 pt and clamps; it measures nothing.

    Until v0.5.39 it sized each label against a per-character Palatino
    advance-width table and stepped the over-wide ones (e.g. 'Geographic
    Footprint', 1.50" at 11 pt in a 1.457" column) down individually, because a
    wrapped label re-grows its row — one of the two mechanisms behind PRL17's
    5.91" table. The converge loop now measures the rendered table and caps the
    label column only when it actually overruns, which also keeps the column
    uniform instead of mixing 11 pt and 10 pt labels.
    """
    deck_path = _assemble(tmp_path, _sample_content())
    table_shape = next(
        s for s in Presentation(deck_path).slides[13].shapes if getattr(s, "has_table", False)
    )
    table = table_shape.table

    sizes = {
        table.cell(i + 1, 0).text_frame.paragraphs[0].runs[0].font.size
        for i in range(len(_sample_content().market_entry_row_labels))
    }
    assert sizes == {Pt(11)}, f"labels must go in at one uniform 11 pt, got {sizes}"
    # The declared rows sum to exactly the clamp. There is deliberately no per-row
    # content-height floor any more: whether the RENDERED table stays inside the
    # clamp is measured by the converge loop (test_market_entry_long_content_converges).
    assert sum(Emu(r.height).inches for r in table.rows) == pytest.approx(5.71, abs=0.001)
    assert Emu(table_shape.height).inches == pytest.approx(5.71, abs=0.001)


@pytest.mark.skipif(
    find_soffice() is None, reason="the converge loop measures on a LibreOffice render"
)
def test_market_entry_repair_caps_the_label_column_before_the_value_columns(tmp_path: Path):
    """A shrink lands where there is most to give: labels 11 -> 10, values stay 9.

    This is the behaviour that made deleting the Palatino width table safe. The
    repair caps every body run at (largest size - k) rather than subtracting k from
    each, so a defect the 11 pt labels caused does not drag the 9 pt value copy —
    deliberately set below the library's 10 pt for headroom — down with it.
    """
    deck_path = _assemble(tmp_path, _sample_content(), converge=True)
    table = next(
        s for s in Presentation(deck_path).slides[13].shapes if getattr(s, "has_table", False)
    ).table

    labels = {table.cell(i + 1, 0).text_frame.paragraphs[0].runs[0].font.size for i in range(12)}
    values = {table.cell(i + 1, 1).text_frame.paragraphs[0].runs[0].font.size for i in range(12)}
    assert len(labels) == 1 and len(values) == 1, "each column stays uniform"
    label_pt, value_pt = labels.pop().pt, values.pop().pt
    assert label_pt <= 11 and value_pt <= 9
    assert value_pt == 9, (
        f"the value columns must keep their 9 pt while the labels absorb the shrink, "
        f"got labels {label_pt} / values {value_pt}"
    )


def test_market_entry_odd_count_blanks_unused_column_and_logo(tmp_path: Path):
    deck_path = _assemble(tmp_path, _content_with_targets(3))  # -> 2 slides
    prs = Presentation(deck_path)
    assert len(prs.slides) == 17  # 16 base - 1 + 2 market-entry
    last_me = prs.slides[14]  # the second (final) market-entry slide
    table = next(s for s in last_me.shapes if getattr(s, "has_table", False)).table
    assert table.cell(1, 1).text.strip(), "the single target fills the first target column"
    assert table.cell(1, 2).text.strip() == "", "the unused target column is blanked"
    boxes = _logo_boxes(last_me)
    assert len(boxes) == 2
    assert boxes[0].text.strip() == "[Placeholder for Logo]", "the single target's logo box is labelled"
    assert boxes[1].text.strip() == "", "the unused (rightmost) logo box is blanked"


def test_market_entry_logo_box_names_each_target(tmp_path: Path):
    """The box says whose logo goes there, and still reads as a placeholder.

    The delivered Project Open Text deck printed "[Glean Logo]" / "[Pinecone
    Logo]" across four slides. The guidance was deliberate and worth keeping —
    the analyst has to know which column is whose — but bracketed prose reads as
    a caption, not as something to clear. Naming the target INSIDE the library's
    own "[Placeholder for …]" wording keeps both.
    """
    base = _sample_content().model_dump()
    base["market_entry_targets"] = [
        {"name": "Kueski", "cells": base["market_entry_targets"][0]["cells"]},
        {"name": "Nubank", "cells": base["market_entry_targets"][1]["cells"]},
    ]
    deck_path = _assemble(tmp_path, PitchDeckContent.model_validate(base))  # 2 targets -> 1 slide
    boxes = _logo_boxes(Presentation(deck_path).slides[13])
    assert [b.text.strip() for b in boxes] == [
        "[Placeholder for Kueski Logo]",
        "[Placeholder for Nubank Logo]",
    ]


# ─── Slide 10: Considerations/Mitigants font sizes match the library ─────────

def test_slide10_risk_table_uses_library_font_sizes(tmp_path: Path):
    deck_path = _assemble(tmp_path, _sample_content())  # 3 risk/mitigant rows
    table = next(
        s for s in Presentation(deck_path).slides[9].shapes if getattr(s, "has_table", False)
    ).table
    # Header row at 12 pt (was 9 pt), matching the library.
    for col, label in ((0, "Considerations"), (1, "Mitigants")):
        run = table.cell(0, col).text_frame.paragraphs[0].runs[0]
        assert run.text == label
        assert run.font.size == Pt(12), "risk/mitigant header must be 12 pt"
    # Body rows at 10 pt (was 8 pt) — three populated rows in the sample content.
    for row in range(1, 4):
        for col in (0, 1):
            run = table.cell(row, col).text_frame.paragraphs[0].runs[0]
            assert run.font.size == Pt(10), "risk/mitigant body must be 10 pt"


def test_investment_highlights_reject_third_bullet():
    # v0.5.24: at most TWO bullets per KIH quadrant — three crowded the boxes.
    base = _sample_content().model_dump()
    base["investment_highlights"][0]["bullets"] = [
        "High retention across multi-year contracts",
        "Predictable cash conversion through cycles",
        "A third bullet the schema must now reject",
    ]
    with pytest.raises(ValidationError):
        PitchDeckContent.model_validate(base)


def _library_risk_table_height() -> int:
    """The library's shipped slide-10 table height (5.18" — the analyst-locked
    render height for the Considerations/Mitigants table)."""
    lib = Presentation(TEMPLATE)
    frame = next(s for s in lib.slides[10].shapes if getattr(s, "has_table", False))
    return frame.height


def test_slide10_risk_table_clamped_to_library_height(tmp_path: Path):
    # v0.5.24: the filled table is clamped back to the library's 5.18" — long
    # mitigant copy re-grew rows past the declared heights (a stored row height
    # is only a render-time minimum) and a live run rendered 5.36".
    target = _library_risk_table_height()
    assert abs(Emu(target).inches - 5.18) < 0.02, "library slide-10 table must ship at ~5.18 in"
    deck_path = _assemble(tmp_path, _sample_content())
    frame = next(
        s for s in Presentation(deck_path).slides[9].shapes if getattr(s, "has_table", False)
    )
    assert frame.height == target, "risk table frame must land on the library height"
    assert sum(r.height for r in frame.table.rows) == target, "row heights must sum to the target"


def _over_tall_risk_content() -> PitchDeckContent:
    """Five rows of three near-cap (160-char) mitigants — PRL18's shape of defect."""
    long_mitigant = (
        "Demonstrate the durability of the platform through multi-year cohort retention, "
        "audited unit economics and a fully documented regulatory compliance record"
    )
    base = _sample_content().model_dump()
    base["risk_mitigants"] = [
        {"risk": f"Detailed acquiror consideration number {i} on diligence depth", "mitigants": [long_mitigant] * 3}
        for i in range(1, 6)
    ]
    return PitchDeckContent.model_validate(base)


def test_slide10_risk_table_is_written_at_the_library_sizes(tmp_path: Path):
    """The fill no longer pre-shrinks: over-tall copy still goes in at 12/10 pt.

    Until v0.5.39 the assembler estimated each row's height and stepped the body
    down before writing — and PRL18 still shipped a 5.36" table, because the
    estimate said it fit. Deciding the size is now the converge loop's job, so the
    fill's contract is simply "the library's sizes, clamped to the library's
    height".
    """
    deck_path = _assemble(tmp_path, _over_tall_risk_content())
    frame = next(
        s for s in Presentation(deck_path).slides[9].shapes if getattr(s, "has_table", False)
    )
    assert frame.table.cell(0, 0).text_frame.paragraphs[0].runs[0].font.size == Pt(12)
    assert frame.table.cell(1, 1).text_frame.paragraphs[0].runs[0].font.size == Pt(10)
    assert frame.height == _library_risk_table_height()


@pytest.mark.skipif(
    find_soffice() is None, reason="the converge loop measures on a LibreOffice render"
)
def test_slide10_risk_table_steps_font_down_when_content_over_tall(tmp_path: Path):
    """Same outcome as before, decided by measurement instead of estimation.

    Over-tall copy must end up below the template's 10 pt with the header left at
    12 pt, and the table must still land on the library height — but now because
    the rendered table was measured overrunning, not because a character-width
    model predicted it would.
    """
    deck_path = _assemble(tmp_path, _over_tall_risk_content(), converge=True)
    frame = next(
        s for s in Presentation(deck_path).slides[9].shapes if getattr(s, "has_table", False)
    )
    table = frame.table

    header_run = table.cell(0, 0).text_frame.paragraphs[0].runs[0]
    assert header_run.font.size == Pt(12), "header must stay 12 pt"
    body_size = table.cell(1, 1).text_frame.paragraphs[0].runs[0].font.size
    assert Pt(6) <= body_size < Pt(10), (
        f"over-tall body copy must step below 10 pt, got {body_size.pt}"
    )
    target = _library_risk_table_height()
    assert frame.height == target, "the stepped-down table still lands on the library height"
    assert sum(r.height for r in table.rows) == target

    # And the rendered table now agrees with the declaration, which is the point.
    findings = verify_deck(deck_path, vision=False, out_dir=tmp_path / "qa")
    assert [
        f for f in findings if f.blocking and f.shape == frame.name
    ] == [], "\n".join(str(f) for f in findings if f.shape == frame.name)


def test_pitch_wireframe_expands_market_entry_slides():
    plan = build_pitch_deck_slide_plan(
        company=Company(legal_name="SampleCo Ltd.", ticker="TSX:SMP"),
        market_entry_target_count=8,
    )
    me = [s for s in plan.slides if s.library_entry_id == "market-entry-targets"]
    assert len(me) == 4
    assert len(plan.slides) == 19
    assert me[0].title.endswith("(1 of 4)")
    assert me[3].title.endswith("(4 of 4)")
    # Count unspecified -> default 8 targets (4 market-entry slides, 19-slide plan).
    default_plan = build_pitch_deck_slide_plan(company=Company(legal_name="X Co", ticker="T:X"))
    assert sum(1 for s in default_plan.slides if s.library_entry_id == "market-entry-targets") == 4
    assert len(default_plan.slides) == 19


def test_pitch_deck_inserts_cap_table_into_slide7(tmp_path: Path):
    workbook = _write_sample_cap_table(tmp_path, currency="USD")
    deck_path = _assemble(tmp_path, _sample_content(), captable_workbook_path=workbook)
    prs = Presentation(deck_path)
    slide7 = prs.slides[6]
    assert next((s for s in slide7.shapes if s.name == "Rectangle 3"), None) is None, (
        "slide 7 cap-table placeholder should be replaced by the picture"
    )
    pictures = [s for s in slide7.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert pictures, "expected a cap-table picture on slide 7"
    note = find_shape(slide7, "Text Placeholder 1").text
    assert "US$MM" in note and "[x]" not in note, "footnote currency derived from the cap table (USD)"
    # The figure-footnote currency token is resolved on the highlights slide too (no stray [x]).
    note_hl = find_shape(prs.slides[12], "Text Placeholder 13").text
    assert "US$MM" in note_hl and "[x]" not in note_hl, "highlights-slide footnote currency derived from the cap table"


@pytest.mark.skipif(
    find_soffice() is None, reason="the range renderer needs LibreOffice (see excel_to_powerpoint)"
)
def test_pitch_deck_inserts_ownership_into_slide(tmp_path: Path):
    from deal_workbook import init_deal_workbook
    from ownership_workbook import build_ownership_workbook, InsiderHolding

    own_wb = build_ownership_workbook(
        deal_workbook=init_deal_workbook(
            deal_dir=tmp_path, deliverable_type="pitch", deal_name="Project Test"
        ),
        insiders=[
            InsiderHolding("Barrenechea, Mark James", "Mark Barrenechea (CEO & Director)", 1219092, "2025-03-31"),
            InsiderHolding("Fowlie, Randy", "Randy Fowlie (Director)", [193000, 0], "2025-12-01"),
        ],
        total_shares_outstanding=261_000_000,
    )
    # No `except RuntimeError: pytest.skip` here any more. `_render_range_to_png`
    # raises RuntimeError for a MISSING LibreOffice *and* for a conversion that
    # failed or produced no PDF, so catching it turned a genuine render defect
    # into a green skip. The skipif above states the one condition that is really
    # environmental; everything else must fail.
    deck_path = _assemble(tmp_path, _sample_content(), ownership_workbook_path=own_wb)

    prs = Presentation(deck_path)
    own_slide = prs.slides[8]  # ownership follows Financial Summary (fixed index 8)
    assert next((s for s in own_slide.shapes if s.name == "Rectangle 1"), None) is None, (
        "ownership insider placeholder (Rectangle 1) should be replaced by the picture"
    )
    assert [s for s in own_slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE], (
        "expected an insider-ownership picture on the ownership slide"
    )
    text = _all_slides_text(prs)
    assert "[Placeholder for Insider Ownership]" not in text, "insider placeholder replaced"
    # No Bloomberg export was ingested, so the institutional side stays a placeholder.
    assert "[Placeholder for Institutional Ownership]" in text
@pytest.mark.skipif(
    find_soffice() is None, reason="the range renderer needs LibreOffice (see excel_to_powerpoint)"
)
def test_pitch_deck_inserts_institutions_with_bloomberg(tmp_path: Path):
    """With a Bloomberg export ingested, the ownership slide's right
    "Institutions" placeholder (Rectangle 3) is replaced by the
    Select-Institutions block picture alongside the insider side."""
    from openpyxl import Workbook

    from ownership_workbook import build_ownership_workbook, InsiderHolding

    # Minimal BBG Summary View export: one SEDI-duplicate insider + two institutions.
    bbg_path = tmp_path / "bbg.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.title = "Summary View"
    ws["E7"] = "Test Co Inc."
    ws["E9"] = "TST CN EQUITY"
    for col, header in {"C": "Holder Name", "L": "Position", "N": "Filing Date", "R": "Insider Status"}.items():
        ws[f"{col}13"] = header
    for i, (name, position, status) in enumerate(
        [("Barrenechea Mark James", 1_219_092, "Y"), ("T Rowe Price Group Inc", 150_311, "N-P"), ("Kelso & Co LP", 14_983, "N-P")]
    ):
        ws[f"C{14 + i}"] = name
        ws[f"L{14 + i}"] = position
        ws[f"R{14 + i}"] = status
    wb.save(bbg_path)

    from deal_workbook import init_deal_workbook

    own_wb = build_ownership_workbook(
        deal_workbook=init_deal_workbook(
            deal_dir=tmp_path, deliverable_type="pitch", deal_name="Project Test"
        ),
        insiders=[
            InsiderHolding("Barrenechea, Mark James", "Mark Barrenechea (CEO & Director)", 1219092, "2025-03-31"),
            InsiderHolding("Fowlie, Randy", "Randy Fowlie (Director)", [193000, 0], "2025-12-01"),
        ],
        total_shares_outstanding=261_000_000,
        bloomberg_export_path=bbg_path,
    )
    deck_path = _assemble(tmp_path, _sample_content(), ownership_workbook_path=own_wb)

    prs = Presentation(deck_path)
    own_slide = prs.slides[8]
    for placeholder in ("Rectangle 1", "Rectangle 3"):
        assert next((s for s in own_slide.shapes if s.name == placeholder), None) is None, (
            f"ownership placeholder {placeholder} should be replaced by a picture"
        )
    pictures = [s for s in own_slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pictures) >= 2, "expected insider + institutions pictures on the ownership slide"
    text = _all_slides_text(prs)
    assert "[Placeholder for Insider Ownership]" not in text
    assert "[Placeholder for Institutional Ownership]" not in text


# ─── Overview bullets stay above the LTM revenue section ─────────────────────

def _overview_font_scale(shape):
    from pptx.oxml.ns import qn

    bodyPr = shape.text_frame._txBody.find(qn("a:bodyPr"))
    autofit = bodyPr.find(qn("a:normAutofit")) if bodyPr is not None else None
    if autofit is None or autofit.get("fontScale") is None:
        return None
    return int(autofit.get("fontScale")) / 1000.0


def _overview_content(n_bullets: int, chars: int) -> PitchDeckContent:
    base = _sample_content().model_dump()
    filler = "funding programs, growth rates and guidance in enough detail to overflow "
    base["company_overview_bullets"] = [
        {"text": (f"Bullet {i + 1}: " + filler * 10)[:chars], "level": 0}
        for i in range(n_bullets)
    ]
    return PitchDeckContent.model_validate(base)


def test_slide7_overview_bullets_fitted_above_ltm_band(tmp_path: Path):
    """An over-long overview block must be boxed to the band above the 'LTM
    Revenue Breakdown' header and carry an explicit autofit fontScale —
    PowerPoint ignores a scale-less <a:normAutofit/> on open, which is how live
    runs rendered overview copy straight through the pie section.

    Runs the converge loop: the box sizing is the assembler's, but the *scale* is
    now measured off the render rather than estimated from em constants, so this
    only means anything with the loop on.
    """
    # ~1,050 characters: over the band at full size, inside it once shrunk.
    content = _overview_content(7, 150)
    deck_path = _assemble(tmp_path, content, converge=True)

    slide7 = Presentation(deck_path).slides[6]
    box = find_shape(slide7, "TextBox 9")
    header = next(s for s in slide7.shapes if getattr(s, "has_text_frame", False)
                  and "LTM Revenue" in s.text_frame.text)
    band_avail = Emu(header.top).inches - 0.12 - Emu(box.top).inches
    assert abs(Emu(box.height).inches - band_avail) < 0.02, (
        "overview box must be sized to the band above the LTM revenue header"
    )
    scale = _overview_font_scale(box)
    assert scale is not None and scale < 100.0, (
        "over-long overview copy must ship an explicit normAutofit fontScale"
    )


def test_overview_copy_past_every_shrink_fails_the_stage(tmp_path: Path):
    """Content the loop cannot fit must fail loudly, not ship shrunk-and-spilling.

    The retired estimator floored at 70% and shipped whatever that produced, so a
    deck whose copy did not fit at 70% went out overflowing. 1,850 characters is
    roughly double pitch-content's documented ~950-char overview budget, and no
    scale on the ladder fits it into the band — so the assembler raises.

    This is the behaviour change that makes the whole phase worth having: the
    remedy for over-budget copy is fewer words, and now something says so.
    """
    content = _overview_content(10, 185)
    with pytest.raises(DeckNotConvergedError) as excinfo:
        _assemble(tmp_path, content, converge=True)

    message = str(excinfo.value)
    assert "TextBox 9" in message, f"the failure must name the offending shape: {message}"
    assert "70%" in message, (
        f"the failure must show that the shrink ladder was exhausted: {message}"
    )


# ─── The converge loop's scratch does not land in the deal directory ─────────


def test_a_converged_assembly_leaves_no_qa_scratch_in_the_deal_directory(tmp_path: Path):
    """The output directory is the analyst's, and it is cloud-synced.

    The assembler used to hand the converge loop `out_dir / ".qa"`, so every deck
    left ~170 render files and ~10 MB there permanently — and paid the mount's
    per-file sync overhead while doing it, which is what turned a 23.5 s transform
    into eight consecutive killed attempts. A successful assembly must leave the
    deck and nothing else.
    """
    deck_path = _assemble(tmp_path, _sample_content(), converge=True)

    assert deck_path.is_file()
    assert not (tmp_path / ".qa").exists(), "the QA scratch tree is back in the deal directory"
    strays = sorted(path.name for path in tmp_path.iterdir() if path.is_dir())
    assert strays == [], f"the assembler left directories in the deal directory: {strays}"


def test_a_failed_assembly_keeps_the_failing_renders_where_the_error_says(tmp_path: Path):
    """Failure is the one case that leaves something: the evidence for why.

    Same over-budget copy as `test_overview_copy_past_every_shrink_fails_the_stage`
    — the point here is where the renders end up and that the error names it.
    """
    with pytest.raises(DeckNotConvergedError) as excinfo:
        _assemble(tmp_path, _overview_content(10, 185), converge=True)

    qa = tmp_path / ".qa"
    assert qa.is_dir(), "a failing converge must keep the failing pass for the analyst"
    assert list(qa.rglob("*.png")), "the kept pass carries no renders to look at"
    assert str(qa) in str(excinfo.value), (
        f"the error must name where the renders were kept: {excinfo.value}"
    )


def test_market_entry_long_content_converges(tmp_path: Path):
    """A filled market-entry table must render inside its clamp, measured.

    PRL17's 5.91"-rendered table is the historical version of this. The table is
    clamped to 5.71" and the loop steps the body font down until the *rendered*
    height agrees — no character-width table, no per-row height estimate.
    """
    base = _sample_content().model_dump()
    cells = base["market_entry_targets"][0]["cells"]
    long_cells = [
        (c + " with materially more descriptive copy than the sample carries")[:120]
        for c in cells
    ]
    base["market_entry_targets"] = [{"cells": long_cells} for _ in range(3)]
    content = PitchDeckContent.model_validate(base)

    deck_path = _assemble(tmp_path, content, market_entry_target_count=3, converge=True)

    findings = verify_deck(deck_path, vision=False, out_dir=tmp_path / "qa")
    geometric = [
        f
        for f in findings
        if f.blocking
        and f.kind in {"rendered-overflow", "masked-overflow", "table-taller-than-library"}
    ]
    assert geometric == [], "\n".join(str(f) for f in geometric)


def test_slide7_short_overview_keeps_full_size(tmp_path: Path):
    deck_path = _assemble(tmp_path, _sample_content())
    box = find_shape(Presentation(deck_path).slides[6], "TextBox 9")
    assert _overview_font_scale(box) is None  # short copy — no downscale

# ─── Configurable slide mix (v0.5.26): FS slide count + KIH toggle ───────────

_EIGHT_LABELS = [
    "Revenue", "Gross Profit", "Adjusted EBITDA", "Net Income",
    "Operating Income", "Free Cash Flow", "Gross Margin", "Return on Equity",
]


def _slide_text(slide) -> str:
    parts = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            parts.append(shape.text)
    return "\n".join(parts)


def test_pitch_wireframe_two_financial_summary_slides():
    plan = build_pitch_deck_slide_plan(
        company=Company(legal_name="SampleCo Ltd.", ticker="TSX:SMP"),
        financial_metric_count=8,
    )
    fs = [s for s in plan.slides if s.library_entry_id == "financial-summary"]
    assert [s.title for s in fs] == ["Financial Summary (1 of 2)", "Financial Summary (2 of 2)"]
    assert fs[0].content_block["financial_summary_slide"] == 1
    assert fs[1].content_block["financial_summary_slide"] == 2
    assert fs[0].content_block["financial_summary_slide_count"] == 2
    # 16 base + 3 extra market-entry (default 8 targets) + 1 extra FS = 20.
    assert len(plan.slides) == 20
    assert [s.order for s in plan.slides] == list(range(20))


def test_pitch_wireframe_excludes_investment_highlights():
    plan = build_pitch_deck_slide_plan(
        company=Company(legal_name="SampleCo Ltd.", ticker="TSX:SMP"),
        include_investment_highlights=False,
    )
    ids = [s.library_entry_id for s in plan.slides]
    assert "key-investment-highlights" not in ids
    # 16 base + 3 extra market-entry − 1 KIH = 18.
    assert len(plan.slides) == 18
    assert [s.order for s in plan.slides] == list(range(18))


def test_pitch_wireframe_rejects_non_multiple_of_four_metric_count():
    with pytest.raises(ValueError):
        build_pitch_deck_slide_plan(
            company=Company(legal_name="SampleCo Ltd.", ticker="TSX:SMP"),
            financial_metric_count=6,
        )


def test_assemble_two_financial_summary_slides_shifts_downstream(tmp_path: Path):
    content = _content_with_targets(8)
    deck_path = _assemble(
        tmp_path,
        content,
        market_entry_target_count=8,
        financial_metric_count=8,
        financial_metric_labels=_EIGHT_LABELS,
    )
    prs = Presentation(deck_path)
    # 16 base + 1 extra FS + 3 extra market-entry = 20 slides.
    assert len(prs.slides) == 20
    # Both FS slides retitled and tiled four labels each, in order.
    assert find_shape(prs.slides[7], "Title 6").text == "Financial Summary (1 of 2)"
    assert find_shape(prs.slides[8], "Title 6").text == "Financial Summary (2 of 2)"
    assert find_shape(prs.slides[7], "Rectangle 13").text == "Revenue"
    assert find_shape(prs.slides[7], "Rectangle 14").text == "Net Income"
    assert find_shape(prs.slides[8], "Rectangle 13").text == "Operating Income"
    assert find_shape(prs.slides[8], "Rectangle 14").text == "Return on Equity"
    # Every slide after the FS section shifted by one.
    assert "[Placeholder for Insider Ownership]" in _slide_text(prs.slides[9])
    assert find_shape(prs.slides[10], "Text Placeholder 6").text == content.risks_tagline
    assert find_shape(prs.slides[11], "Text Placeholder 5").text == content.comps_takeaway
    assert find_shape(prs.slides[12], "Text Placeholder 5").text == content.precedents_takeaway
    assert find_shape(prs.slides[13], "Title 1").text == "Key Investment Highlights"
    titles = [find_shape(prs.slides[14 + j], "Title 1").text for j in range(4)]
    assert titles == [
        f"Potential Canada Market Entry Targets ({j + 1} of 4)" for j in range(4)
    ]
    # Static disclaimer + contact preserved at the tail.
    tail = _all_slides_text(prs)
    assert "These materials are confidential and proprietary" in tail
    assert "Neil Selfe, Managing Principal" in tail


def test_assemble_two_fs_slides_requires_eight_labels(tmp_path: Path):
    with pytest.raises(ValueError):
        _assemble(
            tmp_path,
            _content_with_targets(8),
            financial_metric_count=8,  # plan carries 2 FS slides…
            # …but the default helper labels are only four.
        )


def test_assemble_excludes_investment_highlights_slide(tmp_path: Path):
    raw = _content_with_targets(8).model_dump()
    raw["investment_highlights"] = []
    raw["investment_highlights_tagline"] = None
    content = PitchDeckContent.model_validate(raw)
    deck_path = _assemble(tmp_path, content, include_investment_highlights=False)
    prs = Presentation(deck_path)
    # 16 base − 1 KIH + 3 extra market-entry = 18 slides.
    assert len(prs.slides) == 18
    assert "Key Investment Highlights" not in _all_slides_text(prs)
    # Slides before KIH keep their default positions…
    assert find_shape(prs.slides[9], "Text Placeholder 6").text == content.risks_tagline
    assert find_shape(prs.slides[11], "Text Placeholder 5").text == content.precedents_takeaway
    # …and the market-entry section starts one earlier (12 instead of 13).
    titles = [find_shape(prs.slides[12 + j], "Title 1").text for j in range(4)]
    assert titles == [
        f"Potential Canada Market Entry Targets ({j + 1} of 4)" for j in range(4)
    ]
    tail = _all_slides_text(prs)
    assert "These materials are confidential and proprietary" in tail
    assert "Neil Selfe, Managing Principal" in tail


# ─── Footnote currency mapping from the cap table (v0.5.34) ──────────────────


def _cap_table_with_currency(tmp_path: Path, code: str) -> Path:
    """A deal workbook whose `captable` tab carries `code` as its output currency.

    This was a bare `Workbook()` titled `Cap with Links` with **no defined
    names** — so it exercised neither the tab lookup nor
    `infor_cap_output_ccy`, and passed entirely on the two fallbacks the reader
    used to carry. That is why the mapping below stayed green while every real
    build read `F5` of whatever tab was active (`precedents` — the string
    `'Target'`). It delegates to the deal-workbook fixture now.
    """
    return _write_sample_cap_table(tmp_path / code, currency=code)


def test_output_currency_letter_maps_usd_and_cad_explicitly(tmp_path: Path):
    assert _output_currency_letter(_cap_table_with_currency(tmp_path, "USD")) == "US"
    assert _output_currency_letter(_cap_table_with_currency(tmp_path, "CAD")) == "C"


def test_output_currency_letter_returns_iso_code_for_non_dollar_currency(tmp_path: Path):
    # A non-dollar filer renders its ISO code in the footnote — never a silent
    # 'C' (pre-v0.5.34, GBP mapped to C$MM and CHF matched the startswith('C')).
    assert _output_currency_letter(_cap_table_with_currency(tmp_path, "GBP")) == "GBP"
    assert _output_currency_letter(_cap_table_with_currency(tmp_path, "CHF")) == "CHF"


def test_output_currency_letter_defaults_to_c_only_for_an_empty_cell(tmp_path: Path):
    """An empty output-currency cell is the ONE fallback left: the template default.

    A footnote must never ship the literal `[x]`, so a blank cell reads 'C'. That
    is the whole of it — v0.5.45 deleted the sheet fallback, and a workbook that
    cannot be read at all now raises rather than labelling a client deck 'C$MM'
    on no evidence.
    """
    assert _output_currency_letter(_cap_table_with_currency(tmp_path, "")) == "C"

    with pytest.raises(FileNotFoundError):
        _output_currency_letter(tmp_path / "missing.xlsx")


# ─── B4: the currency label is a property of the slide's CONTENT ─────────────
#
# The delivered Project Open Text deck failed this four ways at once — see the
# comment block above `_SecondCurrency` in the assembler. These tests cover each
# way, plus the guard that turns the class into a build failure.


def _figure_slides(prs) -> dict[int, str]:
    """Every slide carrying figures, mapped to its footnote text.

    Keyed off the footnote shape rather than a written-down slide list, because
    the deck's slide mix is a deck-spec option. The two figure-bearing slides
    with no currency footnote — ownership (shares and percentages) and
    considerations/mitigants (prose) — carry no such shape and so are absent by
    construction, which is the distinction the assembler makes too.
    """
    notes: dict[int, str] = {}
    for index, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            if "All figures in" in shape.text:
                notes[index] = shape.text
    return notes


def test_no_unfilled_placeholder_survives_on_any_slide(tmp_path: Path):
    """Not one literal '[x]' anywhere on a delivered deck — swept, not spot-checked.

    The three known cases (the market-entry footnotes, the ownership takeaway,
    the Contact cards) were each found by reading a shipped deck. Asserting on
    those three would re-run the same inspection; the point is the sweep, over
    every slide, every group and every table cell.
    """
    deck_path = _assemble(
        tmp_path,
        _sample_content(),
        captable_workbook_path=_write_sample_cap_table(tmp_path / "wb", currency="CAD"),
    )
    prs = Presentation(deck_path)

    offenders = [
        f"slide {index}: {shape_name}"
        for index, slide in enumerate(prs.slides, start=1)
        for shape_name, text in _iter_slide_text(slide)
        if "[x]" in text
    ]
    assert offenders == []

    # The deferred placeholders are a different thing and DO ship — the check
    # must not have passed by scrubbing them too.
    text = _all_slides_text(prs)
    assert "[Pie Chart Placeholder]" in text
    assert "[Placeholder for Comps Chart]" in text


def test_an_unfilled_placeholder_fails_the_build(tmp_path: Path):
    """The sweep is a build failure, not a review finding.

    Without this the guard could rot to a no-op and
    `test_no_unfilled_placeholder_survives_on_any_slide` would still pass. It
    re-introduces an `[x]` into a deck that already verified clean, so what is
    being tested is the check and nothing else.
    """
    deck_path = _assemble(tmp_path, _sample_content())

    prs = Presentation(deck_path)
    find_shape(prs.slides[9], "Text Placeholder 6").text_frame.paragraphs[0].runs[0].text = "[x]"
    prs.save(deck_path)

    with pytest.raises(ValueError, match=r"unfilled '\[x\]' placeholders"):
        _verify_pitch_output(deck_path)


def test_every_figure_bearing_slide_carries_a_currency_note(tmp_path: Path):
    """Including the Financial Summary slide, which the library ships without one.

    Slide 8 of the delivered deck read "Sources: Company filings" and stopped
    there, over a tab locked to millions — while the overview slide restated the
    same LTM revenue in a different currency two slides earlier.
    """
    deck_path = _assemble(
        tmp_path,
        _sample_content(),
        captable_workbook_path=_write_sample_cap_table(tmp_path / "wb", currency="USD"),
    )
    prs = Presentation(deck_path)
    notes = _figure_slides(prs)

    # overview, financial summary, comps, precedents, highlights, market entry
    assert len(notes) == 6, f"expected six currency-labelled slides, got {sorted(notes)}"
    for index, note in notes.items():
        assert "US$MM" in note, f"slide {index} is not labelled in the content's currency: {note!r}"
        assert "[x]" not in note

    # The Financial Summary footnote is the appended case: the library's own
    # source line survives alongside the note this module wrote.
    fs_note = notes[8]
    assert "Sources: Company filings" in fs_note
    assert "Note: All figures in US$MM" in fs_note


def test_a_mixed_currency_slide_names_both_currencies_and_the_rate(tmp_path: Path):
    """Slide 7: US$ overview bullets above a C$ cap table, and the rate between.

    The deck stated LTM revenue as C$7,344.2 here and $5,207.9 on the Financial
    Summary slide. Both were right; the deck said so nowhere, and 7,344.2 /
    5,207.9 = 1.4102 was only ever in the workbook.
    """
    workbook = _write_sample_cap_table(tmp_path / "wb", currency="CAD")
    _set_fx_rate(workbook, 1.4102)

    deck_path = _assemble(
        tmp_path, _sample_content(), captable_workbook_path=workbook  # content is USD
    )
    note = find_shape(Presentation(deck_path).slides[6], "Text Placeholder 1").text

    assert "US$MM" in note, "the bullets are the content bundle's, and it is USD"
    assert "capitalization table in C$MM" in note, "the pasted picture is CAD — say so"
    assert "US$1.00 = C$1.4102" in note, "the rate that reconciles the two slides"


def test_a_single_currency_slide_does_not_mention_a_second(tmp_path: Path):
    """No clause when the cap table converts to the currency the content is in."""
    deck_path = _assemble(
        tmp_path,
        _sample_content(),  # USD
        captable_workbook_path=_write_sample_cap_table(tmp_path / "wb", currency="USD"),
    )
    note = find_shape(Presentation(deck_path).slides[6], "Text Placeholder 1").text

    assert "US$MM" in note
    assert "capitalization table in" not in note


def test_a_missing_fx_rate_still_names_both_currencies(tmp_path: Path):
    """An un-entered FX cell loses the rate, not the disclosure.

    The rate is read from the cap table, so there is nothing to print when the
    cell is empty. Naming both currencies without it beats printing a number
    this module made up, and beats silently picking one.
    """
    deck_path = _assemble(
        tmp_path,
        _sample_content(),  # USD
        captable_workbook_path=_write_sample_cap_table(tmp_path / "wb", currency="CAD"),
    )
    note = find_shape(Presentation(deck_path).slides[6], "Text Placeholder 1").text

    assert "capitalization table in C$MM" in note
    assert "=" not in note.split("capitalization table")[1]


def test_the_deck_currency_no_longer_comes_from_the_cap_table(tmp_path: Path):
    """The regression proper: a C$ cap table must not relabel US$ content slides.

    One letter read off `captable!F5` used to label the highlights slide and
    every market-entry slide, whose every box is the content bundle's. Five
    slides of the delivered deck said C$MM over US$ figures.
    """
    deck_path = _assemble(
        tmp_path,
        _sample_content(),  # USD
        captable_workbook_path=_write_sample_cap_table(tmp_path / "wb", currency="CAD"),
    )
    prs = Presentation(deck_path)

    highlights = find_shape(prs.slides[12], "Text Placeholder 13").text
    assert "US$MM" in highlights and "C$MM" not in highlights

    market_entry = find_shape(prs.slides[13], "Text Placeholder 3").text
    assert "US$MM" in market_entry and "C$MM" not in market_entry


def test_a_non_dollar_content_currency_renders_its_iso_code(tmp_path: Path):
    """A GBP filer is labelled 'GBP MM', never mislabelled with a dollar sign."""
    content = _sample_content().model_dump()
    content["figure_currency"] = "GBP"
    deck_path = _assemble(tmp_path, PitchDeckContent.model_validate(content))

    note = find_shape(Presentation(deck_path).slides[6], "Text Placeholder 1").text
    assert "GBP MM" in note and "$" not in note.split("All figures in")[1]


def test_figure_currency_rejects_a_rendered_token():
    """The bundle states a CODE. A token is the thing the deck renders FROM it."""
    base = _sample_content().model_dump()
    for bad in ("US$MM", "us dollars", "USDD", ""):
        base["figure_currency"] = bad
        with pytest.raises(ValidationError):
            PitchDeckContent.model_validate(base)

    base["figure_currency"] = "cad"  # case is normalised, not rejected
    assert PitchDeckContent.model_validate(base).figure_currency == "CAD"


# ─── B10: an ordered list of data is paired to shapes in VISUAL order ────────


def test_section_divider_boxes_read_in_order_top_to_bottom(tmp_path: Path):
    """The divider's boxes, read the way a human reads them, match section_labels.

    The delivered deck read Overview / Valuation / Financial Summary / Process
    against a content bundle holding Overview / Financial Summary / Valuation /
    Process. Nothing upstream disagreed — `slide.shapes` is XML document order,
    the four boxes are all named "Rounded Rectangle 19", and document order on
    this library is the 1st, 3rd, 2nd and 4th box down the slide.

    The library's own document order is what makes this test mean something, so
    it is asserted rather than assumed: a re-save that happened to align the two
    orders would turn every assertion below into a tautology, and that is exactly
    how the bug shipped.
    """
    library_boxes = [
        s for s in Presentation(TEMPLATE).slides[5].shapes if s.name == "Rounded Rectangle 19"
    ]
    assert [s.top for s in library_boxes] != sorted(s.top for s in library_boxes), (
        "the library's divider no longer authors its boxes out of visual order, so "
        "this test can no longer distinguish a geometric sort from a document-order "
        "one — rebuild the fixture with a scrambled slide before trusting it"
    )

    labels = ["Overview", "Financial Summary", "Valuation", "Process"]
    content = _sample_content().model_dump()
    content["section_labels"] = labels
    content["current_section"] = "Overview"
    deck_path = _assemble(tmp_path, PitchDeckContent.model_validate(content))

    boxes = [s for s in Presentation(deck_path).slides[5].shapes if s.name == "Rounded Rectangle 19"]
    top_to_bottom = [s.text.strip() for s in sorted(boxes, key=lambda s: s.top)]
    assert top_to_bottom == labels


def test_the_financial_summary_tiles_are_named_in_reading_order():
    """`_FS_METRIC_TILES` is a name list, so its ORDER is an unchecked claim.

    The tiles are addressed by name (Phase C: never by position), which is right
    — but the `financial-summary` stage emits its labels in reading order, and
    nothing else in the assembler would notice if the name→position mapping went
    stale. The names read out of order alphabetically, so this is not obvious
    from the list itself.
    """
    from pitch_deck_assembler import _FS_METRIC_TILES

    slide = Presentation(TEMPLATE).slides[8]
    tiles = [find_shape(slide, name) for name in _FS_METRIC_TILES]
    assert [s.text.strip() for s in tiles] == ["Metric #1", "Metric #2", "Metric #3", "Metric #4"]
    assert tiles == shapes_in_reading_order(tiles), (
        "the named tile order no longer matches the library's geometry — the FS "
        "labels would render top-left, top-right, bottom-left, bottom-right in the "
        "wrong slots"
    )


def test_the_highlight_quadrants_are_numbered_in_reading_order():
    """The KIH fill sorts by the printed number; that must agree with geometry.

    Sorting on what the reader sees is the strongest ordering available here, but
    it is only correct if the library numbers its quadrants down the slide. The
    four groups arrive in document order 1, 4, 3, 2, so the sort is load-bearing.
    """
    slide = Presentation(TEMPLATE).slides[13]
    groups = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.GROUP]
    numbered = []
    for group in groups:
        oval = next(
            s for s in group.shapes
            if s.name.startswith("Oval") and s.text.strip().isdigit()
        )
        numbered.append((int(oval.text.strip()), group))

    assert [n for n, _ in numbered] != sorted(n for n, _ in numbered), (
        "the library no longer authors the quadrants out of numeric order, so this "
        "test can no longer tell a sorted fill from an unsorted one"
    )
    by_number = [g for _, g in sorted(numbered, key=lambda pair: pair[0])]
    assert by_number == shapes_in_reading_order(by_number)


# ─── B13b: the Contact slide's cards ─────────────────────────────────────────


def test_unfilled_contact_cards_are_removed_not_left_as_placeholders(tmp_path: Path):
    """Three of the four cards ship as '[x]'; they are deleted, not blanked.

    A deck reaching a client with three empty contact cards reads as unfinished —
    and blanking them would leave three empty boxes in a 2x2 grid, which reads
    barely better. The library's own filled card survives untouched, which is
    what makes "no contacts supplied" a sensible default rather than a gap.
    """
    deck_path = _assemble(tmp_path, _sample_content())
    contact = Presentation(deck_path).slides[-1]

    cards = [
        s for s in contact.shapes
        if getattr(s, "has_table", False) and len(s.table.rows) == 3
    ]
    assert len(cards) == 1, "the three placeholder cards are gone, not emptied"
    assert cards[0].table.cell(0, 0).text == "Neil Selfe, Managing Principal"
    # The address table beside them is not a card and must survive.
    assert "200 Bay Street" in _all_slides_text(Presentation(deck_path))


def test_supplied_contacts_fill_cards_in_reading_order(tmp_path: Path):
    """Declared contacts replace the library's cards; surplus cards are deleted."""
    content = _sample_content().model_dump()
    content["contacts"] = [
        {"name": "A. Banker", "title": "Managing Principal",
         "phone": "(416) 555-0101", "email": "abanker@example.com"},
        {"name": "B. Banker", "title": "Vice President",
         "phone": "(416) 555-0102", "email": "bbanker@example.com"},
    ]
    deck_path = _assemble(tmp_path, PitchDeckContent.model_validate(content))
    contact = Presentation(deck_path).slides[-1]

    cards = shapes_in_reading_order(
        s for s in contact.shapes if getattr(s, "has_table", False) and len(s.table.rows) == 3
    )
    assert len(cards) == 2, "the two cards with no contact behind them are deleted"
    assert [c.table.cell(0, 0).text for c in cards] == [
        "A. Banker, Managing Principal",
        "B. Banker, Vice President",
    ]
    assert [c.table.cell(2, 0).text for c in cards] == [
        "abanker@example.com",
        "bbanker@example.com",
    ]
    # The library's own card was replaced, not kept alongside.
    assert "nselfe@inforfg.com" not in _all_slides_text(Presentation(deck_path))


# ─── B13c: the ownership slide's takeaway ────────────────────────────────────


def test_ownership_takeaway_is_written_onto_the_slide(tmp_path: Path):
    """Its two siblings were filled; this one had no field at all to fill from."""
    deck_path = _assemble(tmp_path, _sample_content())
    slide = Presentation(deck_path).slides[8]

    assert find_shape(slide, "Text Placeholder 4").text == (
        "Institutions hold a clear majority of the register, with insider "
        "ownership concentrated in long-tenured management."
    )


def test_ownership_takeaway_is_required():
    base = _sample_content().model_dump()
    del base["ownership_takeaway"]
    with pytest.raises(ValidationError):
        PitchDeckContent.model_validate(base)


def test_a_mixed_currency_footnote_converges(tmp_path: Path):
    """The longest note this module can write still fits its box, MEASURED.

    The overview footnote is a 0.45"-tall two-column placeholder whose bottom
    edge is 0.02" from the slide edge, and it already ships two lines. The mixed
    -currency clause is the most text the assembler ever adds to it, so whether
    it fits is a rendering question — and the repo's answer to a rendering
    question is a render, never an estimate. Runs the converge loop, which
    measures the rendered ink and fails the stage if it lands outside the box.
    """
    workbook = _write_sample_cap_table(tmp_path / "wb", currency="CAD")
    _set_fx_rate(workbook, 1.4102)

    deck_path = _assemble(
        tmp_path, _sample_content(), captable_workbook_path=workbook, converge=True
    )

    note = find_shape(Presentation(deck_path).slides[6], "Text Placeholder 1").text
    assert "US$1.00 = C$1.4102" in note


def test_the_logo_placeholder_is_advisory_not_a_blocking_finding(tmp_path: Path):
    """Naming the target inside the placeholder wording must not create review noise.

    `deck_contract` reports any `[… placeholder …]` it does not recognise as a
    BLOCKING `unfilled-placeholder`. The old `[Glean Logo]` matched nothing and
    so was reported as nothing — invisible to the contract while rendering to a
    client. Moving it into the library's placeholder vocabulary makes it visible,
    which is the point, but it has to land in the advisory tier the deck's other
    deliberate placeholders live in: re-flagging expected things is how a review
    gets ignored.
    """
    from deck_contract import SEVERITY_ADVISORY

    base = _sample_content().model_dump()
    base["market_entry_targets"] = [
        {"name": "Kueski", "cells": base["market_entry_targets"][0]["cells"]},
        {"name": "Nubank", "cells": base["market_entry_targets"][1]["cells"]},
    ]
    deck_path = _assemble(tmp_path, PitchDeckContent.model_validate(base))

    findings = verify_deck(deck_path, render=False, vision=False, out_dir=tmp_path / "qa")
    logos = [f for f in findings if "Logo]" in f.detail]
    assert logos, "the logo boxes should be reported at all — they are work to finish"
    assert all(f.kind == "expected-placeholder" for f in logos), [f.kind for f in logos]
    assert all(f.severity == SEVERITY_ADVISORY for f in logos)


def test_the_chart_placeholders_are_named_in_reading_order():
    """`financial_charts._PLACEHOLDER_NAMES` pairs charts to the tab's metric rows.

    The other half of the Financial Summary slide: the label tiles are checked
    above, and these are the chart boxes beneath them. A stale mapping would put
    every chart under the wrong label — four correct charts, four correct labels,
    and a wrong slide.
    """
    from financial_charts import _PLACEHOLDER_NAMES

    slide = Presentation(TEMPLATE).slides[8]
    boxes = [find_shape(slide, name) for name in _PLACEHOLDER_NAMES]
    assert [s.text.strip() for s in boxes] == [
        f"[Placeholder for Metric #{n} Chart]" for n in (1, 2, 3, 4)
    ]
    assert boxes == shapes_in_reading_order(boxes)


def test_the_metric_groups_are_named_in_reading_order():
    """`earnings_update_assembler._METRIC_GROUPS` is zipped against `kpi_rows`.

    Same class as the pitch deck's divider — an ordered list of data paired to
    shapes — resolved by name rather than by document order, which is correct.
    This is the check that keeps the name order honest.
    """
    from earnings_update_assembler import _METRIC_GROUPS

    slide = Presentation(TEMPLATE).slides[7]  # the library's earnings-summary slide
    groups = [find_shape(slide, name) for name, *_ in _METRIC_GROUPS]
    assert groups == shapes_in_reading_order(groups)
