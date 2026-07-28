"""Tests for the template_layout runtime-verification map.

Two halves, per the design contract:

(a) every declared defined name / slide marker resolves against the SHIPPED
    templates (verification is a no-op change for an unmodified plugin), and
(b) a template that lost its ``infor_`` names — or a re-ordered slide library —
    makes the wired-in writers raise ``TemplateLayoutError`` with a message
    naming the template, what was missing, and the remedy.

The shipped templates are never modified — every mutation happens on a copy
under ``tmp_path``.

Note what (b) is and is not. Verification is name-based: it catches a workbook
that no longer carries the names its writers resolve through. It deliberately
does NOT re-check that a name still points at a particular address — that is the
thing the names exist to stop mattering, and Excel moves a name with its cell, so
a re-saved template is a non-event (``test_a_moved_block_needs_no_code_change``).
The sentinel tables that used to cross-check name against address were deleted
once Phase C had shipped; ``test_the_sentinel_tables_are_gone`` keeps them gone.
"""

from __future__ import annotations

from pathlib import Path

import pytest
from openpyxl import Workbook, load_workbook
from openpyxl.workbook.defined_name import DefinedName
from pptx import Presentation

import template_layout as tl
from template_layout import TemplateLayoutError

PLUGIN_ROOT = Path(__file__).resolve().parents[2]
TEMPLATES = PLUGIN_ROOT / "templates"
CAP_TABLE = TEMPLATES / "INFOR Cap Table Template.xlsx"
OWNERSHIP = TEMPLATES / "INFOR Ownership Template.xlsx"
COMPS = TEMPLATES / "INFOR Comps Template.xlsx"
PRECEDENTS = TEMPLATES / "INFOR Precedents Template.xlsx"
LIBRARY = TEMPLATES / "INFOR Slide Library.pptx"


# ─── defined-name helpers ─────────────────────────────────────────────────────


def _drop_name(ws, name: str) -> None:
    """Remove a sheet-scoped defined name, as a re-created template would."""
    if name in ws.defined_names:
        del ws.defined_names[name]


def _repoint(ws, name: str, ref: str) -> None:
    """Re-point a sheet-scoped defined name — what Excel does when rows move."""
    _drop_name(ws, name)
    ws.defined_names.add(DefinedName(name, attr_text=f"'{ws.title}'!{ref}"))


def _named_sheet(name: str, ref: str):
    """A one-sheet workbook carrying a single sheet-scoped defined name."""
    wb = Workbook()
    ws = wb.active
    ws.title = "S"
    _repoint(ws, name, ref)
    return ws


def _deal_workbook(tmp_path: Path) -> Path:
    from deal_workbook import init_deal_workbook

    return init_deal_workbook(
        deal_dir=tmp_path, deliverable_type="pitch", deal_name="Project Test"
    )


def _deal_workbook_without(tmp_path: Path, tab: str, *names: str) -> Path:
    """A deal workbook whose ``tab`` has had ``names`` stripped."""
    deal = _deal_workbook(tmp_path)
    wb = load_workbook(deal)
    for name in names:
        _drop_name(wb[tab], name)
    wb.save(deal)
    return deal


# ─── (a0) the shipped templates carry every declared defined name ────────────
# The contract in both directions: each `infor_` name exists on the sheet the
# registry says, and resolves to the address the registry says.


@pytest.mark.parametrize("template", sorted(tl.TEMPLATE_NAMED_RANGES))
def test_shipped_template_carries_every_declared_name(template: str):
    wb = load_workbook(TEMPLATES / template)
    try:
        for sheet, targets in tl.TEMPLATE_NAMED_RANGES[template].items():
            assert sheet in wb.sheetnames, f"{template} has no sheet {sheet!r}"
            ws = wb[sheet]
            for name, target in sorted(targets.items()):
                assert name.startswith("infor_"), f"{name} is missing the collision prefix"
                assert tl.defined_name_ref(ws, name) == tl.normalize_ref(target), (
                    f"{template} {sheet}!{name}"
                )
    finally:
        wb.close()


def test_the_registry_covers_the_names_phase_c_promised():
    # The plan named eight cells explicitly; this pins that each is reachable.
    # (`precedents_input_ccy` ships as `infor_prec_output_ccy` — the cell is
    # labelled "Output:" and carries the cap table's OUTPUT currency, so the
    # name follows the artefact.)
    cap = tl.TEMPLATE_NAMED_RANGES[tl.CAP_TABLE_TEMPLATE][tl.CAP_TABLE_SOURCE_SHEET]
    assert cap[tl.NAME_FX_RATE] == "F7"
    assert cap[tl.NAME_SHARE_PRICE] == "F16"
    assert cap[tl.NAME_LTM_REVENUE_VALUATION] == "D47"
    assert cap[tl.NAME_LTM_EBITDA_VALUATION] == "D48"
    assert cap[tl.NAME_CAP_PICTURE_RANGE] == "B15:F40"
    assert cap[tl.NAME_BASIC_SHARES] == "F17"
    assert tl.TEMPLATE_NAMED_RANGES[tl.COMPS_TEMPLATE][tl.COMPS_SOURCE_SHEET][
        tl.NAME_COMPS_OUTPUT_CCY
    ] == "F3"
    assert tl.TEMPLATE_NAMED_RANGES[tl.PRECEDENTS_TEMPLATE][tl.PRECEDENTS_SOURCE_SHEET][
        tl.NAME_PREC_OUTPUT_CCY
    ] == "C2"


def test_every_verified_group_only_names_registered_names():
    # The groups the writers pass to `verify_names` must be drawn from the
    # registry: a group naming something the prep tool never stamps would fail
    # every run, and a name in neither place would be verified by nothing.
    registered = {
        name
        for sheets in tl.TEMPLATE_NAMED_RANGES.values()
        for targets in sheets.values()
        for name in targets
    }
    groups = {
        "CAP_TABLE_WRITE_NAMES": tl.CAP_TABLE_WRITE_NAMES,
        "CAP_TABLE_PICTURE_NAMES": tl.CAP_TABLE_PICTURE_NAMES,
        "CAP_TABLE_SECTION_VII_NAMES": tl.CAP_TABLE_SECTION_VII_NAMES,
        "CAP_TABLE_OUTPUT_CCY_NAMES": tl.CAP_TABLE_OUTPUT_CCY_NAMES,
        "OWNERSHIP_INSIDER_WRITE_NAMES": tl.OWNERSHIP_INSIDER_WRITE_NAMES,
        "OWNERSHIP_BBG_HOLDER_NAMES": tl.OWNERSHIP_BBG_HOLDER_NAMES,
        "OWNERSHIP_BBG_LINK_NAMES": tl.OWNERSHIP_BBG_LINK_NAMES,
        "OWNERSHIP_INSIDERS_PICTURE_NAMES": tl.OWNERSHIP_INSIDERS_PICTURE_NAMES,
        "OWNERSHIP_INSTITUTIONS_PICTURE_NAMES": tl.OWNERSHIP_INSTITUTIONS_PICTURE_NAMES,
        "COMPS_WRITE_NAMES": tl.COMPS_WRITE_NAMES,
        "PRECEDENTS_WRITE_NAMES": tl.PRECEDENTS_WRITE_NAMES,
    }
    for group, names in groups.items():
        assert names, f"{group} is empty"
        assert set(names) <= registered, f"{group} names nothing the registry stamps"


def test_new_names_did_not_disturb_the_capiq_and_legacy_namespaces():
    # The reconnaissance rule: the cap table's Capital IQ names identify the
    # workbook to the add-in, and the comps template carries 1,245 legacy
    # artefacts. Phase C adds alongside them; it must never tidy them away.
    cap = load_workbook(CAP_TABLE)
    assert {"CIQWBGuid", "CIQWBInfo", "IQ_LTM", "CAD_USD"} <= set(cap.defined_names)
    assert len(cap.defined_names) == 33  # unchanged; the infor_ names are sheet-scoped
    comps = load_workbook(COMPS)
    assert len(comps.defined_names) == 1245
    assert not [n for n in comps.defined_names if n.startswith("infor_")]


# ─── (a) every verification passes against the shipped templates ─────────────


def test_cap_table_template_passes_every_verified_group():
    ws = load_workbook(CAP_TABLE)[tl.CAP_TABLE_SOURCE_SHEET]
    tl.verify_cap_table_before_write(ws)  # header + LTM + Section VII
    tl.verify_names(ws, tl.CAP_TABLE_PICTURE_NAMES, template=tl.CAP_TABLE_TEMPLATE)
    tl.verify_names(ws, tl.CAP_TABLE_OUTPUT_CCY_NAMES, template=tl.CAP_TABLE_TEMPLATE)


def test_ownership_template_passes_every_verified_group():
    wb = load_workbook(OWNERSHIP)
    ws = wb[tl.OWNERSHIP_SOURCE_SHEET]
    tl.verify_names(
        ws,
        tl.OWNERSHIP_INSIDER_WRITE_NAMES
        + tl.OWNERSHIP_BBG_LINK_NAMES
        + tl.OWNERSHIP_INSIDERS_PICTURE_NAMES
        + tl.OWNERSHIP_INSTITUTIONS_PICTURE_NAMES,
        template=tl.OWNERSHIP_TEMPLATE,
    )
    tl.verify_names(
        wb[tl.OWNERSHIP_BBG_SOURCE_SHEET],
        tl.OWNERSHIP_BBG_HOLDER_NAMES,
        template=tl.OWNERSHIP_TEMPLATE,
    )


def test_comps_template_passes_every_verified_group():
    ws = load_workbook(COMPS)[tl.COMPS_SOURCE_SHEET]
    tl.verify_names(ws, tl.COMPS_WRITE_NAMES, template=tl.COMPS_TEMPLATE)


def test_precedents_template_passes_every_verified_group():
    ws = load_workbook(PRECEDENTS)[tl.PRECEDENTS_SOURCE_SHEET]
    tl.verify_names(ws, tl.PRECEDENTS_WRITE_NAMES, template=tl.PRECEDENTS_TEMPLATE)


def test_slide_library_passes_every_marker():
    prs = Presentation(str(LIBRARY))
    assert len(prs.slides) == tl.SLIDE_LIBRARY_SLIDE_COUNT
    for index in sorted(tl.LIBRARY_SLIDE_MARKERS):
        tl.verify_library_slide(prs, index)


def test_every_library_marker_identifies_exactly_one_slide():
    # The Phase C precondition. `find_slide_by_marker` replaced the assemblers'
    # hardcoded indices, and it is only sound if each marker is unique in the
    # library: two matches and the assembler's choice would be arbitrary.
    prs = Presentation(str(LIBRARY))
    for index, marker in sorted(tl.LIBRARY_SLIDE_MARKERS.items()):
        assert tl.find_slides_by_marker(prs, marker) == [index], marker.description


def test_earnings_assembler_finds_its_five_slides_by_marker():
    # `_KEEP_LIBRARY_INDICES = (0, 6, 7, 15, 16)` is gone; the five entries are
    # discovered. This pins that they still resolve to the historical indices —
    # so the migration is faithful — without the assembler depending on them.
    from earnings_update_assembler import _KEEP_MARKERS

    prs = Presentation(str(LIBRARY))
    found = [tl.find_slide_by_marker(prs, marker) for marker in _KEEP_MARKERS]
    assert found == [0, 6, 7, 15, 16]


def test_pitch_assembler_clone_and_delete_targets_resolve_by_marker():
    prs = Presentation(str(LIBRARY))
    assert tl.find_slide_by_marker(prs, tl.MARKER_EARNINGS_SUMMARY) == 7
    assert tl.find_slide_by_marker(prs, tl.MARKER_FINANCIAL_SUMMARY) == 8
    assert tl.find_slide_by_marker(prs, tl.MARKER_MARKET_ENTRY) == 14


def test_built_overview_marker_finds_the_overview_slide_in_an_assembled_deck():
    # `OVERVIEW_SLIDE_INDEX = 6` is gone. In a BUILT deck the library's overview
    # marker no longer matches — its title has been filled with the client name
    # — so the pie placeholder is the marker, which is also the shape
    # `financial_charts` is about to replace.
    #
    # Measured on the earnings-update fixture, whose overview slide sits at
    # index 1: the point is that the marker finds it wherever it is. The pitch
    # fixture is deliberately NOT used — it is a finished deck, so
    # `financial-charts` has already swapped the placeholder for a picture,
    # which is the same "present until this stage replaces it" contract the FS
    # self-discovery has always had.
    deck = Path(__file__).parent / "fixtures" / "earnings-update-deck.pptx"
    prs = Presentation(str(deck))
    assert tl.find_slides_by_marker(prs, tl.MARKER_OVERVIEW) == []  # title was filled
    assert tl.find_slide_by_marker(prs, tl.MARKER_BUILT_OVERVIEW) == 1


def test_the_pitch_assembler_always_leaves_the_pie_placeholder_for_the_marker():
    # `financial_charts` locates the overview slide by the pie placeholder, so
    # the assembler must always leave one. It does: `_verify_pitch_output`
    # fails the stage if the placeholder went missing.
    import pitch_deck_assembler

    source = Path(pitch_deck_assembler.__file__).read_text(encoding="utf-8")
    assert tl.MARKER_BUILT_OVERVIEW.contains in source


def test_find_slide_by_marker_rejects_zero_and_multiple_matches():
    prs = Presentation(str(LIBRARY))
    absent = tl.SlideMarker("Title 1", "no such text anywhere", "nonexistent")
    with pytest.raises(tl.TemplateLayoutError, match="no slide carries"):
        tl.find_slide_by_marker(prs, absent)
    assert tl.find_optional_slide_by_marker(prs, absent) is None

    # 'Title 1' with an empty fragment matches every slide that has the shape.
    ambiguous = tl.SlideMarker("Title 1", "", "ambiguous")
    with pytest.raises(tl.TemplateLayoutError, match="must identify exactly one"):
        tl.find_slide_by_marker(prs, ambiguous)
    with pytest.raises(tl.TemplateLayoutError, match="must identify exactly one"):
        tl.find_optional_slide_by_marker(prs, ambiguous)


def test_fs_marker_matches_financial_charts_self_discovery():
    # The FS slide marker must stay in lock-step with the Rectangle 17 pattern
    # financial_charts uses to self-discover FS slides in the assembled deck.
    import financial_charts

    assert tl.MARKER_FINANCIAL_SUMMARY.shape_name == financial_charts._FS_MARKER_PLACEHOLDER


# ─── verify helper semantics ──────────────────────────────────────────────────


def test_verify_names_passes_when_every_name_resolves():
    ws = _named_sheet("infor_thing", "$F$1")
    tl.verify_names(ws, ("infor_thing",), template="T.xlsx")
    assert tl.resolve_name_cell(ws, "infor_thing") == "F1"


def test_verify_names_reports_every_missing_name_at_once():
    # A writer that resolved names one at a time would half-fill the tab and
    # name only the first casualty; the pre-flight lists all of them.
    ws = _named_sheet("infor_present", "$F$1")
    with pytest.raises(TemplateLayoutError) as exc:
        tl.verify_names(
            ws,
            ("infor_present", "infor_absent_one", "infor_absent_two"),
            template="Some Template.xlsx",
        )
    msg = str(exc.value)
    assert "Some Template.xlsx" in msg
    assert "infor_absent_one" in msg and "infor_absent_two" in msg
    assert "infor_present" not in msg  # the resolving name is not reported
    assert "add_template_named_ranges.py" in msg  # the remedy


def test_verify_names_rejects_a_name_pointing_at_another_sheet():
    # A name that survived a sheet copy but still points at the source is not
    # usable here, and must read as missing rather than resolve elsewhere.
    wb = Workbook()
    ws = wb.active
    ws.title = "S"
    other = wb.create_sheet("Other")
    other.defined_names.add(DefinedName("infor_thing", attr_text="'Other'!$F$1"))
    with pytest.raises(TemplateLayoutError, match="infor_thing"):
        tl.verify_names(ws, ("infor_thing",), template="T.xlsx")


def test_verify_workbook_names_names_a_missing_sheet(tmp_path: Path):
    path = tmp_path / "wrong.xlsx"
    Workbook().save(path)
    with pytest.raises(TemplateLayoutError, match="expected sheet 'Cap with Links'"):
        tl.verify_workbook_names(
            path, sheet=tl.CAP_TABLE_SOURCE_SHEET, names=tl.CAP_TABLE_PICTURE_NAMES
        )


def test_verify_workbook_names_passes_on_the_shipped_deal_workbook(tmp_path: Path):
    # The pre-flight both assemblers run, against a workbook produced the way a
    # real run produces one.
    from deal_workbook import TAB_OWNERSHIP

    deal = _deal_workbook(tmp_path)
    tl.verify_workbook_names(
        deal, sheet=TAB_OWNERSHIP, names=tl.OWNERSHIP_INSIDERS_PICTURE_NAMES
    )


def test_resolve_helpers_fail_loudly_and_distinguish_cell_from_range():
    ws = _named_sheet("infor_block", "$B$2:$D$4")
    assert tl.resolve_name_range(ws, "infor_block") == "B2:D4"
    with pytest.raises(TemplateLayoutError, match="resolves to the range 'B2:D4'"):
        tl.resolve_name_cell(ws, "infor_block")
    with pytest.raises(TemplateLayoutError, match="defined name 'infor_absent' is missing"):
        tl.resolve_name_range(ws, "infor_absent")


def test_resolve_workbook_range_falls_back_when_the_name_is_absent(tmp_path: Path):
    # The fallback is the shipped address. It is unreachable after a
    # `verify_workbook_names` pre-flight — which requires the name — so this
    # pins the behaviour for a caller that skipped one, not a supported path.
    path = tmp_path / "nameless.xlsx"
    wb = Workbook()
    wb.active.title = tl.CAP_TABLE_SOURCE_SHEET
    wb.save(path)
    assert tl.resolve_workbook_range(
        path, sheet=tl.CAP_TABLE_SOURCE_SHEET, name=tl.NAME_CAP_PICTURE_RANGE, fallback="B15:F40"
    ) == "B15:F40"
    # ...and the shipped template resolves through the name instead.
    assert tl.resolve_workbook_range(
        CAP_TABLE, sheet=tl.CAP_TABLE_SOURCE_SHEET, name=tl.NAME_CAP_PICTURE_RANGE, fallback="ZZ1"
    ) == "B15:F40"


def test_the_sentinel_tables_are_gone():
    """Drift lock: the name↔address cross-check must not come back.

    Sentinels paired every protected address with the caption text beside it, as
    a Phase C cross-check that the names had been mapped to the right cells.
    That shipped (v0.5.40) and the tables were deleted (v0.5.42). Reintroducing
    one would re-couple the writers to an address, which is the coupling the
    names exist to remove — and two tables that can disagree is its own failure
    mode.
    """
    source = Path(tl.__file__).read_text(encoding="utf-8")
    for gone in ("CellAnchor", "label_addr", "verify_anchors", "require_names"):
        assert gone not in source, f"{gone} is back in template_layout.py"
    assert not [n for n in dir(tl) if n.endswith("_ANCHORS") or n.endswith("_ANCHOR")]


# ─── (b) a template that lost its names raises through the wired-in writers ───


def test_cap_table_without_its_share_inputs_name_fails_the_section_vii_read(tmp_path: Path):
    # read_basic_shares_from_cap_table must raise rather than fall back to a
    # hardcoded F168:F185 window that may no longer be Section VII.
    from deal_workbook import TAB_CAPTABLE
    from ownership_workbook import read_basic_shares_from_cap_table

    stripped = _deal_workbook_without(tmp_path, TAB_CAPTABLE, tl.NAME_CAP_SHARE_INPUTS)
    with pytest.raises(TemplateLayoutError, match=tl.NAME_CAP_SHARE_INPUTS):
        read_basic_shares_from_cap_table(stripped)


def test_ownership_tab_without_its_names_fails_the_builder(tmp_path: Path):
    from deal_workbook import TAB_OWNERSHIP
    from ownership_workbook import InsiderHolding, build_ownership_workbook

    stripped = _deal_workbook_without(
        tmp_path, TAB_OWNERSHIP, tl.NAME_OWN_INSIDER_BLOCK, tl.NAME_OWN_TOTAL_SHARES
    )
    with pytest.raises(TemplateLayoutError) as exc:
        build_ownership_workbook(
            insiders=[InsiderHolding("Doe, Jane", "Jane Doe (CFO)", 500)],
            total_shares_outstanding=1_000_000,
            deal_workbook=stripped,
        )
    msg = str(exc.value)
    # Both missing names are reported, not just the first.
    assert tl.NAME_OWN_INSIDER_BLOCK in msg and tl.NAME_OWN_TOTAL_SHARES in msg
    # Two layers check this, and the OUTER one wins: `deal_workbook.write_tab`
    # verifies each `TabSpec`'s names before handing the sheet to the producer,
    # so the message carries the deal-workbook remedy (re-run
    # build_deal_workbook_template.py) rather than the source-template one. The
    # producer's own `verify_names` is what covers a direct invocation.
    assert "deal workbook" in msg
    assert "build_deal_workbook_template.py" in msg


def test_comps_tab_without_its_names_fails_the_builder(tmp_path: Path):
    from comps_workbook import build_comps_workbook
    from deal_workbook import TAB_COMPS

    stripped = _deal_workbook_without(
        tmp_path, TAB_COMPS, tl.NAME_COMPS_GROUP_BLOCKS[0]
    )
    with pytest.raises(TemplateLayoutError, match=tl.NAME_COMPS_GROUP_BLOCKS[0]):
        build_comps_workbook(
            verticals=[{"name": "Vertical A", "companies": [{"ticker": "TSX:RY"}]}],
            deal_workbook=stripped,
        )


def _one_precedent_group(name: str = "Group A") -> list[dict]:
    return [
        {
            "name": name,
            "transactions": [
                {
                    "input_currency": "USD",
                    "announce_date": "2025-01-01",
                    "target": "Example Target Inc.",
                    "acquiror": "Example Acquiror Inc.",
                    "tev": 999.9,
                    "hq_country": "USA",
                    "revenue_ltm": 888.8,
                }
            ],
        }
    ]


def test_precedents_tab_without_its_names_fails_the_builder(tmp_path: Path):
    from deal_workbook import TAB_PRECEDENTS
    from precedents_workbook import build_precedents_workbook

    stripped = _deal_workbook_without(
        tmp_path, TAB_PRECEDENTS, tl.NAME_PREC_GROUP_LABELS[0]
    )
    with pytest.raises(TemplateLayoutError, match=tl.NAME_PREC_GROUP_LABELS[0]):
        build_precedents_workbook(deal_workbook=stripped, groups=_one_precedent_group())


def test_a_moved_block_needs_no_code_change(tmp_path: Path):
    """The payoff, on the Excel path that actually happens.

    An analyst who inserts a row in Excel gets the defined names moved for them —
    that is what a name is. So the writer must follow the name to wherever it now
    points, with no code change and no complaint. Here group 1's label and block
    are re-pointed one row down (``E7`` -> ``E8``, ``B8:AI13`` -> ``B9:AI14``) and
    the transaction lands on row 9 instead of row 8.

    This is the case the deleted sentinel tables got wrong: they pinned the old
    address, so a correctly re-saved template raised.
    """
    from deal_workbook import TAB_PRECEDENTS
    from precedents_workbook import build_precedents_workbook

    deal = _deal_workbook(tmp_path)
    wb = load_workbook(deal)
    ws = wb[TAB_PRECEDENTS]
    ws.insert_rows(8)
    _repoint(ws, tl.NAME_PREC_GROUP_LABELS[0], "$E$8")
    _repoint(ws, tl.NAME_PREC_GROUP_BLOCKS[0], "$B$9:$AI$14")
    wb.save(deal)

    build_precedents_workbook(deal_workbook=deal, groups=_one_precedent_group("Moved Group"))

    check = load_workbook(deal)[TAB_PRECEDENTS]
    assert check["E8"].value == "Moved Group"
    assert check["F9"].value == "Example Target Inc."
    assert check["F8"].value is None  # the shipped row was left alone


def _reordered_library_copy(dest: Path, from_index: int) -> Path:
    """Save a copy of the shipped library with one slide moved to the end."""
    prs = Presentation(str(LIBRARY))
    sld_id_lst = prs.slides._sldIdLst
    element = list(sld_id_lst)[from_index]
    sld_id_lst.remove(element)
    sld_id_lst.append(element)
    prs.save(str(dest))
    return dest


def test_reordered_library_fails_marker_verification(tmp_path: Path):
    # Move the earnings-summary slide (raw index 7) to the end: index 7 is now
    # the Financial Summary slide and every later index shifts by one.
    reordered = _reordered_library_copy(tmp_path / "library.pptx", 7)
    prs = Presentation(str(reordered))
    with pytest.raises(TemplateLayoutError) as exc:
        tl.verify_library_slide(prs, 7)
    msg = str(exc.value)
    assert "slide index 7" in msg
    assert "earnings summary" in msg


def _earnings_inputs(tmp_path: Path):
    """Minimal real SlidePlan + content artefacts for an end-to-end assembly."""
    from earnings_update_wireframe import build_earnings_update_slide_plan, write_slide_plan
    from schemas import Company
    from tests.test_earnings_update_assembler import _sample_content

    slide_plan = build_earnings_update_slide_plan(
        company=Company(legal_name="SampleCo", ticker="TSX:SMPL"),
        reporting_quarter="Q4 2025",
        comparison_quarter="Q4 2024",
    )
    slide_plan_path = write_slide_plan(slide_plan, tmp_path / "slide_plan.json")
    content_path = tmp_path / "content.json"
    content_path.write_text(_sample_content().model_dump_json(indent=2), encoding="utf-8")
    return slide_plan_path, content_path


def test_reordered_library_still_assembles_the_earnings_deck(tmp_path: Path):
    # The Phase C payoff, end to end. Through v0.5.39 the earnings assembler
    # kept library indices (0, 6, 7, 15, 16) and this same re-order raised
    # TemplateLayoutError — correct, but it meant every library edit was a code
    # migration. The five entries are now located by marker, so moving one is a
    # non-event: the deck still assembles, with the right five slides.
    from earnings_update_assembler import assemble_earnings_update_deck

    reordered = _reordered_library_copy(tmp_path / "library.pptx", 6)  # overview -> end
    slide_plan_path, content_path = _earnings_inputs(tmp_path)

    out = assemble_earnings_update_deck(
        slide_plan_path=slide_plan_path,
        content_path=content_path,
        template_path=reordered,
        output_dir=tmp_path,
        converge=False,  # geometry is deck_repair's business, not this test's
    )
    prs = Presentation(str(out))
    assert len(prs.slides) == 5
    # The overview slide moved to the end of the library, so it is now the last
    # of the five kept entries rather than the second — and it is still filled.
    texts = ["\n".join(s.text for s in slide.shapes if s.has_text_frame) for slide in prs.slides]
    assert any("Introduction to SampleCo" in t for t in texts)
    assert any("Earnings Summary" in t for t in texts)
    assert any("Disclaimer" in t for t in texts)


def test_inserting_a_library_slide_needs_no_code_change(tmp_path: Path):
    """The Phase C payoff, on the edit that historically forced a migration.

    Inserting a slide ahead of the disclaimer/contact closers is precisely what
    v0.5.8 and v0.5.14 did, and each time the earnings assembler shipped the
    wrong slides until someone hand-bumped `_KEEP_LIBRARY_INDICES` — once
    dropping the Contact slide entirely from client decks. Here the closers move
    15/16 -> 16/17 and both assemblers simply follow.
    """
    from earnings_update_assembler import _KEEP_MARKERS, assemble_earnings_update_deck
    from pitch_deck_assembler import _PitchLayout
    from pptx_helpers import clone_slide_after, delete_slide

    shipped = Presentation(str(LIBRARY))
    before = {
        "disclaimer": tl.find_slide_by_marker(shipped, tl.MARKER_DISCLAIMER),
        "contact": tl.find_slide_by_marker(shipped, tl.MARKER_CONTACT),
        "precedents": tl.find_slide_by_marker(shipped, tl.MARKER_PRECEDENTS),
    }

    # Insert one slide just after Comparable Companies, and retitle it so it is a
    # distinct concept rather than a second comps slide (which would make the
    # comps marker ambiguous — itself a caught error, see the finder tests).
    grown = tmp_path / "library.pptx"
    prs = Presentation(str(LIBRARY))
    comps_at = tl.find_slide_by_marker(prs, tl.MARKER_COMPS)
    clone_slide_after(prs, comps_at)
    inserted = comps_at + 1
    for shape in prs.slides[inserted].shapes:
        if shape.name == tl.MARKER_COMPS.shape_name and shape.has_text_frame:
            shape.text_frame.paragraphs[0].runs[0].text = "Trading History Analysis"
    prs.save(str(grown))

    prs = Presentation(str(grown))
    assert len(prs.slides) == len(shipped.slides) + 1
    # Everything after the insertion shifted by one.
    assert tl.find_slide_by_marker(prs, tl.MARKER_DISCLAIMER) == before["disclaimer"] + 1
    assert tl.find_slide_by_marker(prs, tl.MARKER_CONTACT) == before["contact"] + 1
    assert tl.find_slide_by_marker(prs, tl.MARKER_PRECEDENTS) == before["precedents"] + 1

    # The earnings assembler follows the shift and still builds the right deck.
    found = [tl.find_slide_by_marker(prs, marker) for marker in _KEEP_MARKERS]
    assert found == [0, 6, 7, before["disclaimer"] + 1, before["contact"] + 1]

    slide_plan_path, content_path = _earnings_inputs(tmp_path)
    out = assemble_earnings_update_deck(
        slide_plan_path=slide_plan_path,
        content_path=content_path,
        template_path=grown,
        output_dir=tmp_path,
        converge=False,
    )
    built = Presentation(str(out))
    assert len(built.slides) == 5
    texts = ["\n".join(s.text for s in slide.shapes if s.has_text_frame) for slide in built.slides]
    assert any("Introduction to SampleCo" in t for t in texts)
    assert any("Disclaimer" in t for t in texts), "the closers must still be the last two"
    assert any("Contact" in t for t in texts)

    # The pitch layout discovers its indices over the grown library too.
    pitch = Presentation(str(grown))
    delete_slide(pitch, tl.find_slide_by_marker(pitch, tl.MARKER_EARNINGS_SUMMARY))
    layout = _PitchLayout(pitch, template_name=grown.name, expected_total=len(pitch.slides))
    assert layout.overview == 6
    assert layout.precedents == before["precedents"]  # 11 -> 12 raw, -1 for the earnings delete
    assert layout.investment_highlights == layout.precedents + 1
    assert layout.market_entry == [layout.precedents + 2]


def test_library_missing_a_kept_entry_fails_the_earnings_assembler(tmp_path: Path):
    # The other direction: discovery must still be a check. A library that lost
    # one of the five entries halts the run rather than shipping a short deck.
    from earnings_update_assembler import assemble_earnings_update_deck
    from pptx_helpers import delete_slide

    prs = Presentation(str(LIBRARY))
    delete_slide(prs, tl.find_slide_by_marker(prs, tl.MARKER_CONTACT))
    broken = tmp_path / "library.pptx"
    prs.save(str(broken))

    slide_plan_path, content_path = _earnings_inputs(tmp_path)
    with pytest.raises(TemplateLayoutError, match="contact"):
        assemble_earnings_update_deck(
            slide_plan_path=slide_plan_path,
            content_path=content_path,
            template_path=broken,
            output_dir=tmp_path,
            converge=False,
        )
