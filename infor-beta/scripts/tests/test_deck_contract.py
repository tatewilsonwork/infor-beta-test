"""Phase B step 2: the contract, wired as a regression test.

Three historical bugs anchor this file. Each is replayed from a **real artefact**
— the analyst's actual run output, stripped to the one defective slide, with the
slide's shape XML preserved byte-for-byte (see `fixtures/regressions/README.md`).
None is synthesized, and each replay asserts the measurement, not just that
*something* was reported:

  1. PRL17's market-entry table rendering 5.91" against a 5.71" declaration —
     `rendered-overflow`. Invisible in the XML, which declares 5.710".
  2. PRL18's Considerations/Mitigants table declaring 5.360" against the
     library's 5.1715" — `table-taller-than-library`. Visible in the XML.
  3. The v0.5.23 overview bullets running through the LTM Revenue Breakdown band
     — `rendered-overflow`, naming the band it runs into.

The split between 1 and 2 is the point of the two checks, and a contract with
only the XML tier would catch just one of the three.

A fourth defect anchors the **attribution** tier, and it is not historical — it is
live in the earnings-update fixture. Its broker table renders under the summary
box that sits 0.037" below it, so counting unclaimed ink can only see 0.037" of an
overflow that is really 0.15". `masked-overflow` measures it by rendering the
table alone.

**The vision tier is deliberately not asserted for verdicts.** It carries no OCR;
it renders slides, extracts picture crops, and returns an agenda for the
checkpoint agent. What is asserted here is that the agenda points at the right
things and that nothing in it can ever block — never what a reviewer concluded.
"""

from __future__ import annotations

import collections
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from deck_contract import (
    SEVERITY_ADVISORY,
    SEVERITY_BLOCKING,
    default_library_path,
    match_library_slide,
    library_baseline,
    measure_attributed_overflow,
    verify_deck,
    vision_pass,
)
from excel_to_powerpoint import find_soffice

_FIXTURES = Path(__file__).parent / "fixtures"
_REGRESSIONS = _FIXTURES / "regressions"
_PITCH = _FIXTURES / "pitch-deck.pptx"
_EARNINGS = _FIXTURES / "earnings-update-deck.pptx"

_PRL17 = _REGRESSIONS / "prl17-market-entry-table.pptx"
_PRL18 = _REGRESSIONS / "prl18-risk-table.pptx"
_PRL14 = _REGRESSIONS / "prl14-overview-bullets.pptx"

needs_render = pytest.mark.skipif(
    find_soffice() is None,
    reason="the render-measured tier needs LibreOffice (see slide_render)",
)


def _kinds(findings) -> collections.Counter:
    return collections.Counter(f.kind for f in findings)


def _of_kind(findings, kind: str) -> list:
    return [f for f in findings if f.kind == kind]


# ─── Historical replay 1: PRL17's market-entry table ─────────────────────────


@needs_render
def test_catches_prl17_market_entry_table_rendering_past_its_declaration(tmp_path):
    """The bug the XML cannot show: declared 5.710", rendered 5.91".

    The assembler clamps this table to `_ME_TABLE_HEIGHT` = 5.710", *under* the
    library's 5.7197", so `table-taller-than-library` is silent by construction.
    A stored row height is only a render-time minimum, so one wrapped row label
    ('Geographic Footprint') re-grew its row and the table with it.
    """
    findings = verify_deck(_PRL17, vision=False, out_dir=tmp_path)
    overflow = _of_kind(findings, "rendered-overflow")

    table = [f for f in overflow if f.shape and f.shape.startswith("Table")]
    assert table, f"no rendered-overflow on the market-entry table; got {_kinds(findings)}"
    assert table[0].severity == SEVERITY_BLOCKING
    assert table[0].measured_in > 0.15, (
        f"overflow measured {table[0].measured_in:.3f}\" — expected the ~0.18\" that "
        f"PowerPoint reported as a 5.91\" table"
    )

    # The XML tier must NOT be what catches this one.
    assert not _of_kind(findings, "table-taller-than-library"), (
        "the declared height is under the library's, so this bug is only visible "
        "on the render — a table-height finding here means the fixture changed"
    )


# ─── Historical replay 2: PRL18's risk table ─────────────────────────────────


def test_catches_prl18_risk_table_taller_than_the_library(tmp_path):
    """Declared 5.360" against the library's 5.1715" — visible without rendering."""
    findings = verify_deck(_PRL18, render=False, out_dir=tmp_path)
    taller = _of_kind(findings, "table-taller-than-library")

    assert taller, f"no table-height finding; got {_kinds(findings)}"
    assert taller[0].severity == SEVERITY_BLOCKING
    assert taller[0].measured_in == pytest.approx(5.360, abs=0.01)
    assert taller[0].limit_in == pytest.approx(5.1715, abs=0.01)
    assert "considerations" in taller[0].detail.lower()


# ─── Historical replay 3: the v0.5.23 overview overflow ──────────────────────


@needs_render
def test_catches_overview_bullets_running_into_the_ltm_band(tmp_path):
    """Pre-v0.5.23 overview bullets: no autofit, a 0.58" box, 1,126 characters.

    The box declares a bottom of 2.05" and the LTM Revenue Breakdown band header
    sits at 4.747", so the copy renders straight through it. The finding must name
    the band — 'something overflowed' is not actionable, 'it is in the LTM band' is.
    """
    findings = verify_deck(_PRL14, vision=False, out_dir=tmp_path)
    overflow = [f for f in _of_kind(findings, "rendered-overflow") if f.shape == "TextBox 9"]

    assert overflow, f"no rendered-overflow on the overview bullets; got {_kinds(findings)}"
    assert overflow[0].severity == SEVERITY_BLOCKING
    assert overflow[0].measured_in > 2.0, (
        f"overflow measured {overflow[0].measured_in:.3f}\" — the bullets run ~2.6\" "
        f"past a box that declares 0.58\" of height"
    )
    # 'Rectangle 15' is the LTM Revenue Breakdown band on the library overview slide.
    assert "Rectangle 15" in overflow[0].detail, (
        f"the finding does not name the band it overruns: {overflow[0].detail}"
    )


# ─── Per-shape render attribution ────────────────────────────────────────────


@needs_render
def test_catches_the_broker_table_masked_by_the_summary_box(tmp_path):
    """The blind spot `rendered-overflow` structurally cannot see.

    The earnings-update fixture's broker table declares a bottom of 6.184" and the
    summary box starts at 6.221", so the table's re-grown last row (`EPS, Adj.`)
    renders underneath the summary box's text. Only 0.037" of that lands outside
    every declared box — under tolerance — so counting unclaimed ink says nothing.

    Rendering the table alone measures its own ink at ~0.15". Before this check the
    slide reached a reviewer only indirectly, via the unrelated `TextBox 6`
    finding on the same slide.
    """
    findings = verify_deck(_EARNINGS, vision=False, out_dir=tmp_path)
    masked = [f for f in _of_kind(findings, "masked-overflow") if f.shape == "Table 4"]

    assert masked, f"the broker table's masked overflow was not caught; got {_kinds(findings)}"
    assert masked[0].slide == 2  # zero-based; slide 3 in PowerPoint
    assert masked[0].severity == SEVERITY_BLOCKING
    assert masked[0].measured_in == pytest.approx(0.153, abs=0.03), (
        f"attributed overflow measured {masked[0].measured_in:.3f}\" — expected the "
        f"~0.15\" the re-grown 'EPS, Adj.' row renders past 6.184\""
    )
    # Naming the masking shape is the actionable half: a repair step has to know
    # which of the two shapes to shrink.
    assert "Rectangle 1111" in masked[0].detail, (
        f"the finding does not name the shape it renders into: {masked[0].detail}"
    )

    # And it is NOT the unclaimed-ink check that found it.
    assert not [f for f in _of_kind(findings, "rendered-overflow") if f.shape == "Table 4"], (
        "unclaimed ink cannot see this overflow — a rendered-overflow finding on "
        "the broker table means the masking geometry changed"
    )


@needs_render
def test_attribution_does_not_double_report_an_already_visible_overflow(tmp_path):
    """One finding per shape, so the repair loop cannot double-count.

    PRL17's market-entry table is both masked from below (`Text Placeholder 3`
    sits 0.118" under it) and overflowing into unclaimed space, so both checks
    measure it. `rendered-overflow` wins and `masked-overflow` stays quiet.
    """
    findings = verify_deck(_PRL17, vision=False, out_dir=tmp_path)
    tables = [
        f.kind
        for f in findings
        if f.shape == "Table 1215" and f.kind in {"rendered-overflow", "masked-overflow"}
    ]
    assert tables == ["rendered-overflow"], f"expected exactly one overflow finding, got {tables}"


@needs_render
def test_attribution_measures_only_the_probed_shape(tmp_path):
    """The probe subtracts the layout, so a shape's own ink is all that is measured.

    Asserted against the library rather than a built deck: every masked candidate
    there is an unfilled placeholder sitting on a decorated layout, so anything
    beyond render noise (one pixel row = 0.007" at 150 dpi) would mean layout or
    neighbour ink is leaking into the measurement.
    """
    attributed = measure_attributed_overflow(default_library_path(), tmp_path)
    assert attributed, "the library has masked candidates, so the pass must measure something"
    worst = max(depth for shapes in attributed.values() for depth in shapes.values())
    assert worst < 0.02, (
        f"library attribution measured {worst:.3f}\" of own-ink overflow on a blank "
        f"placeholder — the layout subtraction is leaking ink into the measurement"
    )


def test_attribute_false_skips_the_attribution_pass(tmp_path):
    """The extra conversion is opt-out-able without losing the ordinary render tier."""
    findings = verify_deck(_EARNINGS, vision=False, attribute=False, out_dir=tmp_path)
    assert not _of_kind(findings, "masked-overflow")
    assert _of_kind(findings, "rendered-overflow"), "the ordinary render tier must still run"


# ─── The library is the geometric baseline ───────────────────────────────────


@needs_render
def test_blank_library_reports_no_geometric_findings(tmp_path):
    """The reference cannot violate itself.

    This is what keeps the geometric tier usable: the library's own footnote
    placeholders render past their boxes and its tombstone slides park shapes
    off-canvas, and none of that is a defect. If this test fails, the baseline
    subtraction has broken and every deck is about to report template noise.
    """
    library = default_library_path()
    assert library is not None, "the shipped slide library must be resolvable"

    findings = verify_deck(library, vision=False, out_dir=tmp_path)
    geometric = [
        f
        for f in findings
        if f.kind in {"rendered-overflow", "shape-outside-slide", "table-taller-than-library"}
    ]
    assert geometric == [], f"the library violates its own geometry: {[str(f) for f in geometric]}"


@needs_render
def test_library_baseline_is_non_trivial():
    """Guards against the baseline silently becoming all-zeros.

    An empty baseline would still pass the test above while quietly restoring the
    template noise it exists to cancel, so assert it actually measured something.
    """
    baseline = library_baseline(default_library_path(), render=True)
    assert baseline.rendered
    assert baseline.table_height, "no library table heights were measured"
    assert any(baseline.overflow.values()), (
        "the library baseline recorded no overflow anywhere — the reference is empty, "
        "so template-inherent overflow is no longer being cancelled"
    )


def test_match_library_slide_identifies_built_slides():
    """Built slides match their library origin despite filled titles and clones."""
    library = [
        {s.name for s in slide.shapes} for slide in Presentation(default_library_path()).slides
    ]
    pitch = Presentation(_PITCH)

    # Overview: library index 6, title filled to 'Introduction to Propel Holdings Inc.'
    assert match_library_slide(pitch.slides[6], library) == 6
    # Every market-entry clone resolves to the one library entry it was cloned from.
    assert match_library_slide(pitch.slides[13], library) == 14


# ─── String tier ─────────────────────────────────────────────────────────────


def test_flags_the_three_unsubstituted_currency_tokens(tmp_path):
    """The Phase A `[x]` defect: a client deck shipped 'All figures in [x]$MM'.

    Slide 9 (ownership) carries a bare `[x]`; slides 11 and 12 (comps, precedents)
    carry the `[x]$MM` footnote form. `fill_footnote_currency` exists and both
    assemblers call it, but not on these shapes.
    """
    findings = verify_deck(_PITCH, render=False, out_dir=tmp_path)
    currency = _of_kind(findings, "unsubstituted-currency-token")

    assert {f.slide for f in currency} == {10, 11}, (
        f"expected the comps/precedents footnotes (0-based slides 10, 11); got "
        f"{[(f.slide, f.shape) for f in currency]}"
    )
    assert all(f.severity == SEVERITY_BLOCKING for f in currency)

    bare = _of_kind(findings, "unfilled-token")
    assert any(f.slide == 8 and f.shape == "Text Placeholder 4" for f in bare), (
        "the ownership slide's bare '[x]' was not flagged"
    )


def test_flags_the_contact_slide_unfilled_tokens(tmp_path):
    """Nine `[x]` cells ship on the client-facing contact slide.

    The library's contact slide carries three spare banker blocks (name / phone /
    email) that no assembler fills, so they render as a literal '[x]' in a navy
    header and two body rows apiece. Present in both frozen fixtures.

    If this is ever judged intentional it belongs in `EXPECTED_PLACEHOLDERS`, not
    in a lowered severity — the point of the check is that nothing renders '[x]'
    to a client without someone having decided so.
    """
    for deck, contact_slide in ((_PITCH, 17), (_EARNINGS, 4)):
        findings = verify_deck(deck, render=False)
        tokens = [f for f in _of_kind(findings, "unfilled-token") if f.slide == contact_slide]
        assert len(tokens) == 9, (
            f"{deck.name}: expected 9 '[x]' cells on the contact slide, got "
            f"{[f.shape for f in tokens]}"
        )
        assert all(f.severity == SEVERITY_BLOCKING for f in tokens)


def test_known_intentional_placeholders_are_advisory(tmp_path):
    """Comps / precedents / pie placeholders stay visible but never block.

    Their workbook tabs carry un-evaluated CapIQ array formulas that the analyst
    refreshes in Excel, so there is nothing to render at build time.
    """
    pitch = verify_deck(_PITCH, render=False, out_dir=tmp_path)
    expected = _of_kind(pitch, "expected-placeholder")
    assert {f.slide for f in expected} == {10, 11}
    assert all(f.severity == SEVERITY_ADVISORY for f in expected)
    assert not _of_kind(pitch, "unfilled-placeholder")

    earnings = verify_deck(_EARNINGS, render=False)
    pie = _of_kind(earnings, "expected-placeholder")
    assert [(f.slide, f.shape) for f in pie] == [(1, "Rectangle 4")]
    assert pie[0].severity == SEVERITY_ADVISORY


def test_error_token_scan_covers_text_and_spares_rasterised_ranges(tmp_path):
    """`#VALUE!` in a text frame is a defect; inside a pasted range it is not.

    The cap table's forward-estimate columns are CapIQ UDF calls wrapped in
    ``IFERROR(..., "n/a ")``, and the EV/metric rows divide by that text, so an
    un-refreshed CapIQ estimate propagates ``#VALUE!`` into the pasted picture **by
    design** — CapIQ cannot be refreshed in this environment. Both frozen fixtures
    carry that state on their overview slide and must stay clean here.

    A value written into a text frame has no such excuse.
    """
    for deck in (_PITCH, _EARNINGS):
        assert not _of_kind(verify_deck(deck, render=False), "forbidden-string"), (
            f"{deck.name}: the text scan must not reach into rasterised range pictures"
        )

    # A real error value in a text frame is caught.
    prs = Presentation(default_library_path())
    box = prs.slides[0].shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(0.4))
    box.text_frame.text = "EV / Revenue #VALUE!"
    dirty = tmp_path / "error-in-text.pptx"
    prs.save(dirty)

    forbidden = _of_kind(verify_deck(dirty, render=False, out_dir=tmp_path), "forbidden-string")
    assert forbidden, "an error value written into a text frame was not flagged"
    assert forbidden[0].severity == SEVERITY_BLOCKING
    assert "#VALUE!" in forbidden[0].detail


# ─── Vision tier: an agenda, never a verdict ─────────────────────────────────


@needs_render
def test_vision_pass_writes_evidence_and_targets_the_rasterised_pictures(tmp_path):
    """The vision tier's job is to put the right pixels in front of a reviewer.

    No OCR, no verdict. What is checkable is that the evidence exists and the
    agenda points at the content a string scan cannot reach — the overview slide's
    cap-table picture above all, which is a flat PNG with no text layer and no alt
    text.
    """
    result = vision_pass(_PITCH, out_dir=tmp_path)

    assert len(result.review_images) == len(Presentation(_PITCH).slides)
    assert all(p.is_file() and p.stat().st_size > 0 for p in result.review_images.values())
    assert all(p.is_file() for _, _, p in result.picture_crops)

    # The overview cap-table picture must be on the agenda, with a crop to read.
    overview = [t for t in result.targets if t.slide == 6 and t.shape == "Picture 16"]
    assert overview, (
        f"the overview cap-table picture is not on the review agenda; targets on "
        f"slide 7: {[(t.shape, t.crop) for t in result.targets if t.slide == 6]}"
    )
    assert overview[0].crop is not None and overview[0].crop.is_file()
    assert overview[0].render is not None and overview[0].render.is_file()


def test_clearance_hint_is_baseline_relative(tmp_path):
    """Only a gap the *fill* tightened is worth mentioning.

    The INFOR layout seats a header bar directly on its table, so raw clearance
    flags 17 of the pitch fixture's 18 slides and says nothing. Both frozen
    fixtures must therefore be quiet — including the earnings-update broker table,
    whose 0.037" gap above the summary box is the library's own design — while a
    shape the fill moved closer is reported.
    """
    for deck in (_PITCH, _EARNINGS):
        agenda = " ".join(t.question for t in vision_pass(deck, renders={}).targets)
        assert "clearance" not in agenda and "re-grown row" not in agenda, (
            f"{deck.name}: library-inherent clearance is being reported as a hint"
        )

    # Narrow the 0.037" gap the library leaves between that table and the summary
    # box to 0.007" — tightened by more than `_CLEARANCE_TIGHTENED_IN`, but still
    # a gap. (A bigger nudge would make the boxes overlap, which the overlap hint
    # reports instead.)
    prs = Presentation(default_library_path())
    table = next(s for s in prs.slides[7].shapes if getattr(s, "has_table", False))
    table.top += Inches(0.03)
    tightened = tmp_path / "tightened.pptx"
    prs.save(tightened)

    targets = vision_pass(tightened, out_dir=tmp_path, renders={}).targets
    hints = [t for t in targets if "re-grown row" in t.question]
    assert hints, f"a tightened clearance was not flagged; targets: {[t.question[:40] for t in targets]}"
    assert hints[0].shape == table.name


def test_vision_findings_never_block(tmp_path):
    """The tier must not be able to fail a deck unattended.

    `renders={}` skips rendering: the agenda's reasons come from the shape tree,
    so this needs no LibreOffice.
    """
    result = vision_pass(_EARNINGS, out_dir=tmp_path, renders={})
    assert result.findings, "the vision tier produced no agenda at all"
    assert all(f.severity == SEVERITY_ADVISORY for f in result.findings)
    assert all(f.kind in {"vision-review", "render-unavailable"} for f in result.findings)


# ─── Degradation and plumbing ────────────────────────────────────────────────


def test_render_false_skips_the_render_measured_tier(tmp_path):
    """A fast XML-only pass reports no render-measured findings and no vision agenda."""
    findings = verify_deck(_PITCH, render=False, out_dir=tmp_path)
    assert not _of_kind(findings, "rendered-overflow")
    assert not _of_kind(findings, "vision-review")
    # The string tier still ran.
    assert _of_kind(findings, "unsubstituted-currency-token")


def test_missing_library_blocks_rather_than_passing_quietly(tmp_path):
    """No baseline means the geometric tier did not run, and that must be loud."""
    findings = verify_deck(_PITCH, library=tmp_path / "absent.pptx", render=False)
    unavailable = _of_kind(findings, "library-unavailable")
    assert unavailable and unavailable[0].severity == SEVERITY_BLOCKING


def test_missing_deck_raises():
    with pytest.raises(FileNotFoundError):
        verify_deck(_FIXTURES / "does-not-exist.pptx")


def test_findings_are_ordered_blocking_first(tmp_path):
    findings = verify_deck(_PITCH, render=False, out_dir=tmp_path)
    severities = [f.severity for f in findings]
    assert severities == sorted(severities, key=lambda s: s != SEVERITY_BLOCKING)


# ─── Known state of the frozen fixtures ──────────────────────────────────────


@needs_render
@pytest.mark.parametrize(
    "deck, expected",
    [
        (
            "pitch-deck.pptx",
            {
                "unfilled-token": 10,  # ownership bare [x] + 9 contact cells
                "unsubstituted-currency-token": 2,  # comps + precedents footnotes
                "expected-placeholder": 2,  # comps + precedents chart placeholders
                "vision-review": 10,  # slides carrying rasterised pictures
            },
        ),
        (
            "earnings-update-deck.pptx",
            {
                "unfilled-token": 9,  # 9 contact cells
                "rendered-overflow": 1,  # slide 3 Business Updates block
                "masked-overflow": 1,  # slide 3 broker table, under the summary box
                "expected-placeholder": 1,  # pie chart placeholder
                "vision-review": 3,  # cover, overview, contact
            },
        ),
    ],
)
def test_frozen_fixture_findings_are_pinned(deck, expected, tmp_path):
    """Locks the fixtures' known state so a new finding class cannot slip in.

    The fixtures are frozen (sha256 in `fixtures/README.md`), so these counts are
    stable. A change here is either a real regression or a deliberate tolerance
    change — in both cases it should be a conscious edit, not a silent drift.
    """
    findings = verify_deck(_FIXTURES / deck, out_dir=tmp_path)
    assert dict(_kinds(findings)) == expected, "\n".join(str(f) for f in findings)
