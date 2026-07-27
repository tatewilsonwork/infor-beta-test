"""Unit tests for the Financial Summary chart builder + placeholder insertion.

These cover the cross-platform pieces — period-axis detection, the openpyxl
chart formatting, and PNG-into-placeholder insertion. The LibreOffice render
backend is environmental (per repo convention) and is exercised by the skill's
mandatory QA render, not here.

The Excel COM builders' *cleanup ordering* and error normalization ARE covered,
via the fake COM stack in ``tests/fake_com.py`` — no Excel is spawned, which is
the point: spawning it is what leaks. Their render *fidelity* remains
environmental.
"""

import zipfile
from pathlib import Path

import pytest

import financial_charts
from openpyxl import load_workbook
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from financial_summary_workbook import MetricSeries, build_financial_summary_workbook
from ltm_metrics import RevenueSegment, build_ltm_metrics_workbook
from pptx_helpers import INFOR_ACCENTS
from financial_charts import (
    _make_openpyxl_chart,
    _make_openpyxl_pie,
    _make_single_value_chart,
    insert_png_into_placeholder,
    ltm_revenue_overview_range,
    period_axis_columns,
    render_financial_summary_charts_into_deck,
    render_ltm_revenue_pie_into_deck,
)


def _chart_part_count(xlsx_path: Path) -> int:
    """Count the native chart XML parts persisted inside an .xlsx (zip) file."""
    with zipfile.ZipFile(xlsx_path) as z:
        return sum(
            1
            for n in z.namelist()
            if n.startswith("xl/charts/chart") and n.endswith(".xml")
        )


def _raise_runtime(*_args, **_kwargs):
    """A stand-in backend that mimics 'Excel COM unavailable' (RuntimeError)."""
    raise RuntimeError("excel COM unavailable in test")


def _no_libreoffice(monkeypatch):
    """Pretend LibreOffice is absent.

    Patches the shared locator, NOT ``shutil.which``: ``find_soffice`` also probes
    the standard install locations, so a which()-only patch leaves a real Windows
    MSI install visible and the "degrades gracefully" assertions never run.
    The renderers import the locator lazily, so patching the source module's
    attribute reaches them at call time.
    """
    monkeypatch.setattr("excel_to_powerpoint.find_soffice", lambda: None)

# A valid 1x1 transparent PNG — avoids a PIL dependency in the test.
_PNG_1X1 = (
    b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01\x08\x06"
    b"\x00\x00\x00\x1f\x15\xc4\x89\x00\x00\x00\nIDATx\x9cc\x00\x01\x00\x00\x05\x00"
    b"\x01\r\n-\xb4\x00\x00\x00\x00IEND\xaeB`\x82"
)

_FISCAL = ["FY2021", "FY2022", "FY2023", "FY2024", "FY2025"]


def _metrics() -> list[MetricSeries]:
    return [
        MetricSeries("Revenue", "US$MM", [3100.0, 3450.0, 3820.0, 4180.0, 4520.0],
                     result_label="LTM Revenue"),
        MetricSeries("Gross Profit", "US$MM", [1240.0, 1410.0, 1570.0, 1740.0, 1900.0],
                     result_label="LTM Gross Profit"),
        MetricSeries("Adjusted EBITDA", "US$MM", [820.0, 940.0, 1080.0, 1210.0, 1330.0],
                     result_label="LTM Adj. EBITDA"),
        MetricSeries("Net Income", "US$MM", [410.0, 470.0, 540.0, 610.0, 680.0],
                     result_label="LTM Net Income"),
    ]


def _fs_workbook(tmp_path: Path, **overrides) -> Path:
    kwargs = dict(
        company_name="SampleCo",
        currency_note="Figures in US$MM unless noted",
        period_note="FY = fiscal year; LTM = trailing twelve months as of Q3 2026",
        fiscal_labels=_FISCAL,
        metrics=_metrics(),
        output_dir=tmp_path,
    )
    kwargs.update(overrides)
    return build_financial_summary_workbook(**kwargs)


def test_period_axis_columns_with_ltm(tmp_path: Path):
    ws = load_workbook(_fs_workbook(tmp_path)).active
    # Metric | FY1..FY5 | LTM | Units -> period axis B..G.
    assert period_axis_columns(ws) == (2, 7)


def test_period_axis_columns_suppressed_ltm(tmp_path: Path):
    ws = load_workbook(_fs_workbook(tmp_path, show_ltm=False)).active
    # Metric | FY1..FY5 | Units -> period axis B..F.
    assert period_axis_columns(ws) == (2, 6)


def test_make_openpyxl_chart_applies_infor_formatting(tmp_path: Path):
    ws = load_workbook(_fs_workbook(tmp_path)).active
    chart = _make_openpyxl_chart(ws, data_row=6, first_col=2, last_col=7)

    assert chart.type == "col"
    assert chart.grouping == "clustered"
    assert chart.gapWidth == 50
    assert chart.title is None
    assert chart.legend is None
    # One series, filled INFOR blue (openpyxl wraps a hex string in a ColorChoice).
    assert len(chart.series) == 1
    fill = chart.series[0].graphicalProperties.solidFill
    assert getattr(fill, "srgbClr", fill) == "46566E"
    # Data labels on every bar, Outside End, in the SAME currency format as the
    # tab's value cells (so a bar reads "$102.7" exactly like its cell).
    assert chart.dataLabels.showVal is True
    assert chart.dataLabels.dLblPos == "outEnd"
    import financial_summary_workbook

    assert chart.dataLabels.numFmt == '$#,##0.0_);($#,##0.0);"--"'
    assert chart.dataLabels.numFmt == financial_summary_workbook._VALUE_FORMAT
    # P1.2: VALUE only — category/series/legend-key/percent flags off, so the
    # LibreOffice render shows "589.8", not "FY2025; Row 2; 589.808".
    assert chart.dataLabels.showCatName in (False, None)
    assert chart.dataLabels.showSerName in (False, None)
    assert chart.dataLabels.showLegendKey in (False, None)
    assert chart.dataLabels.showPercent in (False, None)
    # Value axis hidden (no line/label), no gridlines; category axis kept.
    assert chart.y_axis.delete is True
    assert chart.y_axis.majorGridlines is None
    assert chart.x_axis.delete is False
    # v0.5.17: no chart-area border, black category-axis line.
    assert chart.graphical_properties.line.noFill is True
    assert chart.x_axis.spPr.line.solidFill.srgbClr == "000000"
    # v0.5.18: the black axis line carries a visible (non-hairline) width so the
    # LibreOffice render shows a baseline beneath the bars (Issue 2).
    assert chart.x_axis.spPr.line.width == financial_charts._AXIS_LINE_WIDTH_EMU
    assert chart.x_axis.spPr.line.width > 0


def test_add_openpyxl_charts_yields_four_charts(tmp_path: Path):
    # Building one chart per metric row produces four independent charts.
    wb = load_workbook(_fs_workbook(tmp_path))
    ws = wb.active
    charts = [_make_openpyxl_chart(ws, data_row=r, first_col=2, last_col=7) for r in (6, 7, 8, 9)]
    assert len(charts) == 4
    assert all(c.gapWidth == 50 and c.y_axis.delete is True for c in charts)


def _deck_with_placeholder(tmp_path: Path) -> Path:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    box = slide.shapes.add_textbox(Inches(0.35), Inches(1.51), Inches(4.53), Inches(2.51))
    box.name = "Rectangle 17"
    box.text_frame.text = "[Placeholder for Metric #1 Chart]"
    path = tmp_path / "deck.pptx"
    prs.save(path)
    return path


def test_insert_png_replaces_placeholder_with_picture(tmp_path: Path):
    deck = _deck_with_placeholder(tmp_path)
    out = insert_png_into_placeholder(
        deck_path=deck,
        slide_index=0,
        placeholder_name="Rectangle 17",
        png_bytes=_PNG_1X1,
        output_path=deck,
    )
    prs = Presentation(out)
    slide = prs.slides[0]
    names = [s.name for s in slide.shapes]
    assert "Rectangle 17" not in names
    pictures = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pictures) == 1
    pic = pictures[0]
    # Picture lands at the placeholder's exact box.
    assert pic.left == Inches(0.35)
    assert pic.top == Inches(1.51)
    assert pic.width == Inches(4.53)
    assert pic.height == Inches(2.51)


def test_insert_missing_placeholder_raises(tmp_path: Path):
    deck = _deck_with_placeholder(tmp_path)
    import pytest

    with pytest.raises(KeyError):
        insert_png_into_placeholder(
            deck_path=deck,
            slide_index=0,
            placeholder_name="Rectangle 999",
            png_bytes=_PNG_1X1,
            output_path=deck,
        )


def test_single_value_chart_has_black_axis_and_no_border(tmp_path: Path):
    # v0.5.17 border/axis fix also applies to the LibreOffice single-value renderer.
    ws = load_workbook(_fs_workbook(tmp_path)).active
    chart = _make_single_value_chart(ws, 6)
    assert chart.graphical_properties.line.noFill is True
    assert chart.x_axis.spPr.line.solidFill.srgbClr == "000000"
    # v0.5.18: visible baseline width in the LibreOffice single-value renderer too.
    assert chart.x_axis.spPr.line.width == financial_charts._AXIS_LINE_WIDTH_EMU
    # P1.2: the LibreOffice single-value renderer also shows VALUE only.
    assert chart.dataLabels.showVal is True
    assert chart.dataLabels.showCatName in (False, None)
    assert chart.dataLabels.showSerName in (False, None)
    assert chart.dataLabels.showLegendKey in (False, None)
    assert chart.dataLabels.showPercent in (False, None)
    # ... and in the tab's currency format.
    assert chart.dataLabels.numFmt == financial_charts._VALUE_FORMAT


# --- LTM revenue pie (overview slide) ---------------------------------------


def _ltm_workbook(tmp_path: Path, segments=None, sheet_name: str = "ltm-metrics") -> Path:
    """Build an ltm-metrics workbook, renaming its tab to the combined-workbook name."""
    segs = segments if segments is not None else [
        RevenueSegment("Cloud Services", 1932.0),
        RevenueSegment("Customer Support", 1480.0),
        RevenueSegment("License", 360.0),
        RevenueSegment("Professional Services", 290.0),
    ]
    path = build_ltm_metrics_workbook(
        company_name="SampleCo",
        period_label="LTM ended March 31, 2026",
        currency="US$MM",
        segmentation_basis="Service line",
        segments=segs,
        revenue_bridge=None,
        ebitda_bridge=None,
        output_dir=tmp_path,
    )
    wb = load_workbook(path)
    wb.active.title = sheet_name
    wb.save(path)
    return path


def test_ltm_revenue_overview_range_locates_block(tmp_path: Path):
    ws = load_workbook(_ltm_workbook(tmp_path)).active
    rng = ltm_revenue_overview_range(ws)
    # Section row 6, header row 7, four segments at rows 8-11; Total (row 12) excluded.
    assert rng == (8, 11)
    first, last = rng
    assert [ws.cell(row=r, column=1).value for r in range(first, last + 1)] == [
        "Cloud Services",
        "Customer Support",
        "License",
        "Professional Services",
    ]
    assert [ws.cell(row=r, column=2).value for r in range(first, last + 1)] == [
        1932.0,
        1480.0,
        360.0,
        290.0,
    ]
    assert ws.cell(row=last + 1, column=1).value == "Total"


def test_ltm_revenue_overview_range_none_when_absent(tmp_path: Path):
    from openpyxl import Workbook

    wb = Workbook()
    wb.active["A1"] = "Something Else"
    assert ltm_revenue_overview_range(wb.active) is None


def test_make_openpyxl_pie_applies_infor_formatting(tmp_path: Path):
    from openpyxl.chart import PieChart

    ws = load_workbook(_ltm_workbook(tmp_path)).active
    first, last = ltm_revenue_overview_range(ws)
    n = last - first + 1
    pie = _make_openpyxl_pie(ws, first, last)

    assert isinstance(pie, PieChart)
    assert pie.title is None
    # v0.5.22: legend docked on the RIGHT (was top), not overlaying the plot.
    assert pie.legend.position == "r"
    assert pie.legend.overlay is False
    # v0.5.22: no chart-area border (the openpyxl mirror of the COM strip — the
    # default outline used to frame the pie picture on the slide).
    assert pie.graphical_properties.line.noFill is True
    fills = [dp.graphicalProperties.solidFill.srgbClr for dp in pie.series[0].data_points]
    assert fills == INFOR_ACCENTS[:n]  # INFOR theme accents, in order
    # v0.5.23: the series charts the Top-4+Other source block's fraction column
    # (G), not the raw segment rows, and the data labels show VALUE only (not
    # the recomputed percentage). The label config lives on the SERIES-level
    # dLbls (v0.5.22 — the only level where per-point suppression is honored).
    assert "$G$" in pie.series[0].val.numRef.f
    cat = pie.series[0].cat
    cat_ref = cat.numRef.f if cat.numRef is not None else cat.strRef.f
    assert "$E$" in cat_ref
    assert pie.series[0].dLbls.showVal is True
    assert pie.series[0].dLbls.showPercent in (False, None)


def test_make_openpyxl_pie_pins_plot_area_left_of_legend(tmp_path: Path):
    # v0.5.22: the pie's plot area is pinned to the LEFT of the chart box via a
    # manual layout so it sits clear of the right-docked legend.
    ws = load_workbook(_ltm_workbook(tmp_path)).active
    first, last = ltm_revenue_overview_range(ws)
    pie = _make_openpyxl_pie(ws, first, last)
    manual = pie.layout.manualLayout
    assert manual.x == financial_charts._PIE_PLOT_X
    assert manual.y == financial_charts._PIE_PLOT_Y
    assert manual.w == financial_charts._PIE_PLOT_W
    assert manual.h == financial_charts._PIE_PLOT_H
    # The pie stays in the left portion of the box, clear of the right legend.
    assert manual.x + manual.w < 0.7
    assert manual.xMode == "edge" and manual.yMode == "edge"


def test_pie_slice_fills_cycle_past_six():
    # The pie itself now tops out at 5 slices, but the shared style helper still
    # cycles the accent fills for any hypothetical larger chart.
    from openpyxl import Workbook
    from openpyxl.chart import PieChart, Reference

    ws = Workbook().active
    for i in range(7):
        ws.cell(row=i + 1, column=1, value=f"Segment {i}")
        ws.cell(row=i + 1, column=2, value=100.0 + i)
    chart = PieChart()
    chart.add_data(Reference(ws, min_col=2, max_col=2, min_row=1, max_row=7), titles_from_data=False)
    financial_charts._style_openpyxl_pie(chart, 7)
    fills = [dp.graphicalProperties.solidFill.srgbClr for dp in chart.series[0].data_points]
    assert fills[6] == INFOR_ACCENTS[0]  # the 7th slice cycles back to accent1


def test_pie_grouping_and_grouped_labels():
    # ≤5 segments chart as-is, descending by $; >5 keep the 4 largest + "Other".
    from financial_charts import _grouped_pie_labels_amounts, _pie_grouping

    assert _pie_grouping([1.0, 3.0, 2.0]) == ([1, 2, 0], False)
    kept, has_other = _pie_grouping([10.0, 60.0, 5.0, 20.0, 4.0, 3.0, 2.0])
    assert (kept, has_other) == ([1, 3, 0, 2], True)

    # The PRL17 live-run shape: 7 segments -> top 4 + Other = remainder.
    names = ["CreditFresh", "MoneyKey Bank", "QuidMarket", "MoneyKey Direct",
             "LaaS Fees", "Fora", "Other Revenue"]
    amounts = [392.5, 99.8, 54.5, 31.9, 21.2, 13.1, 4.0]
    labels, grouped = _grouped_pie_labels_amounts(names, amounts)
    assert labels == ["CreditFresh", "MoneyKey Bank", "QuidMarket", "MoneyKey Direct", "Other"]
    assert grouped[:4] == [392.5, 99.8, 54.5, 31.9]
    assert grouped[4] == pytest.approx(21.2 + 13.1 + 4.0)


def test_pie_groups_more_than_five_segments_into_top4_other(tmp_path: Path):
    """>5 segments must chart exactly 5 slices — the 4 largest plus 'Other' —
    via a live source block in columns E:G (the legend overflowed the pie when
    every segment charted; analyst asked for 4 + Other)."""
    segs = [
        RevenueSegment("CreditFresh Bank Program", 392.5),
        RevenueSegment("MoneyKey Bank Service Program", 99.8),
        RevenueSegment("QuidMarket Direct Lending", 54.5),
        RevenueSegment("MoneyKey Direct Lending & CSO", 31.9),
        RevenueSegment("Lending-as-a-Service Fees", 21.2),
        RevenueSegment("Fora Direct Lending", 13.1),
        RevenueSegment("Other Revenue", 4.0),
    ]
    ws = load_workbook(_ltm_workbook(tmp_path, segments=segs)).active
    first, last = ltm_revenue_overview_range(ws)
    assert last - first + 1 == 7
    total_row = last + 1
    pie = _make_openpyxl_pie(ws, first, last)

    # Exactly 5 slices, spanning the source block rows.
    src_ref = pie.series[0].val.numRef.f
    assert src_ref.endswith(f"$G${first}:$G${first + 4}")
    assert len(pie.series[0].data_points) == 5

    # Source block: title + headers + 4 live segment refs (descending $) + Other.
    assert ws.cell(row=first - 2, column=5).value == "Pie Chart Source (Top 4 + Other)"
    assert ws.cell(row=first - 1, column=5).value == "Segment"
    assert ws.cell(row=first - 1, column=7).value == "% of Total"
    assert [ws.cell(row=first + j, column=5).value for j in range(4)] == [
        f"=A{first}", f"=A{first + 1}", f"=A{first + 2}", f"=A{first + 3}"
    ]
    assert ws.cell(row=first + 4, column=5).value == "Other"
    assert ws.cell(row=first, column=6).value == f"=B{first}"
    assert ws.cell(row=first + 4, column=6).value == f"=B{total_row}-SUM(F{first}:F{first + 3})"
    assert ws.cell(row=first, column=7).value == f"=F{first}/$B${total_row}"


def test_pie_source_block_rerun_replaces_stale_rows(tmp_path: Path):
    # A re-run with fewer segments must clear the old block's extra slice rows.
    segs7 = [RevenueSegment(f"Segment {i}", 100.0 - i) for i in range(7)]
    path = _ltm_workbook(tmp_path, segments=segs7)
    wb = load_workbook(path)
    ws = wb.active
    first, last = ltm_revenue_overview_range(ws)
    _make_openpyxl_pie(ws, first, last)
    assert ws.cell(row=first + 4, column=5).value == "Other"

    segs3 = [RevenueSegment(f"Segment {i}", 100.0 - i) for i in range(3)]
    path3 = _ltm_workbook(tmp_path / "smaller", segments=segs3)
    ws3 = load_workbook(path3).active
    f3, l3 = ltm_revenue_overview_range(ws3)
    _make_openpyxl_pie(ws3, f3, l3)
    _make_openpyxl_pie(ws3, f3, l3)  # idempotent re-run
    assert ws3.cell(row=f3 - 2, column=5).value == "Pie Chart Source"  # no "(Top 4 + Other)"
    assert [ws3.cell(row=f3 + j, column=5).value for j in range(3)] == [
        f"=A{f3}", f"=A{f3 + 1}", f"=A{f3 + 2}"
    ]
    assert ws3.cell(row=f3 + 3, column=5).value is None  # nothing stale below the block


def test_pie_data_labels_show_value_only_with_percent_format(tmp_path: Path):
    # v0.5.19: the pie's data labels show VALUE only (category/series/legend-key and
    # percentage flags all off) and carry the "%" number format, so the source
    # block's fraction renders e.g. 0.452 as "45.2%". Since v0.5.22 the config
    # lives on the series-level dLbls.
    ws = load_workbook(_ltm_workbook(tmp_path)).active
    first, last = ltm_revenue_overview_range(ws)
    pie = _make_openpyxl_pie(ws, first, last)
    labels = pie.series[0].dLbls
    assert labels.showVal is True
    assert labels.showPercent in (False, None)
    assert labels.showCatName in (False, None)
    assert labels.showSerName in (False, None)
    assert labels.showLegendKey in (False, None)
    assert labels.numFmt == '#,##0.0%_);(#,##0.0%);"--"'
    assert labels.numFmt == financial_charts._PIE_LABEL_FORMAT
    # The default segments (47.6/36.4/8.9/7.1%) are all above the 3% threshold —
    # no per-point suppression entries.
    assert not labels.dLbl


def test_pie_labels_inside_end_and_white(tmp_path: Path):
    # v0.5.24: every pie label sits INSIDE its slice (Best Fit floated the
    # small-slice labels outside the pie in the live run) and renders white —
    # black is unreadable inside the dark accent fills. The legend keeps its
    # black Palatino 8.
    ws = load_workbook(_ltm_workbook(tmp_path)).active
    first, last = ltm_revenue_overview_range(ws)
    pie = _make_openpyxl_pie(ws, first, last)
    labels = pie.series[0].dLbls
    assert labels.dLblPos == "inEnd"
    label_font = labels.txPr.p[0].pPr.defRPr
    assert label_font.solidFill.srgbClr == "FFFFFF"
    assert label_font.latin.typeface == "Palatino Linotype"
    legend_font = pie.legend.txPr.p[0].pPr.defRPr
    assert legend_font.solidFill.srgbClr == "000000"


def _assert_label_suppressed(label):
    """A per-point all-show-flags-off override — nothing rendered for the slice."""
    assert label.showVal is False
    assert label.showPercent is False
    assert label.showCatName is False
    assert label.showSerName is False
    assert label.showLegendKey is False


def test_pie_labels_only_on_slices_above_3pct(tmp_path: Path):
    # v0.5.22: slices at or below 3% of the total carry NO data label (they overlap
    # each other in the short overview box). 600/270/100/30 of 1000 → 60% / 27% /
    # 10% / 3.0% — the 3.0% slice is not ABOVE the threshold, so it is suppressed.
    segs = [
        RevenueSegment("Cloud", 600.0),
        RevenueSegment("Support", 270.0),
        RevenueSegment("License", 100.0),
        RevenueSegment("Other", 30.0),
    ]
    ws = load_workbook(_ltm_workbook(tmp_path, segments=segs)).active
    first, last = ltm_revenue_overview_range(ws)
    pie = _make_openpyxl_pie(ws, first, last)
    suppressed = pie.series[0].dLbls.dLbl
    assert [d.idx for d in suppressed] == [3]
    _assert_label_suppressed(suppressed[0])


def test_single_pie_chart_suppresses_small_slices():
    # The throwaway LibreOffice render workbook charts literal fractions; the same
    # >3% rule applies to it directly.
    from openpyxl import Workbook

    from financial_charts import _make_single_pie_chart

    wb = Workbook()
    ws = wb.active
    fractions = [0.60, 0.27, 0.105, 0.025]
    for i, frac in enumerate(fractions, start=1):
        ws.cell(row=i, column=1, value=f"Segment {i}")
        ws.cell(row=i, column=2, value=frac)
    pie = _make_single_pie_chart(ws, len(fractions))
    suppressed = pie.series[0].dLbls.dLbl
    assert [d.idx for d in suppressed] == [3]
    _assert_label_suppressed(suppressed[0])


def test_suppressed_pie_label_indices_threshold():
    # Strictly ABOVE 3% keeps a label; 3% exactly, below, and non-numeric do not.
    assert financial_charts._suppressed_pie_label_indices(
        [0.5, 0.03, 0.031, 0.029, None]
    ) == [1, 3, 4]


def _deck_with_pie_placeholder(tmp_path: Path) -> Path:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    box = slide.shapes.add_textbox(Inches(2.62), Inches(3.40), Inches(4.51), Inches(1.77))
    box.name = "Rectangle 4"
    box.text_frame.text = "[Pie Chart Placeholder]"
    path = tmp_path / "overview_deck.pptx"
    prs.save(path)
    return path


def test_insert_png_replaces_pie_placeholder(tmp_path: Path):
    deck = _deck_with_pie_placeholder(tmp_path)
    out = insert_png_into_placeholder(
        deck_path=deck,
        slide_index=0,
        placeholder_name="Rectangle 4",
        png_bytes=_PNG_1X1,
        output_path=deck,
    )
    prs = Presentation(out)
    slide = prs.slides[0]
    assert "Rectangle 4" not in [s.name for s in slide.shapes]
    pictures = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pictures) == 1


def test_render_pie_returns_none_without_ltm_tab(tmp_path: Path):
    from openpyxl import Workbook

    wb = Workbook()
    wb.active.title = "financial-summary"  # no ltm-metrics tab
    combined = tmp_path / "pitch-Test.xlsx"
    wb.save(combined)

    deck = _deck_with_pie_placeholder(tmp_path)
    out = render_ltm_revenue_pie_into_deck(
        deck_path=deck,
        combined_workbook_path=combined,
        slide_index=0,
        output_path=deck,
    )
    assert out is None
    # The placeholder is left intact (the null path, like ownership).
    prs = Presentation(deck)
    text = " ".join(
        s.text_frame.text for s in prs.slides[0].shapes if s.has_text_frame
    )
    assert "[Pie Chart Placeholder]" in text


# --- LibreOffice is resolved through the shared locator ------------------------


def test_recalc_values_resolves_libreoffice_through_find_soffice(tmp_path: Path, monkeypatch):
    """v0.5.36: the chart renderers must resolve LibreOffice via ``find_soffice``.

    Phase A (v0.5.35) made LibreOffice the default renderer everywhere but left
    these sites on a bare ``shutil.which``, which the Windows MSI never satisfies —
    so dev either failed or silently degraded. Locked platform-independently: every
    PATH lookup returns None and only the locator knows where the binary is, so a
    site that goes back to ``which`` cannot find it.
    """
    path = _fs_workbook(tmp_path)
    wb = load_workbook(path)
    wb.active.title = "financial-summary"
    wb.save(path)
    monkeypatch.setattr("shutil.which", lambda *_a, **_k: None)
    monkeypatch.setattr("excel_to_powerpoint.find_soffice", lambda: "/opt/libreoffice/soffice")
    seen: list[str] = []

    def fake_convert(soffice, src, _out_fmt, out_dir):
        seen.append(soffice)
        (Path(out_dir) / f"{Path(src).stem}.xlsx").write_bytes(Path(src).read_bytes())

    monkeypatch.setattr("excel_to_powerpoint._soffice_convert", fake_convert)

    financial_charts._libreoffice_recalc_values(path, "financial-summary", 2, 7, [6])

    assert seen == ["/opt/libreoffice/soffice"]


# --- graceful degradation when LibreOffice is unavailable (Issue 1b / Issue 3) ---


def test_fs_charts_persist_when_libreoffice_missing(tmp_path: Path, monkeypatch):
    """Issue 1b: the four native charts are saved on the `financial-summary` tab even
    when LibreOffice is absent. The backend returns ``{}`` (no deck PNGs) instead of
    raising, so the stage degrades gracefully rather than aborting."""
    path = _fs_workbook(tmp_path)
    wb = load_workbook(path)
    wb.active.title = "financial-summary"
    wb.save(path)
    _no_libreoffice(monkeypatch)

    pngs = financial_charts._build_charts_openpyxl_libreoffice(
        path, "financial-summary", 2, 7, [6, 7, 8, 9]
    )

    assert pngs == {}  # degraded: no deck PNGs, but no exception
    assert _chart_part_count(path) == 4  # native charts persisted regardless


def test_fs_orchestrator_degrades_to_none_when_render_unavailable(tmp_path: Path, monkeypatch):
    """Issue 1b glue: when the backend persists the workbook charts but cannot render
    their PNGs (LibreOffice missing), the orchestrator returns None and leaves the deck
    untouched rather than crashing on a missing PNG."""
    path = _fs_workbook(tmp_path)
    wb = load_workbook(path)
    wb.active.title = "financial-summary"
    wb.save(path)
    deck = _deck_with_placeholder(tmp_path)
    before = deck.read_bytes()
    # Force both backends to the "charts saved, no PNGs" outcome on any platform.
    monkeypatch.setattr(financial_charts, "_build_charts_com", _raise_runtime)
    monkeypatch.setattr(financial_charts, "_build_charts_openpyxl_libreoffice", lambda *a, **k: {})

    out = render_financial_summary_charts_into_deck(deck_path=deck, combined_workbook_path=path)

    assert out is None
    assert deck.read_bytes() == before  # deck untouched — placeholders preserved


def test_pie_persists_when_libreoffice_missing(tmp_path: Path, monkeypatch):
    """Issue 3 parity: the native pie is saved on the `ltm-metrics` tab even when
    LibreOffice is absent. The backend returns None (no PNG) instead of raising."""
    path = _ltm_workbook(tmp_path)
    first, last = ltm_revenue_overview_range(load_workbook(path).active)
    _no_libreoffice(monkeypatch)

    png = financial_charts._build_pie_openpyxl_libreoffice(path, "ltm-metrics", first, last)

    assert png is None  # degraded: no PNG, but no exception
    assert _chart_part_count(path) == 1  # native pie persisted regardless


def test_pie_orchestrator_degrades_to_none_when_render_unavailable(tmp_path: Path, monkeypatch):
    """Issue 3 glue: when the pie backend persists the workbook but cannot render the
    PNG, the orchestrator returns None and leaves the overview placeholder."""
    path = _ltm_workbook(tmp_path)
    deck = _deck_with_pie_placeholder(tmp_path)
    before = deck.read_bytes()
    monkeypatch.setattr(financial_charts, "_build_pie_com", _raise_runtime)
    monkeypatch.setattr(financial_charts, "_build_pie_openpyxl_libreoffice", lambda *a, **k: None)

    out = render_ltm_revenue_pie_into_deck(
        deck_path=deck, combined_workbook_path=path, slide_index=0
    )

    assert out is None
    assert deck.read_bytes() == before  # overview placeholder preserved


# --- both chart steps on ONE combined workbook (v0.5.20 regression) ----------


def _combined_workbook(dir_path: Path) -> Path:
    """A combined pitch workbook carrying BOTH the `financial-summary` and
    `ltm-metrics` tabs — the real financial-charts input, where the FS-chart and
    pie steps run back-to-back against the same file."""
    dir_path.mkdir(parents=True, exist_ok=True)
    fs = load_workbook(_fs_workbook(dir_path))
    fs.active.title = "financial-summary"
    ltm_ws = load_workbook(_ltm_workbook(dir_path)).active
    dst = fs.create_sheet("ltm-metrics")
    for row in ltm_ws.iter_rows():
        for cell in row:
            dst.cell(row=cell.row, column=cell.column, value=cell.value)
    combined = dir_path / "pitch-SampleCo.xlsx"
    fs.save(combined)
    return combined


def test_fs_charts_and_pie_coexist_on_combined_workbook(tmp_path: Path, monkeypatch):
    """Regression (v0.5.21): both chart steps persist via a load→save of the SAME
    combined workbook. The steps must compose — exactly five charts (4 FS + 1
    pie) after both run, in either order, with neither step wiping the other's
    charts (an openpyxl build that drops chart parts on load) nor a re-run
    accumulating a duplicate set next to the stale one (openpyxl 3.x, which
    round-trips chart parts — pre-fix a re-run grew 5 → 10)."""
    _no_libreoffice(monkeypatch)

    # SKILL.md order: FS charts first, pie second.
    combined = _combined_workbook(tmp_path / "fs_then_pie")
    first, last = ltm_revenue_overview_range(load_workbook(combined)["ltm-metrics"])
    financial_charts._build_charts_openpyxl_libreoffice(
        combined, "financial-summary", 2, 7, [6, 7, 8, 9]
    )
    financial_charts._build_pie_openpyxl_libreoffice(combined, "ltm-metrics", first, last)
    assert _chart_part_count(combined) == 5  # 4 FS charts + 1 pie

    # Re-running either step replaces its charts instead of accumulating.
    financial_charts._build_charts_openpyxl_libreoffice(
        combined, "financial-summary", 2, 7, [6, 7, 8, 9]
    )
    financial_charts._build_pie_openpyxl_libreoffice(combined, "ltm-metrics", first, last)
    assert _chart_part_count(combined) == 5

    # Reverse order loses nothing either.
    combined = _combined_workbook(tmp_path / "pie_then_fs")
    first, last = ltm_revenue_overview_range(load_workbook(combined)["ltm-metrics"])
    financial_charts._build_pie_openpyxl_libreoffice(combined, "ltm-metrics", first, last)
    financial_charts._build_charts_openpyxl_libreoffice(
        combined, "financial-summary", 2, 7, [6, 7, 8, 9]
    )
    assert _chart_part_count(combined) == 5


def test_pie_legend_pinned_full_right_side_at_8pt(tmp_path: Path):
    """The legend is pinned to the full remaining right side of the chart box at
    Palatino 8 — Excel's auto legend wrapped every entry and silently dropped
    the ones that no longer fit (the "Other" entry vanished in the live run)."""
    ws = load_workbook(_ltm_workbook(tmp_path)).active
    first, last = ltm_revenue_overview_range(ws)
    pie = _make_openpyxl_pie(ws, first, last)
    manual = pie.legend.layout.manualLayout
    assert manual.x == financial_charts._PIE_LEGEND_X
    assert manual.y == financial_charts._PIE_LEGEND_Y
    assert manual.w == financial_charts._PIE_LEGEND_W
    assert manual.h == financial_charts._PIE_LEGEND_H
    # Legend starts where the pinned plot area ends — no overlap either way.
    assert manual.x >= financial_charts._PIE_PLOT_X + financial_charts._PIE_PLOT_W
    legend_pr = pie.legend.txPr.p[0].pPr.defRPr
    assert legend_pr.sz == financial_charts._PIE_LEGEND_FONT_SIZE_PT * 100
    assert legend_pr.latin.typeface == "Palatino Linotype"

# ─── Configurable slide mix (v0.5.26): N metric rows across N FS slides ──────


def _metrics8() -> list[MetricSeries]:
    return _metrics() + [
        MetricSeries("Operating Income", "US$MM", [620.0, 710.0, 815.0, 920.0, 1010.0],
                     result_label="LTM Operating Income"),
        MetricSeries("Free Cash Flow", "US$MM", [350.0, 400.0, 460.0, 520.0, 580.0],
                     result_label="LTM Free Cash Flow"),
        MetricSeries("Gross Margin", "%", [0.40, 0.41, 0.41, 0.42, 0.42],
                     ltm_value=0.42),
        MetricSeries("Return on Equity", "%", [0.12, 0.13, 0.14, 0.15, 0.16],
                     ltm_value=0.16),
    ]


def test_metric_data_rows_four_and_eight(tmp_path: Path):
    ws4 = load_workbook(_fs_workbook(tmp_path)).active
    assert financial_charts.metric_data_rows(ws4) == [6, 7, 8, 9]
    ws8 = load_workbook(
        _fs_workbook(tmp_path, metrics=_metrics8(), metric_count=8, file_stem="EightMetrics")
    ).active
    assert financial_charts.metric_data_rows(ws8) == [6, 7, 8, 9, 10, 11, 12, 13]


def test_fs_chart_anchor_grid_scales_with_data_rows():
    # 4 metrics (last data row 9): the historical B11/I11/B27/I27 grid.
    assert [financial_charts._fs_chart_anchor(i, 9) for i in range(4)] == [
        "B11", "I11", "B27", "I27",
    ]
    # 8 metrics (last data row 13): grid starts below the data at row 15.
    assert [financial_charts._fs_chart_anchor(i, 13) for i in range(8)] == [
        "B15", "I15", "B31", "I31", "B47", "I47", "B63", "I63",
    ]


def test_persist_eight_charts_when_libreoffice_missing(tmp_path: Path, monkeypatch):
    path = _fs_workbook(tmp_path, metrics=_metrics8(), metric_count=8)
    wb = load_workbook(path)
    wb.active.title = "financial-summary"
    wb.save(path)
    _no_libreoffice(monkeypatch)

    pngs = financial_charts._build_charts_openpyxl_libreoffice(
        path, "financial-summary", 2, 7, list(range(6, 14))
    )

    assert pngs == {}  # degraded: no deck PNGs, but no exception
    assert _chart_part_count(path) == 8  # one native chart per metric row


def _deck_with_fs_slides(tmp_path: Path, n: int) -> Path:
    """A deck with one non-FS slide followed by `n` Financial Summary slides,
    each carrying the four named chart placeholders (as the assembler's FS
    clones do)."""
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])  # non-FS slide — discovery must skip it
    placeholder_texts = [
        ("Rectangle 17", "[Placeholder for Metric #1 Chart]"),
        ("Rectangle 7", "[Placeholder for Metric #2 Chart]"),
        ("Rectangle 19", "[Placeholder for Metric #3 Chart]"),
        ("Rectangle 18", "[Placeholder for Metric #4 Chart]"),
    ]
    for _ in range(n):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        for name, text in placeholder_texts:
            box = slide.shapes.add_textbox(Inches(0.35), Inches(1.51), Inches(4.53), Inches(2.51))
            box.name = name
            box.text_frame.text = text
    path = tmp_path / "deck.pptx"
    prs.save(path)
    return path


def test_find_financial_summary_slides_skips_non_fs(tmp_path: Path):
    deck = _deck_with_fs_slides(tmp_path, 2)
    assert financial_charts._find_financial_summary_slides(deck) == [1, 2]


def test_find_financial_summary_slides_raises_when_absent(tmp_path: Path):
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    path = tmp_path / "plain.pptx"
    prs.save(path)
    with pytest.raises(ValueError):
        financial_charts._find_financial_summary_slides(path)


def test_orchestrator_inserts_eight_charts_across_two_fs_slides(tmp_path: Path, monkeypatch):
    path = _fs_workbook(tmp_path, metrics=_metrics8(), metric_count=8)
    wb = load_workbook(path)
    wb.active.title = "financial-summary"
    wb.save(path)
    deck = _deck_with_fs_slides(tmp_path, 2)
    rows = list(range(6, 14))
    monkeypatch.setattr(financial_charts, "_build_charts_com", _raise_runtime)
    monkeypatch.setattr(
        financial_charts,
        "_build_charts_openpyxl_libreoffice",
        lambda *a, **k: {r: _PNG_1X1 for r in rows},
    )

    out = render_financial_summary_charts_into_deck(deck_path=deck, combined_workbook_path=path)

    prs = Presentation(out)
    # The non-FS slide is untouched; each FS slide got its four pictures.
    assert not [s for s in prs.slides[0].shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    for idx in (1, 2):
        slide = prs.slides[idx]
        pictures = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        assert len(pictures) == 4
        assert not any(s.name == "Rectangle 17" for s in slide.shapes)


def test_orchestrator_raises_on_row_slide_mismatch(tmp_path: Path, monkeypatch):
    # A 4-metric tab against a 2-FS-slide deck is a wireframe/financial-summary
    # stage mismatch — fail loudly instead of charting the wrong rows.
    path = _fs_workbook(tmp_path)
    wb = load_workbook(path)
    wb.active.title = "financial-summary"
    wb.save(path)
    deck = _deck_with_fs_slides(tmp_path, 2)
    monkeypatch.setattr(financial_charts, "_build_charts_com", _raise_runtime)
    monkeypatch.setattr(
        financial_charts,
        "_build_charts_openpyxl_libreoffice",
        lambda *a, **k: {r: _PNG_1X1 for r in (6, 7, 8, 9)},
    )
    with pytest.raises(ValueError):
        render_financial_summary_charts_into_deck(deck_path=deck, combined_workbook_path=path)


# ─── COM chart builders: cleanup ordering (v0.5.38) ───────────────────────────
#
# `_build_charts_com` / `_build_pie_com` had NO pytest coverage before v0.5.38 —
# the fallback tests above stub them out entirely — so their restructuring onto
# `excel_com_app` would otherwise have shipped unverified. The fake COM stack
# spawns no Excel, which is the point: spawning it is what leaks.


class _FakeChartCells:
    Left = 10.0
    Top = 20.0


class _FakeSeriesCollection:
    def __init__(self):
        self.Count = 0

    def __call__(self, _i=None):
        return self

    def NewSeries(self):
        return _FakeSeries()

    def Delete(self):
        pass


class _FakeSeries:
    Values = None
    XValues = None


class _FakeComChart:
    ChartType = None

    def __init__(self, log):
        self._log = log
        self._sc = _FakeSeriesCollection()

    def SeriesCollection(self, _i=None):
        return self._sc

    def Export(self, Filename, FilterName):  # noqa: N803
        self._log.append("Export")
        Path(Filename).write_bytes(b"\x89PNG\r\n\x1a\n" + b"0" * 16)


class _FakeChartObject:
    Left = 0.0
    Top = 0.0

    def __init__(self, log):
        self.Chart = _FakeComChart(log)


class _FakeChartObjects:
    def __init__(self, log):
        self._log = log

    def Delete(self):
        self._log.append("ChartObjects.Delete")

    def Add(self, **_kwargs):
        return _FakeChartObject(self._log)


class _FakeChartWorksheet:
    def __init__(self, log):
        self._log = log

    def ChartObjects(self):
        return _FakeChartObjects(self._log)

    def Cells(self, _r, _c=None):
        return _FakeChartCells()

    def Range(self, _a, _b=None):
        return object()


class _FakeChartWorkbook:
    def __init__(self, log):
        self._log = log

    def Worksheets(self, _name):
        return _FakeChartWorksheet(self._log)

    def Save(self):
        self._log.append("wb.Save")

    def Close(self, SaveChanges):  # noqa: N803
        self._log.append("wb.Close")


def _install_chart_com_fakes(monkeypatch, log):
    from tests.fake_com import FakeExcelApp, install_fake_com

    class _Workbooks:
        def Open(self, _path, **_kwargs):
            return _FakeChartWorkbook(log)

    class _App(FakeExcelApp):
        Workbooks = _Workbooks()

        def CalculateFull(self):
            pass

        def Calculate(self):
            pass

    install_fake_com(monkeypatch, log, _App(log))


def test_com_chart_build_closes_workbook_before_quit(tmp_path: Path, monkeypatch):
    from tests.fake_com import assert_released_before_quit

    log: list[str] = []
    _install_chart_com_fakes(monkeypatch, log)
    monkeypatch.setattr(financial_charts, "_format_com_chart", lambda *_a: None)

    pngs = financial_charts._build_charts_com(
        tmp_path / "combined.xlsx", "financial-summary", 2, 7, [6, 7, 8, 9]
    )

    assert sorted(pngs) == [6, 7, 8, 9]
    assert all(v.startswith(b"\x89PNG") for v in pngs.values())
    assert log.count("Export") == 4, "one chart exported per metric row"
    assert_released_before_quit(log, "wb.Save", "wb.Close")


def test_com_pie_build_closes_workbook_before_quit(tmp_path: Path, monkeypatch):
    from tests.fake_com import assert_released_before_quit

    log: list[str] = []
    _install_chart_com_fakes(monkeypatch, log)
    monkeypatch.setattr(financial_charts, "_format_com_pie", lambda *_a: None)
    monkeypatch.setattr(
        financial_charts, "_write_pie_source_block_com", lambda _ws, _f, _l: (30, 34, [1.0, 2.0])
    )

    png = financial_charts._build_pie_com(tmp_path / "combined.xlsx", "ltm-metrics", 10, 14)

    assert png.startswith(b"\x89PNG")
    assert_released_before_quit(log, "wb.Save", "wb.Close")


@pytest.mark.parametrize(
    "builder, args",
    [
        ("_build_charts_com", ("financial-summary", 2, 7, [6])),
        ("_build_pie_com", ("ltm-metrics", 10, 14)),
    ],
)
def test_com_chart_startup_failure_is_runtime_error(tmp_path: Path, monkeypatch, builder, args):
    """A failed startup must read as RuntimeError so the documented fallback to
    openpyxl + LibreOffice engages — and must not be double-wrapped."""
    from tests.fake_com import install_fake_com

    log: list[str] = []
    install_fake_com(monkeypatch, log, None)  # DispatchEx raises

    with pytest.raises(RuntimeError, match="Excel COM unavailable") as exc_info:
        getattr(financial_charts, builder)(tmp_path / "combined.xlsx", *args)

    assert "chart build failed" not in str(exc_info.value), "no double wrapping"
    assert "pie build failed" not in str(exc_info.value)
    assert log[-1] == "CoUninitialize"
