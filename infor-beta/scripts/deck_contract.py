"""Executable contract for a built INFOR deck — the Phase B visual oracle.

The dominant failure mode of this plugin is **visual fidelity that nothing in the
system can see until a human opens the artefact**: valid XML, plausible content,
and a table that renders 0.2" taller than it declares or a bullet block that runs
straight through the section header below it. `verify_deck(path)` is the machine
that looks.

Two tiers, deliberately separated so non-determinism stays out of the gate.

**Deterministic — blocking.** Reproducible measurement, no model in the loop:

  1. ``forbidden-string`` / ``unsubstituted-currency-token`` / ``unfilled-token``
     / ``unfilled-placeholder`` — a scan of every shape, group child, and table cell.
  2. ``shape-outside-slide`` — declared shape extents vs. the slide box.
  3. ``table-taller-than-library`` — a filled table's declared height vs. the
     height the shipped library ships that same table at.
  4. ``rendered-overflow`` — **rendered** ink below a shape's declared bottom,
     counted where no other shape has claimed the territory.
  5. ``masked-overflow`` — the same overflow **attributed to one shape** by
     rendering it alone, for the shapes check 4 is structurally blind to.

Checks 4 and 5 are why this module renders. Checks 1–3 read the XML, and the XML
is exactly what the historical bugs were invisible to: PRL17's market-entry table
declared 5.710" — *under* the library's 5.720" — and rendered 5.91", because **a
stored row height is only a render-time minimum** and the layout engine re-grows
any row whose text needs more. Nothing in the file is wrong. You have to look at
the pixels.

**Vision — advisory, agent-inspected, surfaced at the `deck` checkpoint.**
Overlap, collision, unreadable contrast, chart-label pileup. These need a model
looking at pixels, not text extracted from them, so `vision_pass` deliberately
carries **no OCR**: it renders every slide, extracts every embedded picture at
native resolution, and returns structured `vision-review` findings naming what to
adjudicate on which slide. The reviewer is the checkpoint agent.

They are advisory and never produce an unattended pass/fail — there is no CI in
this repo, and pytest is deliberately not gated on them. Mechanising this tier
into an assertion was tried and abandoned: OCR recall on 8-9 pt figures is not
good enough to block a client deliverable on, and a test that asserts what an OCR
engine happened to read is a test of the OCR engine.

What the error-value scan does and does not cover
-------------------------------------------------
`ERROR_TOKENS` scans **text shapes and table cells only** — never rasterised
picture content — and that scope is deliberate, not a limitation to fix later.

An Excel error value inside a pasted range picture is usually **correct**. The cap
table's forward-estimate columns are CapIQ UDF calls (`SP_REV_EST` /
`SP_EBITDA_EST` at `E47:F48`), wrapped by rows 34/35 in
``IFERROR(..., "n/a ")``; the EV/metric rows then divide ``$F$31`` by that text,
so an un-refreshed CapIQ estimate propagates ``#VALUE!`` into `E39:F40` **by
design**. CapIQ cannot be refreshed in this environment, so that is the normal
state of a shipped cap table — and the same holds for the comps and precedents
tabs (array formulas ship un-evaluated for the analyst to refresh) and for the
financial-summary LTM link (``#N/A`` until the combined workbook exists).

A workbook-side scan is also not the missing piece. Measured on
`fixtures/pitch-workbook.xlsx`: of the 121 formula cells on the `captable` tab,
**zero** carry a cached value under ``data_only=True`` — including ``=TODAY()-1``.
openpyxl only ever returns the cache, and nothing populated it, so such a scan
sees ``None`` everywhere and looks green while proving nothing. Real values exist
only on the COM path (live Excel) or after a LibreOffice recalc-on-load. Anything
built there must fail loudly on an empty cache rather than pass quietly.

An error token in deck **text**, by contrast, has no such excuse: nothing writes
a CapIQ formula into a text frame, so it means a value was written from an
unresolved computation. That stays blocking.

The shipped library is the geometric baseline
---------------------------------------------
"Does content render below its declared box?" has a non-zero answer on a
perfectly good INFOR slide: the library's own footnote placeholders render ~0.07"
past their boxes, and its tombstone slides park decorative shapes off-canvas
entirely. Those are design, not defects, and a contract that reported them would
be ignored within a week.

So the reference is not zero — it is **the same shape on the matching slide of the
blank library**. Every geometric check measures the built deck, measures the
library, and reports the *excess*. That makes the contract self-maintaining in
the Phase C sense: an analyst who re-heights a library table or nudges a
placeholder moves the reference with it, and only what the *fill* did is
attributed to the fill.

Built slides are matched to library slides by shape-name signature rather than by
index or title text: a built deck's titles are filled (``Introduction to Acme``
vs. the library's ``Introduction to [Client Name]``), the pitch flow deletes and
clones slides, and the market-entry section repeats one library slide N times.

String checks get **no** baseline. A ``[x]`` in a built deck is a defect whether
or not the library is where it came from — that token exists precisely to be
substituted.

Why LibreOffice is the oracle
-----------------------------
`slide_render` renders through LibreOffice on every platform (v0.5.35), and this
module measures on that render deliberately, **without correcting toward
PowerPoint's metrics**. The two engines do not lay out text identically even with
the same Palatino Linotype file installed: measured on the pitch fixture's
overview slide, LibreOffice wraps earlier and produces one extra line and a text
block ~6% taller. That makes LibreOffice the *conservative* renderer — "fits
under LibreOffice" implies "fits in PowerPoint", but not the reverse. Measuring
on the conservative side is the whole point; a contract that passed on
LibreOffice and failed in PowerPoint would be worse than no contract.

Read a `rendered-overflow` depth accordingly: it is a LibreOffice measurement,
and may run slightly larger than what PowerPoint shows. It says content does not
fit its declared box; it is not a PowerPoint pixel prediction.

`masked-overflow` — the same measurement, per shape
--------------------------------------------------
`rendered-overflow` counts ink that lands in **nobody's** declared box. That is
what keeps it free of false positives, and it is also structurally blind to an
overflow running *straight into the shape below*, which the lower shape's own
claim masks. No threshold fixes it: inside that box the pixels belong to two
shapes at once.

The earnings-update fixture is the instance. Its broker table declares a bottom of
6.184" with the summary box starting at 6.221", and the table's re-grown last row
(`EPS, Adj.`) renders underneath the summary box's text. Only **0.037"** of that
overflow is unclaimed — under tolerance, so `rendered-overflow` says nothing.

`masked-overflow` measures it by **rendering the shape alone**. For every
growth-prone shape with a masking neighbour, a probe deck carries the shape on its
own layout plus an empty slide of that same layout; subtracting the second from
the first leaves nothing but that shape's ink, and the deepest inked row below its
declared bottom is the answer. On the broker table that is **0.153"** against a
0.000" library baseline. Absolute positioning is what makes it sound — removing a
neighbour cannot move or re-wrap the shape under test — and every probe rides in
one deck, so the pass costs one extra conversion (2.5-7.7 s measured).

This is a prerequisite for the converge loop, not a refinement: a repair step
that cannot attribute overflowing ink to a shape does not know which shape to
shrink. A shape `rendered-overflow` already reported is skipped, so the two checks
never double-count the same defect.

Relationship to the estimation code
-----------------------------------
This module measures; it does not estimate. That was a deliberate separation while
the estimation code still existed — a contract built on the same character-width
table and em constants could not have caught those constants being wrong — and as
of v0.5.39 the estimation code is gone: `palatino_text_width_in`,
`estimate_text_height_in`, `fit_text_scale`, `set_table_height`'s row minimums and
the assemblers' font ladders were all replaced by `deck_repair`, which decides
sizes from this module's measurements.

Usage
-----
    from deck_contract import verify_deck, SEVERITY_BLOCKING

    findings = verify_deck("Pitch Deck - Acme.pptx")
    blocking = [f for f in findings if f.blocking]

`render=False` gives a fast XML-only pass. When LibreOffice is absent the
render-measured tier degrades to one loud ``render-unavailable`` finding — a deck
is never reported clean because nothing looked at it.
"""

from __future__ import annotations

import os
import re
import tempfile
from copy import deepcopy
from dataclasses import dataclass, field, replace
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.util import Emu

from pptx_helpers import apply_text_scale, normal_autofit_scale, strip_autofit
from template_layout import LIBRARY_SLIDE_MARKERS

# ─── Findings ────────────────────────────────────────────────────────────────

SEVERITY_BLOCKING = "blocking"
SEVERITY_ADVISORY = "advisory"


@dataclass(frozen=True)
class Finding:
    """One contract violation.

    ``slide`` is zero-based; add 1 for the number the analyst sees in PowerPoint.
    ``measured_in`` / ``limit_in`` carry the measurement behind a geometric
    finding, so the repair step in Phase B step 3 knows how far it has to move
    rather than only that something is wrong.
    """

    kind: str
    severity: str
    slide: int
    detail: str
    shape: str | None = None
    measured_in: float | None = None
    limit_in: float | None = None

    @property
    def blocking(self) -> bool:
        return self.severity == SEVERITY_BLOCKING

    def __str__(self) -> str:
        where = f"slide {self.slide + 1}"
        if self.shape:
            where += f" {self.shape!r}"
        return f"[{self.severity}] {self.kind}: {where} — {self.detail}"


# ─── Forbidden strings ───────────────────────────────────────────────────────

ERROR_TOKENS = ("#N/A", "#REF!", "#VALUE!", "#DIV/0!", "#NAME?", "#NULL!", "#NUM!")
TEMPLATE_TOKENS = ("{{",)

# The library footnotes ship 'All figures in [x]$MM'; `fill_footnote_currency`
# substitutes the letter. v0.5.34 wired that in but not on every slide carrying
# the token, and a client-facing deck shipped a literal '[x]' (Phase A finding).
# The same token is also the library's generic "fill me in" marker — the contact
# slide's three spare banker blocks are nine '[x]' cells that no assembler fills.
CURRENCY_TOKEN = "[x]"
_CURRENCY_CONTEXT = "[x]$"

_PLACEHOLDER_RE = re.compile(r"\[[^\[\]]*placeholder[^\[\]]*\]", re.IGNORECASE)

# Placeholders that are legitimately left unfilled: the comps and precedents
# slides ship as placeholders on purpose (their workbook tabs carry un-evaluated
# CapIQ array formulas the analyst refreshes in Excel), and the earnings-update
# plan has no `financial-charts` stage so its pie placeholder is correct. Kept
# visible as advisory rather than suppressed.
EXPECTED_PLACEHOLDERS = {
    "[Placeholder for Comps Chart]",
    "[Placeholder for Precedents Chart]",
    "[Pie Chart Placeholder]",
    "[Placeholder for Institutional Ownership]",
}


# ─── Tolerances ──────────────────────────────────────────────────────────────
# Every tolerance is set from measurement against the shipped library and the
# frozen fixtures. Each comment records what it was set from.

# Excess over the library baseline before a shape's off-slide overhang is
# reported. The library's own footnote placeholder overhangs the right edge by
# 0.48" and its tombstone slides park shapes fully off-canvas.
_SLIDE_BOUNDS_TOL_IN = 0.05

# Excess over the library's shipped table height. The library ships the
# market-entry frame at 5.7197" and the assembler clamps to 5.710", so a correct
# deck lands *under* the reference; PRL18's risk table declared 5.360" against a
# 5.1715" library height (+0.19").
_TABLE_HEIGHT_TOL_IN = 0.02

# Excess rendered overflow over the library baseline. Observed: library footnote
# placeholders 0.07", library tombstone shapes up to 0.75" — all cancelling. Real
# excesses: PRL17 market-entry +0.08", the pitch fixture's own market-entry and
# risk tables ~+0.10", PRL14's overview bullets +2.5".
_OVERFLOW_TOL_IN = 0.05

# Claimed boxes are padded outward by this much so a neighbour's antialiased
# border is not read as the shape under test overflowing into it.
_CLAIM_PAD_IN = 0.02

# Vertical white gap bridged while growing the overflow region downward. Palatino
# body copy at 10.5 pt leaves ~0.06" of white between lines, so the growth has to
# bridge more than that to follow a spilling bullet block; the gap from a table's
# last border to the footnote below it is ~0.10", so this must stay under the
# smallest legitimate inter-shape gap. 0.12" bridges every text line on the
# fixtures and no inter-shape gap.
_OVERFLOW_GAP_IN = 0.12

# Grayscale below this is ink. Matches the threshold the retired render-parity
# test used (deleted in Phase D with the second backend it compared against), so a
# geometric finding here means the same thing as a parity measurement there.
_INK_MAX_LUMA = 250

# A probe row needs this fraction of its width in unclaimed ink (min 3 px) before
# it counts, so antialias speckle cannot manufacture an overflow.
_ROW_INK_FRACTION = 0.01
_ROW_INK_MIN_PX = 3

# A layout/master shape covering more than this fraction of the slide is a
# background, not a content box, and may not claim territory. Layout/master only:
# a slide shape always claims (the market-entry table alone covers 71% of its
# slide and must claim against its neighbours).
_BACKGROUND_AREA_FRACTION = 0.40

# Minimum shape-name overlap for a built slide to be matched to a library slide.
_MATCH_MIN_SCORE = 0.30

# Two content-bearing boxes must overlap by more than this in BOTH axes before the
# vision tier puts the slide on the review agenda — enough to skip shapes that
# merely abut or share a border.
_OVERLAP_MIN_IN = 0.05

# A vertical clearance must be this much tighter than the library's for the same
# shape pair before the vision tier mentions it. Raw clearance carries no signal:
# the INFOR layout seats a header bar directly on its table, so 17 of the pitch
# fixture's 18 slides have a sub-0.12" gap somewhere by design.
_CLEARANCE_TIGHTENED_IN = 0.02

RENDER_DPI = 150
LIBRARY_TEMPLATE_NAME = "INFOR Slide Library.pptx"


# ─── Shape traversal ─────────────────────────────────────────────────────────


def _iter_shapes(shapes):
    """Every shape, recursing into groups depth-first."""
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes)


def _shape_texts(slide):
    """(label, text) for every text-bearing shape and table cell on the slide."""
    for shape in _iter_shapes(slide.shapes):
        if getattr(shape, "has_text_frame", False) and shape.text_frame.text:
            yield shape.name, shape.text_frame.text
        if getattr(shape, "has_table", False):
            for r, row in enumerate(shape.table.rows):
                for c, cell in enumerate(row.cells):
                    if cell.text:
                        yield f"{shape.name}[{r},{c}]", cell.text


@dataclass(frozen=True)
class _Box:
    """A shape's declared extents, in inches."""

    left: float
    top: float
    right: float
    bottom: float
    name: str = ""

    @property
    def area(self) -> float:
        return max(0.0, self.right - self.left) * max(0.0, self.bottom - self.top)


def _box_of(shape) -> _Box | None:
    if shape.left is None or shape.top is None or shape.width is None or shape.height is None:
        return None
    left, top = Emu(shape.left).inches, Emu(shape.top).inches
    return _Box(
        left, top, left + Emu(shape.width).inches, top + Emu(shape.height).inches, shape.name
    )


def _has_content(shape) -> bool:
    """True when the shape puts ink of its own on the slide (text or a table)."""
    return (
        getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip()
    ) or getattr(shape, "has_table", False)


def _content_boxes(slide) -> list[_Box]:
    """Declared boxes of every content-bearing shape on the slide itself."""
    boxes = []
    for shape in slide.shapes:
        box = _box_of(shape)
        if box is not None and _has_content(shape):
            boxes.append(box)
    return boxes


def _growth_prone(slide):
    """Shapes whose rendered height depends on content: tables and filled text."""
    for shape in slide.shapes:
        if getattr(shape, "has_table", False):
            yield shape
        elif getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
            yield shape


# ─── Matching a built slide to its library slide ─────────────────────────────


def _signature(slide) -> set[str]:
    return {shape.name for shape in _iter_shapes(slide.shapes)}


def match_library_slide(slide, library_signatures: list[set[str]]) -> int | None:
    """Index of the library slide a built slide came from, or None.

    Scored by shape-name overlap (intersection over union). Robust to the pitch
    flow deleting and cloning slides, to filled titles, and to the market-entry
    section repeating one library slide — all of which defeat matching by index
    or by title text.
    """
    signature = _signature(slide)
    if not signature:
        return None
    best_index, best_score = None, 0.0
    for index, reference in enumerate(library_signatures):
        union = signature | reference
        if not union:
            continue
        score = len(signature & reference) / len(union)
        if score > best_score:
            best_index, best_score = index, score
    return best_index if best_score >= _MATCH_MIN_SCORE else None


def _concept_name(library_index: int | None) -> str:
    if library_index is None:
        return "unmatched slide"
    marker = LIBRARY_SLIDE_MARKERS.get(library_index)
    return marker.description if marker else f"library slide {library_index + 1}"


def default_library_path() -> Path | None:
    """The shipped slide library, resolved the way the skills resolve templates."""
    env = os.environ.get("CLAUDE_PLUGIN_ROOT")
    roots = [Path(env)] if env else []
    roots.append(Path(__file__).resolve().parent.parent)  # scripts/.. == infor-beta/
    for root in roots:
        candidate = root / "templates" / LIBRARY_TEMPLATE_NAME
        if candidate.is_file():
            return candidate
    return None


# ─── Measurement primitives (run over both the deck and the library) ─────────


def _measure_outside(prs) -> dict[int, dict[str, float]]:
    """Per slide, per shape: how far the declared box falls outside the slide."""
    width, height = Emu(prs.slide_width).inches, Emu(prs.slide_height).inches
    out: dict[int, dict[str, float]] = {}
    for index, slide in enumerate(prs.slides):
        per_shape: dict[str, float] = {}
        for shape in _iter_shapes(slide.shapes):
            box = _box_of(shape)
            if box is None:
                continue
            excess = max(-box.left, -box.top, box.right - width, box.bottom - height)
            if excess > 0:
                per_shape[shape.name] = max(per_shape.get(shape.name, 0.0), excess)
        out[index] = per_shape
    return out


def _clearance_pairs(slide) -> list[tuple[str, str, float]]:
    """(shape, shape_below, clearance_in) for growth-prone shapes and what sits under them.

    Only pairs that overlap horizontally and are within `_OVERFLOW_GAP_IN`
    vertically — i.e. close enough that a re-grown row or a wrapped line would
    land on the lower shape.
    """
    below = _content_boxes(slide)
    pairs: list[tuple[str, str, float]] = []
    for shape in _growth_prone(slide):
        box = _box_of(shape)
        if box is None:
            continue
        for other in below:
            if other.name == box.name or other.top < box.bottom:
                continue
            if other.right <= box.left or other.left >= box.right:
                continue
            clearance = other.top - box.bottom
            if clearance <= _OVERFLOW_GAP_IN:
                pairs.append((box.name, other.name, clearance))
    return pairs


def _measure_clearances(prs) -> dict[int, dict[tuple[str, str], float]]:
    """Per slide, the tightest clearance recorded for each (shape, shape_below) pair."""
    out: dict[int, dict[tuple[str, str], float]] = {}
    for index, slide in enumerate(prs.slides):
        per_pair: dict[tuple[str, str], float] = {}
        for shape, other, clearance in _clearance_pairs(slide):
            key = (shape, other)
            per_pair[key] = min(per_pair.get(key, clearance), clearance)
        out[index] = per_pair
    return out


def _table_heights(prs) -> dict[int, float]:
    """Per slide, the tallest table's height in inches (frame vs. row sum)."""
    out: dict[int, float] = {}
    for index, slide in enumerate(prs.slides):
        tallest = 0.0
        for shape in slide.shapes:
            if not getattr(shape, "has_table", False):
                continue
            frame = Emu(shape.height).inches
            rows = sum(Emu(row.height).inches for row in shape.table.rows)
            tallest = max(tallest, frame, rows)
        if tallest:
            out[index] = tallest
    return out


def _claim_boxes(slide, exclude_name: str, slide_area: float) -> list[_Box]:
    """Boxes that legitimately own territory, for masking the overflow probe.

    Slide shapes always claim. Layout and master shapes claim only when they are
    not backgrounds — a full-slide decorative rectangle would otherwise mask the
    page and make every overflow invisible.
    """
    boxes: list[_Box] = []
    for shape in _iter_shapes(slide.shapes):
        box = _box_of(shape)
        if box is not None and box.name != exclude_name:
            boxes.append(box)
    for source in (slide.slide_layout, slide.slide_layout.slide_master):
        for shape in _iter_shapes(source.shapes):
            box = _box_of(shape)
            if box is None:
                continue
            if slide_area > 0 and box.area / slide_area > _BACKGROUND_AREA_FRACTION:
                continue
            boxes.append(box)
    return boxes


def _overflow_depth_in(ink, box: _Box, claims: list[_Box], page_w: int, page_h: int, dpi: float):
    """Depth (inches) that unclaimed rendered ink extends below ``box``'s bottom.

    Walks down from the declared bottom row by row. A row counts when it holds
    enough ink that no other shape's box accounts for, and the walk bridges up to
    ``_OVERFLOW_GAP_IN`` of blank rows so it follows a spilling bullet block
    across its inter-line white space without stepping across the gap to the next
    shape below.
    """

    def px_x(v: float) -> int:
        return max(0, min(page_w, int(round(v * dpi))))

    def px_y(v: float) -> int:
        return max(0, min(page_h, int(round(v * dpi))))

    x0, x1 = px_x(box.left), px_x(box.right)
    y_start = px_y(box.bottom)
    if x1 - x0 < 2 or page_h - y_start < 2:
        return 0.0

    band = ink[y_start:page_h, x0:x1].copy()
    for claim in claims:
        cy0 = max(0, px_y(claim.top - _CLAIM_PAD_IN) - y_start)
        cy1 = min(band.shape[0], px_y(claim.bottom + _CLAIM_PAD_IN) - y_start)
        cx0 = max(0, px_x(claim.left - _CLAIM_PAD_IN) - x0)
        cx1 = min(band.shape[1], px_x(claim.right + _CLAIM_PAD_IN) - x0)
        if cy1 > cy0 and cx1 > cx0:
            band[cy0:cy1, cx0:cx1] = False

    per_row = band.sum(axis=1)
    threshold = max(_ROW_INK_MIN_PX, int(round((x1 - x0) * _ROW_INK_FRACTION)))
    gap_rows = int(round(_OVERFLOW_GAP_IN * dpi))
    deepest, blank = -1, 0
    for row, count in enumerate(per_row):
        if count >= threshold:
            deepest, blank = row, 0
        else:
            blank += 1
            if blank > gap_rows:
                break
    return 0.0 if deepest < 0 else (deepest + 1) / dpi


def _measure_overflow(prs, renders: dict[int, Path]) -> dict[int, dict[str, float]]:
    """Per slide, per growth-prone shape: rendered overflow depth in inches."""
    import numpy as np
    from PIL import Image

    slide_w, slide_h = Emu(prs.slide_width).inches, Emu(prs.slide_height).inches
    slide_area = slide_w * slide_h
    out: dict[int, dict[str, float]] = {}
    for index, png in sorted(renders.items()):
        slide = prs.slides[index]
        image = Image.open(png).convert("L")
        page_w, page_h = image.size
        dpi = page_h / slide_h if slide_h else RENDER_DPI
        ink = np.array(image) < _INK_MAX_LUMA

        per_shape: dict[str, float] = {}
        for shape in _growth_prone(slide):
            box = _box_of(shape)
            if box is None or box.bottom >= slide_h:
                continue
            claims = _claim_boxes(slide, shape.name, slide_area)
            depth = _overflow_depth_in(ink, box, claims, page_w, page_h, dpi)
            if depth > 0:
                per_shape[shape.name] = max(per_shape.get(shape.name, 0.0), depth)
        out[index] = per_shape
    return out


def _overrun_slide_shapes(slide, box: _Box, reach: float) -> list[str]:
    """Slide shapes whose declared box the overflowing content runs into.

    Restricted to shapes on the slide itself that carry content, so the message
    names 'the LTM Revenue Breakdown band' rather than the layout's decorative
    rules. ``reach`` is the *unmasked* ink extent (see `_measure_overflow`): the
    masked depth stops at the first claimed box, which is precisely the box we
    want to name.
    """
    hit: list[str] = []
    for shape in _iter_shapes(slide.shapes):
        other = _box_of(shape)
        if other is None or other.name == box.name or not _has_content(shape):
            continue
        if (
            other.top < reach
            and other.bottom > box.bottom
            and other.right > box.left
            and other.left < box.right
        ):
            hit.append(other.name)
    return sorted(set(hit))


# ─── Per-shape render attribution ────────────────────────────────────────────
# `_measure_overflow` counts ink that lands in NOBODY's declared box, which is
# what keeps it free of false positives — and is also its blind spot: an overflow
# running straight into the shape below is masked by that shape's own claim.
# Pixels in that region belong to two shapes at once and no threshold can
# separate them.
#
# Attribution separates them by rendering the shape alone. For each masked
# candidate the probe deck carries two extra slides: one holding only that shape
# on its own layout, and one holding nothing at all. Subtracting the second from
# the first leaves exactly the shape's own ink — every neighbour, and all
# layout/master decoration, gone. Absolute positioning is what makes this sound:
# removing a neighbour cannot move or re-wrap the shape under test.
#
# Every probe goes into ONE deck, so the whole pass costs one extra LibreOffice
# conversion (measured 2.5-7.7 s per artefact).


def _masking_neighbours(slide) -> dict[str, list[str]]:
    """Growth-prone shape name -> content shapes that mask the territory below it.

    A neighbour masks when it overlaps horizontally, extends below the shape's
    declared bottom, and starts no more than `_OVERFLOW_GAP_IN` below it (a
    negative gap — an overlapping box — masks from the first row down). Those are
    exactly the shapes whose overflow `_measure_overflow` cannot see.
    """
    below = _content_boxes(slide)
    masked: dict[str, list[str]] = {}
    for shape in _growth_prone(slide):
        box = _box_of(shape)
        if box is None:
            continue
        for other in below:
            if other.name == box.name or other.right <= box.left or other.left >= box.right:
                continue
            if other.bottom > box.bottom and other.top <= box.bottom + _OVERFLOW_GAP_IN:
                masked.setdefault(shape.name, []).append(other.name)
    return masked


# Relationship-bearing XML dropped from a probed shape. A probe slide is built by
# deep-copying shape XML onto a fresh slide (the technique `clone_slide_after`
# already ships), which does NOT carry slide-level relationships — so a surviving
# r:id would dangle and could fail the render. None of this affects the
# measurement: a shape's own fill and hyperlink styling stay inside its box, and
# the probe only ever measures ink BELOW the box.
_PROBE_STRIP_TAGS = ("a:hlinkClick", "a:hlinkHover", "a:blipFill")
_PROBE_STRIP_ATTRS = ("r:embed", "r:link", "r:id")


def _strip_relationships(element):
    for tag in _PROBE_STRIP_TAGS:
        for node in element.findall(f".//{qn(tag)}"):
            node.getparent().remove(node)
    for node in element.iter():
        for attr in _PROBE_STRIP_ATTRS:
            node.attrib.pop(qn(attr), None)
    return element


@dataclass(frozen=True)
class ProbeRequest:
    """One shape to render in isolation, optionally altered first.

    ``key`` is whatever the caller wants the measurement filed under — the
    attribution pass uses ``(slide, shape)``; the repair loop's fit ladders use
    ``(slide, shape, candidate_size)`` so one conversion measures a whole ladder.
    ``mutate`` receives the *copied* shape on the probe slide, so it can alter the
    probe without touching the deck.

    ``bottom`` overrides the y (inches) below which ink counts as overflow. Needed
    when the mutation resizes the shape — a repair probe that re-clamps a table to
    a smaller height must still be measured against the height it is aiming for,
    not against the height the deck currently declares.
    """

    slide: int
    shape: str
    key: object = None
    mutate: object = None  # Callable[[Shape], None] | None
    bottom: float | None = None

    def filed_under(self):
        return self.key if self.key is not None else (self.slide, self.shape)


@dataclass
class ProbePlan:
    """Where each probe landed in the probe deck."""

    path: Path
    shape: dict[object, int] = field(default_factory=dict)     # key -> probe slide
    layout: dict[int, int] = field(default_factory=dict)       # deck slide -> empty probe
    box: dict[object, "_Box"] = field(default_factory=dict)    # key -> declared box
    slide: dict[object, int] = field(default_factory=dict)     # key -> deck slide

    @property
    def indices(self) -> list[int]:
        return sorted({*self.shape.values(), *self.layout.values()})


def _append_probe_slide(prs, source_slide, keep_name: str | None, mutate=None) -> int:
    """Append a slide carrying only `keep_name` from `source_slide` (or nothing)."""
    new = prs.slides.add_slide(source_slide.slide_layout)
    for shape in list(new.shapes):  # drop the placeholders add_slide seeds
        shape._element.getparent().remove(shape._element)
    if keep_name is not None:
        for shape in source_slide.shapes:
            if shape.name == keep_name:
                new.shapes._spTree.append(_strip_relationships(deepcopy(shape._element)))
                break
        if mutate is not None:
            for shape in new.shapes:
                if shape.name == keep_name:
                    mutate(shape)
                    break
    return len(prs.slides._sldIdLst) - 1


def build_probe_deck(deck: Path, out_path: Path, requests) -> ProbePlan | None:
    """Write one deck holding every requested isolation probe; None when empty.

    Each request contributes one slide holding just that shape, and each distinct
    layout contributes one empty slide so the layout's own ink can be subtracted.
    Everything rides in a single deck because the LibreOffice conversion — not the
    page count — is what costs.
    """
    requests = list(requests)
    if not requests:
        return None
    prs = Presentation(deck)
    originals = list(prs.slides)
    plan = ProbePlan(path=out_path)
    layout_probe: dict[int, int] = {}
    for request in requests:
        if request.slide >= len(originals):
            continue
        slide = originals[request.slide]
        shape = next((s for s in slide.shapes if s.name == request.shape), None)
        box = _box_of(shape) if shape is not None else None
        if box is None:
            continue
        layout_key = id(slide.slide_layout)
        if layout_key not in layout_probe:
            layout_probe[layout_key] = _append_probe_slide(prs, slide, None)
        plan.layout[request.slide] = layout_probe[layout_key]
        key = request.filed_under()
        plan.shape[key] = _append_probe_slide(prs, slide, request.shape, request.mutate)
        plan.box[key] = box if request.bottom is None else replace(box, bottom=request.bottom)
        plan.slide[key] = request.slide
    if not plan.shape:
        return None
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    return plan


def measure_probe_overflow(deck: Path | str, requests, work_dir: Path | str) -> dict:
    """Render each requested shape alone and report ink below its declared bottom.

    Returns ``{request.filed_under(): depth_in}``, omitting probes that measured
    nothing. Raises RuntimeError when the probe deck cannot be rendered, so a
    caller reports the tier unavailable rather than reporting a deck clean.

    This is the one measurement primitive behind both `masked-overflow` and the
    repair loop's fit ladders. It is deliberately estimate-free: no character
    widths, no line-height constants, no wrap model — LibreOffice lays the text
    out and the pixels are read back.
    """
    import numpy as np
    from PIL import Image

    deck = Path(deck)
    work = Path(work_dir)
    plan = build_probe_deck(deck, work / "probe.pptx", requests)
    if plan is None:
        return {}

    renders = _render_slides(plan.path, work / "png", plan.indices)
    ink = {
        index: np.array(Image.open(png).convert("L")) < _INK_MAX_LUMA
        for index, png in renders.items()
    }
    slide_h = Emu(Presentation(deck).slide_height).inches

    out: dict = {}
    for key, probe in plan.shape.items():
        if probe not in ink:
            continue
        box = plan.box[key]
        if box.bottom >= slide_h:
            continue
        layout = plan.layout.get(plan.slide[key])
        own = ink[probe] & ~ink[layout] if layout is not None and layout in ink else ink[probe]
        page_h, page_w = own.shape
        dpi = page_h / slide_h if slide_h else RENDER_DPI
        depth = _own_ink_depth_in(own, box, page_w, page_h, dpi)
        if depth > 0:
            out[key] = depth
    return out


def _own_ink_depth_in(own, box: _Box, page_w: int, page_h: int, dpi: float) -> float:
    """Depth (inches) the shape's own ink extends below its declared bottom.

    No claim masking and no gap bridging: `own` holds nothing but this shape's
    ink, so the deepest inked row IS the answer — a blank band cannot mean "the
    next shape starts here".
    """
    x0 = max(0, min(page_w, int(round(box.left * dpi))))
    x1 = max(0, min(page_w, int(round(box.right * dpi))))
    y_start = max(0, min(page_h, int(round(box.bottom * dpi))))
    if x1 - x0 < 2 or page_h - y_start < 2:
        return 0.0
    per_row = own[y_start:page_h, x0:x1].sum(axis=1)
    threshold = max(_ROW_INK_MIN_PX, int(round((x1 - x0) * _ROW_INK_FRACTION)))
    inked = [row for row, count in enumerate(per_row) if count >= threshold]
    return (inked[-1] + 1) / dpi if inked else 0.0


def bake_stored_autofit(shape) -> None:
    """Render probe mutation: show the shape at the size PowerPoint will draw it.

    LibreOffice treats any `<a:normAutofit>` as "shrink to fit" and recomputes its
    own scale, so an autofit shape's rendered overflow always measures zero — the
    renderer squeezes the text in no matter how much there is. **PowerPoint does
    not**: it applies the stored `fontScale` and nothing more, which is precisely
    the v0.5.23 defect (a scale-less autofit renders full size and spills).

    So a probe of an autofit shape bakes the stored scale into the run sizes and
    replaces the autofit with `<a:noAutofit/>`. The result is a shape whose
    rendered height LibreOffice reports honestly — and reports ~one line taller
    than PowerPoint, keeping it the conservative oracle.

    `lnSpcReduction` is deliberately not modelled: dropping it renders the text
    slightly taller than it will ship, which is the safe direction.
    """
    if not getattr(shape, "has_text_frame", False):
        return  # a table has no autofit to bake
    scale = normal_autofit_scale(shape)
    if scale is None:
        return
    if scale != 100.0:
        apply_text_scale(shape, scale)
    strip_autofit(shape)


def _attribution_requests(prs) -> list[ProbeRequest]:
    """Every shape whose overflow the unclaimed-ink measurement cannot see.

    Two reasons it cannot, and both need the same isolation probe:

      - **A neighbour masks it.** The shape below claims the territory the
        overflow lands in (`_masking_neighbours`).
      - **An autofit hides it.** The shape carries `<a:normAutofit>`, so
        LibreOffice shrinks the text to fit and renders no overflow at all, while
        PowerPoint applies only the stored scale. Every INFOR overview bullet
        block is in this class, which means without this the contract was blind to
        the PRL14 defect on any freshly built deck.
    """
    requests: list[ProbeRequest] = []
    for index, slide in enumerate(prs.slides):
        masked = _masking_neighbours(slide)
        for shape in _growth_prone(slide):
            autofit = (
                normal_autofit_scale(shape) is not None
                if getattr(shape, "has_text_frame", False)
                else False
            )
            if shape.name in masked or autofit:
                requests.append(
                    ProbeRequest(index, shape.name, mutate=bake_stored_autofit)
                )
    return requests


def measure_attributed_overflow(deck: Path | str, work_dir: Path | str) -> dict[int, dict[str, float]]:
    """Per slide, per shape: how far the shape's OWN ink runs below its box.

    Raises RuntimeError when the probe deck cannot be rendered, so the caller
    reports the tier as unavailable rather than reporting a deck clean.
    """
    deck = Path(deck)
    requests = _attribution_requests(Presentation(deck))
    measured = measure_probe_overflow(deck, requests, work_dir)
    out: dict[int, dict[str, float]] = {}
    for (index, name), depth in measured.items():
        out.setdefault(index, {})[name] = depth
    return out


# ─── Library baseline ────────────────────────────────────────────────────────


@dataclass
class LibraryBaseline:
    """What the blank library does, so only what the fill did is attributed to it."""

    signatures: list[set[str]] = field(default_factory=list)
    table_height: dict[int, float] = field(default_factory=dict)
    overflow: dict[int, dict[str, float]] = field(default_factory=dict)
    attributed: dict[int, dict[str, float]] = field(default_factory=dict)
    outside: dict[int, dict[str, float]] = field(default_factory=dict)
    clearance: dict[int, dict[tuple[str, str], float]] = field(default_factory=dict)
    rendered: bool = False
    attributed_ok: bool = False

    def clearances_for(self, slide) -> dict[tuple[str, str], float]:
        """The matching library slide's clearance map for a built slide."""
        matched = match_library_slide(slide, self.signatures)
        return self.clearance.get(matched, {}) if matched is not None else {}


_BASELINE_CACHE: dict[tuple[str, int, bool, bool], LibraryBaseline] = {}


def library_baseline(
    library_path: Path, *, render: bool = True, attribute: bool = True
) -> LibraryBaseline:
    """Measure the blank library. Cached per (path, mtime, render, attribute)."""
    library_path = Path(library_path)
    key = (str(library_path.resolve()), int(library_path.stat().st_mtime), render, attribute)
    cached = _BASELINE_CACHE.get(key)
    if cached is not None:
        return cached

    prs = Presentation(library_path)
    baseline = LibraryBaseline(
        signatures=[_signature(slide) for slide in prs.slides],
        table_height=_table_heights(prs),
        outside=_measure_outside(prs),
        clearance=_measure_clearances(prs),
    )
    if render:
        tmp = Path(tempfile.mkdtemp(prefix="deck-contract-library-"))
        try:
            renders = _render_slides(library_path, tmp / "slides", list(range(len(prs.slides))))
            baseline.overflow = _measure_overflow(prs, renders)
            baseline.rendered = True
        except RuntimeError:
            baseline.rendered = False  # caller reports render-unavailable
        if attribute and baseline.rendered:
            try:
                baseline.attributed = measure_attributed_overflow(
                    library_path, tmp / "attribution"
                )
                baseline.attributed_ok = True
            except RuntimeError:
                baseline.attributed_ok = False
    _BASELINE_CACHE[key] = baseline
    return baseline


# ─── Deterministic checks ────────────────────────────────────────────────────


def _check_forbidden_strings(prs) -> list[Finding]:
    findings: list[Finding] = []
    for index, slide in enumerate(prs.slides):
        for name, text in _shape_texts(slide):
            for token in ERROR_TOKENS:
                if token in text:
                    findings.append(
                        Finding(
                            "forbidden-string",
                            SEVERITY_BLOCKING,
                            index,
                            f"contains the spreadsheet error value {token!r} — a formula "
                            f"did not resolve before the value was written",
                            shape=name,
                        )
                    )
            for token in TEMPLATE_TOKENS:
                if token in text:
                    findings.append(
                        Finding(
                            "forbidden-string",
                            SEVERITY_BLOCKING,
                            index,
                            f"contains the unrendered template token {token!r}",
                            shape=name,
                        )
                    )
            if CURRENCY_TOKEN in text:
                currency = _CURRENCY_CONTEXT in text
                findings.append(
                    Finding(
                        "unsubstituted-currency-token" if currency else "unfilled-token",
                        SEVERITY_BLOCKING,
                        index,
                        (
                            f"carries the literal currency token {_CURRENCY_CONTEXT!r} — "
                            f"`pptx_helpers.fill_footnote_currency` was not called for this "
                            f"shape, so the deck ships '[x]$MM' to the client"
                            if currency
                            else f"reads the bare library fill-me-in token "
                            f"{CURRENCY_TOKEN!r}, which renders verbatim on the slide"
                        ),
                        shape=name,
                    )
                )
            for match in _PLACEHOLDER_RE.findall(text):
                expected = match in EXPECTED_PLACEHOLDERS
                findings.append(
                    Finding(
                        "expected-placeholder" if expected else "unfilled-placeholder",
                        SEVERITY_ADVISORY if expected else SEVERITY_BLOCKING,
                        index,
                        (
                            f"still reads {match!r} (known-intentional — nothing is "
                            f"rendered into it at build time)"
                            if expected
                            else f"still reads {match!r} — the assembler left a library "
                            f"placeholder unfilled"
                        ),
                        shape=name,
                    )
                )
    return findings


def _check_shapes_within_slide(prs, baseline: LibraryBaseline) -> list[Finding]:
    """A shape may not hang further off the slide than the library's does."""
    width, height = Emu(prs.slide_width).inches, Emu(prs.slide_height).inches
    measured = _measure_outside(prs)
    findings: list[Finding] = []
    for index, slide in enumerate(prs.slides):
        matched = match_library_slide(slide, baseline.signatures)
        reference = baseline.outside.get(matched, {}) if matched is not None else {}
        for name, excess in sorted(measured.get(index, {}).items()):
            allowed = reference.get(name, 0.0)
            if excess - allowed > _SLIDE_BOUNDS_TOL_IN:
                findings.append(
                    Finding(
                        "shape-outside-slide",
                        SEVERITY_BLOCKING,
                        index,
                        f"declared box extends {excess:.2f}\" outside the "
                        f"{width:.0f}x{height:.0f}\" slide; the library's {_concept_name(matched)} "
                        f"slide allows {allowed:.2f}\"",
                        shape=name,
                        measured_in=excess,
                        limit_in=allowed + _SLIDE_BOUNDS_TOL_IN,
                    )
                )
    return findings


def _check_table_heights(prs, baseline: LibraryBaseline) -> list[Finding]:
    """A filled table must not declare more height than the library ships.

    Catches PRL18's Considerations/Mitigants table: long mitigant copy grew the
    declared rows to 5.360" against the library's 5.1715". Does NOT catch PRL17's
    market-entry table, which declared 5.710" and rendered 5.91" — that is
    `rendered-overflow`'s job, and the split is the point: one defect is in the
    file, the other only exists once something lays the text out.
    """
    findings: list[Finding] = []
    for index, slide in enumerate(prs.slides):
        matched = match_library_slide(slide, baseline.signatures)
        limit = baseline.table_height.get(matched) if matched is not None else None
        if limit is None:
            continue
        for shape in slide.shapes:
            if not getattr(shape, "has_table", False):
                continue
            frame = Emu(shape.height).inches
            rows = sum(Emu(row.height).inches for row in shape.table.rows)
            measured = max(frame, rows)
            if measured - limit > _TABLE_HEIGHT_TOL_IN:
                which = "frame is" if frame >= rows else "row heights sum to"
                findings.append(
                    Finding(
                        "table-taller-than-library",
                        SEVERITY_BLOCKING,
                        index,
                        f"{which} {measured:.3f}\" but the library ships the "
                        f"{_concept_name(matched)} table at {limit:.3f}\" — filled content "
                        f"grew the declared rows, so the table runs past where it was "
                        f"designed to end",
                        shape=shape.name,
                        measured_in=measured,
                        limit_in=limit,
                    )
                )
    return findings


def _check_rendered_overflow(
    prs, renders: dict[int, Path], baseline: LibraryBaseline
) -> list[Finding]:
    """Rendered ink below a shape's declared bottom, over the library's own.

    **Autofit shapes are excluded and handled by `masked-overflow` instead.** This
    check measures the deck as LibreOffice draws it, and LibreOffice honours any
    `<a:normAutofit>` by inventing its own shrink — so what it reports for such a
    shape is neither what PowerPoint will draw nor a bound on it. On the
    earnings-update fixture's Business Updates block the two disagree by an order
    of magnitude: 0.14" measured here (LibreOffice's squeeze residual) against
    1.51" once the stored scale is baked in and the autofit removed. Reporting the
    smaller number would understate a real defect and mislead the repair loop
    about how far it has to move.
    """
    import numpy as np
    from PIL import Image

    slide_h = Emu(prs.slide_height).inches
    measured = _measure_overflow(prs, renders)
    findings: list[Finding] = []
    for index, slide in enumerate(prs.slides):
        matched = match_library_slide(slide, baseline.signatures)
        reference = baseline.overflow.get(matched, {}) if matched is not None else {}
        autofit_shapes = {
            shape.name
            for shape in _growth_prone(slide)
            if getattr(shape, "has_text_frame", False)
            and normal_autofit_scale(shape) is not None
        }
        firing = {
            name: depth
            for name, depth in sorted(measured.get(index, {}).items())
            if name not in autofit_shapes
            and depth - reference.get(name, 0.0) > _OVERFLOW_TOL_IN
        }
        if not firing:
            continue
        image = Image.open(renders[index]).convert("L")
        page_w, page_h = image.size
        dpi = page_h / slide_h if slide_h else RENDER_DPI
        ink = np.array(image) < _INK_MAX_LUMA

        for name, depth in firing.items():
            allowed = reference.get(name, 0.0)
            excess = depth - allowed
            shape = next((s for s in slide.shapes if s.name == name), None)
            box = _box_of(shape) if shape is not None else None
            if box is None:
                continue
            # Unmasked extent: how far the content actually runs, for naming what
            # it runs into (the masked depth stops at that very box).
            raw = _overflow_depth_in(ink, box, [], page_w, page_h, dpi)
            overrun = _overrun_slide_shapes(slide, box, box.bottom + max(depth, raw))
            into = (
                f" and continues into the declared box of {', '.join(overrun)}"
                if overrun
                else ""
            )
            baseline_note = (
                f" (the library's {_concept_name(matched)} slide renders {allowed:.2f}\" past"
                f" the same box, so {excess:.2f}\" of this is the fill's)"
                if allowed
                else ""
            )
            findings.append(
                Finding(
                    "rendered-overflow",
                    SEVERITY_BLOCKING,
                    index,
                    f"renders {depth:.2f}\" of content below its declared bottom edge "
                    f"({box.bottom:.2f}\"){into}{baseline_note} — the declared box is only "
                    f"a minimum, and the layout engine grew past it",
                    shape=name,
                    measured_in=depth,
                    limit_in=allowed + _OVERFLOW_TOL_IN,
                )
            )
    return findings


def _check_masked_overflow(
    prs, attributed: dict[int, dict[str, float]], baseline: LibraryBaseline, already: set
) -> list[Finding]:
    """Attributed overflow for shapes `rendered-overflow` structurally cannot see.

    This is the check that closes the documented blind spot. Its measurement is
    the shape's OWN ink (see `measure_attributed_overflow`), so the neighbour that
    was masking the overflow no longer hides it — and the finding names that
    neighbour, which is the actionable part: a repair step has to know which of
    the two shapes to shrink.

    Shapes `rendered-overflow` already reported are skipped: both checks measure
    the same defect from different sides, and one finding per shape keeps the
    repair loop from double-counting. That is why PRL17's market-entry table —
    masked below AND overflowing into unclaimed space — reports once.
    """
    findings: list[Finding] = []
    for index, slide in enumerate(prs.slides):
        measured = attributed.get(index)
        if not measured:
            continue
        matched = match_library_slide(slide, baseline.signatures)
        reference = baseline.attributed.get(matched, {}) if matched is not None else {}
        neighbours = _masking_neighbours(slide)
        for name, depth in sorted(measured.items()):
            if (index, name) in already:
                continue
            allowed = reference.get(name, 0.0)
            excess = depth - allowed
            if excess <= _OVERFLOW_TOL_IN:
                continue
            box = _box_of(next((s for s in slide.shapes if s.name == name), None))
            if box is None:  # unreachable: attribution needs a box to measure from
                continue
            reached = [
                other
                for other in neighbours.get(name, [])
                if (b := _box_of(next((s for s in slide.shapes if s.name == other), None)))
                is not None
                and b.top < box.bottom + depth
            ]
            into = (
                f", landing inside the declared box of {', '.join(sorted(set(reached)))}"
                if reached
                else ""
            )
            shape_obj = next((s for s in slide.shapes if s.name == name), None)
            autofit = (
                normal_autofit_scale(shape_obj)
                if shape_obj is not None and getattr(shape_obj, "has_text_frame", False)
                else None
            )
            why = (
                f"measured with the shape's {autofit:.0f}% autofit baked into its run "
                f"sizes and the autofit removed — LibreOffice would otherwise squeeze "
                f"the text to fit and show no overflow, where PowerPoint applies only "
                f"the stored scale"
                if autofit is not None
                else "measured by rendering this shape alone, because the shape(s) below "
                "claim that territory and pixel counting cannot separate the two "
                "shapes' ink"
            )
            baseline_note = (
                f" (the library's {_concept_name(matched)} slide renders {allowed:.2f}\" past"
                f" the same box, so {excess:.2f}\" of this is the fill's)"
                if allowed
                else ""
            )
            findings.append(
                Finding(
                    "masked-overflow",
                    SEVERITY_BLOCKING,
                    index,
                    f"renders {depth:.2f}\" of its own content below its declared bottom "
                    f"edge ({box.bottom:.2f}\"){into}{baseline_note} — {why}",
                    shape=name,
                    measured_in=depth,
                    limit_in=allowed + _OVERFLOW_TOL_IN,
                )
            )
    return findings


# ─── Rendering ───────────────────────────────────────────────────────────────


def _render_slides(deck: Path, out_dir: Path, indices: list[int]) -> dict[int, Path]:
    """Render slides through `slide_render`; raises RuntimeError without LibreOffice."""
    from slide_render import render_deck_to_png

    paths = render_deck_to_png(deck, out_dir, slide_indices=indices, dpi=RENDER_DPI)
    return dict(zip(indices, paths))


# ─── Vision tier (advisory, agent-inspected) ─────────────────────────────────
# No OCR. Overlap, collision, unreadable contrast and chart-label pileup need a
# model looking at pixels; this tier's job is to put the right pixels in front of
# one and say what to look for. Nothing here ever blocks, and pytest is not gated
# on it.


@dataclass(frozen=True)
class VisionTarget:
    """One thing for the checkpoint agent to look at, and what to look for."""

    slide: int
    question: str
    render: Path | None = None
    shape: str | None = None
    crop: Path | None = None


@dataclass
class VisionPass:
    """Rendered evidence plus the review agenda for the `deck` checkpoint.

    ``review_images`` and ``picture_crops`` are written whether or not anything
    adjudicates them, because the checkpoint shows them to the analyst either
    way. ``targets`` is the machine-readable agenda; ``findings`` is the same
    agenda collapsed to one advisory `vision-review` per slide, so
    `verify_deck`'s return value stays readable.
    """

    findings: list[Finding] = field(default_factory=list)
    targets: list[VisionTarget] = field(default_factory=list)
    review_images: dict[int, Path] = field(default_factory=dict)
    picture_crops: list[tuple[int, str, Path]] = field(default_factory=list)


def _picture_shapes(slide):
    for shape in _iter_shapes(slide.shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            yield shape


def _write_picture_crops(prs, out_dir: Path) -> list[tuple[int, str, Path]]:
    """Extract each embedded picture blob at native resolution.

    The blob beats the slide render for reviewing a pasted range or a chart: the
    cap-table picture is placed at ~4.5x5.4" but carries far more pixels than
    150 dpi of slide gives it, and legibility of 8-9 pt figures depends on that.
    """
    crops: list[tuple[int, str, Path]] = []
    out_dir.mkdir(parents=True, exist_ok=True)
    for index, slide in enumerate(prs.slides):
        for shape in _picture_shapes(slide):
            try:
                image = shape.image
            except (ValueError, AttributeError):
                continue  # linked / OLE picture with no embedded blob
            safe = re.sub(r"[^A-Za-z0-9]+", "_", shape.name).strip("_")
            path = out_dir / f"slide{index + 1:02d}_{safe}.{image.ext}"
            path.write_bytes(image.blob)
            crops.append((index, shape.name, path))
    return crops


def _content_box_overlaps(slide) -> list[tuple[str, str]]:
    """Pairs of content-bearing shapes whose declared boxes overlap.

    A deterministic *hint* for the vision tier, not a finding: overlapping boxes
    are routine and legitimate in PowerPoint (a label over a filled band), so
    this only tells the reviewing agent which slides are worth a close look.
    """
    boxes = _content_boxes(slide)
    pairs: list[tuple[str, str]] = []
    for i, a in enumerate(boxes):
        for b in boxes[i + 1 :]:
            if min(a.right, b.right) - max(a.left, b.left) > _OVERLAP_MIN_IN and (
                min(a.bottom, b.bottom) - max(a.top, b.top) > _OVERLAP_MIN_IN
            ):
                pairs.append((a.name, b.name))
    return pairs


def _tight_clearance(slide, reference: dict[tuple[str, str], float]) -> list[tuple[str, str, float]]:
    """Growth-prone shapes that sit closer to the content below than the library's do.

    (shape, shape_below, clearance_in), reported only when the fill has *tightened*
    the gap relative to the same pair on the matching library slide — a pair the
    library does not have at all counts as tightened.

    Raw clearance is not a signal: the INFOR layout puts a header bar directly
    above its table by design, so 17 of the pitch fixture's 18 slides have a
    sub-0.12" gap somewhere and flagging them all says nothing. Only a gap the
    fill narrowed is interesting.

    Kept as a *hint* now that `masked-overflow` measures the same region properly:
    attribution answers "does this shape's own ink run past its box", while this
    answers "did the fill move two shapes closer together than the library
    designed" — which is worth a look even when neither shape overflows at all.
    It names the pair and never asserts a collision.
    """
    tight: list[tuple[str, str, float]] = []
    for shape, other, clearance in _clearance_pairs(slide):
        allowed = reference.get((shape, other))
        if allowed is None or clearance < allowed - _CLEARANCE_TIGHTENED_IN:
            tight.append((shape, other, clearance))
    return sorted(tight, key=lambda t: t[2])


_VISION_CHECKLIST = (
    "text drawn over other text, shapes colliding, text too faint against its "
    "background, and chart/table labels piling up"
)


def vision_pass(
    deck: Path | str,
    *,
    out_dir: Path | str | None = None,
    renders: dict[int, Path] | None = None,
    baseline: LibraryBaseline | None = None,
) -> VisionPass:
    """Render the deck, extract its pictures, and build the review agenda.

    Returns advisory findings only. There is no OCR and no automated verdict: the
    checkpoint agent (or the analyst) is the reader. A slide earns an agenda entry
    when it carries a picture whose content no string scan can reach, when two
    content-bearing boxes overlap, or when the fill has tightened a vertical
    clearance the library left wider; every slide's render is written either way so
    the reviewer can page through them.

    ``baseline`` supplies the library reference for the clearance hint. It is
    measured (without rendering) when not supplied, and the hint is skipped
    entirely if the library cannot be resolved.
    """
    deck = Path(deck)
    prs = Presentation(deck)
    if baseline is None:
        library = default_library_path()
        baseline = library_baseline(library, render=False) if library else LibraryBaseline()
    root = Path(out_dir) if out_dir else Path(tempfile.mkdtemp(prefix="deck-contract-"))
    root.mkdir(parents=True, exist_ok=True)

    result = VisionPass()
    result.picture_crops = _write_picture_crops(prs, root / "pictures")

    if renders is None:
        try:
            renders = _render_slides(deck, root / "slides", list(range(len(prs.slides))))
        except RuntimeError as exc:
            renders = {}
            result.findings.append(
                Finding(
                    "render-unavailable",
                    SEVERITY_ADVISORY,
                    0,
                    f"the vision tier could not render the deck, so there is nothing "
                    f"for the checkpoint to review ({exc})",
                )
            )
    result.review_images = renders

    crops_by_slide: dict[int, list[tuple[str, Path]]] = {}
    for index, name, path in result.picture_crops:
        crops_by_slide.setdefault(index, []).append((name, path))

    for index, slide in enumerate(prs.slides):
        render = result.review_images.get(index)
        reasons: list[str] = []

        for name, path in crops_by_slide.get(index, []):
            result.targets.append(
                VisionTarget(
                    index,
                    "read what this picture actually renders — it is a rasterised "
                    "range or chart, so no string scan can reach its content; check "
                    "its labels for pileup and its figures for legibility",
                    render=render,
                    shape=name,
                    crop=path,
                )
            )
        pictures = [name for name, _ in crops_by_slide.get(index, [])]
        if pictures:
            reasons.append(f"rasterised picture(s) {', '.join(pictures)}")

        overlaps = _content_box_overlaps(slide)
        if overlaps:
            shown = "; ".join(f"{a} / {b}" for a, b in overlaps[:3])
            reasons.append(f"{len(overlaps)} overlapping content-box pair(s) ({shown})")
            result.targets.append(
                VisionTarget(
                    index,
                    f"these declared boxes overlap — confirm nothing is drawn on top "
                    f"of readable content: {shown}",
                    render=render,
                )
            )

        tight = _tight_clearance(slide, baseline.clearances_for(slide))
        if tight:
            shown = "; ".join(f"{a} over {b} by {gap:.2f}\"" for a, b, gap in tight[:3])
            reasons.append(f"{len(tight)} tight vertical clearance(s) ({shown})")
            result.targets.append(
                VisionTarget(
                    index,
                    f"a re-grown row or wrapped line here lands on the content below, "
                    f"and `rendered-overflow` cannot measure it because the neighbour's "
                    f"own box masks it — check for text under text: {shown}",
                    render=render,
                    shape=tight[0][0],
                )
            )

        if not reasons:
            continue
        result.findings.append(
            Finding(
                "vision-review",
                SEVERITY_ADVISORY,
                index,
                f"needs a look at the render for {_VISION_CHECKLIST}; flagged because "
                f"of {' and '.join(reasons)}"
                + (f". Render: {render}" if render else " (no render available)"),
            )
        )
    return result


# ─── Entry point ─────────────────────────────────────────────────────────────


def _attribution_findings(
    deck: Path, prs, baseline: LibraryBaseline, work: Path, *, already: set
) -> list[Finding]:
    """Run the attribution pass, or report loudly that it could not run."""
    if not baseline.attributed_ok:
        return [
            Finding(
                "render-unavailable",
                SEVERITY_BLOCKING,
                0,
                "could not render the reference library's attribution probes, so "
                "per-shape masked overflow had no baseline and was NOT checked",
            )
        ]
    try:
        attributed = measure_attributed_overflow(deck, work)
    except RuntimeError as exc:
        return [
            Finding(
                "render-unavailable",
                SEVERITY_BLOCKING,
                0,
                f"could not render the deck's attribution probes, so per-shape masked "
                f"overflow did NOT run: {exc}",
            )
        ]
    return _check_masked_overflow(prs, attributed, baseline, already)


def verify_deck(
    path: Path | str,
    *,
    library: Path | str | None = None,
    render: bool = True,
    vision: bool = True,
    attribute: bool = True,
    out_dir: Path | str | None = None,
) -> list[Finding]:
    """Verify a built deck against the contract; return every finding.

    Blocking findings come first, then advisory; within each tier, ordered by
    slide. **A list with no blocking findings is not a clean bill of health on its
    own** — look for ``render-unavailable`` and ``library-unavailable``, which mean
    a whole tier did not run, and read the advisory ``vision-review`` entries,
    which are the agenda for a human or agent rather than a verdict.

    ``render=False`` skips the render-measured tier (and the vision tier with it)
    for a fast XML-only pass. ``attribute=False`` keeps the ordinary render tier
    but skips the per-shape attribution pass and its one extra conversion.
    ``library`` overrides the reference slide library. ``out_dir`` is where
    renders and picture crops are written; pass a durable path when the
    checkpoint needs to show them.
    """
    deck = Path(path)
    if not deck.is_file():
        raise FileNotFoundError(f"deck not found: {deck}")
    prs = Presentation(deck)

    findings: list[Finding] = []
    findings.extend(_check_forbidden_strings(prs))

    library_path = Path(library) if library else default_library_path()
    if library_path is None or not library_path.is_file():
        findings.append(
            Finding(
                "library-unavailable",
                SEVERITY_BLOCKING,
                0,
                f"the reference slide library ({LIBRARY_TEMPLATE_NAME}) could not be "
                f"resolved, so the geometric checks had no baseline and did NOT run",
            )
        )
        return findings

    baseline = library_baseline(library_path, render=render, attribute=attribute)
    findings.extend(_check_shapes_within_slide(prs, baseline))
    findings.extend(_check_table_heights(prs, baseline))

    root = Path(out_dir) if out_dir else Path(tempfile.mkdtemp(prefix="deck-contract-"))
    if render:
        renders: dict[int, Path] = {}
        try:
            renders = _render_slides(deck, root / "slides", list(range(len(prs.slides))))
        except RuntimeError as exc:
            findings.append(
                Finding(
                    "render-unavailable",
                    SEVERITY_BLOCKING,
                    0,
                    f"could not render the deck, so the render-measured check "
                    f"(rendered overflow) did NOT run: {exc}",
                )
            )
        if renders and not baseline.rendered:
            findings.append(
                Finding(
                    "render-unavailable",
                    SEVERITY_BLOCKING,
                    0,
                    "could not render the reference library, so rendered overflow had no "
                    "baseline and was NOT checked",
                )
            )
        elif renders:
            overflow = _check_rendered_overflow(prs, renders, baseline)
            findings.extend(overflow)
            if attribute:
                findings.extend(
                    _attribution_findings(
                        deck,
                        prs,
                        baseline,
                        root / "attribution",
                        already={(f.slide, f.shape) for f in overflow},
                    )
                )
        if vision:
            findings.extend(
                vision_pass(deck, out_dir=root, renders=renders, baseline=baseline).findings
            )

    order = {SEVERITY_BLOCKING: 0, SEVERITY_ADVISORY: 1}
    findings.sort(key=lambda f: (order.get(f.severity, 9), f.slide, f.kind, f.shape or ""))
    return findings
